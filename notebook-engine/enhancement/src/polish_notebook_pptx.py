from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterator

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor.from_string("1F3B5A")
TEAL = RGBColor.from_string("1FA7A8")
GOLD = RGBColor.from_string("F2C14E")
PALE_TEAL = RGBColor.from_string("E7F6F5")
PALE_BLUE = RGBColor.from_string("EEF4FA")
INK = RGBColor.from_string("243647")
MID = RGBColor.from_string("6C8193")
LIGHT_BORDER = RGBColor.from_string("9EC7CC")

GENERIC_LESSON_FOCUS_MARKERS = (
    "interactive student notebook",
    "student notebook",
    "the lesson idea",
    "guided practice",
    "independent practice",
    "collaborative practice",
    "exit ticket",
    "reflection",
    "vocabulary",
)

VOCAB_HEADER_MARKERS = {
    "word",
    "student-friendly meaning",
    "student friendly meaning",
    "example",
    "visual/icon",
    "visual / icon",
    "key math language",
    "key language",
}

GENERIC_FOCUS_MARKERS = (
    "content objective",
    "language objective",
    "why it matters",
    "name:",
    "period:",
    "teacher:",
    "interactive student notebook",
    "session map",
    "today's notebook path",
)

STATISTICS_SIGNAL_TOKENS = {
    "statistical",
    "nonstatistical",
    "variability",
    "data",
    "distribution",
    "survey",
    "responses",
    "question",
    "questions",
    "dot",
    "plot",
}

SORTING_SIGNAL_TOKENS = {
    "classify",
    "compare",
    "favorite",
    "match",
    "nonstatistical",
    "sort",
    "statistical",
    "variability",
}

DATA_VISUAL_SIGNAL_TOKENS = {
    "data",
    "distribution",
    "dot",
    "graph",
    "plot",
    "table",
}

MIN_CHIP_FONT_PT = 9.0
MIN_FOOTER_FONT_PT = 9.0
MIN_META_FONT_PT = 10.5
MIN_BODY_FONT_PT = 12.0
MIN_SECTION_FONT_PT = 16.0
MIN_TITLE_FONT_PT = 20.0
TEXT_MARGIN_LEFT = Inches(0.08)
TEXT_MARGIN_RIGHT = Inches(0.08)
TEXT_MARGIN_TOP = Inches(0.05)
TEXT_MARGIN_BOTTOM = Inches(0.05)
MIN_LINE_HEIGHT = Inches(0.08)


@dataclass
class PolishStats:
    shapes_seen: int = 0
    text_frames_seen: int = 0
    table_cells_seen: int = 0
    font_lifts: int = 0
    margin_repairs: int = 0
    wrap_repairs: int = 0
    heading_styles: int = 0
    panel_tabs: int = 0
    response_zones: int = 0
    title_backplates: int = 0
    prompt_cards: int = 0
    added_slides: int = 0

    def to_dict(self) -> dict[str, int]:
        return {
            "shapesSeen": self.shapes_seen,
            "textFramesSeen": self.text_frames_seen,
            "tableCellsSeen": self.table_cells_seen,
            "fontLifts": self.font_lifts,
            "marginRepairs": self.margin_repairs,
            "wrapRepairs": self.wrap_repairs,
            "headingStyles": self.heading_styles,
            "panelTabs": self.panel_tabs,
            "responseZones": self.response_zones,
            "titleBackplates": self.title_backplates,
            "promptCards": self.prompt_cards,
            "addedSlides": self.added_slides,
        }


@dataclass
class NotebookContext:
    session_title: str
    unit_line: str
    footer_text: str
    standards: list[str]
    lesson_focus: str
    vocabulary: list[str]
    examples: list[str]
    claim_text: str


def iter_text_shapes(shapes) -> Iterator:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(shape.shapes)
            continue
        yield shape


def normalize_text(value: str) -> str:
    return " ".join(part.strip() for part in str(value or "").splitlines() if part.strip()).strip()


def truncate_text(text: str, limit: int) -> str:
    text = normalize_text(text)
    if len(text) <= limit:
        return text
    trimmed = text[: limit - 3].rstrip(" ,;:")
    return trimmed + "..."


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return normalize_text(shape.text)


def unique_in_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for item in items:
        key = item.lower()
        if key in seen:
            continue
        seen.add(key)
        ordered.append(item)
    return ordered


def slide_text_lines(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape_text(shape)
            if text:
                lines.append(text)
    return lines


def slide_text_blob(slide) -> str:
    return " | ".join(slide_text_lines(slide))


def is_standards_chip(shape, text: str) -> bool:
    return bool(re.fullmatch(r"\d+\.[A-Z]{1,3}", text.strip()) and shape.width <= Inches(1.1))


def is_small_chip(shape, text: str) -> bool:
    return bool(text and len(text) <= 12 and shape.width <= Inches(1.2) and shape.height <= Inches(0.42))


def is_footer(shape, slide_height, text: str) -> bool:
    return bool(text and shape.top >= slide_height - Inches(0.65))


def is_page_number(shape, text: str) -> bool:
    return bool(text.isdigit() and len(text) <= 2 and shape.top >= Inches(6.8))


def is_brand_or_meta_line(shape, text: str) -> bool:
    return bool(
        text
        and (
            text.startswith("EduWonderLab")
            or text.startswith("Unit ")
            or "Lesson " in text
            or "Session " in text
        )
        and shape.top <= Inches(1.0)
        and shape.height <= Inches(0.45)
    )


def is_main_cover_title(shape, text: str) -> bool:
    return bool(
        text
        and shape.top <= Inches(0.75)
        and shape.width >= Inches(5.5)
        and ("Notebook" in text or "Session" in text)
        and len(text) <= 100
    )


def is_section_title(shape, text: str) -> bool:
    return bool(
        text
        and shape.top <= Inches(1.65)
        and shape.height >= Inches(0.32)
        and len(text) <= 36
        and "·" not in text
        and not is_brand_or_meta_line(shape, text)
        and not is_standards_chip(shape, text)
    )


def slide_title_text(slide) -> str:
    for idx in range(min(6, len(slide.shapes))):
        shape = slide.shapes[idx]
        if getattr(shape, "has_text_frame", False):
            text = shape_text(shape)
            if text and not is_brand_or_meta_line(shape, text) and not is_standards_chip(shape, text):
                return text
    return ""


def title_shapes(slide) -> tuple[str, str]:
    title = ""
    subtitle = ""
    found = 0
    for idx in range(min(8, len(slide.shapes))):
        shape = slide.shapes[idx]
        if getattr(shape, "has_text_frame", False):
            text = shape_text(shape)
            if not text or is_standards_chip(shape, text) or is_brand_or_meta_line(shape, text):
                continue
            found += 1
            if found == 1:
                title = text
            elif found == 2:
                subtitle = text
                break
    return title, subtitle


def small_heading_candidate(shape, text: str, slide_height) -> bool:
    return bool(
        text
        and 1 <= len(text.split()) <= 5
        and shape.height <= Inches(0.45)
        and shape.top >= Inches(1.0)
        and shape.top <= slide_height - Inches(1.0)
        and shape.width <= Inches(4.4)
        and not is_page_number(shape, text)
        and not is_standards_chip(shape, text)
        and not is_brand_or_meta_line(shape, text)
        and not is_section_title(shape, text)
    )


def heading_style_palette(text: str) -> tuple[RGBColor, RGBColor]:
    lowered = text.lower()
    if any(token in lowered for token in ("error", "fix", "repair")):
        return GOLD, NAVY
    if any(token in lowered for token in ("discuss", "partner", "talk", "wonder", "reason")):
        return TEAL, RGBColor.from_string("FFFFFF")
    return NAVY, RGBColor.from_string("FFFFFF")


def style_heading_box(shape, stats: PolishStats) -> None:
    text = shape_text(shape)
    if not text:
        return
    fill_rgb, font_rgb = heading_style_palette(text)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.03)
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(max(run.font.size.pt if run.font.size else 0, 11.5))
            run.font.bold = True
            run.font.color.rgb = font_rgb
    stats.heading_styles += 1


def panel_tab_labels(slide_title: str) -> tuple[str | None, str | None]:
    key = slide_title.lower()
    mapping = {
        "guided practice": ("Work the lesson", "Write it out"),
        "drag-sort": ("Sort + decide", "Defend your sort"),
        "error analysis": ("Spot the error", "Explain + repair"),
        "independent practice": ("Solve + show", "Check your reasoning"),
        "exit ticket + reflection": ("Reflect", "Exit ticket"),
        "real world transfer": ("Transfer the lesson", "Connect it"),
        "dot plot constructor": ("Build the plot", "Analyze the pattern"),
        "frayer model": ("Build the meaning", "Use the word"),
        "partner a/b — turn and teach": ("Partner A", "Partner B"),
        "be curious": ("Notice the scenario", "Notice + wonder"),
    }
    for token, labels in mapping.items():
        if token in key:
            return labels
    return (None, None)


def add_panel_tab(slide, panel_shape, label: str, fill_rgb: RGBColor, stats: PolishStats) -> None:
    width = min(panel_shape.width * 0.46, Inches(2.4))
    left = panel_shape.left + Inches(0.06)
    top = max(panel_shape.top - Inches(0.18), Inches(0.62))
    tab = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.28))
    tab.fill.solid()
    tab.fill.fore_color.rgb = fill_rgb
    tab.line.fill.background()
    tf = tab.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.11)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    stats.panel_tabs += 1


def large_panel_shapes(slide) -> list:
    panels = []
    for shape in iter_text_shapes(slide.shapes):
        if (
            getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape.width >= Inches(3.0)
            and shape.height >= Inches(2.0)
            and shape.top >= Inches(0.85)
            and shape.top <= Inches(6.2)
        ):
            panels.append(shape)
    return sorted(panels, key=lambda item: (item.top, item.left))


def line_placeholder_shapes(slide) -> list:
    placeholders = []
    for shape in iter_text_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape_text(shape):
            continue
        if shape.width < Inches(1.6):
            continue
        if shape.height > MIN_LINE_HEIGHT:
            continue
        if shape.top < Inches(1.1) or shape.top > Inches(6.7):
            continue
        placeholders.append(shape)
    return sorted(placeholders, key=lambda item: (round(item.left.inches, 1), item.top))


def cluster_line_placeholders(placeholders: list) -> list[list]:
    groups: list[list] = []
    for shape in placeholders:
        placed = False
        for group in groups:
            anchor = group[0]
            same_column = abs(shape.left - anchor.left) <= Inches(0.6) and abs(shape.width - anchor.width) <= Inches(1.0)
            nearby = abs(shape.top - group[-1].top) <= Inches(1.2)
            if same_column and nearby:
                group.append(shape)
                placed = True
                break
        if not placed:
            groups.append([shape])
    return [group for group in groups if len(group) >= 2]


def response_zone_label(slide_title: str, cluster_left: int, slide_width: int) -> str:
    right_side = cluster_left > slide_width // 2
    key = slide_title.lower()
    if "error analysis" in key:
        return "Repair your thinking" if right_side else "Fix the mistake"
    if "drag-sort" in key:
        return "Defend your sort" if right_side else "Sort + record"
    if "guided practice" in key:
        return "Write the frame" if right_side else "Show your work"
    if "partner" in key:
        return "Respond here" if right_side else "Partner notes"
    if "independent practice" in key:
        return "Check + explain" if right_side else "Solve here"
    if "exit ticket" in key:
        return "Exit response" if right_side else "Reflection"
    if "real world transfer" in key:
        return "Connect it" if right_side else "Transfer your thinking"
    return "Write here"


def add_response_zone(slide, group: list, slide_title: str, slide_width: int, stats: PolishStats) -> None:
    left = min(shape.left for shape in group) - Inches(0.12)
    top = min(shape.top for shape in group) - Inches(0.12)
    right = max(shape.left + shape.width for shape in group) + Inches(0.12)
    bottom = max(shape.top + max(shape.height, MIN_LINE_HEIGHT) for shape in group) + Inches(0.18)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, right - left, bottom - top)
    box.fill.background()
    box.line.color.rgb = LIGHT_BORDER
    box.line.width = Pt(1.3)

    label = response_zone_label(slide_title, int(left), slide_width)
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.08),
        max(top - Inches(0.16), Inches(0.68)),
        min(Inches(1.9), right - left - Inches(0.16)),
        Inches(0.24),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = TEAL
    pill.line.fill.background()
    tf = pill.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    stats.response_zones += 1


def add_slide_publisher_polish(slide, slide_height, slide_width, stats: PolishStats) -> None:
    title = slide_title_text(slide)
    if not title:
        return
    for shape in iter_text_shapes(slide.shapes):
        text = shape_text(shape)
        if small_heading_candidate(shape, text, slide_height):
            style_heading_box(shape, stats)

    left_label, right_label = panel_tab_labels(title)
    panels = large_panel_shapes(slide)
    if left_label and panels:
        add_panel_tab(slide, panels[0], left_label, NAVY, stats)
    if right_label and len(panels) >= 2:
        add_panel_tab(slide, panels[1], right_label, TEAL, stats)

    placeholders = line_placeholder_shapes(slide)
    for group in cluster_line_placeholders(placeholders):
        add_response_zone(slide, group, title, slide_width, stats)


def collect_vocabulary(prs: Presentation) -> list[str]:
    terms: list[str] = []
    for slide in prs.slides:
        title, _subtitle = title_shapes(slide)
        if "vocabulary" not in title.lower():
            continue
        header_bottom = Inches(1.0)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape_text(shape)
            lowered = text.lower()
            if not text:
                continue
            if lowered in VOCAB_HEADER_MARKERS or lowered.startswith("key math language"):
                header_bottom = max(header_bottom, shape.top + shape.height)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape_text(shape)
            lowered = text.lower()
            if not text or len(text) > 40:
                continue
            if lowered in VOCAB_HEADER_MARKERS:
                continue
            if shape.top < header_bottom - Inches(0.02):
                continue
            if shape.left > Inches(3.35):
                continue
            if shape.width > Inches(2.75):
                continue
            if not (1 <= len(text.split()) <= 4):
                continue
            if is_brand_or_meta_line(shape, text) or is_standards_chip(shape, text) or is_page_number(shape, text):
                continue
            if lowered.startswith("unit ") or lowered.startswith("lesson ") or lowered in {"vocabulary", "statistics"}:
                continue
            if text not in terms:
                terms.append(text)
    return unique_in_order(terms)[:6]


def collect_examples(prs: Presentation) -> list[str]:
    examples: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape_text(shape)
            if not text:
                continue
            if "?" not in text and "because" not in text and "\"" not in text:
                continue
            if len(text) < 18 or len(text) > 180:
                continue
            if is_brand_or_meta_line(shape, text) or is_standards_chip(shape, text):
                continue
            if text not in examples:
                examples.append(text)
    return examples[:8]


def detect_claim_text(prs: Presentation, examples: list[str]) -> str:
    for slide in prs.slides:
        title, _subtitle = title_shapes(slide)
        if "error analysis" in title.lower():
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape_text(shape)
                if "\"" in text or "student says" in text.lower():
                    return text
    return examples[0] if examples else ""


def cover_focus_candidate(shape, text: str) -> bool:
    lowered = text.lower()
    return bool(
        text
        and shape.top >= Inches(1.2)
        and shape.top <= Inches(3.3)
        and shape.height >= Inches(0.4)
        and shape.width >= Inches(2.4)
        and len(text.split()) >= 2
        and len(text) <= 60
        and not is_brand_or_meta_line(shape, text)
        and not is_standards_chip(shape, text)
        and not is_page_number(shape, text)
        and lowered not in VOCAB_HEADER_MARKERS
        and not any(marker in lowered for marker in GENERIC_FOCUS_MARKERS)
        and not any(marker in lowered for marker in GENERIC_LESSON_FOCUS_MARKERS)
    )


def detect_cover_focus(prs: Presentation) -> str:
    if not prs.slides:
        return ""
    candidates: list[tuple[float, str]] = []
    for shape in prs.slides[0].shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape_text(shape)
        if cover_focus_candidate(shape, text):
            candidates.append((shape.top / 914400, text))
    if not candidates:
        return ""
    candidates.sort(key=lambda item: item[0])
    return candidates[0][1]


def detect_lesson_focus(prs: Presentation, vocabulary: list[str]) -> str:
    cover_focus = detect_cover_focus(prs)
    if cover_focus:
        return cover_focus
    generic_tokens = (
        "notebook",
        "session",
        "learning target",
        "key math language",
        "guided notes",
        "guided practice",
        "independent practice",
        "collaborative practice",
        "reflection",
        "exit ticket",
        "be curious",
        "partner",
        "turn and teach",
        "cover",
    )
    for slide in prs.slides:
        title, subtitle = title_shapes(slide)
        for candidate in (title, subtitle):
            cleaned = normalize_text(candidate)
            lowered = cleaned.lower()
            if not cleaned or any(token in lowered for token in generic_tokens):
                continue
            if len(cleaned.split()) >= 3:
                return cleaned
    if vocabulary:
        return vocabulary[0]
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape_text(shape)
        if text.startswith("Lesson ") or text.startswith("Unit "):
            if ":" in text:
                return text.split(":", 1)[1].strip()
            return text
    title, subtitle = title_shapes(first_slide)
    return subtitle or title or "the lesson idea"


def build_notebook_context(prs: Presentation) -> NotebookContext:
    first_slide = prs.slides[0]
    session_title, unit_line = title_shapes(first_slide)
    footer_text = ""
    standards: list[str] = []
    for shape in first_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape_text(shape)
        if is_standards_chip(shape, text):
            standards.append(text)
        elif shape.top >= Inches(7.0) and text and not text.isdigit():
            footer_text = text
    vocabulary = collect_vocabulary(prs)
    examples = collect_examples(prs)
    lesson_focus = detect_lesson_focus(prs, vocabulary)
    claim_text = detect_claim_text(prs, examples)
    return NotebookContext(
        session_title=session_title or "Interactive Student Notebook",
        unit_line=unit_line or lesson_focus,
        footer_text=footer_text or session_title or lesson_focus,
        standards=standards[:5],
        lesson_focus=lesson_focus,
        vocabulary=vocabulary,
        examples=examples,
        claim_text=claim_text,
    )


def lesson_terms(context: NotebookContext) -> list[str]:
    terms = unique_in_order([term for term in context.vocabulary if normalize_text(term)])
    signal_text = " ".join([context.lesson_focus, *context.examples[:2], context.claim_text]).lower()
    if any(token in signal_text for token in STATISTICS_SIGNAL_TOKENS):
        existing = {term.lower() for term in terms}
        for fallback in ("Statistical Question", "Variability", "Data", "Distribution"):
            if fallback.lower() not in existing:
                terms.append(fallback)
                existing.add(fallback.lower())
    return terms[:6]


def first_matching_term(terms: list[str], keywords: tuple[str, ...], fallback: str) -> str:
    for term in terms:
        lowered = term.lower()
        if any(keyword in lowered for keyword in keywords):
            return term
    return fallback


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.3)
    return shape


def set_shape_paragraphs(
    shape,
    paragraphs: list[dict[str, object]],
    *,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margin_left=Inches(0.08),
    margin_right=Inches(0.08),
    margin_top=Inches(0.04),
    margin_bottom=Inches(0.04),
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    tf.margin_left = margin_left
    tf.margin_right = margin_right
    tf.margin_top = margin_top
    tf.margin_bottom = margin_bottom
    for index, spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = spec.get("align", align)
        run = paragraph.add_run()
        run.text = str(spec.get("text", ""))
        run.font.size = Pt(float(spec.get("size", 11)))
        run.font.bold = bool(spec.get("bold", False))
        run.font.color.rgb = spec.get("color", INK)
        if "space_before" in spec:
            paragraph.space_before = Pt(float(spec["space_before"]))
        if "space_after" in spec:
            paragraph.space_after = Pt(float(spec["space_after"]))


def set_shape_text(
    shape,
    text: str,
    size: float,
    *,
    bold: bool = False,
    color=INK,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.MIDDLE,
    margin_left=Inches(0.08),
    margin_right=Inches(0.08),
    margin_top=Inches(0.04),
    margin_bottom=Inches(0.04),
) -> None:
    set_shape_paragraphs(
        shape,
        [{"text": text, "size": size, "bold": bold, "color": color, "align": align}],
        align=align,
        vertical=vertical,
        margin_left=margin_left,
        margin_right=margin_right,
        margin_top=margin_top,
        margin_bottom=margin_bottom,
    )


def move_shape_behind(shape, target_shape) -> None:
    shape_element = getattr(shape, "_element", None)
    target_element = getattr(target_shape, "_element", None)
    if shape_element is None or target_element is None:
        return
    parent = shape_element.getparent()
    if parent is None:
        return
    parent.remove(shape_element)
    target_element.addprevious(shape_element)


def add_text(slide, left, top, width, height, text, size, *, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_pill_chip(slide, left, top, width, text, fill_rgb, *, font_rgb=RGBColor.from_string("FFFFFF")):
    pill = add_rect(slide, left, top, width, Inches(0.28), fill_rgb, line_rgb=fill_rgb, radius=True)
    set_shape_text(
        pill,
        text,
        9.5,
        bold=True,
        color=font_rgb,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margin_left=Inches(0.04),
        margin_right=Inches(0.04),
        margin_top=Inches(0.01),
        margin_bottom=Inches(0.01),
    )
    return pill


def add_step_card(slide, left, top, width, step_number: int, title: str, body: str, fill_rgb) -> None:
    card = add_rect(slide, left, top, width, Inches(1.02), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    badge = add_rect(slide, left + Inches(0.12), top + Inches(0.12), Inches(0.42), Inches(0.42), fill_rgb, line_rgb=fill_rgb, radius=True)
    set_shape_text(
        badge,
        str(step_number),
        11,
        bold=True,
        color=RGBColor.from_string("FFFFFF"),
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    set_shape_paragraphs(
        card,
        [
            {"text": title, "size": 11.2, "bold": True, "color": NAVY},
            {"text": body, "size": 10.4, "color": INK, "space_before": 3},
        ],
        vertical=MSO_ANCHOR.TOP,
        margin_left=Inches(0.68),
        margin_right=Inches(0.12),
        margin_top=Inches(0.09),
        margin_bottom=Inches(0.08),
    )


def extract_question_prompts(text: str) -> list[str]:
    prompts: list[str] = []
    for quoted in re.findall(r'"([^"]+\?)"', text):
        cleaned = normalize_text(quoted)
        if cleaned:
            prompts.append(cleaned)
    if prompts:
        return prompts
    for candidate in re.findall(r"[^.?!]*\?", text):
        cleaned = normalize_text(candidate)
        if len(cleaned) >= 12:
            prompts.append(cleaned)
    return prompts


def lesson_question_examples(context: NotebookContext) -> list[str]:
    prompts: list[str] = []
    for source in [*context.examples, context.claim_text, context.lesson_focus]:
        prompts.extend(extract_question_prompts(source))
    prompts = unique_in_order([prompt for prompt in prompts if prompt])
    if len(prompts) < 2 and is_statistics_context(context):
        prompts.extend(
            [
                "What is your favorite park?",
                "In what state is Acadia National Park found?",
                "How many books do students read in a month?",
            ]
        )
        prompts = unique_in_order(prompts)
    return prompts[:3]


def add_question_choice_card(
    slide,
    left,
    top,
    width,
    label: str,
    prompt: str,
    accent_rgb,
    *,
    hint_text: str,
) -> None:
    card = add_rect(slide, left, top, width, Inches(2.1), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    header = add_rect(slide, left, top, width, Inches(0.34), accent_rgb, line_rgb=accent_rgb, radius=True)
    set_shape_text(
        header,
        label,
        10.4,
        bold=True,
        color=RGBColor.from_string("FFFFFF"),
        margin_left=Inches(0.08),
        margin_right=Inches(0.08),
        margin_top=Inches(0.02),
        margin_bottom=Inches(0.01),
    )
    add_text(slide, left + Inches(0.12), top + Inches(0.48), width - Inches(0.24), Inches(0.54), truncate_text(prompt, 95), 11.2, color=NAVY)
    many_tile = add_rect(slide, left + Inches(0.12), top + Inches(1.18), Inches(1.62), Inches(0.62), TEAL, line_rgb=TEAL, radius=True)
    set_shape_paragraphs(
        many_tile,
        [
            {"text": "o o o", "size": 10.8, "bold": True, "color": RGBColor.from_string("FFFFFF"), "align": PP_ALIGN.CENTER},
            {"text": "many answers", "size": 8.8, "bold": True, "color": RGBColor.from_string("FFFFFF"), "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margin_left=Inches(0.03),
        margin_right=Inches(0.03),
        margin_top=0,
        margin_bottom=0,
    )
    one_tile = add_rect(slide, left + Inches(1.92), top + Inches(1.18), Inches(1.48), Inches(0.62), NAVY, line_rgb=NAVY, radius=True)
    set_shape_paragraphs(
        one_tile,
        [
            {"text": "o", "size": 10.8, "bold": True, "color": RGBColor.from_string("FFFFFF"), "align": PP_ALIGN.CENTER},
            {"text": "one answer", "size": 8.8, "bold": True, "color": RGBColor.from_string("FFFFFF"), "align": PP_ALIGN.CENTER},
        ],
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margin_left=Inches(0.03),
        margin_right=Inches(0.03),
        margin_top=0,
        margin_bottom=0,
    )
    add_text(slide, left + Inches(3.58), top + Inches(1.26), width - Inches(3.7), Inches(0.36), hint_text, 9.4, color=MID)


def prompt_card_label(slide_title: str) -> str:
    key = slide_title.lower()
    if "error analysis" in key:
        return "Fix the mistake"
    if "guided practice" in key or "independent practice" in key:
        return "Show your thinking"
    if "exit ticket" in key or "reflection" in key:
        return "Write your takeaway"
    if "real world" in key:
        return "Transfer the idea"
    if "be curious" in key:
        return "Notice + wonder"
    return "Think it through"


def is_prompt_card_candidate(shape, text: str) -> bool:
    lowered = text.lower()
    starts_like_prompt = lowered.startswith(
        (
            "use ",
            "write ",
            "compare ",
            "explain ",
            "complete ",
            "record ",
            "show ",
            "solve ",
            "sort ",
            "look ",
            "read ",
            "notice ",
            "wonder ",
            "discuss ",
            "justify ",
            "decide ",
        )
    )
    upper_prompt_band = shape.top <= Inches(1.95) and len(text.split()) >= 6
    return bool(
        text
        and len(text) >= 24
        and len(text) <= 180
        and shape.top >= Inches(1.05)
        and shape.top <= Inches(2.65)
        and shape.width >= Inches(4.4)
        and shape.height <= Inches(1.35)
        and not is_section_title(shape, text)
        and not is_brand_or_meta_line(shape, text)
        and not is_standards_chip(shape, text)
        and (text.endswith("?") or starts_like_prompt or upper_prompt_band)
    )


def add_title_backplate(slide, shape, stats: PolishStats) -> None:
    backplate = add_rect(
        slide,
        max(shape.left - Inches(0.12), Inches(0.22)),
        max(shape.top - Inches(0.08), Inches(0.25)),
        min(shape.width + Inches(0.3), Inches(12.0)),
        shape.height + Inches(0.18),
        RGBColor.from_string("FFFFFF"),
        line_rgb=LIGHT_BORDER,
        radius=True,
    )
    move_shape_behind(backplate, shape)
    accent = add_rect(
        slide,
        backplate.left + Inches(0.14),
        backplate.top + backplate.height - Inches(0.09),
        min(Inches(2.2), backplate.width - Inches(0.28)),
        Inches(0.05),
        GOLD,
        line_rgb=GOLD,
        radius=True,
    )
    move_shape_behind(accent, shape)
    tf = shape.text_frame
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = NAVY
    stats.title_backplates += 1


def add_prompt_card(slide, shape, slide_title: str, stats: PolishStats) -> None:
    card = add_rect(
        slide,
        max(shape.left - Inches(0.18), Inches(0.24)),
        max(shape.top - Inches(0.12), Inches(0.9)),
        min(shape.width + Inches(0.36), Inches(12.2)),
        shape.height + Inches(0.24),
        RGBColor.from_string("FFFDF8"),
        line_rgb=LIGHT_BORDER,
        radius=True,
    )
    move_shape_behind(card, shape)
    bar = add_rect(
        slide,
        card.left,
        card.top,
        Inches(0.14),
        card.height,
        TEAL,
        line_rgb=TEAL,
        radius=False,
    )
    move_shape_behind(bar, shape)
    pill = add_rect(
        slide,
        card.left + Inches(0.22),
        max(card.top - Inches(0.16), Inches(0.72)),
        Inches(1.7),
        Inches(0.26),
        NAVY,
        line_rgb=NAVY,
        radius=True,
    )
    set_shape_text(
        pill,
        prompt_card_label(slide_title),
        9.8,
        bold=True,
        color=RGBColor.from_string("FFFFFF"),
        align=PP_ALIGN.CENTER,
        margin_left=Inches(0.04),
        margin_right=Inches(0.04),
        margin_top=Inches(0.01),
        margin_bottom=Inches(0.01),
    )
    tf = shape.text_frame
    tf.margin_left = Inches(0.14)
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = INK
    stats.prompt_cards += 1


def add_header_and_footer(slide, prs: Presentation, context: NotebookContext, page_number: int, slide_title: str, slide_subtitle: str) -> None:
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor.from_string("F9F4EA"))
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.55), NAVY)
    add_text(slide, Inches(0.25), Inches(0.0), Inches(8.0), Inches(0.34), slide_title, 20, bold=True, color=RGBColor.from_string("FFFFFF"))
    add_text(slide, Inches(0.25), Inches(0.34), Inches(8.0), Inches(0.21), slide_subtitle, 10.5, color=RGBColor.from_string("FFFFFF"))
    chip_left = Inches(8.4)
    for idx, standard in enumerate(context.standards[:5]):
        fill = GOLD if idx == min(4, len(context.standards) - 1) else RGBColor.from_string("355B78")
        chip = add_rect(slide, chip_left + Inches(0.87 * idx), Inches(0.1), Inches(0.82), Inches(0.34), fill, radius=True)
        set_shape_text(
            chip,
            standard,
            10,
            bold=True,
            color=NAVY if fill == GOLD else RGBColor.from_string("FFFFFF"),
            align=PP_ALIGN.CENTER,
            margin_left=0,
            margin_right=0,
            margin_top=Inches(0.01),
            margin_bottom=Inches(0.01),
        )
    add_rect(slide, 0, Inches(7.12), prs.slide_width, Inches(0.38), NAVY)
    add_text(slide, Inches(0.25), Inches(7.12), Inches(10.0), Inches(0.38), context.footer_text, 9, bold=True, color=RGBColor.from_string("FFFFFF"))
    add_text(slide, Inches(12.53), Inches(7.12), Inches(0.45), Inches(0.38), str(page_number), 10, bold=True, color=RGBColor.from_string("FFFFFF"), align=PP_ALIGN.CENTER)


def blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = normalize_text(getattr(layout, "name", ""))
        if name and "blank" in name.lower():
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def vocabulary_visual_cue(term: str) -> tuple[str, str]:
    lowered = term.lower()
    if "statistical" in lowered:
        return "o o o", "many answers"
    if "variability" in lowered:
        return "o  o   o", "answers vary"
    if "distribution" in lowered:
        return "o o o o", "look for spread"
    if "data" in lowered:
        return "o | o | o", "collected answers"
    return "o + label", "say it, show it"


def add_vocabulary_in_action_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    terms = lesson_terms(context)
    add_header_and_footer(slide, prs, context, page_number, "Vocabulary in Action", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.0), Inches(6.15), Inches(5.95), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(7.2), Inches(1.0), Inches(5.75), Inches(5.95), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Use the words well", NAVY, stats)
    add_panel_tab(slide, right_panel, "Write like a mathematician", TEAL, stats)
    add_text(slide, Inches(1.0), Inches(1.28), Inches(5.75), Inches(0.32), "Point to the cue, say the word, then add a quick label or sketch.", 11.2, bold=True, color=NAVY)

    card_left = Inches(1.0)
    card_top = Inches(1.8)
    card_width = Inches(1.74)
    gap = Inches(0.18)
    for index, term in enumerate(terms[:3], start=1):
        left = card_left + (card_width + gap) * (index - 1)
        cue_visual, cue_label = vocabulary_visual_cue(term)
        card = add_rect(slide, left, card_top, card_width, Inches(2.18), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
        set_shape_paragraphs(
            card,
            [
                {"text": f"{index}. {term}", "size": 11.6, "bold": True, "color": NAVY},
                {"text": cue_label, "size": 10.1, "bold": True, "color": TEAL, "space_before": 4},
                {"text": "Say it. Label it. Sketch it.", "size": 9.6, "color": INK, "space_before": 4},
            ],
            vertical=MSO_ANCHOR.TOP,
            margin_left=Inches(0.12),
            margin_right=Inches(0.12),
            margin_top=Inches(0.12),
            margin_bottom=Inches(0.08),
        )
        cue_tile = add_rect(slide, left + Inches(0.16), card_top + Inches(1.5), card_width - Inches(0.32), Inches(0.42), GOLD, line_rgb=GOLD, radius=True)
        set_shape_paragraphs(
            cue_tile,
            [
                {"text": cue_visual, "size": 10.6, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
                {"text": "visual cue", "size": 8.6, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER},
            ],
            align=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE,
            margin_left=Inches(0.03),
            margin_right=Inches(0.03),
            margin_top=0,
            margin_bottom=0,
        )

    coach_box = add_rect(slide, Inches(1.0), Inches(4.22), Inches(5.75), Inches(1.62), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    set_shape_paragraphs(
        coach_box,
        [
            {"text": "Word coach", "size": 11.2, "bold": True, "color": NAVY},
            {"text": "Point: Which card matches the page?", "size": 10.0, "color": INK, "space_before": 4},
            {"text": "Say: This word helps explain ___.", "size": 10.0, "color": INK, "space_before": 2},
            {"text": "Show: Add one label or tiny sketch.", "size": 10.0, "color": INK, "space_before": 2},
        ],
        vertical=MSO_ANCHOR.TOP,
        margin_left=Inches(0.14),
        margin_right=Inches(0.12),
        margin_top=Inches(0.12),
        margin_bottom=Inches(0.08),
    )

    add_text(slide, Inches(7.45), Inches(1.28), Inches(5.25), Inches(0.28), "Say it. Write it. Sketch it.", 11.2, bold=True, color=NAVY)
    add_step_card(slide, Inches(7.45), Inches(1.72), Inches(5.05), 1, "Say it aloud", "Tell a partner what the word helps you notice.", NAVY)
    add_step_card(slide, Inches(7.45), Inches(2.88), Inches(5.05), 2, "Write it clearly", "Finish the frame so the word explains the math.", TEAL)
    add_step_card(slide, Inches(7.45), Inches(4.04), Inches(5.05), 3, "Show it visually", "Add one symbol, label, or tiny sketch.", GOLD)
    frame = add_rect(slide, Inches(7.45), Inches(5.18), Inches(5.05), Inches(1.0), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    if is_statistics_context(context):
        statistical_term = first_matching_term(terms, ("statistical",), "statistical question")
        variability_term = first_matching_term(terms, ("variability",), "variability")
        frame_text = (
            f"In this lesson, a {statistical_term.lower()} helps me predict ___. I can use {variability_term.lower()} to explain why ___."
        )
    else:
        frame_text = (
            f"In this lesson, {terms[0]} helps me see ___. I can use {terms[min(1, len(terms) - 1)]} to explain ___."
        )
    add_text(
        slide,
        frame.left + Inches(0.14),
        frame.top + Inches(0.14),
        Inches(4.75),
        Inches(0.62),
        frame_text,
        10.8,
        color=INK,
    )
    add_pill_chip(slide, Inches(7.45), Inches(6.34), Inches(1.58), "partner talk", NAVY)
    add_pill_chip(slide, Inches(9.18), Inches(6.34), Inches(1.58), "quick sketch", TEAL)
    add_text(slide, Inches(10.9), Inches(6.34), Inches(1.6), Inches(0.28), "Goal: make the vocabulary explain the math.", 9.8, color=MID)
    stats.added_slides += 1


def sentence_starter(context: NotebookContext) -> str:
    terms = lesson_terms(context)
    if is_statistics_context(context):
        statistical_term = first_matching_term(terms, ("statistical",), "statistical question")
        variability_term = first_matching_term(terms, ("variability",), "variability")
        return f"This is a {statistical_term.lower()} because ___. I know there will be {variability_term.lower()} because ___."
    if terms:
        return f"Use {terms[0]} and {terms[min(1, len(terms)-1)]} in your explanation."
    return f"Use the lesson language from {context.lesson_focus} in your explanation."


def add_vocab_chips(slide, context: NotebookContext, left, top, width) -> None:
    terms = lesson_terms(context)
    if not terms:
        return
    add_text(slide, left, top, width, Inches(0.28), "Lesson language", 11.5, bold=True, color=NAVY)
    chip_top = top + Inches(0.38)
    chip_left = left
    row_height = Inches(0.34)
    for term in terms[:6]:
        chip = add_rect(slide, chip_left, chip_top, min(Inches(1.95), width - Inches(0.12)), row_height, PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
        set_shape_text(
            chip,
            term,
            10.2,
            bold=True,
            color=NAVY,
            margin_left=Inches(0.06),
            margin_right=Inches(0.06),
            margin_top=Inches(0.02),
            margin_bottom=Inches(0.02),
        )
        chip_top += Inches(0.42)


def lesson_signal_tokens(context: NotebookContext) -> set[str]:
    tokens = set(re.findall(r"[a-z0-9]+", context.lesson_focus.lower()))
    for vocab in lesson_terms(context):
        tokens.update(re.findall(r"[a-z0-9]+", vocab.lower()))
    for example in context.examples:
        tokens.update(re.findall(r"[a-z0-9]+", example.lower()))
    if context.claim_text:
        tokens.update(re.findall(r"[a-z0-9]+", context.claim_text.lower()))
    return tokens


def is_generic_lesson_focus(text: str) -> bool:
    lowered = normalize_text(text).lower()
    if not lowered:
        return True
    return any(marker == lowered or marker in lowered for marker in GENERIC_LESSON_FOCUS_MARKERS)


def context_signal_score(context: NotebookContext) -> int:
    score = 0
    if context.lesson_focus and not is_generic_lesson_focus(context.lesson_focus):
        score += 2
    if len(context.vocabulary) >= 2:
        score += 2
    elif context.vocabulary:
        score += 1
    if len(context.examples) >= 2:
        score += 2
    elif context.examples:
        score += 1
    if context.claim_text:
        score += 1
    if len(lesson_signal_tokens(context)) >= 8:
        score += 1
    return score


def context_supports_extensions(context: NotebookContext) -> bool:
    return context_signal_score(context) >= 4


def is_statistics_context(context: NotebookContext) -> bool:
    return bool(lesson_signal_tokens(context) & STATISTICS_SIGNAL_TOKENS)


def is_sorting_lesson(context: NotebookContext) -> bool:
    return bool(lesson_signal_tokens(context) & SORTING_SIGNAL_TOKENS)


def has_data_visual_signal(context: NotebookContext) -> bool:
    return bool(lesson_signal_tokens(context) & DATA_VISUAL_SIGNAL_TOKENS)


def mode_title(mode: str) -> str:
    return {
        "challenge_lab": "Challenge Lab",
        "compare_and_justify": "Compare and Justify",
        "example_nonexample": "Example vs. Non-Example",
        "fix_the_mistake": "Fix the Mistake",
        "mini_debate": "Mini Debate",
        "representation_connection": "Representation Connection",
        "vocabulary_in_action": "Vocabulary in Action",
    }.get(mode, mode.replace("_", " ").title())


def select_extension_modes(context: NotebookContext) -> list[str]:
    if not context.lesson_focus or not context_supports_extensions(context):
        return []
    focus_tokens = lesson_signal_tokens(context)
    focus_text = context.lesson_focus.lower()
    signal_score = context_signal_score(context)
    vocabulary_terms = lesson_terms(context)
    statistics_context = is_statistics_context(context)
    sorting_context = is_sorting_lesson(context)
    data_visual_context = has_data_visual_signal(context)
    modes: list[str] = []

    extra_scores: list[tuple[int, str]] = []
    if len(vocabulary_terms) >= 2:
        extra_scores.append((5, "vocabulary_in_action"))
    if sorting_context:
        extra_scores.append((6, "example_nonexample"))
    if len(context.examples) >= 2 or any(token in focus_text for token in ("compare", "justify", "strategy", "explain", "reason")) or sorting_context:
        extra_scores.append((5, "compare_and_justify"))
    if any(token in focus_text for token in ("error", "mistake", "wrong", "fix", "misconception")) or context.claim_text:
        extra_scores.append((4, "fix_the_mistake"))
    if statistics_context:
        if data_visual_context:
            extra_scores.append((3, "representation_connection"))
    elif any(token in focus_tokens for token in ("table", "graph", "equation", "model", "representation", "diagram", "data")):
        extra_scores.append((4, "representation_connection"))
    if any(token in focus_text for token in ("example", "nonexample", "non-example", "sort", "match", "classify")) and not sorting_context:
        extra_scores.append((3, "example_nonexample"))
    if signal_score >= 5:
        extra_scores.append((3, "mini_debate"))
    if signal_score >= 6 or len(context.examples) >= 2:
        extra_scores.append((2, "challenge_lab"))

    seen: set[str] = set()
    for _score, mode in sorted(extra_scores, key=lambda item: (-item[0], item[1])):
        if mode in seen:
            continue
        modes.append(mode)
        seen.add(mode)
        if len(modes) >= 4:
            break
    if signal_score >= 5 and len(modes) < 2:
        for fallback_mode in ("challenge_lab", "mini_debate"):
            if fallback_mode in seen:
                continue
            modes.append(fallback_mode)
            seen.add(fallback_mode)
            if len(modes) >= 2:
                break
    return modes


def add_challenge_lab_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    add_header_and_footer(slide, prs, context, page_number, "Challenge Lab", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.0), Inches(7.35), Inches(5.9), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(8.35), Inches(1.0), Inches(4.6), Inches(5.9), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Create + justify", NAVY, stats)
    add_panel_tab(slide, right_panel, "Use the lesson language", TEAL, stats)

    cards = [
        ("1. Build one", f"Create one new example that fits {context.lesson_focus.lower()}."),
        ("2. Make a near miss", "Change one part so it almost fits, but does not quite work anymore."),
        ("3. Defend your choice", f"Compare your idea to: {context.examples[0] if context.examples else context.lesson_focus}."),
    ]
    top = Inches(1.5)
    for title, body in cards:
        card = add_rect(slide, Inches(1.0), top, Inches(6.95), Inches(1.25), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
        add_text(slide, card.left + Inches(0.12), card.top + Inches(0.08), Inches(2.3), Inches(0.24), title, 11.5, bold=True, color=NAVY)
        add_text(slide, card.left + Inches(0.12), card.top + Inches(0.42), Inches(6.65), Inches(0.42), body, 12, color=INK)
        add_text(slide, card.left + Inches(0.12), card.top + Inches(0.86), Inches(6.65), Inches(0.22), "Write directly in the workspace below each move.", 10, color=MID)
        top += Inches(1.48)

    workspace = add_rect(slide, Inches(1.0), Inches(5.98), Inches(6.95), Inches(0.78), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    add_text(slide, workspace.left + Inches(0.12), workspace.top + Inches(0.06), Inches(2.1), Inches(0.2), "Notebook extension", 10, bold=True, color=TEAL)
    add_text(slide, workspace.left + Inches(0.12), workspace.top + Inches(0.28), Inches(6.6), Inches(0.32), sentence_starter(context), 10.5, color=INK)

    add_vocab_chips(slide, context, Inches(8.55), Inches(1.5), Inches(4.15))
    add_text(slide, Inches(8.55), Inches(4.2), Inches(4.05), Inches(0.26), "Success criteria", 11.5, bold=True, color=NAVY)
    criteria = [
        "uses lesson vocabulary",
        "matches the lesson idea",
        "shows one clear justification",
    ]
    criteria_top = Inches(4.58)
    for item in criteria:
        add_rect(slide, Inches(8.58), criteria_top + Inches(0.02), Inches(0.18), Inches(0.18), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=False)
        add_text(slide, Inches(8.85), criteria_top, Inches(3.75), Inches(0.2), item, 10.5, color=INK)
        criteria_top += Inches(0.34)
    add_text(slide, Inches(8.55), Inches(5.88), Inches(4.05), Inches(0.34), "Sentence starter", 11.5, bold=True, color=NAVY)
    add_text(slide, Inches(8.55), Inches(6.24), Inches(4.05), Inches(0.42), f"My example fits because ___. {sentence_starter(context)}", 10.5, color=INK)
    stats.added_slides += 1


def add_compare_and_justify_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    terms = lesson_terms(context)
    statistics_context = is_statistics_context(context)
    prompts = lesson_question_examples(context)
    add_header_and_footer(slide, prs, context, page_number, "Compare and Justify", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.05), Inches(5.85), Inches(5.8), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(6.95), Inches(1.05), Inches(5.85), Inches(5.8), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Question A" if statistics_context else "Model A", NAVY, stats)
    add_panel_tab(slide, right_panel, "Question B" if statistics_context else "Model B", TEAL, stats)
    if statistics_context:
        add_text(slide, Inches(1.0), Inches(1.34), Inches(5.45), Inches(0.28), "Compare the two question cards. Decide which one would create variable data.", 11.1, bold=True, color=NAVY)
        add_text(slide, Inches(7.15), Inches(1.34), Inches(5.45), Inches(0.28), "Use statistical question, nonstatistical question, and variability in your explanation.", 10.9, bold=True, color=NAVY)
        add_question_choice_card(
            slide,
            Inches(1.0),
            Inches(1.8),
            Inches(5.35),
            "Question Card A",
            prompts[0] if prompts else "What is your favorite park?",
            NAVY,
            hint_text="Circle the chip that matches this question.",
        )
        add_question_choice_card(
            slide,
            Inches(7.15),
            Inches(1.8),
            Inches(5.35),
            "Question Card B",
            prompts[1] if len(prompts) > 1 else "In what state is Acadia National Park found?",
            TEAL,
            hint_text="Look for one fixed answer or many possible answers.",
        )
    else:
        add_text(slide, Inches(1.0), Inches(1.38), Inches(5.45), Inches(0.28), "Compare the two representations, then justify which one best fits the lesson.", 11.3, bold=True, color=NAVY)
        add_text(slide, Inches(7.15), Inches(1.38), Inches(5.45), Inches(0.28), "Name one similarity and one difference before you choose.", 11.3, bold=True, color=NAVY)
        for top in [Inches(1.95), Inches(2.7), Inches(3.45)]:
            add_rect(slide, Inches(1.0), top, Inches(5.35), Inches(0.58), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
            add_rect(slide, Inches(7.15), top, Inches(5.35), Inches(0.58), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    add_pill_chip(slide, Inches(1.0), Inches(4.18), Inches(1.35), "look for", NAVY)
    add_pill_chip(slide, Inches(2.5), Inches(4.18), Inches(1.55), "many answers", TEAL)
    add_pill_chip(slide, Inches(4.18), Inches(4.18), Inches(1.3), "one answer", GOLD, font_rgb=NAVY)
    add_pill_chip(slide, Inches(7.15), Inches(4.18), Inches(1.4), "talk move", NAVY)
    add_pill_chip(slide, Inches(8.72), Inches(4.18), Inches(1.7), "evidence first", TEAL)
    add_text(slide, Inches(1.0), Inches(4.35), Inches(5.35), Inches(0.26), f"Use these words: {', '.join(terms[:3])}", 10.3, color=MID)
    add_text(slide, Inches(7.15), Inches(4.35), Inches(5.35), Inches(0.26), "Write a claim and back it up with evidence.", 10.3, color=MID)
    final_box = add_rect(slide, Inches(1.0), Inches(4.76), Inches(11.5), Inches(1.5), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    final_justification = (
        "Question ___ is statistical/nonstatistical because ___. I know the responses will ___ because ___."
        if statistics_context
        else f"I chose ___ because ___. {sentence_starter(context)}"
    )
    set_shape_paragraphs(
        final_box,
        [
            {"text": "Final justification", "size": 11.5, "bold": True, "color": NAVY},
            {"text": final_justification, "size": 11, "color": INK, "space_before": 4},
        ],
        vertical=MSO_ANCHOR.TOP,
        margin_left=Inches(0.12),
        margin_right=Inches(0.12),
        margin_top=Inches(0.12),
        margin_bottom=Inches(0.08),
    )
    stats.added_slides += 1


def add_fix_the_mistake_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    terms = lesson_terms(context)
    add_header_and_footer(slide, prs, context, page_number, "Fix the Mistake", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.05), Inches(7.1), Inches(5.85), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(8.15), Inches(1.05), Inches(4.65), Inches(5.85), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Spot the error", NAVY, stats)
    add_panel_tab(slide, right_panel, "Repair it", TEAL, stats)
    claim = context.claim_text or f"This answer matches {context.lesson_focus.lower()}."
    add_text(slide, Inches(1.0), Inches(1.42), Inches(6.65), Inches(0.42), "A student says:", 11.4, bold=True, color=NAVY)
    add_rect(slide, Inches(1.0), Inches(1.9), Inches(6.55), Inches(1.2), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    add_text(slide, Inches(1.15), Inches(2.08), Inches(6.15), Inches(0.54), truncate_text(claim, 130), 12, color=INK)
    add_text(slide, Inches(1.0), Inches(3.38), Inches(6.6), Inches(0.26), "What is wrong or incomplete?", 11.4, bold=True, color=NAVY)
    for top in [Inches(3.76), Inches(4.52), Inches(5.28)]:
        add_rect(slide, Inches(1.0), top, Inches(6.55), Inches(0.48), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    add_pill_chip(slide, Inches(1.0), Inches(6.0), Inches(1.55), "look for", NAVY)
    add_pill_chip(slide, Inches(2.72), Inches(6.0), Inches(1.7), "one fixed answer", GOLD, font_rgb=NAVY)
    add_pill_chip(slide, Inches(4.6), Inches(6.0), Inches(1.72), "many possible answers", TEAL)
    add_text(slide, Inches(8.35), Inches(1.42), Inches(4.25), Inches(0.3), "Use the lesson evidence to repair the claim.", 11.2, bold=True, color=NAVY)
    add_step_card(slide, Inches(8.35), Inches(1.86), Inches(4.15), 1, "Name the mistake", "Point to the word or idea that makes the claim break down.", NAVY)
    add_step_card(slide, Inches(8.35), Inches(3.02), Inches(4.15), 2, "Repair it with evidence", f"Use {terms[0] if terms else 'the lesson evidence'} in your fix.", TEAL)
    add_step_card(slide, Inches(8.35), Inches(4.18), Inches(4.15), 3, "Defend your repair", "Tell why your new version fits the lesson better.", GOLD)
    add_rect(slide, Inches(8.35), Inches(5.44), Inches(4.15), Inches(0.82), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    add_text(slide, Inches(8.48), Inches(5.62), Inches(3.9), Inches(0.35), sentence_starter(context), 10.6, color=INK)
    stats.added_slides += 1


def add_representation_connection_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    statistics_context = is_statistics_context(context)
    add_header_and_footer(slide, prs, context, page_number, "Representation Connection", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.05), Inches(12.0), Inches(5.85), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Track the idea three ways" if statistics_context else "Model the idea three ways", NAVY, stats)
    if statistics_context:
        add_text(slide, Inches(1.0), Inches(1.42), Inches(11.55), Inches(0.28), "Connect the lesson across the question, possible responses, and the reason it does or does not fit.", 11.1, bold=True, color=NAVY)
        columns = [
            ("Question / prompt", "Write or sketch the survey question you are analyzing.", Inches(1.0)),
            ("Possible responses", "List 3 or 4 answers you might collect from different people.", Inches(4.25)),
            ("Why it fits", "Explain how the responses show variability or one fixed answer.", Inches(7.5)),
        ]
    else:
        add_text(slide, Inches(1.0), Inches(1.42), Inches(11.55), Inches(0.28), "Connect the lesson across a model, an equation, and a written explanation.", 11.4, bold=True, color=NAVY)
        columns = [
            ("Model / diagram", "Show the idea with shapes, marks, or a visual model.", Inches(1.0)),
            ("Equation / notation", "Translate the model into symbols or a number sentence.", Inches(4.25)),
            ("Words / reasoning", "Explain how the model and equation match.", Inches(7.5)),
        ]
    for title, body, left in columns:
        card = add_rect(slide, left, Inches(1.95), Inches(2.8), Inches(2.45), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
        add_text(slide, card.left + Inches(0.12), card.top + Inches(0.12), Inches(2.45), Inches(0.24), title, 11.2, bold=True, color=NAVY)
        add_text(slide, card.left + Inches(0.12), card.top + Inches(0.48), Inches(2.45), Inches(0.95), body, 10.6, color=INK)
        add_rect(slide, card.left + Inches(0.12), card.top + Inches(1.56), Inches(2.55), Inches(0.62), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)

    add_text(slide, Inches(1.0), Inches(4.78), Inches(11.6), Inches(0.24), "Bridge the three columns with a sentence starter and the lesson vocabulary.", 10.8, color=MID)
    add_vocab_chips(slide, context, Inches(1.0), Inches(5.12), Inches(11.55))
    stats.added_slides += 1


def add_example_nonexample_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    prompts = lesson_question_examples(context)
    add_header_and_footer(slide, prs, context, page_number, "Example vs. Non-Example", f"Bonus · {context.lesson_focus}")
    left_panel = add_rect(slide, Inches(0.8), Inches(1.05), Inches(5.85), Inches(5.85), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(6.95), Inches(1.05), Inches(5.85), Inches(5.85), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Example", NAVY, stats)
    add_panel_tab(slide, right_panel, "Non-example", TEAL, stats)
    add_text(slide, Inches(1.0), Inches(1.34), Inches(5.45), Inches(0.28), "Card sort: which question belongs in the Example column?", 11.1, bold=True, color=NAVY)
    add_text(slide, Inches(7.15), Inches(1.34), Inches(5.45), Inches(0.28), "Use the lesson vocabulary to defend each choice.", 11.1, bold=True, color=NAVY)
    add_question_choice_card(
        slide,
        Inches(1.0),
        Inches(1.76),
        Inches(5.35),
        "Card A",
        prompts[0] if prompts else "What is your favorite park?",
        NAVY,
        hint_text="Drag or rewrite this card into the best column.",
    )
    add_question_choice_card(
        slide,
        Inches(7.15),
        Inches(1.76),
        Inches(5.35),
        "Card B",
        prompts[1] if len(prompts) > 1 else "In what state is Acadia National Park found?",
        TEAL,
        hint_text="Use the chips to help you decide.",
    )
    add_pill_chip(slide, Inches(1.0), Inches(4.08), Inches(1.7), "expect variability", TEAL)
    add_pill_chip(slide, Inches(2.9), Inches(4.08), Inches(1.55), "many answers", NAVY)
    add_pill_chip(slide, Inches(7.15), Inches(4.08), Inches(1.55), "one answer", GOLD, font_rgb=NAVY)
    add_pill_chip(slide, Inches(8.9), Inches(4.08), Inches(1.65), "fixed fact", NAVY)
    add_text(slide, Inches(1.0), Inches(4.4), Inches(5.35), Inches(0.26), f"Example because ___. {sentence_starter(context)}", 10.2, color=MID)
    add_text(slide, Inches(7.15), Inches(4.4), Inches(5.35), Inches(0.26), "Non-example because ___.", 10.2, color=MID)
    quick_write_box = add_rect(slide, Inches(1.0), Inches(4.76), Inches(11.5), Inches(1.5), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
    set_shape_paragraphs(
        quick_write_box,
        [
            {"text": "Quick write", "size": 11.5, "bold": True, "color": NAVY},
            {"text": "The clearest example is ___ because ___. The non-example misses ___.", "size": 11, "color": INK, "space_before": 4},
        ],
        vertical=MSO_ANCHOR.TOP,
        margin_left=Inches(0.12),
        margin_right=Inches(0.12),
        margin_top=Inches(0.12),
        margin_bottom=Inches(0.08),
    )
    add_pill_chip(slide, Inches(1.0), Inches(6.34), Inches(1.55), "partner talk", NAVY)
    add_text(slide, Inches(2.72), Inches(6.34), Inches(4.2), Inches(0.26), "Say: I chose this card because the responses would ___.", 9.8, color=MID)
    stats.added_slides += 1


def add_mini_debate_slide(prs: Presentation, context: NotebookContext, stats: PolishStats) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    page_number = len(prs.slides)
    add_header_and_footer(slide, prs, context, page_number, "Mini Debate", f"Bonus · {context.lesson_focus}")
    claim_card = add_rect(slide, Inches(0.8), Inches(1.0), Inches(12.15), Inches(1.25), GOLD, line_rgb=GOLD, radius=True)
    add_text(slide, claim_card.left + Inches(0.16), claim_card.top + Inches(0.08), Inches(1.5), Inches(0.24), "Claim", 12, bold=True, color=NAVY)
    claim_text = context.claim_text or (context.examples[0] if context.examples else f"This example matches {context.lesson_focus.lower()}.")
    add_text(slide, claim_card.left + Inches(0.16), claim_card.top + Inches(0.42), Inches(11.7), Inches(0.56), claim_text, 12.5, color=NAVY)

    left_panel = add_rect(slide, Inches(0.8), Inches(2.55), Inches(5.85), Inches(4.1), PALE_BLUE, line_rgb=LIGHT_BORDER, radius=True)
    right_panel = add_rect(slide, Inches(6.95), Inches(2.55), Inches(5.85), Inches(4.1), PALE_TEAL, line_rgb=LIGHT_BORDER, radius=True)
    add_panel_tab(slide, left_panel, "Agree", NAVY, stats)
    add_panel_tab(slide, right_panel, "Disagree", TEAL, stats)

    add_text(slide, Inches(1.0), Inches(2.98), Inches(5.45), Inches(0.34), "What evidence supports the claim?", 11.5, bold=True, color=NAVY)
    add_text(slide, Inches(7.15), Inches(2.98), Inches(5.45), Inches(0.34), "What evidence pushes back?", 11.5, bold=True, color=NAVY)

    for top in [Inches(3.5), Inches(4.28), Inches(5.06)]:
        add_rect(slide, Inches(1.0), top, Inches(5.35), Inches(0.56), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)
        add_rect(slide, Inches(7.15), top, Inches(5.35), Inches(0.56), RGBColor.from_string("FFFFFF"), line_rgb=LIGHT_BORDER, radius=True)

    add_text(slide, Inches(1.08), Inches(5.88), Inches(5.2), Inches(0.42), f"I agree because ___. {sentence_starter(context)}", 10.5, color=INK)
    add_text(slide, Inches(7.23), Inches(5.88), Inches(5.2), Inches(0.42), f"I disagree because ___. {sentence_starter(context)}", 10.5, color=INK)

    verdict = add_rect(slide, Inches(0.8), Inches(6.78), Inches(12.15), Inches(0.26), NAVY, line_rgb=NAVY, radius=True)
    add_text(slide, verdict.left + Inches(0.12), verdict.top + Inches(0.02), Inches(5.0), Inches(0.18), "Final verdict: Which side made the stronger case?", 10.5, bold=True, color=RGBColor.from_string("FFFFFF"))
    stats.added_slides += 1


def add_extension_slides(
    prs: Presentation,
    context: NotebookContext,
    stats: PolishStats,
    requested_modes: list[str] | None = None,
) -> list[str]:
    if not context.lesson_focus:
        return []
    renderers = {
        "challenge_lab": add_challenge_lab_slide,
        "compare_and_justify": add_compare_and_justify_slide,
        "example_nonexample": add_example_nonexample_slide,
        "fix_the_mistake": add_fix_the_mistake_slide,
        "mini_debate": add_mini_debate_slide,
        "representation_connection": add_representation_connection_slide,
        "vocabulary_in_action": add_vocabulary_in_action_slide,
    }
    applied_modes: list[str] = []
    for mode in requested_modes if requested_modes is not None else select_extension_modes(context):
        renderer = renderers.get(mode)
        if renderer is None:
            continue
        renderer(prs, context, stats)
        applied_modes.append(mode)
    return applied_modes


def lesson_overlap_count(text: str, context: NotebookContext) -> int:
    text_tokens = set(re.findall(r"[a-z0-9]+", text.lower()))
    return len(text_tokens & lesson_signal_tokens(context))


def audit_added_slide_text(mode: str, text: str, context: NotebookContext) -> list[str]:
    issues: list[str] = []
    lowered = text.lower()
    overlap = lesson_overlap_count(text, context)
    terms = lesson_terms(context)
    if "the lesson idea" in lowered or "source example" in lowered:
        issues.append("generic wording remained in the added slide")
    if mode == "vocabulary_in_action":
        if any(marker in lowered for marker in ("student-friendly meaning", "visual/icon")):
            issues.append("table headers leaked into the vocabulary extension slide")
        if terms and not any(term.lower() in lowered for term in terms[:2]):
            issues.append("vocabulary extension slide is not using the extracted lesson terms")
    if is_statistics_context(context):
        if mode == "representation_connection":
            issues.append("statistics lesson received a generic representation-connection slide")
        if mode == "compare_and_justify" and not ("statistical" in lowered and "variability" in lowered):
            issues.append("statistics compare slide is missing core lesson language")
        if mode == "example_nonexample" and not ("statistical" in lowered or "variability" in lowered):
            issues.append("example/non-example slide is too generic for the statistics lesson")
    elif overlap < 3 and mode in {"compare_and_justify", "example_nonexample", "representation_connection", "challenge_lab"}:
        issues.append("added slide is not anchored tightly enough to the lesson signal")
    return issues


def audit_category_payload(status: str, score: int, notes: list[str] | None = None) -> dict[str, object]:
    return {"status": status, "score": score, "notes": notes or []}


def scaffold_marker_count(text: str) -> int:
    markers = (
        "card a",
        "card b",
        "question card",
        "many answers?",
        "many answers",
        "one answer?",
        "one answer",
        "expect variability",
        "fixed fact",
        "partner talk",
        "talk move",
        "say it",
        "show it",
        "visual cue",
        "word coach",
        "quick sketch",
        "look for",
        "name the mistake",
        "repair it with evidence",
        "defend your repair",
        "quick write",
        "final justification",
        "use these words",
        "sentence starter",
    )
    lowered = text.lower()
    return sum(1 for marker in markers if marker in lowered)


def support_marker_count(text: str) -> int:
    markers = (
        "partner talk",
        "talk move",
        "sentence starter",
        "say it",
        "show it",
        "word coach",
        "quick sketch",
        "use these words",
        "look for",
        "repair it with evidence",
        "defend your repair",
    )
    lowered = text.lower()
    return sum(1 for marker in markers if marker in lowered)


def is_content_text_shape(shape, slide_height: int, text: str) -> bool:
    if not text:
        return False
    if shape.top < Inches(0.72):
        return False
    if is_standards_chip(shape, text) or is_page_number(shape, text):
        return False
    if is_footer(shape, slide_height, text) or is_brand_or_meta_line(shape, text):
        return False
    return True


def text_shape_count(slide, slide_height: int) -> int:
    count = 0
    for shape in slide.shapes:
        text = shape_text(shape) if getattr(shape, "has_text_frame", False) else ""
        if is_content_text_shape(shape, slide_height, text):
            count += 1
    return count


def slide_text_character_count(slide, slide_height: int) -> int:
    total = 0
    for shape in slide.shapes:
        text = shape_text(shape) if getattr(shape, "has_text_frame", False) else ""
        if is_content_text_shape(shape, slide_height, text):
            total += len(text)
    return total


def build_deck_audit_categories(
    context: NotebookContext,
    original_slide_count: int,
    applied_modes: list[str],
    original_shape_deltas: list[int],
    reviewed_slides: list[dict[str, object]],
) -> tuple[dict[str, dict[str, object]], list[str], dict[str, object]]:
    total_original_shape_delta = sum(max(delta, 0) for delta in original_shape_deltas)
    upgraded_original_slides = sum(1 for delta in original_shape_deltas if delta >= 2)
    strongly_upgraded_original_slides = sum(1 for delta in original_shape_deltas if delta >= 4)
    scaffold_counts = [int(item.get("scaffoldMarkers", 0)) for item in reviewed_slides]
    support_counts = [int(item.get("supportMarkers", 0)) for item in reviewed_slides]
    overlap_counts = [int(item.get("lesson_overlap", 0)) for item in reviewed_slides]
    shape_counts = [int(item.get("shapeCount", 0)) for item in reviewed_slides]
    text_shape_counts = [int(item.get("textShapeCount", 0)) for item in reviewed_slides]
    text_char_counts = [int(item.get("textCharacterCount", 0)) for item in reviewed_slides]
    added_slide_count = len(reviewed_slides)
    min_overlap = min(overlap_counts) if overlap_counts else 0
    min_scaffold = min(scaffold_counts) if scaffold_counts else 0
    total_support_markers = sum(support_counts)
    max_shape_count = max(shape_counts) if shape_counts else 0
    max_text_shape_count = max(text_shape_counts) if text_shape_counts else 0
    max_text_char_count = max(text_char_counts) if text_char_counts else 0

    issues: list[str] = []
    categories: dict[str, dict[str, object]] = {}

    change_pass = total_original_shape_delta >= max(18, original_slide_count + 8)
    change_weak = total_original_shape_delta >= max(10, original_slide_count // 2 + 4)
    categories["changeMagnitude"] = audit_category_payload(
        "pass" if change_pass else "weak" if change_weak else "fail",
        2 if change_pass else 1 if change_weak else 0,
        [f"Original-slide shape delta: {total_original_shape_delta} across {original_slide_count} source slides."],
    )
    if not change_weak:
        issues.append("deck-wide enhancement delta is too small to count as a rigorous upgrade")

    required_upgraded = min(original_slide_count, max(6, original_slide_count // 2))
    required_strong = min(original_slide_count, max(1, original_slide_count // 4, 3 if original_slide_count >= 6 else 1))
    weak_upgraded = min(original_slide_count, max(4, original_slide_count // 3))
    coverage_pass = upgraded_original_slides >= required_upgraded and strongly_upgraded_original_slides >= required_strong
    coverage_weak = upgraded_original_slides >= weak_upgraded
    categories["existingSlideCoverage"] = audit_category_payload(
        "pass" if coverage_pass else "weak" if coverage_weak else "fail",
        2 if coverage_pass else 1 if coverage_weak else 0,
        [
            f"{upgraded_original_slides} source slides gained 2+ shapes.",
            f"{strongly_upgraded_original_slides} source slides gained 4+ shapes.",
        ],
    )
    if not coverage_weak:
        issues.append("too few original slides show visible enhancement coverage")

    scaffold_pass = added_slide_count > 0 and min_scaffold >= 4 and sum(scaffold_counts) >= added_slide_count * 5
    scaffold_weak = added_slide_count == 0 or (min_scaffold >= 3 and sum(scaffold_counts) >= max(6, added_slide_count * 3))
    categories["addedSlideScaffolds"] = audit_category_payload(
        "pass" if scaffold_pass else "weak" if scaffold_weak else "fail",
        2 if scaffold_pass else 1 if scaffold_weak else 0,
        [f"Added-slide scaffold markers: {sum(scaffold_counts)} total, minimum {min_scaffold} on any added slide."],
    )
    if added_slide_count > 0 and not scaffold_weak:
        issues.append("added slides still look too under-scaffolded to qualify as rigorous enhancement")

    overlap_target = 10 if is_statistics_context(context) else 6
    anchoring_pass = added_slide_count == 0 or min_overlap >= overlap_target
    anchoring_weak = added_slide_count == 0 or min_overlap >= max(4, overlap_target - 2)
    categories["lessonAnchoring"] = audit_category_payload(
        "pass" if anchoring_pass else "weak" if anchoring_weak else "fail",
        2 if anchoring_pass else 1 if anchoring_weak else 0,
        [f"Lowest added-slide lesson overlap: {min_overlap} tokens."],
    )
    if added_slide_count > 0 and not anchoring_weak:
        issues.append("at least one added slide is too weakly anchored to the lesson to pass a rigorous audit")

    supports_pass = added_slide_count > 0 and total_support_markers >= added_slide_count * 3
    supports_weak = total_support_markers >= max(2, added_slide_count * 2)
    categories["studentSupports"] = audit_category_payload(
        "pass" if supports_pass else "weak" if supports_weak else "fail",
        2 if supports_pass else 1 if supports_weak else 0,
        [f"Student-support markers across added slides: {total_support_markers}."],
    )
    if added_slide_count > 0 and not supports_weak:
        issues.append("added slides do not include enough explicit student-facing supports")

    layout_pass = (
        added_slide_count == 0
        or (max_shape_count <= 40 and max_text_shape_count <= 24 and max_text_char_count <= 850)
    )
    layout_weak = (
        added_slide_count == 0
        or (max_shape_count <= 48 and max_text_shape_count <= 28 and max_text_char_count <= 980)
    )
    categories["layoutPressure"] = audit_category_payload(
        "pass" if layout_pass else "weak" if layout_weak else "fail",
        2 if layout_pass else 1 if layout_weak else 0,
        [
            f"Max added-slide shape count: {max_shape_count}.",
            f"Max added-slide text boxes: {max_text_shape_count}.",
            f"Max added-slide text characters: {max_text_char_count}.",
        ],
    )
    if added_slide_count > 0 and not layout_weak:
        issues.append("at least one added slide is too crowded to count as a polished classroom-ready layout")

    metrics = {
        "originalSlideCount": original_slide_count,
        "addedSlideCount": added_slide_count,
        "totalOriginalShapeDelta": total_original_shape_delta,
        "upgradedOriginalSlides": upgraded_original_slides,
        "stronglyUpgradedOriginalSlides": strongly_upgraded_original_slides,
        "minAddedSlideLessonOverlap": min_overlap,
        "minAddedSlideScaffoldMarkers": min_scaffold,
        "totalAddedSlideSupportMarkers": total_support_markers,
        "maxAddedSlideShapeCount": max_shape_count,
        "maxAddedSlideTextShapeCount": max_text_shape_count,
        "maxAddedSlideTextCharacterCount": max_text_char_count,
        "requestedModes": applied_modes,
    }
    return categories, issues, metrics


def audit_polished_output(
    source_path: Path,
    output_path: Path,
    context: NotebookContext,
    applied_modes: list[str],
    original_slide_count: int,
) -> dict[str, object]:
    source_review = Presentation(str(source_path))
    review = Presentation(str(output_path))
    slide_height = review.slide_height
    issues: list[str] = []
    reviewed_slides: list[dict[str, object]] = []
    flagged_modes: list[str] = []
    if len(review.slides) < original_slide_count:
        issues.append("polish output lost source slides")
    original_shape_deltas: list[int] = []
    for idx in range(min(original_slide_count, len(source_review.slides), len(review.slides))):
        original_shape_deltas.append(len(review.slides[idx].shapes) - len(source_review.slides[idx].shapes))
    added_slides = list(review.slides)[original_slide_count:]
    if len(added_slides) != len(applied_modes):
        issues.append("added slide count does not match the requested extension slide count")
    for idx, slide in enumerate(added_slides):
        mode = applied_modes[idx] if idx < len(applied_modes) else "unknown"
        text = slide_text_blob(slide)
        slide_issues = audit_added_slide_text(mode, text, context)
        if slide_issues and mode not in flagged_modes:
            flagged_modes.append(mode)
        issues.extend(f"{mode_title(mode)}: {item}" for item in slide_issues)
        reviewed_slides.append(
            {
                "mode": mode,
                "title": mode_title(mode),
                "shapeCount": len(slide.shapes),
                "textShapeCount": text_shape_count(slide, slide_height),
                "textCharacterCount": slide_text_character_count(slide, slide_height),
                "lesson_overlap": lesson_overlap_count(text, context),
                "scaffoldMarkers": scaffold_marker_count(text),
                "supportMarkers": support_marker_count(text),
                "issues": slide_issues,
            }
        )
    categories, category_issues, metrics = build_deck_audit_categories(
        context,
        original_slide_count,
        applied_modes,
        original_shape_deltas,
        reviewed_slides,
    )
    issues.extend(category_issues)
    critical_categories = (
        "addedSlideScaffolds",
        "lessonAnchoring",
        "studentSupports",
        "layoutPressure",
    )
    return {
        "passed": not issues
        and all(item["status"] != "fail" for item in categories.values())
        and all(categories[name]["status"] == "pass" for name in critical_categories),
        "issues": issues,
        "flaggedModes": flagged_modes,
        "categories": categories,
        "deckMetrics": metrics,
        "reviewedSlides": reviewed_slides,
    }


def run_polish_pass(
    source_path: Path,
    output_path: Path,
    requested_modes: list[str] | None = None,
) -> tuple[NotebookContext, PolishStats, list[str], int]:
    prs = Presentation(str(source_path))
    stats = PolishStats()
    context = build_notebook_context(prs)
    original_slide_count = len(prs.slides)
    slide_height = prs.slide_height
    slide_width = prs.slide_width
    for slide in prs.slides:
        title_candidates: list[tuple[object, str]] = []
        prompt_candidates: list[tuple[object, str]] = []
        for shape in iter_text_shapes(slide.shapes):
            stats.shapes_seen += 1
            if getattr(shape, "has_table", False):
                repair_table(shape.table, stats)
            if getattr(shape, "has_text_frame", False):
                text = shape_text(shape)
                if not text:
                    continue
                min_pt = target_min_font_pt(shape, slide_height, text)
                repair_text_frame(shape.text_frame, min_pt, stats)
                if is_section_title(shape, text):
                    title_candidates.append((shape, text))
                elif is_prompt_card_candidate(shape, text):
                    prompt_candidates.append((shape, text))
        if title_candidates:
            add_title_backplate(slide, title_candidates[0][0], stats)
        if prompt_candidates:
            add_prompt_card(slide, prompt_candidates[0][0], slide_title_text(slide) or prompt_candidates[0][1], stats)
        add_slide_publisher_polish(slide, slide_height, slide_width, stats)
    applied_modes = add_extension_slides(prs, context, stats, requested_modes=requested_modes)
    prs.save(str(output_path))
    return context, stats, applied_modes, original_slide_count


def determine_quality_tier(context: NotebookContext, stats: PolishStats, applied_modes: list[str]) -> tuple[str, dict[str, int], list[str]]:
    signal_score = context_signal_score(context)
    quality_signals = {
        "lessonSignal": signal_score,
        "styleRepairs": int(stats.heading_styles > 0) + int(stats.prompt_cards > 0) + int(stats.response_zones > 0) + int(stats.title_backplates > 0),
        "extensionSlides": stats.added_slides,
        "distinctModes": len(applied_modes),
    }
    warnings: list[str] = []
    if signal_score < 4:
        warnings.append("Lesson signal was too thin for flagship extension slides, so the enhancer stayed in repair-only mode.")
    if signal_score >= 4 and stats.added_slides < 2:
        warnings.append("Lesson signal supported extensions, but fewer than two premium extension slides were produced.")
    if quality_signals["styleRepairs"] < 2:
        warnings.append("The source deck needed more visual hierarchy repair than the enhancer could fully add in one pass.")
    if signal_score >= 5 and stats.added_slides >= 2 and quality_signals["styleRepairs"] >= 3 and len(applied_modes) >= 2:
        return "premium", quality_signals, warnings
    if quality_signals["styleRepairs"] >= 1 or stats.added_slides >= 1:
        return "enhanced", quality_signals, warnings
    return "fail", quality_signals, warnings


def target_min_font_pt(shape, slide_height, text: str) -> float:
    if is_page_number(shape, text):
        return MIN_FOOTER_FONT_PT
    if is_standards_chip(shape, text):
        return MIN_CHIP_FONT_PT
    if is_small_chip(shape, text):
        return MIN_CHIP_FONT_PT
    if is_footer(shape, slide_height, text):
        return MIN_FOOTER_FONT_PT
    if is_brand_or_meta_line(shape, text):
        return MIN_META_FONT_PT
    if is_main_cover_title(shape, text):
        return MIN_TITLE_FONT_PT
    if is_section_title(shape, text):
        return MIN_SECTION_FONT_PT
    return MIN_BODY_FONT_PT


def raise_text_frame_fonts(text_frame, min_pt: float, stats: PolishStats) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is None:
                continue
            current_pt = run.font.size.pt
            if current_pt < min_pt:
                run.font.size = Pt(min_pt)
                stats.font_lifts += 1


def repair_text_frame(text_frame, min_pt: float, stats: PolishStats) -> None:
    stats.text_frames_seen += 1
    if not text_frame.word_wrap:
        text_frame.word_wrap = True
        stats.wrap_repairs += 1
    if text_frame.margin_left < TEXT_MARGIN_LEFT:
        text_frame.margin_left = TEXT_MARGIN_LEFT
        stats.margin_repairs += 1
    if text_frame.margin_right < TEXT_MARGIN_RIGHT:
        text_frame.margin_right = TEXT_MARGIN_RIGHT
        stats.margin_repairs += 1
    if text_frame.margin_top < TEXT_MARGIN_TOP:
        text_frame.margin_top = TEXT_MARGIN_TOP
        stats.margin_repairs += 1
    if text_frame.margin_bottom < TEXT_MARGIN_BOTTOM:
        text_frame.margin_bottom = TEXT_MARGIN_BOTTOM
        stats.margin_repairs += 1
    raise_text_frame_fonts(text_frame, min_pt, stats)


def repair_table(table, stats: PolishStats) -> None:
    for row in table.rows:
        for cell in row.cells:
            stats.table_cells_seen += 1
            repair_text_frame(cell.text_frame, MIN_BODY_FONT_PT, stats)


def polish_notebook_pptx(source_path: Path, output_path: Path) -> dict[str, object]:
    source_path = source_path.resolve()
    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    context, stats, applied_modes, original_slide_count = run_polish_pass(source_path, output_path)
    output_audit = audit_polished_output(source_path, output_path, context, applied_modes, original_slide_count)
    audit_passes = 1
    if output_audit["issues"]:
        filtered_modes = [mode for mode in applied_modes if mode not in set(output_audit["flaggedModes"])]
        if filtered_modes != applied_modes:
            context, stats, applied_modes, original_slide_count = run_polish_pass(
                source_path,
                output_path,
                requested_modes=filtered_modes,
            )
            output_audit = audit_polished_output(source_path, output_path, context, applied_modes, original_slide_count)
            audit_passes = 2
    quality_tier, quality_signals, warnings = determine_quality_tier(context, stats, applied_modes)
    if not output_audit["passed"]:
        warnings = warnings + [f"Output audit: {issue}" for issue in output_audit["issues"]]
        quality_tier = "enhanced" if quality_tier == "premium" else quality_tier
    return {
        "source": str(source_path),
        "output": str(output_path),
        "context": {
            "lessonFocus": context.lesson_focus,
            "vocabulary": context.vocabulary,
            "exampleCount": len(context.examples),
        },
        "qualityTier": quality_tier,
        "qualitySignals": quality_signals,
        "warnings": warnings,
        "auditPasses": audit_passes,
        "outputAudit": output_audit,
        "extensionModes": applied_modes,
        "stats": stats.to_dict(),
    }


def write_pptx_polish_report(path: Path, payload: dict[str, object]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=True) + "\n", encoding="utf-8")
