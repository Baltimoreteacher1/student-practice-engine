import unittest

from pptx import Presentation

from notebook_engine import (
    AUTO_FIT_MIN_FONT_PT,
    apply_flagship_activity_layer,
    build_flagship_activity_specs,
    build_flagship_session_data,
    choose_activity_candidates,
    choose_flagship_activity_mix,
    ensure_source_aligned_be_curious_activity,
    fit_text_style,
    flagship_activity_has_generic_markers,
    FORMAL_REVIEW_MIN_FONT_PT,
    HIGH_AGENCY_PREMIUM_FEATURES,
    ensure_peer_discussion_support,
    ensure_source_anchored_problem_activity,
    generated_discussion_questions,
    has_shipping_box_context,
    inferred_problem_focus_phrase,
    is_specific_discussion_question,
    is_low_value_vocabulary_term,
    math_profile_for_text,
    normalize_i_can_objective,
    PUBLISHER_STYLE_VERSION,
    problem_workbook_content,
    publisher_copyedit_issues,
    publisher_copyedit_text,
    load_activity_library,
    problem_activity_has_generic_markers,
    prompt_stack_layout_mode,
    premium_layout_uses_full_spread,
    refresh_problem_activity_for_context,
    render_practice_slide,
    rendered_problem_target_status,
    retarget_template_role_slide_to_source,
    select_session_premium_features,
    session_engagement_modes,
    session_engagement_slide_count,
    short_problem_solving_subtitle,
    source_problem_candidates,
    source_problem_response_prompt,
    validate_flagship_activities,
)


class NotebookEnginePublisherPolishTests(unittest.TestCase):
    def test_style_version_marks_next_generation_release(self) -> None:
        self.assertTrue(PUBLISHER_STYLE_VERSION.endswith("-vnext"))

    def test_vnext_premium_spreads_use_full_page_layouts(self) -> None:
        self.assertTrue(premium_layout_uses_full_spread({"premium_layout": "evidence_ladder"}))
        self.assertTrue(premium_layout_uses_full_spread({"premium_layout": "real_world_transfer"}))
        self.assertFalse(premium_layout_uses_full_spread({"premium_layout": "strategy_comparison"}))

    def test_render_practice_slide_uses_evidence_ladder_spread_copy(self) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        plan_slide = {
            "section": "Session 1",
            "title": "Independent Practice",
            "subtitle": "Use a claim, evidence, and reasoning ladder.",
            "kind": "practice",
            "premium_layout": "evidence_ladder",
            "primary_text": "Which strategy matches the model?",
            "secondary_text": "Use the source example to defend your thinking.",
            "response_prompt": "Write the explanation that proves your answer.",
            "tasks": ["Which strategy matches the model?", "What evidence proves it?"],
            "bullets": [],
            "discussion_questions": [],
        }

        render_practice_slide(
            slide,
            plan_slide=plan_slide,
            page=3,
            footer_text="Notebook Footer",
            image_lookup={},
        )

        all_text = " ".join(
            shape.text_frame.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Evidence-Based Explanation", all_text)
        self.assertIn("Use the Ladder", all_text)
        self.assertNotIn("Try a Similar One", all_text)

    def test_formal_review_min_font_floor_is_publisher_readable(self) -> None:
        self.assertGreaterEqual(FORMAL_REVIEW_MIN_FONT_PT, 10.4)

    def test_auto_fit_floor_keeps_body_copy_larger_than_micro_labels(self) -> None:
        self.assertGreaterEqual(AUTO_FIT_MIN_FONT_PT, 11.2)

    def test_fit_text_style_keeps_standard_body_copy_above_review_floor(self) -> None:
        size, _spacing = fit_text_style(
            "Model This: Some middle school students are organizing a canned food drive. "
            "Their goal is to collect 100 cans each week for the local food pantry.",
            6190488,
            420624,
            base_size=11.6,
            min_size=10.4,
        )

        self.assertGreaterEqual(size, AUTO_FIT_MIN_FONT_PT)

    def test_inferred_problem_focus_phrase_handles_statistics_lessons(self) -> None:
        slide = {
            "title": "Collaborative Practice",
            "subtitle": "Use ordered data and a dot plot to justify the median.",
            "primary_text": "Describe data using the median.",
            "tasks": ["Explain how the dot plot shows the middle value."],
            "bullets": [],
        }

        self.assertEqual(inferred_problem_focus_phrase(slide), "data analysis problem")

    def test_math_profile_for_text_identifies_data_analysis(self) -> None:
        profile = math_profile_for_text("Describe data using the median, mean, and a dot plot.")

        self.assertEqual(profile["topic"], "data_analysis")
        self.assertIn("median", profile["default_terms"])
        self.assertEqual(profile["answer_label"], "Median")

    def test_low_value_vocab_rejects_context_phrases(self) -> None:
        self.assertTrue(is_low_value_vocabulary_term("Teaching Experience"))
        self.assertTrue(is_low_value_vocabulary_term("Describe Data"))
        self.assertFalse(is_low_value_vocabulary_term("Median"))

    def test_problem_workbook_content_includes_explanation_workspace(self) -> None:
        content = problem_workbook_content(
            {
                "primary_text": "What is the median of the data set?",
                "tasks": ["Explain why the median is a measure of center."],
                "response_prompt": "Explain why your answer makes sense.",
            },
            variant="practice",
        )

        self.assertEqual(content["panel_title"], "Solve, Check, Explain")
        self.assertEqual(content["left_workspace"][1], "Show the setup, labels, and calculations neatly.")
        self.assertEqual(content["right_workspaces"][1][0], "Explain + Check")

    def test_publisher_copyedit_rewrites_internal_generator_jargon(self) -> None:
        edited = publisher_copyedit_text("Workbook baseline adapted from the uploaded model.")

        self.assertEqual(edited, "Student notebook adapted from the lesson structure.")

    def test_publisher_copyedit_issues_flag_jargon_and_generic_assessment_copy(self) -> None:
        issues = publisher_copyedit_issues("Use source-faithful practice to show what you know.")

        self.assertIn("internal generator jargon", issues)
        self.assertIn("generic assessment phrasing", issues)

    def test_short_problem_solving_subtitle_uses_clean_generic_fallback(self) -> None:
        subtitle = short_problem_solving_subtitle(
            {
                "kind": "practice",
                "practice_phase": "Solve Together",
            }
        )

        self.assertEqual(subtitle, "Solve the source problem together, then solve a similar one.")
        self.assertNotIn("lesson problem", subtitle)

    def test_normalize_i_can_objective_converts_source_objective_stems(self) -> None:
        self.assertEqual(normalize_i_can_objective("We will summarize the data set"), "I can summarize the data set.")
        self.assertEqual(normalize_i_can_objective("Students will explain the median"), "I can explain the median.")
        self.assertEqual(normalize_i_can_objective("Determine the area of the figure"), "I can determine the area of the figure.")

    def test_generated_discussion_questions_are_specific_to_data_analysis(self) -> None:
        questions = generated_discussion_questions(
            {
                "kind": "practice",
                "title": "Collaborative Practice",
                "subtitle": "Use ordered data and a dot plot to justify the median.",
                "primary_text": "Describe the data set using the median.",
                "tasks": ["Explain how the dot plot shows the middle value."],
                "source_problem_cards": ["What is the median of the data set shown on the dot plot?"],
            }
        )

        self.assertGreaterEqual(len(questions), 2)
        self.assertIn("ordered values", questions[0].lower())
        self.assertIn("dot plot", questions[1].lower())
        self.assertTrue(all(is_specific_discussion_question(question) for question in questions[:2]))

    def test_ensure_peer_discussion_support_replaces_generic_compare_prompt(self) -> None:
        slide = {
            "kind": "practice",
            "template_role": "collaborative_practice",
            "title": "Collaborative Practice",
            "subtitle": "Partners solve and compare.",
            "primary_text": "Find the volume of the rectangular prism.",
            "partner_a_problem": "A rectangular prism has length 4 units, width 3 units, and height 2 units. Find the volume.",
            "partner_b_problem": "A rectangular prism has length 5 units, width 3 units, and height 2 units. Find the volume.",
            "partner_prompt": "Partner A solves the first problem. Partner B solves the second problem. Then compare.",
            "discussion_questions": ["Q1: ___", "Q2: The ___ changes because ___.", "Q3: ___"],
        }

        ensure_peer_discussion_support(slide)

        self.assertGreaterEqual(len(slide["discussion_questions"]), 2)
        self.assertTrue(all(is_specific_discussion_question(question) for question in slide["discussion_questions"][:2]))
        self.assertIn("what stayed the same", slide["partner_prompt"].lower())

    def test_box_plot_is_not_classified_as_shipping_box_problem(self) -> None:
        slide = {
            "title": "Practice",
            "subtitle": "Use the box plot to compare the quartiles.",
            "primary_text": "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            "tasks": ["Describe what the box plot shows about the spread of the data."],
            "bullets": ["A box-and-whisker plot is a type of data display."],
        }

        self.assertEqual(inferred_problem_focus_phrase(slide), "data analysis problem")
        self.assertFalse(has_shipping_box_context("Use the box plot to compare quartiles."))

    def test_rendered_problem_target_status_accepts_any_preserved_target(self) -> None:
        slide = {
            "kind": "practice",
            "title": "Error Analysis Activity",
            "subtitle": "Solve the data-analysis problem, then solve a similar one.",
            "primary_text": "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            "tasks": [
                "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
                "How do the monthly temperatures between the two cities compare?",
            ],
            "source_problem_cards": [
                "How do the monthly temperatures between the two cities compare?",
                "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            ],
        }

        preserved, matched_target = rendered_problem_target_status(
            "Student Work How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            slide,
        )

        self.assertTrue(preserved)
        self.assertIn("slowest and fastest finish times", matched_target)

    def test_source_problem_response_prompt_uses_data_display_language(self) -> None:
        source_records = [
            {
                "title": "Running a Road Race",
                "text": "How can you use a data display to show and explain the distribution of the values in the data set?",
            }
        ]

        prompt = source_problem_response_prompt(source_records, "practice")

        self.assertIn("data display", prompt.lower())
        self.assertNotIn("formula", prompt.lower())

    def test_source_problem_candidates_prefer_shorter_data_analysis_questions(self) -> None:
        slides = [
            {
                "title": "Between which two values does the upper quarter of the data fall?",
                "text": "The lower quarter of the data is between 29 minutes and 34.5 minutes. One quarter of the data lies between the third quartile and the upper extreme. Between which two values does the upper quarter of the data fall?",
                "text_items": [],
                "problem_texts": [],
            },
            {
                "title": "Let’s Explore More",
                "text": "What does the length of the box and whiskers tell you about the spread of the data in the box plot?\u200b?",
                "text_items": [],
                "problem_texts": ["What does the length of the box and whiskers tell you about the spread of the data in the box plot?\u200b?"],
            },
        ]

        candidates = source_problem_candidates(slides, limit=3)

        self.assertTrue(candidates)
        self.assertTrue(candidates[0].startswith("Between which two values"))
        self.assertNotIn("\u200b", " ".join(candidates))

    def test_prompt_stack_layout_mode_switches_to_focus_before_three_long_prompts_shrink_too_far(self) -> None:
        problems = [
            "P1: Which interval of the box plot contains the middle half of the data, and what does that tell you about the spread?",
            "P2: How do the median and upper quartile compare between the sixth and seventh graders in the display?",
            "P3: Which values show the slowest and fastest finish times, and how do you know from the whiskers?",
        ]

        self.assertEqual(prompt_stack_layout_mode(problems), "focus")

    def test_refresh_problem_activity_for_context_rewrites_formula_copy_on_data_analysis_slide(self) -> None:
        slide = {
            "kind": "practice",
            "title": "Practice Extension",
            "subtitle": "Solve the data-analysis problem, then solve a similar one.",
            "primary_text": "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            "tasks": ["How do the slowest and fastest finish times compare between the sixth and seventh graders?"],
            "activity_name": "table-to-equation matching",
            "activity_family": "match_pair",
            "activity_instructions": "Match each table, figure, or value set to the equation or statement that represents it.",
            "movable_pieces": ["Known information", "Volume formula", "Missing value", "Check the answer"],
            "answer_check": "Check that the relationship stays true across all matched parts.",
        }

        refresh_problem_activity_for_context(slide)

        self.assertTrue(slide["activity_name"].startswith("data evidence matching"))
        self.assertNotIn("formula", " ".join(slide["movable_pieces"]).lower())
        self.assertIn("evidence", slide["answer_check"].lower())

    def test_problem_activity_has_generic_markers_detects_template_filler(self) -> None:
        slide = {
            "activity_name": "strategy ranking activity",
            "activity_instructions": "Choose the path that feels like the best fit, then defend your choice with evidence.",
            "movable_pieces": ["Path A", "Path B", "Path C", "Best fit"],
            "answer_check": "Choose a path you can solve and explain clearly.",
        }

        self.assertTrue(problem_activity_has_generic_markers(slide))

    def test_ensure_source_aligned_be_curious_activity_uses_source_terms(self) -> None:
        slide = {
            "kind": "be_curious",
            "source_slide_numbers": [1],
            "vocabulary": [
                {"word": "Median"},
                {"word": "Dot Plot"},
                {"word": "Ordered Data"},
            ],
        }
        deck = {
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Be Curious",
                    "text": "What does the dot plot show about the median of the data set?",
                    "text_items": [],
                    "problem_texts": [],
                    "notes": "",
                }
            ]
        }

        ensure_source_aligned_be_curious_activity(slide, deck=deck)

        self.assertTrue(slide["activity_name"].startswith("source clue reveal"))
        self.assertEqual(slide["activity_family"], "reveal_discuss")
        self.assertIn("notice", slide["activity_instructions"].lower())
        self.assertTrue(slide["movable_pieces"])

    def test_retarget_template_role_slide_to_source_rewrites_choice_board(self) -> None:
        slide = {
            "kind": "practice",
            "template_role": "choice_board",
            "title": "Choice Board",
            "subtitle": "Choose one path and solve it.",
            "primary_text": "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            "response_prompt": "I chose Path ___ because ___.",
            "source_slide_numbers": [40, 41, 42],
            "source_problem_cards": [
                "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
                "What does the length of the box and whiskers tell you about the spread of the data in the box plot?",
                "Between which two values does the upper quarter of the data fall?",
            ],
            "activity_name": "strategy ranking activity",
            "activity_family": "compare_rank",
            "activity_instructions": "Choose the path that feels like the best fit, then defend your choice with evidence.",
            "movable_pieces": ["Path A", "Path B", "Path C", "Best fit"],
            "answer_check": "Choose a path you can solve and explain clearly.",
        }

        retarget_template_role_slide_to_source(slide)

        self.assertEqual(slide["title"], "Choose a Source Problem")
        self.assertFalse(slide["activity_name"])
        self.assertTrue(slide["choice_paths"])
        self.assertTrue(all(not item["label"].startswith("Path ") for item in slide["choice_paths"]))
        self.assertIn("source-aligned problem", slide["subtitle"].lower())

    def test_ensure_source_anchored_problem_activity_replaces_generic_problem_activity(self) -> None:
        slide = {
            "kind": "practice",
            "title": "Independent Practice",
            "subtitle": "Use the box plot to compare the quartiles.",
            "primary_text": "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
            "tasks": ["Describe what the box plot shows about the spread of the data."],
            "source_slide_numbers": [40, 41],
            "source_problem_cards": [
                "How do the slowest and fastest finish times compare between the sixth and seventh graders?",
                "What does the length of the box and whiskers tell you about the spread of the data in the box plot?",
            ],
            "activity_name": "strategy ranking activity",
            "activity_family": "compare_rank",
            "activity_instructions": "Choose the path that feels like the best fit, then defend your choice with evidence.",
            "movable_pieces": ["Path A", "Path B", "Path C", "Best fit"],
            "answer_check": "Choose a path you can solve and explain clearly.",
        }

        ensure_source_anchored_problem_activity(slide)

        self.assertIn(slide["activity_family"], {"build_construct", "sort_classify", "match_pair", "sequence_order", "detect_justify"})
        self.assertNotIn("path a", " ".join(slide["movable_pieces"]).lower())
        self.assertTrue(
            "analysis" in slide["activity_instructions"].lower()
            or "evidence" in slide["activity_instructions"].lower()
        )

    def test_session_engagement_slide_count_counts_activity_roles_and_premium_layouts(self) -> None:
        session = {
            "slides": [
                {"kind": "be_curious"},
                {"kind": "practice", "activity_name": "sort the cards", "activity_family": "sort_classify"},
                {"kind": "practice", "premium_layout": "strategy_comparison"},
                {"kind": "challenge", "partner_prompt": "Turn and teach your partner how you solved it."},
                {"kind": "reflection"},
            ]
        }

        self.assertEqual(session_engagement_slide_count(session), 4)
        self.assertEqual(session_engagement_modes(session), ["compare", "notice", "partner", "sort"])

    def test_session_engagement_modes_include_vnext_premium_layouts(self) -> None:
        session = {
            "slides": [
                {"kind": "practice", "premium_layout": "evidence_ladder"},
                {"kind": "challenge", "premium_layout": "real_world_transfer"},
            ]
        }

        self.assertEqual(session_engagement_slide_count(session), 2)
        self.assertEqual(session_engagement_modes(session), ["apply", "create", "discuss", "justify"])

    def test_select_session_premium_features_guarantees_multiple_interactive_features(self) -> None:
        deck = {
            "lesson_title": "Compare Tables and Graphs",
            "summary": "Students compare representations, explain mistakes, and justify which strategy fits best.",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Compare the table and graph",
                    "text": "Explain how the table and graph show the same relationship.",
                    "text_items": ["Compare the table and graph.", "Explain how you know."],
                    "problem_texts": ["How does the graph match the table values?"],
                    "notes": "",
                },
                {
                    "slide_number": 2,
                    "title": "A student made a mistake",
                    "text": "Find the error in the graph and justify the correction.",
                    "text_items": ["Find the mistake.", "Explain the correction."],
                    "problem_texts": ["What is the error in the graph?"],
                    "notes": "",
                },
                {
                    "slide_number": 3,
                    "title": "Session 2",
                    "text": "",
                    "text_items": [],
                    "problem_texts": [],
                    "notes": "",
                },
                {
                    "slide_number": 4,
                    "title": "Which strategy fits best?",
                    "text": "Compare models, choose a strategy, and explain why it fits.",
                    "text_items": ["Which strategy fits best?", "Compare the model and graph."],
                    "problem_texts": ["Which strategy should you use to explain the graph?"],
                    "notes": "",
                },
                {
                    "slide_number": 5,
                    "title": "Create your own representation",
                    "text": "Create your own table or graph that matches the relationship.",
                    "text_items": ["Create your own table.", "Create your own graph."],
                    "problem_texts": ["Create your own representation and explain why it works."],
                    "notes": "",
                },
            ],
        }
        session = {
            "slides": [
                {"kind": "cover", "title": "Session 2"},
                {"kind": "worked_example", "title": "Guided Practice", "primary_text": "Compare the table and graph."},
                {"kind": "practice", "title": "Collaborative Practice", "primary_text": "Find the error and explain it."},
                {"kind": "challenge", "title": "Independent Practice", "primary_text": "Choose the best strategy and justify it."},
            ]
        }

        features, _context_anchor = select_session_premium_features(deck, "session_2", session)
        interactive_features = [feature for feature in features if feature in HIGH_AGENCY_PREMIUM_FEATURES]

        self.assertGreaterEqual(len(interactive_features), 2)
        self.assertIn(features[0], HIGH_AGENCY_PREMIUM_FEATURES)

    def test_select_session_premium_features_can_choose_vnext_argument_and_transfer_moves(self) -> None:
        deck = {
            "lesson_title": "Use Evidence to Defend a Survey Claim",
            "summary": "Students make a claim, cite evidence, explain their reasoning, and transfer the lesson to a school survey context.",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "School survey claim",
                    "text": "Make a claim about the school survey and defend it with evidence and reasoning.",
                    "text_items": ["Make a claim.", "Use evidence and reasoning."],
                    "problem_texts": ["Which claim can you defend with the survey evidence?"],
                    "notes": "",
                },
                {
                    "slide_number": 2,
                    "title": "Evidence from the data",
                    "text": "Use the data evidence to justify your claim about the school survey.",
                    "text_items": ["Use evidence from the survey."],
                    "problem_texts": ["What evidence best supports the claim?"],
                    "notes": "",
                },
                {
                    "slide_number": 3,
                    "title": "Transfer the idea",
                    "text": "Apply the same reasoning to a new real-world survey situation.",
                    "text_items": ["Transfer the reasoning."],
                    "problem_texts": ["How would you apply the same reasoning to a new survey context?"],
                    "notes": "",
                },
            ],
        }
        session = {
            "slides": [
                {"kind": "cover", "title": "Session 2"},
                {"kind": "worked_example", "title": "Guided Practice", "primary_text": "Make a claim and support it with evidence."},
                {"kind": "practice", "title": "Collaborative Practice", "primary_text": "Explain which evidence best supports the claim."},
                {"kind": "challenge", "title": "Independent Practice", "primary_text": "Transfer the reasoning to a new survey situation."},
            ]
        }

        features, _context_anchor = select_session_premium_features(deck, "session_2", session)

        self.assertIn("evidence_ladder", features)
        self.assertIn("real_world_transfer", features)

    def test_activity_library_exposes_broad_bank_with_inferred_keywords(self) -> None:
        library = load_activity_library()

        self.assertGreaterEqual(len(library), 150)
        coordinate_plotting = next(item for item in library if item["name"] == "coordinate point plotting")
        self.assertIn("coordinate", coordinate_plotting["keywords"])
        self.assertIn("graph", coordinate_plotting["keywords"])

    def test_choose_activity_candidates_surfaces_topic_fit_from_large_bank(self) -> None:
        deck = {
            "lesson_title": "Compare tables and graphs",
            "summary": "Students analyze tables, graphs, coordinate pairs, and explain patterns.",
            "slides": [
                {
                    "title": "Compare the graph and table",
                    "text": "Use the graph, table, and coordinates to explain the relationship.",
                    "notes": "",
                    "text_items": [],
                }
            ],
        }

        names = [item["name"] for item in choose_activity_candidates(deck, load_activity_library(), limit=12)]

        self.assertTrue(any("graph" in name or "coordinate" in name or "table" in name for name in names))
        self.assertTrue(any(name == "coordinate point plotting" for name in names))

    def test_choose_flagship_activity_mix_prioritizes_middle_and_end_lesson_moves(self) -> None:
        deck = {
            "lesson_title": "Connect tables, graphs, and equations",
            "summary": "Students compare representations, justify mistakes, and explain which strategy fits.",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Session 1",
                    "text": "Notice how the table and graph connect.",
                    "text_items": ["Notice the table.", "Notice the graph."],
                    "problem_texts": [],
                    "notes": "",
                },
                {
                    "slide_number": 2,
                    "title": "Representation Match",
                    "text": "Match the table, graph, and equation that show the same relationship.",
                    "text_items": ["table", "graph", "equation"],
                    "problem_texts": ["How does the graph match the table values and equation?"],
                    "notes": "",
                },
                {
                    "slide_number": 3,
                    "title": "Fix the graph mistake",
                    "text": "A student matched the wrong point to the equation. Explain the error.",
                    "text_items": ["Fix the mistake.", "Use evidence."],
                    "problem_texts": ["What is the error in the graph, and how do you know?"],
                    "notes": "",
                },
                {
                    "slide_number": 4,
                    "title": "Compare strategies",
                    "text": "Compare which strategy best explains the relationship.",
                    "text_items": ["Compare the strategies.", "Justify the best fit."],
                    "problem_texts": ["Which strategy best explains the graph and why?"],
                    "notes": "",
                },
            ],
        }
        session = {
            "slides": [
                {"kind": "cover", "title": "Session 1"},
                {"kind": "be_curious", "title": "Notice + Wonder", "source_slide_numbers": [1]},
                {"kind": "guided_notes", "title": "Guided Notes", "source_slide_numbers": [2]},
                {"kind": "practice", "title": "Practice", "source_slide_numbers": [2, 3]},
                {"kind": "challenge", "title": "Challenge", "source_slide_numbers": [3, 4]},
                {"kind": "exit_ticket", "title": "Exit Ticket", "source_slide_numbers": [4]},
            ]
        }

        session_data = build_flagship_session_data(deck, "session_1", session)
        mix = choose_flagship_activity_mix(session_data)
        specs = build_flagship_activity_specs(session_data, mix)

        self.assertGreaterEqual(session_data["target_count"], 2)
        self.assertGreaterEqual(len(specs), 2)
        self.assertIn("middle", {spec["placement"] for spec in specs})
        self.assertIn("end", {spec["placement"] for spec in specs})
        self.assertTrue(all(spec["sourceAnchors"] for spec in specs))
        self.assertTrue(all(not flagship_activity_has_generic_markers(spec) for spec in specs))

    def test_apply_flagship_activity_layer_attaches_source_anchored_specs_to_slides(self) -> None:
        deck = {
            "lesson_title": "Volume of rectangular prisms",
            "summary": "Students use dimensions, formulas, and units to solve volume problems and explain their reasoning.",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Launch",
                    "text": "What do you notice about the prism model?",
                    "text_items": ["Notice the prism.", "Notice the dimensions."],
                    "problem_texts": [],
                    "notes": "",
                },
                {
                    "slide_number": 2,
                    "title": "Worked example",
                    "text": "Use V = l x w x h to find the volume of the rectangular prism.",
                    "text_items": ["length", "width", "height", "cubic units"],
                    "problem_texts": ["What is the volume of the rectangular prism shown?"],
                    "notes": "",
                },
                {
                    "slide_number": 3,
                    "title": "Shipping rule",
                    "text": "Use the known volume and dimensions to find the missing length and compare it to the 12-inch rule.",
                    "text_items": ["missing length", "shipping rule", "formula"],
                    "problem_texts": ["What is the missing length of the box, and does it meet the shipping rule?"],
                    "notes": "",
                },
                {
                    "slide_number": 4,
                    "title": "Challenge",
                    "text": "Compare which solve path is strongest and explain why.",
                    "text_items": ["compare", "justify", "evidence"],
                    "problem_texts": ["Which solve path best proves the box meets the condition?"],
                    "notes": "",
                },
            ],
        }
        session = {
            "slides": [
                {"kind": "cover", "title": "Session 1"},
                {"kind": "be_curious", "title": "Notice + Wonder", "source_slide_numbers": [1]},
                {"kind": "worked_example", "title": "Worked Example", "source_slide_numbers": [2]},
                {"kind": "practice", "title": "Practice", "source_slide_numbers": [2, 3]},
                {"kind": "challenge", "title": "Challenge", "source_slide_numbers": [3, 4]},
                {"kind": "exit_ticket", "title": "Exit Ticket", "source_slide_numbers": [4]},
            ]
        }

        apply_flagship_activity_layer(deck, "session_1", session)
        flagship_slides = [slide for slide in session["slides"] if slide.get("flagship_activity")]

        self.assertGreaterEqual(len(session["flagship_activities"]), 2)
        self.assertEqual(len(flagship_slides), len(session["flagship_activities"]))
        self.assertTrue(all(slide["activity_name"] == slide["flagship_activity"]["title"] for slide in flagship_slides))
        self.assertEqual(validate_flagship_activities(session, deck=deck), [])

    def test_apply_flagship_activity_layer_falls_back_without_breaking_thin_sessions(self) -> None:
        deck = {
            "lesson_title": "Lesson",
            "summary": "Students reflect on the lesson.",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Session 1",
                    "text": "Lesson overview.",
                    "text_items": ["Lesson overview."],
                    "problem_texts": [],
                    "notes": "",
                }
            ],
        }
        session = {
            "slides": [
                {"kind": "cover", "title": "Session 1", "source_slide_numbers": [1]},
                {"kind": "reflection", "title": "Reflection", "source_slide_numbers": [1]},
            ]
        }

        apply_flagship_activity_layer(deck, "session_1", session)

        self.assertEqual(session["flagship_activities"], [])
        self.assertEqual(session["flagship_activity_target"], 0)
        self.assertTrue(session["flagship_activity_fallbacks"])

    def test_flagship_activity_has_generic_markers_rejects_placeholder_copy(self) -> None:
        spec = {
            "title": "Sort / Match / Classify",
            "purpose": "Move a piece and explain.",
            "directions": "Sort the cards into Category A and Category B.",
            "prompts": ["Card 1", "Card 2", "Category A", "Category B"],
            "supports": ["Sentence starter: ___"],
        }

        self.assertTrue(flagship_activity_has_generic_markers(spec))


if __name__ == "__main__":
    unittest.main()
