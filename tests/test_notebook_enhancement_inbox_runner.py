import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
ENHANCEMENT_SRC = ROOT / "notebook-engine" / "enhancement" / "src"
if str(ENHANCEMENT_SRC) not in sys.path:
    sys.path.insert(0, str(ENHANCEMENT_SRC))

from run_enhancement_inbox import (  # type: ignore[import-not-found]
    discover_bundle_dirs,
    discover_pptx_files,
    process_enhancement_inbox,
)


def sample_plan() -> dict:
    return {
        "session_1": {
            "slides": [
                {
                    "kind": "practice",
                    "title": "Practice",
                    "subtitle": "",
                    "primary_text": "Students should use the lesson idea.",
                    "secondary_text": "Use the space below to do the work and complete the activity.",
                    "bullets": [],
                    "tasks": [],
                    "response_prompt": "Do the work.",
                    "sentence_starters": [],
                    "vocabulary": [],
                    "activity_name": "",
                    "activity_instructions": "Complete the activity.",
                    "source_slide_numbers": [1],
                }
            ]
        },
        "session_2": {"slides": []},
    }


def sample_deck() -> dict:
    return {
        "slides": [
            {
                "slide_number": 1,
                "title": "Compare the area model and equation",
                "text": "Use the area model to explain which strategy matches the equation.",
                "text_items": ["Compare the area model."],
                "problem_texts": ["Which strategy matches the area model?"],
                "notes": "",
            }
        ]
    }


def write_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(0, 0, prs.slide_width // 2, prs.slide_height // 10)
    title_box.text_frame.paragraphs[0].add_run().text = "Vocabulary"
    term_box = slide.shapes.add_textbox(0, prs.slide_height // 5, prs.slide_width // 4, prs.slide_height // 12)
    term_box.text_frame.paragraphs[0].add_run().text = "equation"
    prompt_box = slide.shapes.add_textbox(0, prs.slide_height // 3, prs.slide_width // 2, prs.slide_height // 10)
    prompt_box.text_frame.paragraphs[0].add_run().text = "Which strategy matches the equation, and how do you know?"
    prs.save(str(path))


class NotebookEnhancementInboxRunnerTests(unittest.TestCase):
    def test_discover_bundle_dirs_finds_job_folders(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            job_dir = inbox_dir / "job-a"
            job_dir.mkdir(parents=True)
            (job_dir / "notebook_plan.json").write_text("{}", encoding="utf-8")

            bundles = discover_bundle_dirs(inbox_dir)

            self.assertEqual(bundles, [job_dir.resolve()])

    def test_process_enhancement_inbox_polishes_bundles_without_render(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            output_dir = Path(tmp_dir) / "OUTPUT"
            job_dir = inbox_dir / "job-a"
            job_dir.mkdir(parents=True)
            (job_dir / "notebook_plan.json").write_text(__import__("json").dumps(sample_plan()), encoding="utf-8")
            (job_dir / "source_deck.json").write_text(__import__("json").dumps(sample_deck()), encoding="utf-8")

            summary = process_enhancement_inbox(inbox_dir, output_dir, render_outputs=False)

            self.assertEqual(summary["status"], "completed_with_review_needed")
            self.assertEqual(summary["processed_count"], 1)
            self.assertEqual(summary["results"][0]["status"], "polished")
            self.assertTrue(summary["results"][0]["manual_review_required"])
            self.assertTrue((output_dir / "job-a" / "notebook_plan.json").exists())
            self.assertTrue((output_dir / "enhancement_inbox_last_run.json").exists())

    def test_discover_pptx_files_finds_raw_notebook_inputs(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            inbox_dir.mkdir()
            pptx_path = inbox_dir / "sample-notebook.pptx"
            pptx_path.write_bytes(b"placeholder")

            pptx_files = discover_pptx_files(inbox_dir)

            self.assertEqual(pptx_files, [pptx_path.resolve()])

    def test_process_enhancement_inbox_continues_after_failed_bundle(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            output_dir = Path(tmp_dir) / "OUTPUT"
            bad_job = inbox_dir / "job-bad"
            bad_job.mkdir(parents=True)
            (bad_job / "notebook_plan.json").write_text("{not valid json", encoding="utf-8")
            pptx_path = inbox_dir / "sample-notebook.pptx"
            write_sample_pptx(pptx_path)

            summary = process_enhancement_inbox(inbox_dir, output_dir, render_outputs=False)

            self.assertEqual(summary["status"], "completed_with_errors")
            statuses = {result["status"] for result in summary["results"]}
            self.assertIn("failed", statuses)
            self.assertIn("polished_pptx", statuses)
