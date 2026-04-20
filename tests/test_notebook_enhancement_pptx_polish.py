import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ENHANCEMENT_SRC = ROOT / "notebook-engine" / "enhancement" / "src"
if str(ENHANCEMENT_SRC) not in sys.path:
    sys.path.insert(0, str(ENHANCEMENT_SRC))

from polish_notebook_pptx import polish_notebook_pptx  # type: ignore[import-not-found]
from run_enhancement_inbox import process_enhancement_inbox  # type: ignore[import-not-found]


def build_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "Independent Practice"
    title_run.font.size = Pt(14)
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.82), Inches(3.5), Inches(0.22))
    subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
    subtitle_run.text = "Unit 2 · Lesson 2.1"
    subtitle_run.font.size = Pt(9.5)
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.0), Inches(1.2))
    body_box.text_frame.word_wrap = False
    body_box.text_frame.margin_left = 0
    body_box.text_frame.margin_right = 0
    body_run = body_box.text_frame.paragraphs[0].add_run()
    body_run.text = "Which strategy matches the equation, and how do you know this notebook draft needs stronger readability?"
    body_run.font.size = Pt(9)
    claim_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(8.0), Inches(0.42))
    claim_run = claim_box.text_frame.paragraphs[0].add_run()
    claim_run.text = 'A student says, "Both representations work because they show the same relationship."'
    claim_run.font.size = Pt(9.5)
    label_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(1.8), Inches(0.28))
    label_run = label_box.text_frame.paragraphs[0].add_run()
    label_run.text = "Notebook move"
    label_run.font.size = Pt(10)
    for top in [3.55, 3.95]:
        line_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(4.8), Inches(0.01))
        line_box.text_frame.clear()
    vocab_slide = prs.slides.add_slide(prs.slide_layouts[6])
    vocab_title = vocab_slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    vocab_title_run = vocab_title.text_frame.paragraphs[0].add_run()
    vocab_title_run.text = "Vocabulary"
    vocab_title_run.font.size = Pt(14)
    vocab_term_a = vocab_slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(2.4), Inches(0.32))
    vocab_term_a.text_frame.paragraphs[0].add_run().text = "equation"
    vocab_term_a.text_frame.paragraphs[0].runs[0].font.size = Pt(10.5)
    vocab_term_b = vocab_slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(2.4), Inches(0.32))
    vocab_term_b.text_frame.paragraphs[0].add_run().text = "strategy"
    vocab_term_b.text_frame.paragraphs[0].runs[0].font.size = Pt(10.5)
    prs.save(str(path))


def build_thin_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "Independent Practice"
    title_run.font.size = Pt(14)
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.2), Inches(1.0))
    body_run = body_box.text_frame.paragraphs[0].add_run()
    body_run.text = "This notebook draft needs stronger readability."
    body_run.font.size = Pt(9)
    prs.save(str(path))


def build_statistics_sample_pptx(path: Path) -> None:
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    title = cover.shapes.add_textbox(Inches(0.25), Inches(0.02), Inches(7.8), Inches(0.3))
    title.text_frame.paragraphs[0].add_run().text = "Statistical Questions - Session 1"
    subtitle = cover.shapes.add_textbox(Inches(0.25), Inches(0.3), Inches(7.8), Inches(0.22))
    subtitle.text_frame.paragraphs[0].add_run().text = "Grade 6 | Unit 2 Lesson 1"
    focus_box = cover.shapes.add_textbox(Inches(1.05), Inches(1.98), Inches(3.1), Inches(1.6))
    focus_run = focus_box.text_frame.paragraphs[0].add_run()
    focus_run.text = "Statistical Questions & Variability"
    focus_run.font.size = Pt(20)
    notebook_label = cover.shapes.add_textbox(Inches(4.6), Inches(1.28), Inches(8.3), Inches(0.55))
    notebook_label.text_frame.paragraphs[0].add_run().text = "Interactive Student Notebook"

    vocab = prs.slides.add_slide(prs.slide_layouts[6])
    vocab_title = vocab.shapes.add_textbox(Inches(0.25), Inches(0.02), Inches(7.8), Inches(0.3))
    vocab_title.text_frame.paragraphs[0].add_run().text = "Vocabulary"
    vocab_header = vocab.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.45))
    vocab_header.text_frame.paragraphs[0].add_run().text = "Key Math Language"
    headers = [
        (Inches(0.85), "Word"),
        (Inches(2.6), "Student-Friendly Meaning"),
        (Inches(6.1), "Example"),
        (Inches(10.4), "Visual/Icon"),
    ]
    for left, text in headers:
        box = vocab.shapes.add_textbox(left, Inches(1.1), Inches(1.95), Inches(0.45))
        box.text_frame.paragraphs[0].add_run().text = text
    rows = [
        ("Statistical Question", "A question with different answers.", '"What is your favorite park?"', "chart"),
        ("Variability", "Answers are not all the same.", "Students may name many parks.", "updown"),
        ("Data", "Information we collect.", "Survey responses from classmates.", "list"),
    ]
    top = Inches(1.6)
    for word, meaning, example, icon in rows:
        for left, width, text in [
            (Inches(0.9), Inches(1.6), word),
            (Inches(2.65), Inches(3.35), meaning),
            (Inches(6.15), Inches(4.15), example),
            (Inches(10.45), Inches(1.85), icon),
        ]:
            box = vocab.shapes.add_textbox(left, top, width, Inches(1.18))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(11)
        top += Inches(1.28)

    compare = prs.slides.add_slide(prs.slide_layouts[6])
    compare_title = compare.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    compare_title.text_frame.paragraphs[0].add_run().text = "Sort It"
    compare_prompt = compare.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.8), Inches(0.8))
    compare_prompt.text_frame.paragraphs[0].add_run().text = 'Which prompt is statistical: "What is your favorite park?" or "In what state is Acadia National Park found?"'

    error_slide = prs.slides.add_slide(prs.slide_layouts[6])
    error_title = error_slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    error_title.text_frame.paragraphs[0].add_run().text = "Error Analysis"
    claim = error_slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.8), Inches(0.9))
    claim.text_frame.paragraphs[0].add_run().text = 'A student says, "In what state is Acadia National Park found?" is statistical because people like different parks.'

    prs.save(str(path))


def slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            texts.append(shape.text_frame.text)
    return texts


def find_shape_with_text(slide, target: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text == target:
            return shape
    raise AssertionError(f"Could not find shape with text: {target}")


class NotebookEnhancementPptxPolishTests(unittest.TestCase):
    def test_polish_notebook_pptx_lifts_tiny_fonts_and_wrap(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            source_path = Path(tmp_dir) / "draft.pptx"
            output_path = Path(tmp_dir) / "draft-polished.pptx"
            build_sample_pptx(source_path)

            report = polish_notebook_pptx(source_path, output_path)

            self.assertTrue(output_path.exists())
            original = Presentation(str(source_path))
            polished = Presentation(str(output_path))
            title_shape = find_shape_with_text(polished.slides[0], "Independent Practice")
            subtitle_shape = find_shape_with_text(polished.slides[0], "Unit 2 · Lesson 2.1")
            body_shape = find_shape_with_text(polished.slides[0], "Which strategy matches the equation, and how do you know this notebook draft needs stronger readability?")
            self.assertGreaterEqual(title_shape.text_frame.paragraphs[0].runs[0].font.size.pt, 16.0)
            self.assertLessEqual(subtitle_shape.text_frame.paragraphs[0].runs[0].font.size.pt, 11.0)
            self.assertGreaterEqual(body_shape.text_frame.paragraphs[0].runs[0].font.size.pt, 12.0)
            self.assertTrue(body_shape.text_frame.word_wrap)
            self.assertGreater(report["stats"]["fontLifts"], 0)
            self.assertGreater(report["stats"]["headingStyles"], 0)
            self.assertGreater(report["stats"]["responseZones"], 0)
            self.assertGreater(report["stats"]["titleBackplates"], 0)
            self.assertGreater(report["stats"]["promptCards"], 0)
            self.assertGreaterEqual(report["stats"]["addedSlides"], 4)
            self.assertEqual(report["qualityTier"], "enhanced")
            self.assertEqual(report["auditPasses"], 1)
            self.assertFalse(report["outputAudit"]["passed"])
            self.assertGreater(len(polished.slides[0].shapes), len(original.slides[0].shapes))
            self.assertEqual(len(polished.slides), len(original.slides) + report["stats"]["addedSlides"])
            all_slide_texts = [" ".join(slide_texts(slide)) for slide in polished.slides]
            self.assertTrue(any("Compare and Justify" in text for text in all_slide_texts))
            self.assertTrue(any("Fix the Mistake" in text for text in all_slide_texts))
            self.assertTrue(any("Vocabulary in Action" in text for text in all_slide_texts))
            self.assertTrue(any("Representation Connection" in text for text in all_slide_texts))

    def test_process_enhancement_inbox_polishes_raw_pptx_files(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            output_dir = Path(tmp_dir) / "OUTPUT"
            inbox_dir.mkdir()
            source_path = inbox_dir / "draft-notebook.pptx"
            build_sample_pptx(source_path)

            summary = process_enhancement_inbox(inbox_dir, output_dir, render_outputs=False)

            job_dir = output_dir / "draft-notebook"
            self.assertEqual(summary["status"], "completed_with_review_needed")
            self.assertEqual(summary["results"][0]["status"], "polished_pptx")
            self.assertFalse(summary["results"][0]["premium_passed"])
            self.assertTrue(summary["results"][0]["manual_review_required"])
            self.assertEqual(summary["results"][0]["quality_tier"], "enhanced")
            self.assertTrue((job_dir / "draft-notebook - Polished.pptx").exists())
            self.assertTrue((job_dir / "pptx_polish_report.json").exists())

    def test_polish_notebook_pptx_keeps_thin_source_out_of_premium_extension_mode(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            source_path = Path(tmp_dir) / "thin-draft.pptx"
            output_path = Path(tmp_dir) / "thin-draft-polished.pptx"
            build_thin_sample_pptx(source_path)

            report = polish_notebook_pptx(source_path, output_path)

            self.assertTrue(output_path.exists())
            self.assertEqual(report["qualityTier"], "enhanced")
            self.assertEqual(report["stats"]["addedSlides"], 0)
            self.assertEqual(report["extensionModes"], [])
            self.assertTrue(report["warnings"])

    def test_statistics_context_uses_real_vocab_and_sorting_extensions(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            source_path = Path(tmp_dir) / "stats-draft.pptx"
            output_path = Path(tmp_dir) / "stats-draft-polished.pptx"
            build_statistics_sample_pptx(source_path)

            report = polish_notebook_pptx(source_path, output_path)
            polished = Presentation(str(output_path))
            all_slide_texts = [" ".join(slide_texts(slide)) for slide in polished.slides]

            self.assertEqual(report["context"]["lessonFocus"], "Statistical Questions & Variability")
            self.assertEqual(report["context"]["vocabulary"][:3], ["Statistical Question", "Variability", "Data"])
            self.assertNotIn("Student-Friendly Meaning", report["context"]["vocabulary"])
            self.assertIn("example_nonexample", report["extensionModes"])
            self.assertNotIn("representation_connection", report["extensionModes"])
            self.assertTrue(report["outputAudit"]["passed"])
            self.assertEqual(report["outputAudit"]["categories"]["layoutPressure"]["status"], "pass")
            self.assertLessEqual(report["outputAudit"]["deckMetrics"]["maxAddedSlideShapeCount"], 40)
            self.assertTrue(any("Example vs. Non-Example" in text for text in all_slide_texts))

    def test_process_enhancement_inbox_marks_review_needed_for_thin_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            inbox_dir = Path(tmp_dir) / "INBOX"
            output_dir = Path(tmp_dir) / "OUTPUT"
            inbox_dir.mkdir()
            source_path = inbox_dir / "thin-draft.pptx"
            build_thin_sample_pptx(source_path)

            summary = process_enhancement_inbox(inbox_dir, output_dir, render_outputs=False)

            self.assertEqual(summary["status"], "completed_with_review_needed")
            self.assertFalse(summary["results"][0]["premium_passed"])
            self.assertTrue(summary["results"][0]["manual_review_required"])
