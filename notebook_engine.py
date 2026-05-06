#!/usr/bin/env python3
"""Reusable PPTX-first engine for compact Session 1 student notebook generation."""

from __future__ import annotations

import argparse
import base64
import json
import mimetypes
import os
import re
import ssl
import sys
import zipfile
from collections import Counter
from html import escape as html_escape
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable
from urllib import error, request

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent
OPENAI_API_BASE = os.getenv("OPENAI_API_BASE", "https://api.openai.com/v1")
DEFAULT_MODEL = os.getenv("OPENAI_MODEL", "gpt-5.4")
OPENAI_TIMEOUT_SECONDS = int(os.getenv("OPENAI_TIMEOUT_SECONDS", "240"))
DEBUG_TRACEBACK_ENV = "NOTEBOOK_ENGINE_DEBUG_TRACEBACK"
ACTIVITY_LIBRARY_PATH = ROOT / "activity_library.txt"
ACTIVITY_DATABASE_PATH = ROOT / "activity_database.json"
SLIDE_W = Emu(12192000)
SLIDE_H = Inches(7.5)
SESSION_KEY_ORDER = ("session_1", "session_2")
SESSION_OUTPUT_KEYS = {"session_1": "session1", "session_2": "session2"}
PUBLISHER_STYLE_VERSION = "reference-classroom-2026-04-vnext"
FONT_HEAD = "Inter"
FONT_BODY = "Inter"
MIN_SESSION_SLIDES = 10
MAX_SESSION_SLIDES = 18
PREMIUM_TARGET_SESSION_SLIDES = 14
DEFAULT_TEXT_LINE_SPACING = 1.18
TABLE_TEXT_LINE_SPACING = 1.15
MIN_ACTIVITY_SLIDES_PER_SESSION = 7
MAX_ACTIVITY_SLIDES_PER_SESSION = 11
COMMON_CA_BUNDLE_PATHS = [
    "/etc/ssl/cert.pem",
    "/private/etc/ssl/cert.pem",
    "/opt/homebrew/etc/openssl@3/cert.pem",
    "/usr/local/etc/openssl@3/cert.pem",
]

BG = RGBColor(255, 255, 255)
PAPER = RGBColor(255, 255, 255)
NAVY = RGBColor(23, 50, 77)
TEAL = RGBColor(31, 166, 162)
GOLD = RGBColor(242, 193, 91)
CORAL = RGBColor(231, 132, 61)
SAGE = RGBColor(90, 142, 157)
INK = RGBColor(23, 50, 77)
MUTED = RGBColor(111, 134, 153)
LINE = RGBColor(214, 222, 227)
SOFT_LINE = RGBColor(227, 234, 238)
SOFT_NAVY = RGBColor(71, 122, 141)
PAPER_WARM = RGBColor(255, 250, 241)
MARGIN_ROSE = RGBColor(244, 222, 186)
PALE_BLUE = RGBColor(236, 247, 252)
PALE_GOLD = RGBColor(248, 221, 164)
PALE_CORAL = RGBColor(253, 241, 223)
PALE_SAGE = RGBColor(239, 247, 249)
PALE_NAVY = RGBColor(235, 244, 247)
EXACT_ESOL_TEMPLATE_FAMILY = "esol_exact_workbook"

SLIDE_KINDS = [
    "cover",
    "be_curious",
    "learning_target",
    "vocabulary",
    "guided_notes",
    "worked_example",
    "practice",
    "quick_review",
    "challenge",
    "reflection",
    "exit_ticket",
]

PROBLEM_SOLVING_KINDS = {
    "worked_example",
    "practice",
    "challenge",
    "exit_ticket",
}

REFERENCE_PROBLEM_TEMPLATE_ROLES = {
    "guided_practice",
    "interactive_activity",
    "best_fit_review",
}

PROBLEM_INTERACTIVE_PAIR_KINDS = {
    "worked_example",
    "practice",
    "challenge",
}

PROBLEM_RENDER_MODE_FOCUS = "problem_focus"
PROBLEM_RENDER_MODE_INTERACTIVE = "interactive_apply"
SOURCE_ANCHORED_PROBLEM_ACTIVITY_FAMILIES = {
    "build_construct",
    "sequence_order",
    "sort_classify",
    "match_pair",
    "detect_justify",
}
SOURCE_ALIGNED_TEMPLATE_ACTIVITY_ROLES = {
    "drag_sort",
    "error_analysis",
    "two_column_compare",
    "choice_board",
    "collaborative_practice",
    "independent_practice",
    "turn_and_teach",
}
GENERIC_PROBLEM_ACTIVITY_MARKERS = (
    "strategy ranking",
    "compare-and-rank",
    "compare and rank",
    "cover-and-reveal",
    "gallery walk",
    "sticky-note sort",
    "sticky note sort",
    "reasoning ladder",
    "independent solve path",
    "error analysis repair",
    "fits you best",
    "best fit",
)
GENERIC_PROBLEM_ACTIVITY_PIECES = {
    "path a",
    "path b",
    "path c",
    "strategy a",
    "strategy b",
    "what is correct?",
    "what is incorrect?",
    "fix it",
    "explain it",
    "read",
    "write",
    "solve",
    "explain",
}

FORMAL_REVIEW_ACTION_STARTERS = {
    "choose",
    "compare",
    "decide",
    "determine",
    "explain",
    "find",
    "identify",
    "list",
    "solve",
    "substitute",
    "use",
    "write",
}

FORMAL_REVIEW_MIN_FONT_PT = 10.4
AUTO_FIT_MIN_FONT_PT = 11.2

PUBLISHER_DISPLAY_REPLACEMENTS = (
    (r"\binteractive apply:\s*", "Practice Extension: "),
    (r"\binteractive apply\b", "Practice Extension"),
    (r"\bsource-faithful\b", "carefully sequenced"),
    (r"\bsource-aligned\b", "lesson-aligned"),
    (r"\bworkbook baseline\b", "student notebook"),
    (r"\buploaded model\b", "lesson structure"),
    (r"\blesson problem\b", "source problem"),
    (r"\bshow what you know\b", "demonstrate what you understand"),
    (r"\bshow what you can do independently\b", "demonstrate what you understand independently"),
    (r"\bshow what you can do on your own\b", "demonstrate what you understand on your own"),
    (r"\buse the same idea\b", "use the same strategy"),
    (r"\bplan your next step\b", "name the next idea to strengthen"),
    (r"\bpartner problem\b", "similar problem"),
)

PUBLISHER_INTERNAL_JARGON_PATTERNS = (
    r"\bsource-faithful\b",
    r"\bsource-aligned\b",
    r"\buploaded model\b",
    r"\bworkbook baseline\b",
    r"\binteractive apply\b",
    r"\bproblem-first\b",
)

PUBLISHER_WEAK_COPY_PATTERNS = (
    (r"\bshow what you know\b", "generic assessment phrasing"),
    (r"\bshow what you can do independently\b", "generic assessment phrasing"),
    (r"\bshow what you can do on your own\b", "generic assessment phrasing"),
    (r"\buse the same idea\b", "generic transfer phrasing"),
    (r"\bplan your next step\b", "generic reflection phrasing"),
    (r"\bmodel the next step before independent work\b", "generic guided-practice phrasing"),
)

KIND_DEFAULT_SECTION = {
    "cover": "Welcome",
    "be_curious": "Launch",
    "learning_target": "Focus",
    "vocabulary": "Support",
    "guided_notes": "Guided Notes",
    "worked_example": "Worked Example",
    "practice": "Practice",
    "quick_review": "Review",
    "challenge": "Challenge",
    "reflection": "Reflect",
    "exit_ticket": "Exit Ticket",
}

KIND_ACCENT = {
    "cover": CORAL,
    "be_curious": CORAL,
    "learning_target": TEAL,
    "vocabulary": SAGE,
    "guided_notes": TEAL,
    "worked_example": CORAL,
    "practice": GOLD,
    "quick_review": GOLD,
    "challenge": TEAL,
    "reflection": SAGE,
    "exit_ticket": CORAL,
}

STOPWORDS = {
    "about",
    "after",
    "again",
    "also",
    "another",
    "because",
    "before",
    "between",
    "being",
    "could",
    "does",
    "from",
    "have",
    "into",
    "lesson",
    "more",
    "much",
    "must",
    "only",
    "over",
    "slide",
    "slides",
    "student",
    "students",
    "their",
    "there",
    "these",
    "this",
    "those",
    "through",
    "today",
    "using",
    "what",
    "when",
    "where",
    "which",
    "while",
    "with",
    "would",
    "your",
}

STUDENT_FRIENDLY_DEFINITIONS = {
    "barrier": "something that makes the task harder to do",
    "area": "the amount of space inside a figure",
    "body weight": "how much a body weighs as a whole amount",
    "claim": "what you say is true or your answer",
    "data": "facts, numbers, or observations you can use as evidence",
    "data set": "a group of values you study together",
    "decimal": "a number that shows part of a whole using place value",
    "decimals": "numbers that show parts of a whole using place value",
    "equivalent": "equal in value even when written in a different form",
    "estimate": "a close answer based on what you see or know",
    "exact": "the precise answer with no estimate",
    "evidence": "facts or data that support a claim",
    "fraction": "a number that names a part of a whole",
    "fractions": "numbers that name parts of a whole",
    "whole": "the complete amount before you find a part of it",
    "part": "one piece or portion of a whole amount",
    "volume": "the amount of space inside a three-dimensional figure",
    "notation": "the way a number is written",
    "notations": "different ways a number can be written",
    "polygon": "a closed shape made from straight line segments",
    "percent": "how many out of every one hundred",
    "percentage": "how many out of every one hundred",
    "percentages": "amounts written out of every one hundred",
    "octagon": "an eight-sided polygon",
    "regular": "having equal sides and equal angles",
    "regular polygon": "a polygon with equal sides and equal angles",
    "triangle": "a shape with three straight sides",
    "prism": "a solid shape with matching ends and flat faces",
    "rectangular prism": "a box-shaped solid with rectangular faces",
    "congruent": "the same size and the same shape",
    "congruent triangle": "a triangle that is the same size and the same shape as another triangle",
    "congruent triangles": "triangles that are the same size and the same shape",
    "compose": "to put parts together to make a whole figure",
    "decompose": "to break a figure into smaller, known shapes",
    "dimension": "a measurement such as a length, width, or height",
    "dimensions": "measurements such as length, width, and height",
    "length": "how long something is from end to end",
    "width": "how wide something is from side to side",
    "height": "how tall something is from bottom to top",
    "formula": "a math rule written with numbers and symbols",
    "lab data": "observations or measurements collected during an investigation",
    "line plot": "a graph that uses marks above a number line to show data values",
    "mean": "the average value in a data set",
    "median": "the middle value in an ordered data set",
    "ordered data": "data values arranged from least to greatest or greatest to least",
    "outlier": "a value much larger or smaller than the rest of the data",
    "reasoning": "the explanation that connects the evidence to the claim",
    "sentence frame": "a sentence pattern that helps you write or speak",
    "sentence frames": "sentence patterns that help you write or speak",
    "sentence starter": "a sentence beginning that helps you explain your thinking",
    "sentence starters": "sentence beginnings that help you explain your thinking",
    "unit cube": "a cube with side lengths of one unit",
    "unit cubes": "cubes with side lengths of one unit",
    "apothem": "the distance from the center of a regular polygon to the midpoint of a side",
    "composite": "made from two or more shapes",
    "figure": "a shape or drawing used in a math problem",
    "budget": "the amount of money planned for a purchase",
    "approximate": "close to the exact value",
    "dot plot": "a graph that shows data values with dots above a number line",
}

SOURCE_TERM_PRIORITY = (
    "median",
    "mean",
    "data set",
    "ordered data",
    "dot plot",
    "line plot",
    "outlier",
    "body weight",
    "rectangular prism",
    "unit cubes",
    "unit cube",
    "volume formula",
    "fractional edge lengths",
    "regular polygon",
    "octagon",
    "congruent triangles",
    "congruent triangle",
    "decompose",
    "compose",
    "congruent",
    "polygon",
    "base",
    "dimensions",
    "dimension",
    "volume",
    "formula",
    "estimate",
    "exact",
    "whole",
    "part",
    "length",
    "width",
    "height",
    "area",
    "prism",
    "claim",
    "evidence",
    "reasoning",
    "sentence starters",
    "sentence starter",
    "sentence frames",
    "sentence frame",
    "lab data",
    "data",
)

EXACT_ESOL_TEMPLATE_SEQUENCE: list[tuple[str, str]] = [
    ("learning_target", "learning_objectives"),
    ("be_curious", "prior_session_review"),
    ("vocabulary", "vocabulary_table"),
    ("worked_example", "guided_practice"),
    ("practice", "interactive_activity"),
    ("practice", "best_fit_review"),
]

SPANISH_TRANSLATIONS = {
    "area": "area",
    "base": "base",
    "compare": "comparar",
    "cube": "cubo",
    "cubes": "cubos",
    "diagonal": "diagonal",
    "diagonals": "diagonales",
    "dimension": "dimension",
    "dimensions": "dimensiones",
    "equation": "ecuacion",
    "explain": "explicar",
    "formula": "formula",
    "graph": "grafica",
    "height": "altura",
    "label": "etiqueta",
    "labels": "etiquetas",
    "length": "longitud",
    "missing dimension": "dimension que falta",
    "multiply": "multiplicar",
    "octagon": "octagono",
    "polygon": "poligono",
    "prism": "prisma",
    "rectangular prism": "prisma rectangular",
    "regular polygon": "poligono regular",
    "rhombus": "rombo",
    "solve": "resolver",
    "table": "tabla",
    "triangle": "triangulo",
    "congruent": "congruente",
    "congruent triangle": "triangulo congruente",
    "congruent triangles": "triangulos congruentes",
    "compose": "componer",
    "decompose": "descomponer",
    "unit": "unidad",
    "unit cube": "cubo unitario",
    "unit cubes": "cubos unitarios",
    "volume": "volumen",
    "width": "ancho",
}

GENERIC_VOCAB_TERMS = {
    "session",
    "student",
    "notebook",
    "collaborate",
    "connect",
    "curious",
    "mindset",
    "reveal",
    "workspace",
    "analyzing",
    "reflect",
    "apply",
}

LOW_VALUE_VOCAB_TERMS = {
    "achieved",
    "case",
    "describe",
    "discussion",
    "doesn",
    "drive",
    "each",
    "experience",
    "experiences",
    "following",
    "family",
    "food",
    "grade",
    "goal",
    "middle",
    "move",
    "movie",
    "number",
    "numbers",
    "planning",
    "question",
    "questions",
    "reasonable",
    "reorganization",
    "school",
    "science",
    "snack",
    "snacks",
    "sentence",
    "starter",
    "starters",
    "still",
    "student",
    "students",
    "study",
    "some",
    "solution",
    "take",
    "teacher",
    "teachers",
    "teaching",
    "theater",
    "they",
    "three",
    "today",
    "voted",
    "week",
    "years",
    "written",
    "writing",
}

PLACEHOLDER_VOCAB_DEFINITION = "a lesson word to use when you explain your math thinking"

PROMPT_STARTERS = (
    "question:",
    "what ",
    "how ",
    "why ",
    "which ",
    "would ",
    "explain ",
    "describe ",
    "determine ",
    "estimate ",
    "compare ",
    "tell whether",
    "solve ",
    "find ",
    "justify ",
)

DANGLING_ENDING_WORDS = {
    "and",
    "as",
    "at",
    "because",
    "by",
    "for",
    "from",
    "in",
    "into",
    "of",
    "on",
    "or",
    "so",
    "that",
    "to",
    "with",
}

DANGLING_ENDING_PAIRS = {
    "for a",
    "for an",
    "for the",
    "from a",
    "from the",
    "in a",
    "in an",
    "in the",
    "into a",
    "into the",
    "of a",
    "of an",
    "of the",
    "on a",
    "on the",
    "such as",
    "to a",
    "to the",
    "with a",
    "with the",
}

NUMBER_WORD_STARTERS = {
    "zero",
    "one",
    "two",
    "three",
    "four",
    "five",
    "six",
    "seven",
    "eight",
    "nine",
    "ten",
    "eleven",
    "twelve",
    "thirteen",
    "fourteen",
    "fifteen",
    "sixteen",
    "seventeen",
    "eighteen",
    "nineteen",
    "twenty",
    "thirty",
    "forty",
    "fifty",
    "sixty",
    "seventy",
    "eighty",
    "ninety",
    "hundred",
}

SAFE_TERMINAL_ABBREVIATIONS = {
    "cm.",
    "ft.",
    "hr.",
    "hrs.",
    "in.",
    "kg.",
    "km.",
    "lb.",
    "lbs.",
    "m.",
    "min.",
    "mm.",
    "oz.",
    "sec.",
    "yd.",
}

LESSON_LEAD_VERBS = (
    "analyze",
    "apply",
    "build",
    "calculate",
    "compare",
    "create",
    "define",
    "describe",
    "determine",
    "develop",
    "evaluate",
    "explain",
    "explore",
    "find",
    "graph",
    "identify",
    "interpret",
    "model",
    "represent",
    "show",
    "simplify",
    "solve",
    "use",
)

ACTIVITY_FAMILIES = [
    "sort_classify",
    "match_pair",
    "sequence_order",
    "build_construct",
    "plot_place",
    "detect_justify",
    "compare_rank",
    "reveal_discuss",
]

ACTIVITY_FAMILY_OPTIONS = [""] + ACTIVITY_FAMILIES

ACTIVITY_FAMILY_LABELS = {
    "sort_classify": "Sort + Classify",
    "match_pair": "Match + Connect",
    "sequence_order": "Sequence + Order",
    "build_construct": "Build + Construct",
    "plot_place": "Plot + Place",
    "detect_justify": "Detect + Justify",
    "compare_rank": "Compare + Rank",
    "reveal_discuss": "Reveal + Discuss",
}

ACTIVITY_FAMILY_SHORT_LABELS = {
    "sort_classify": "Sort",
    "match_pair": "Match",
    "sequence_order": "Sequence",
    "build_construct": "Build",
    "plot_place": "Place",
    "detect_justify": "Justify",
    "compare_rank": "Compare",
    "reveal_discuss": "Reveal",
}

ACTIVITY_FAMILY_PROMPTS = {
    "sort_classify": "Best for categorizing examples, sorting representations, separating cases, or comparing correct vs incorrect thinking.",
    "match_pair": "Best for matching representations, vocabulary, graphs, tables, equations, scenarios, or evidence pairs.",
    "sequence_order": "Best for step sorting, ranking, sequencing solution paths, reconstructing reasoning, or organizing process order.",
    "build_construct": "Best for building equations, tables, models, expressions, diagrams, or visual structures from source content.",
    "plot_place": "Best for placement tasks such as graphing, number lines, coordinate grids, labeling, or locating features.",
    "detect_justify": "Best for spotting errors, defending claims, justifying choices, checking solutions, and evidence-based explanation.",
    "compare_rank": "Best for comparing strategies, rates, quantities, graphs, or ordering options by strength or magnitude.",
    "reveal_discuss": "Best for cover-and-reveal, clue-based discussion, sticky-note response work, and curiosity-building prompts.",
}

HIGH_AGENCY_ACTIVITY_FAMILIES = {
    "sort_classify",
    "match_pair",
    "sequence_order",
    "build_construct",
    "plot_place",
    "detect_justify",
    "compare_rank",
}

CHECKMARK_CHIP = "✓"
REFERENCE_WORKBOOK_BASELINE_GUIDANCE = (
    "Start from the compact Session 1 workbook baseline, then adapt it to the uploaded lesson slides instead of inventing a new notebook structure each time. "
    "The default model should be a premium six-slide notebook: Objectives + Session Map, Be Curious, Vocabulary + Reference Tool, Guided Problem, Interactive Activity, and Best-Fit Interactive Review. "
    "Problem pages should show the actual lesson problem, leave visible room to solve, keep the source wording verbatim when available, and stay tightly anchored to the lesson language, visuals, and setup. "
    "After the notebook pages, students should continue the remaining practice in the book instead of expecting a full second notebook session."
)
EXACT_ESOL_WORKBOOK_GUIDANCE = (
    "Use the compact Session 1 notebook format as the default direction. "
    "Do not invent a new structure and do not reduce source fidelity. Adapt the math content, examples, sequence, visuals, vocabulary, and objectives to the uploaded lesson slides while keeping architecture and renderer stable. "
    "Use this exact six-slide architecture for Session 1 only: Objectives + Session Map, Be Curious / Notice-Wonder, Vocabulary + Reference Tool, Guided Problem, Interactive Activity, Best-Fit Interactive Review. "
    "Keep editable PPTX output with real text boxes, shapes, tables, and movable student-facing elements. "
    "Use larger readable text, clean spacing, strong hierarchy, restrained color, and light notebook-tab section identity. "
    "Avoid worksheet crowding, decorative clutter, flattened image-only pages, and any second-session notebook expansion."
)
PREMIUM_NOTEBOOK_GUIDANCE = (
    "Build flagship, TpT-premium student notebooks that preserve source fidelity and feel sellable as polished classroom resources, not plain worksheets. "
    "Default to the locked six-slide Session 1 notebook architecture: Objectives + Session Map, Be Curious, Vocabulary + Reference Tool, Guided Problem, Interactive Activity, and Best-Fit Interactive Review. "
    "Adapt the notebook structure to the uploaded lesson instead of forcing a longer workbook arc onto every deck. "
    "Lift exact source problems, examples, objective language, and visuals whenever those appear in the slides. "
    "Keep fonts large, layouts spacious, and formatting stable with clear paragraph separation and no blended text lines. "
    "Use embedded text inside cards and prompt boxes, blank boxed response spaces instead of handwriting lines, and no loose floating text. "
    "Keep on-slide copy concise enough to preserve large, highly readable text without crowding the page. "
    f"{REFERENCE_WORKBOOK_BASELINE_GUIDANCE} "
    "Student engagement is non-negotiable: avoid decks that only ask students to read, write, and solve in static boxes. Build a strong launch, a useful vocabulary page, a guided source-problem arc, and two discussion-ready interactive pages instead of spreading attention across a long deck. "
    "Use only database-backed interactive activities, keep them tightly anchored to the exact source slides, and prefer fewer stronger interactions over decorative extras. "
    "Integrate the first written source problem directly into the modeled page and the your-turn page instead of paraphrasing it away. "
    "After those pages, direct students back to the book for the remaining practice load. "
    "Run a professional publisher-style copy edit and formatting review before final output."
)
NON_NEGOTIABLE_QUALITY_LOCK_GUIDANCE = (
    "Non-negotiable output lock: keep every notebook in premium flagship style with high readability, strong source fidelity, and sellable TpT polish. "
    "Never use worksheet-style handwriting lines for student writing. Use blank boxed response workspaces with embedded prompt text."
)
LEGACY_EXACT_PREMIUM_NOTEBOOK_GUIDANCE = (
    f"{EXACT_ESOL_WORKBOOK_GUIDANCE} "
    "Build flagship, TpT-premium student notebooks that preserve source fidelity and feel sellable as polished classroom resources, not plain worksheets. "
    "Lift exact source problems, examples, objective language, and visuals whenever those appear in the slides. "
    "Keep fonts large, layouts spacious, and formatting stable with clear paragraph separation and no blended text lines. Keep on-slide copy concise enough to preserve large, highly readable text without crowding the page. "
    f"{REFERENCE_WORKBOOK_BASELINE_GUIDANCE} "
    "Use only database-backed interactive activities, keep them tightly anchored to the exact source slides, and prefer fewer stronger interactions over decorative extras. "
    "Integrate written source problems directly into worked example, guided practice, collaborate/turn-and-teach, source-fidelity task, and reflection/exit pages instead of paraphrasing them away. "
    "Run a professional publisher-style copy edit and formatting review before final output."
)

PREMIUM_FEATURE_OPTIONS = [
    "error_analysis",
    "multi_representation",
    "scenario_continuity",
    "scaffold_fade",
    "strategy_comparison",
    "evidence_ladder",
    "decision_tree",
    "writing_revolution",
    "turn_and_teach",
    "create_your_own",
    "real_world_transfer",
    "mastery_tracker",
]

HIGH_AGENCY_PREMIUM_FEATURES = [
    "error_analysis",
    "multi_representation",
    "strategy_comparison",
    "evidence_ladder",
    "decision_tree",
    "turn_and_teach",
    "create_your_own",
    "real_world_transfer",
]

HIGH_AGENCY_PREMIUM_LAYOUTS = set(HIGH_AGENCY_PREMIUM_FEATURES)

HIGH_AGENCY_TEMPLATE_ROLES = {
    "drag_sort",
    "error_analysis",
    "two_column_compare",
    "choice_board",
    "collaborative_practice",
    "independent_practice",
    "turn_and_teach",
    "tiered_exit",
}

ENGAGEMENT_MODE_ACTIVITY_FAMILIES = {
    "sort_classify": {"sort"},
    "match_pair": {"match"},
    "sequence_order": {"sequence"},
    "build_construct": {"build"},
    "plot_place": {"build", "place"},
    "detect_justify": {"justify"},
    "compare_rank": {"compare"},
    "reveal_discuss": {"notice", "discuss"},
}

FAMILY_KEYWORD_HINTS = {
    "sort_classify": {"sort", "classify", "group", "category"},
    "match_pair": {"match", "pair", "connect", "table", "graph", "equation", "vocabulary"},
    "sequence_order": {"sequence", "steps", "order", "process"},
    "build_construct": {"build", "construct", "equation", "expression", "table", "model"},
    "plot_place": {"graph", "coordinate", "number line", "point", "grid", "plot"},
    "detect_justify": {"error", "justify", "claim", "evidence", "reasoning"},
    "compare_rank": {"compare", "rank", "strategy", "best", "strongest"},
    "reveal_discuss": {"notice", "wonder", "predict", "discussion", "reveal"},
}

ENGAGEMENT_MODE_TEMPLATE_ROLES = {
    "drag_sort": {"sort", "build"},
    "error_analysis": {"justify"},
    "two_column_compare": {"compare"},
    "choice_board": {"choice"},
    "collaborative_practice": {"partner"},
    "independent_practice": {"apply"},
    "turn_and_teach": {"teach", "partner"},
    "tiered_exit": {"choice", "justify"},
}

ENGAGEMENT_MODE_PREMIUM_LAYOUTS = {
    "error_analysis": {"justify"},
    "multi_representation": {"build"},
    "strategy_comparison": {"compare"},
    "evidence_ladder": {"justify", "discuss"},
    "decision_tree": {"choice"},
    "turn_and_teach": {"teach", "partner"},
    "create_your_own": {"create"},
    "real_world_transfer": {"apply", "create"},
}

ENGAGEMENT_MODE_BY_ACTIVITY_FAMILY = {
    "sort_classify": "sort",
    "match_pair": "match",
    "sequence_order": "sequence",
    "build_construct": "build",
    "plot_place": "place",
    "detect_justify": "justify",
    "compare_rank": "compare",
}

ENGAGEMENT_MODE_BY_PREMIUM_LAYOUT = {
    "error_analysis": "justify",
    "multi_representation": "build",
    "strategy_comparison": "compare",
    "evidence_ladder": "justify",
    "decision_tree": "decision",
    "turn_and_teach": "teach",
    "create_your_own": "create",
    "real_world_transfer": "apply",
}

ENGAGEMENT_MODE_BY_TEMPLATE_ROLE = {
    "drag_sort": "sort",
    "error_analysis": "justify",
    "two_column_compare": "compare",
    "choice_board": "choice",
    "collaborative_practice": "collaborate",
    "independent_practice": "independent_solve",
    "turn_and_teach": "teach",
    "tiered_exit": "choice",
}


def guidance_has_exact_workbook_contract(guidance: str) -> bool:
    lowered = normalize_whitespace(guidance).lower()
    required_markers = (
        "session 1",
        "six-slide",
        "objectives",
        "be curious",
        "vocabulary",
        "guided problem",
    )
    return bool(lowered) and all(marker in lowered for marker in required_markers)


def ensure_exact_workbook_guidance(guidance: str) -> str:
    cleaned = normalize_whitespace(guidance)
    if not cleaned:
        return LEGACY_EXACT_PREMIUM_NOTEBOOK_GUIDANCE
    if guidance_has_exact_workbook_contract(cleaned):
        return cleaned
    if cleaned.startswith("Build flagship, TpT-premium student notebooks"):
        return LEGACY_EXACT_PREMIUM_NOTEBOOK_GUIDANCE
    return normalize_whitespace(f"{LEGACY_EXACT_PREMIUM_NOTEBOOK_GUIDANCE} {cleaned}")


def normalize_saved_guidance(guidance: str) -> str:
    cleaned = normalize_whitespace(guidance)
    if not cleaned:
        return PREMIUM_NOTEBOOK_GUIDANCE
    if cleaned == normalize_whitespace(LEGACY_EXACT_PREMIUM_NOTEBOOK_GUIDANCE):
        return PREMIUM_NOTEBOOK_GUIDANCE
    if cleaned.startswith("Build flagship, TpT-premium student notebooks") and not guidance_has_exact_workbook_contract(cleaned):
        return normalize_whitespace(f"{EXACT_ESOL_WORKBOOK_GUIDANCE} {cleaned}")
    return cleaned


def enforce_runtime_quality_guidance(guidance: str) -> str:
    cleaned = normalize_saved_guidance(guidance)
    if guidance_has_exact_workbook_contract(cleaned):
        base = ensure_exact_workbook_guidance(cleaned)
    elif not cleaned:
        base = ensure_exact_workbook_guidance(PREMIUM_NOTEBOOK_GUIDANCE)
    elif cleaned.startswith("Build flagship, TpT-premium student notebooks"):
        base = ensure_exact_workbook_guidance(cleaned)
    else:
        base = ensure_exact_workbook_guidance(f"{PREMIUM_NOTEBOOK_GUIDANCE} {cleaned}")
    if normalize_whitespace(NON_NEGOTIABLE_QUALITY_LOCK_GUIDANCE).lower() in normalize_whitespace(base).lower():
        return base
    return normalize_whitespace(f"{base} {NON_NEGOTIABLE_QUALITY_LOCK_GUIDANCE}")

PREMIUM_FEATURE_LABELS = {
    "error_analysis": "Error Analysis Block",
    "multi_representation": "Multi-Representation Builder",
    "scenario_continuity": "Scenario Continuity Extension",
    "scaffold_fade": "Scaffold Fade Sequence",
    "strategy_comparison": "Strategy Comparison Panel",
    "evidence_ladder": "Evidence Ladder",
    "decision_tree": "Decision Tree / Math Cheat Code",
    "writing_revolution": "Writing Revolution Explanation Layer",
    "turn_and_teach": "Turn-and-Teach Prompt",
    "create_your_own": "Create-Your-Own Task",
    "real_world_transfer": "Real-World Transfer Studio",
    "mastery_tracker": "Reflection / Mastery Tracker",
}

PREMIUM_BLOCKING_LAYOUTS = {
    "error_analysis",
    "multi_representation",
    "strategy_comparison",
    "evidence_ladder",
    "decision_tree",
    "create_your_own",
    "real_world_transfer",
    "mastery_tracker",
}

PREMIUM_FULL_SPREAD_LAYOUTS = {
    "evidence_ladder",
    "real_world_transfer",
}

FLAGSHIP_ACTIVITY_TYPES = [
    "Error Analysis",
    "Sort / Match / Classify",
    "Notice / Wonder extension",
    "Build the Rule / Find the Pattern",
    "My Turn / Your Turn",
    "Which One Doesn't Belong",
    "Example vs. Non-Example",
    "Complete the Representation",
    "Explain the Strategy",
    "Mini Math Debate",
    "Vocabulary in Action",
    "Real-World Transfer",
    "Fix the Mistake",
    "Compare and Justify",
    "Table / Graph / Equation connection task",
]

FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED = "embedded_section"
FLAGSHIP_ACTIVITY_FOOTPRINT_FULL = "full_slide"

FLAGSHIP_ACTIVITY_TYPE_TO_FAMILY = {
    "Error Analysis": "detect_justify",
    "Sort / Match / Classify": "sort_classify",
    "Notice / Wonder extension": "reveal_discuss",
    "Build the Rule / Find the Pattern": "build_construct",
    "My Turn / Your Turn": "sequence_order",
    "Which One Doesn't Belong": "compare_rank",
    "Example vs. Non-Example": "sort_classify",
    "Complete the Representation": "build_construct",
    "Explain the Strategy": "sequence_order",
    "Mini Math Debate": "compare_rank",
    "Vocabulary in Action": "match_pair",
    "Real-World Transfer": "build_construct",
    "Fix the Mistake": "detect_justify",
    "Compare and Justify": "compare_rank",
    "Table / Graph / Equation connection task": "match_pair",
}

FLAGSHIP_ACTIVITY_TYPE_TO_LAYOUT = {
    "Error Analysis": "mistake_analysis_cards",
    "Sort / Match / Classify": "sort_mat",
    "Notice / Wonder extension": "notice_wonder_board",
    "Build the Rule / Find the Pattern": "pattern_workspace",
    "My Turn / Your Turn": "workspace_duo",
    "Which One Doesn't Belong": "debate_board",
    "Example vs. Non-Example": "example_non_example_board",
    "Complete the Representation": "representation_workspace",
    "Explain the Strategy": "strategy_workspace",
    "Mini Math Debate": "debate_board",
    "Vocabulary in Action": "vocabulary_workspace",
    "Real-World Transfer": "transfer_workspace",
    "Fix the Mistake": "mistake_analysis_cards",
    "Compare and Justify": "comparison_board",
    "Table / Graph / Equation connection task": "representation_connection_board",
}

FLAGSHIP_ACTIVITY_LAYOUT_ZONE_LABELS = {
    "mistake_analysis_cards": ["Notice the issue", "Fix the step", "Use evidence", "Explain why"],
    "sort_mat": ["Sort A", "Sort B", "Sort C", "Explain the rule"],
    "notice_wonder_board": ["Notice", "Wonder", "Predict", "Use a clue"],
    "pattern_workspace": ["Source pattern", "Build the rule", "Test it", "Explain it"],
    "workspace_duo": ["Model first", "My turn", "Check it", "Teach it"],
    "debate_board": ["Claim A", "Claim B", "Evidence", "My conclusion"],
    "example_non_example_board": ["Example", "Non-example", "Rule", "Why"],
    "representation_workspace": ["Source model", "Complete it", "Check the match", "Explain the connection"],
    "strategy_workspace": ["Step 1", "Step 2", "Step 3", "Why it works"],
    "vocabulary_workspace": ["Lesson word", "Source clue", "Use it here", "Say / write it"],
    "transfer_workspace": ["Source idea", "New situation", "Solve it", "Transfer proof"],
    "comparison_board": ["Compare", "Stronger evidence", "Counterpoint", "Justify"],
    "representation_connection_board": ["Table", "Graph / model", "Rule / equation", "Explain the match"],
}

FLAGSHIP_ACTIVITY_RENDERABLE_TEMPLATE_ROLES = {
    "best_fit_review",
    "drag_sort",
    "error_analysis",
    "guided_practice",
    "interactive_activity",
    "turn_and_teach",
    "vocabulary_activity",
}

FLAGSHIP_ACTIVITY_GENERIC_MARKERS = (
    "card 1",
    "card 2",
    "category a",
    "category b",
    "path a",
    "path b",
    "path c",
    "best fit",
    "sort the cards",
    "move a piece",
    "use the pieces",
    "problem a",
    "problem b",
    "example a",
    "example b",
)

TEMPLATE_ROLES_WITH_DRAGGABLES = {
    "best_fit_review",
    "learning_objectives",
    "interactive_activity",
    "vocabulary_activity",
    "drag_sort",
    "error_analysis",
    "two_column_compare",
    "choice_board",
    "collaborative_practice",
    "independent_practice",
    "turn_and_teach",
}

LOCKED_SECTION_LABELS = {
    "cover": "Cover",
    "learning_target": "Learning Target",
    "be_curious": "Notice + Wonder",
    "vocabulary": "Vocabulary",
    "guided_notes": "Guided Notes",
    "worked_example": "Guided Practice",
    "reflection": "Reflection",
    "exit_ticket": "Exit Ticket",
}

CONTEXT_STOPWORDS = STOPWORDS | GENERIC_VOCAB_TERMS | {
    "answer",
    "area",
    "base",
    "box",
    "challenge",
    "collaborative",
    "context",
    "dimension",
    "dimensions",
    "equation",
    "example",
    "exit",
    "find",
    "formula",
    "graph",
    "guided",
    "height",
    "independent",
    "learning",
    "length",
    "math",
    "notes",
    "notice",
    "objective",
    "practice",
    "prism",
    "problem",
    "rectangular",
    "reflection",
    "session",
    "show",
    "surveyed",
    "solve",
    "target",
    "teacher",
    "teachers",
    "teaching",
    "ticket",
    "use",
    "volume",
    "vocabulary",
    "wonder",
    "work",
    "width",
    "years",
}

ACTIVITY_KIND_PRIORITY = {
    "be_curious": 100,
    "vocabulary": 96,
    "practice": 94,
    "challenge": 92,
    "exit_ticket": 90,
    "quick_review": 88,
    "guided_notes": 84,
    "worked_example": 82,
    "reflection": 80,
    "cover": 0,
    "learning_target": 0,
}

ACTIVITY_KIND_CAPS = {
    "cover": 0,
    "learning_target": 0,
    "be_curious": 1,
    "vocabulary": 1,
    "guided_notes": 1,
    "worked_example": 1,
    "practice": 6,
    "quick_review": 1,
    "challenge": 2,
    "reflection": 1,
    "exit_ticket": 1,
}

PRIMARY_TEXT_LIMITS = {
    "cover": 96,
    "be_curious": 96,
    "learning_target": 98,
    "vocabulary": 92,
    "guided_notes": 156,
    "worked_example": 156,
    "practice": 118,
    "quick_review": 108,
    "challenge": 156,
    "reflection": 98,
    "exit_ticket": 126,
}

SECONDARY_TEXT_LIMITS = {
    "cover": 84,
    "be_curious": 88,
    "learning_target": 90,
    "vocabulary": 88,
    "guided_notes": 98,
    "worked_example": 90,
    "practice": 90,
    "quick_review": 86,
    "challenge": 90,
    "reflection": 90,
    "exit_ticket": 86,
}

TASK_TEXT_LIMITS = {
    "guided_notes": 102,
    "worked_example": 96,
    "practice": 92,
    "quick_review": 88,
    "challenge": 94,
    "exit_ticket": 90,
}

TASK_COUNT_LIMITS = {
    "guided_notes": 3,
    "worked_example": 3,
    "practice": 3,
    "quick_review": 2,
    "challenge": 3,
    "exit_ticket": 2,
}

BULLET_TEXT_LIMITS = {
    "be_curious": 70,
    "learning_target": 68,
    "guided_notes": 82,
    "worked_example": 82,
    "practice": 82,
    "quick_review": 78,
    "challenge": 80,
    "reflection": 76,
    "exit_ticket": 78,
}

RESPONSE_PROMPT_LIMITS = {
    "be_curious": 80,
    "guided_notes": 92,
    "worked_example": 88,
    "practice": 96,
    "quick_review": 96,
    "challenge": 96,
    "reflection": 96,
    "exit_ticket": 92,
}

ACTIVITY_KIND_DEFAULTS = {
    "be_curious": [
        "gallery walk sticky-note sort",
        "movable sticky-note responses",
        "cover-and-reveal",
        "estimate-then-reveal prediction",
        "card flip reveal discussion",
    ],
    "learning_target": [
        "movable discussion marker selection",
        "strategy ranking activity",
    ],
    "vocabulary": [
        "visual vocabulary matching",
        "definition matching",
        "math vocabulary categorization",
        "algebra vocabulary placement",
    ],
    "guided_notes": [
        "model annotation placement",
        "concept map building",
        "flowchart completion",
        "label-the-diagram placement",
        "build-a-model representation",
    ],
    "worked_example": [
        "reasoning ladder build",
        "error analysis repair cards",
        "step sequencing",
        "equation builder",
        "step-by-step solution reconstruction",
        "missing-step detective",
        "equation solution path builder",
    ],
    "practice": [
        "strategy sort mat",
        "claim-evidence match",
        "drag-and-sort categorization",
        "matching pairs",
        "table-to-equation matching",
        "representation matching (table/graph/equation)",
        "build-the-table activity",
    ],
    "quick_review": [
        "four-corner claim sort",
        "true/false sorting",
        "cover-and-reveal",
        "concept review matching",
        "peer answer check placement",
    ],
    "challenge": [
        "claim-evidence sort mat",
        "multi-step problem path selection",
        "drag-to-justify reasoning",
        "build-the-argument reasoning cards",
        "mathematical debate claim placement",
    ],
    "reflection": [
        "reflection evidence sort",
        "reasoning ladder build",
        "sentence builder",
        "build-a-math-explanation sentence strips",
        "discussion evidence placement",
        "reasoning ladder building",
    ],
    "exit_ticket": [
        "show-what-you-know self-check",
        "error analysis repair cards",
        "check-the-solution placement",
        "equation error correction drag",
        "mathematical error explanation placement",
    ],
}


def normalize_whitespace(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def is_placeholder_api_key(value: str) -> bool:
    cleaned = normalize_whitespace(value)
    if not cleaned:
        return True
    upper = cleaned.upper()
    return upper.startswith("PASTE_") or "YOUR_API_KEY" in upper or upper == "OPENAI_API_KEY"


def choose_api_key(*candidates: str) -> str:
    for candidate in candidates:
        cleaned = (candidate or "").strip()
        if cleaned and not is_placeholder_api_key(cleaned):
            return cleaned
    return ""


def truncate_text(value: str, limit: int) -> str:
    value = normalize_whitespace(value)
    if len(value) <= limit:
        return value
    cutoff = max(1, limit - 1)
    candidate = value[:cutoff].rstrip()
    boundary = candidate.rfind(" ")
    if boundary >= max(8, int(cutoff * 0.6)):
        candidate = candidate[:boundary].rstrip(" ,;:-")
    candidate = candidate.rstrip(" ,;:-") or value[:cutoff].rstrip()
    return candidate + "…"


def truncate_display_copy(value: str, limit: int) -> str:
    cleaned = normalize_whitespace(value)
    if len(cleaned) <= limit:
        return cleaned
    clauses = [normalize_whitespace(part) for part in re.split(r"(?<=[?.!;:])\s+|,\s+", cleaned) if normalize_whitespace(part)]
    candidate = ""
    for clause in clauses:
        test_value = normalize_whitespace(f"{candidate} {clause}")
        if len(test_value) > limit:
            break
        candidate = test_value
    if candidate:
        sentence_safe_candidate = trim_dangling_display_text(candidate.rstrip(",;:-"))
        if sentence_safe_candidate and (
            len(sentence_safe_candidate) >= max(24, int(limit * 0.55))
            or sentence_safe_candidate.endswith((".", "?", "!"))
            or is_problem_like_text(sentence_safe_candidate)
        ):
            return sentence_safe_candidate
    return trim_dangling_display_text(truncate_text(cleaned, limit))


def unique_nonempty(values: Iterable[str], limit: int | None = None) -> list[str]:
    seen: set[str] = set()
    results: list[str] = []
    for value in values:
        cleaned = normalize_whitespace(value)
        if not cleaned:
            continue
        lowered = cleaned.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        results.append(cleaned)
        if limit and len(results) >= limit:
            break
    return results


def text_tail_words(text: str, count: int) -> list[str]:
    words = normalize_whitespace(text).split()
    tail = words[-count:] if count > 0 else words
    return [re.sub(r"^[^A-Za-z0-9]+|[^A-Za-z0-9]+$", "", word).lower() for word in tail]


def trim_dangling_display_text(value: str) -> str:
    cleaned = normalize_whitespace(value)
    had_ellipsis = cleaned.endswith("…")
    if cleaned.endswith(("?", "!")) and not had_ellipsis:
        return normalize_whitespace(cleaned.rstrip(" ,;:-"))
    cleaned = cleaned.rstrip("…").rstrip(" ,;:-")
    words = cleaned.split()
    if len(words) <= 2 and not had_ellipsis:
        return normalize_whitespace(cleaned)
    while words:
        raw_last = words[-1].lower()
        if raw_last in SAFE_TERMINAL_ABBREVIATIONS:
            break
        if had_ellipsis:
            stripped_last = re.sub(r"^[^A-Za-z0-9]+|[^A-Za-z0-9]+$", "", raw_last)
            if len(words) > 1 and stripped_last.isalpha() and len(stripped_last) <= 3:
                words = words[:-1]
                had_ellipsis = False
                continue
        tail_one = text_tail_words(" ".join(words), 1)
        tail_two = text_tail_words(" ".join(words), 2)
        if len(tail_two) == 2 and " ".join(tail_two) in DANGLING_ENDING_PAIRS:
            words = words[:-2]
            continue
        if tail_one and tail_one[0] in DANGLING_ENDING_WORDS:
            words = words[:-1]
            continue
        break
    return normalize_whitespace(" ".join(words).rstrip(" ,;:-"))


def has_dangling_display_text(value: str) -> bool:
    cleaned = normalize_whitespace(value)
    if not cleaned:
        return False
    if cleaned.endswith(("?", "!")):
        return False
    words = cleaned.rstrip("…").split()
    if len(words) <= 2 and not cleaned.endswith(("…", ",", ";", ":")):
        return False
    if cleaned.endswith(".") and cleaned.split()[-1].lower() in SAFE_TERMINAL_ABBREVIATIONS:
        return False
    comparison = cleaned.rstrip("…").rstrip()
    return trim_dangling_display_text(cleaned) != comparison or comparison.endswith((",", ";", ":"))


def display_text_key(value: str) -> str:
    cleaned = normalize_whitespace(value).lower()
    cleaned = re.sub(r"[^a-z0-9]+", " ", cleaned)
    return normalize_whitespace(cleaned)


def first_distinct_text(candidates: Iterable[str], *, excluded: Iterable[str] = ()) -> str:
    excluded_keys = {display_text_key(item) for item in excluded if display_text_key(item)}
    for candidate in candidates:
        cleaned = trim_dangling_display_text(normalize_whitespace(candidate))
        if not cleaned:
            continue
        key = display_text_key(cleaned)
        if key and key not in excluded_keys:
            return cleaned
    return ""


def lesson_focus_phrase(deck: dict[str, Any], reference_text: str = "") -> str:
    cleaned = normalize_whitespace(reference_text or deck.get("lesson_title", ""))
    cleaned = re.sub(r"^(I can|We will|Students will)\s+", "", cleaned, flags=re.IGNORECASE)
    cleaned = cleaned.rstrip(".")
    lead_verb_pattern = r"^(?:" + "|".join(re.escape(verb) for verb in LESSON_LEAD_VERBS) + r")\s+"
    cleaned = re.sub(lead_verb_pattern, "", cleaned, flags=re.IGNORECASE)
    cleaned = normalize_whitespace(cleaned).rstrip(".")
    if not cleaned:
        cleaned = normalize_whitespace(deck.get("lesson_title", "this lesson")).rstrip(".")
    if not cleaned:
        return "this lesson"
    return cleaned[:1].lower() + cleaned[1:]


def first_standard_text(deck: dict[str, Any]) -> str:
    standards = deck.get("standards", [])
    if standards:
        return normalize_whitespace(str(standards[0]))
    text = normalize_whitespace(
        " ".join(
            [
                deck.get("lesson_title", ""),
                deck.get("summary", ""),
                " ".join(slide.get("title", "") for slide in deck.get("slides", [])[:8]),
                " ".join(slide.get("text", "") for slide in deck.get("slides", [])[:8]),
            ]
        )
    )
    match = re.search(r"\b\d+\.[A-Z]\.[A-Z](?:\.[A-Z])?\.\d+\b", text)
    if match:
        return match.group(0)
    lowered = text.lower()
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return "5.MD.C.5"
    if any(term in lowered for term in ("parallelogram", "triangle", "trapezoid", "rhombus")) and "area" in lowered:
        return "6.G.A.1"
    return ""


def collapse_repeated_opening_phrase(text: str, max_words: int = 8) -> str:
    words = normalize_whitespace(text).split()
    lowered = [word.lower() for word in words]
    upper = min(max_words, len(words) // 2)
    for size in range(upper, 1, -1):
        if lowered[:size] == lowered[size : size * 2]:
            return " ".join(words[size:])
    return normalize_whitespace(text)


def replace_case_insensitive(text: str, pattern: str, replacement: str) -> str:
    def repl(match: re.Match[str]) -> str:
        original = match.group(0)
        if original.isupper():
            return replacement.upper()
        if original[:1].isupper():
            return replacement[:1].upper() + replacement[1:]
        return replacement

    return re.sub(pattern, repl, text, flags=re.IGNORECASE)


def rewrite_publisher_display_text(text: str) -> str:
    rewritten = normalize_whitespace(text)
    for pattern, replacement in PUBLISHER_DISPLAY_REPLACEMENTS:
        rewritten = replace_case_insensitive(rewritten, pattern, replacement)
    return normalize_whitespace(rewritten)


def publisher_copyedit_text(
    value: str,
    *,
    title: str = "",
    preserve_source_prompt: bool = False,
) -> str:
    cleaned = normalize_whitespace((value or "").replace("\u200b", " ").replace("\ufeff", " "))
    if not cleaned:
        return ""
    cleaned = collapse_repeated_opening_phrase(cleaned)
    cleaned = rewrite_publisher_display_text(cleaned)
    if preserve_source_prompt or "reveal:" in cleaned.lower():
        cleaned = clean_source_prompt(cleaned, title)
    else:
        cleaned = re.sub(r"\s*Reveal:\s*$", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s+(Mindset|Workspace|Analyzing)\s*$", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+([,.;:?!])", r"\1", cleaned)
    cleaned = re.sub(r"([(\[])\s+", r"\1", cleaned)
    cleaned = re.sub(r"\s+([)\]])", r"\1", cleaned)
    cleaned = re.sub(r"\s*…\s*", "…", cleaned)
    if not preserve_source_prompt:
        cleaned = trim_dangling_display_text(cleaned)
    return normalize_whitespace(cleaned)


def publisher_copyedit_list(
    values: Iterable[str],
    *,
    title: str = "",
    preserve_source_prompt: bool = False,
    limit: int | None = None,
) -> list[str]:
    edited = (
        publisher_copyedit_text(value, title=title, preserve_source_prompt=preserve_source_prompt)
        for value in values
    )
    return unique_nonempty(edited, limit=limit)


def publisher_copyedit_issues(text: str) -> list[str]:
    cleaned = normalize_whitespace(text)
    if not cleaned:
        return []
    issues: list[str] = []
    if re.search(r"[\u200b\u200c\u200d\ufeff]", text):
        issues.append("hidden zero-width character")
    if collapse_repeated_opening_phrase(cleaned) != cleaned:
        issues.append("repeated opening phrase")
    if re.search(r"\s+[,.!?;:]", cleaned):
        issues.append("space before punctuation")
    if re.search(r"([?!])(?:\s*[\u200b\u200c\u200d\ufeff]*[?!])+$", text):
        issues.append("repeated ending punctuation")
    if "Reveal:" in cleaned:
        issues.append("teacher reveal cue")
    if re.search(r"(?<!/)\s+(Mindset|Workspace|Analyzing)\s*$", cleaned, flags=re.IGNORECASE):
        issues.append("teacher slide artifact")
    if has_dangling_display_text(cleaned):
        issues.append("dangling or incomplete ending")
    if any(re.search(pattern, cleaned, flags=re.IGNORECASE) for pattern in PUBLISHER_INTERNAL_JARGON_PATTERNS):
        issues.append("internal generator jargon")
    for pattern, label in PUBLISHER_WEAK_COPY_PATTERNS:
        if re.search(pattern, cleaned, flags=re.IGNORECASE):
            issues.append(label)
            break
    return issues


def action_prompt_count(prompts: Iterable[str]) -> int:
    count = 0
    for prompt in prompts:
        first_word = normalize_whitespace(prompt).split(" ", 1)[0].lower()
        if first_word in FORMAL_REVIEW_ACTION_STARTERS:
            count += 1
    return count


def ensure_terminal_punctuation(text: str, default: str = ".") -> str:
    cleaned = normalize_whitespace(text)
    if not cleaned:
        return ""
    if cleaned.endswith(("?", "!", ".")):
        return cleaned
    return f"{cleaned}{default}"


def slugify(value: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "-", value.strip()).strip("-").lower()
    return slug or "notebook"


def infer_activity_family(activity_name: str) -> str:
    name = activity_name.lower()
    if any(term in name for term in ("reveal", "mystery", "clue", "sticky-note", "sticky note", "flip", "discussion", "peer answer check", "estimate-and-adjust", "estimate-then-reveal")):
        return "reveal_discuss"
    if any(term in name for term in ("error", "detective", "justify", "evidence", "claim", "argument", "debate", "incorrect", "check-the-solution", "test point", "reasoning")):
        return "detect_justify"
    if any(term in name for term in ("plot", "placement", "number line", "coordinate", "point", "grid", "label", "marking", "mark ", "quadrant", "graph interpretation choice marking", "chooser", "mapping")):
        return "plot_place"
    if any(term in name for term in ("sequence", "sequencing", "ordering", "ranking", "ranking ", "path", "reconstruction", "ladder", "hierarchy", "continuation", "step-by-step", "step sorting", "solution path")):
        return "sequence_order"
    if any(term in name for term in ("builder", "building", "build-", "build a", "build-the", "construction", "assembly", "model", "flowchart", "concept map", "tile", "machine", "construction", "sentence strips")):
        return "build_construct"
    if any(term in name for term in ("match", "matching", "pairs", "pair")):
        return "match_pair"
    if any(term in name for term in ("compare", "comparison", "which-doesn", "vs", "rate comparison", "strategy ranking", "ordering cards", "magnitude", "ranking activity")):
        return "compare_rank"
    if any(term in name for term in ("sort", "categorization", "classification", "classify", "relevance", "feature sorting")):
        return "sort_classify"
    return "build_construct"


def clamp_int(value: Any, low: int, high: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = low
    return max(low, min(high, parsed))


def load_activity_database(path: Path = ACTIVITY_DATABASE_PATH) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    try:
        raw_payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return []
    raw_items = raw_payload.get("activities", raw_payload) if isinstance(raw_payload, dict) else raw_payload
    if not isinstance(raw_items, list):
        return []

    records: list[dict[str, Any]] = []
    for raw_item in raw_items:
        if not isinstance(raw_item, dict):
            continue
        name = normalize_whitespace(str(raw_item.get("name", "")))
        if not name:
            continue
        family = raw_item.get("family", "")
        if family not in ACTIVITY_FAMILY_LABELS:
            family = infer_activity_family(name)
        records.append(
            {
                "name": name,
                "family": family,
                "family_label": ACTIVITY_FAMILY_LABELS[family],
                "family_prompt": ACTIVITY_FAMILY_PROMPTS[family],
                "instructions": truncate_text(normalize_whitespace(str(raw_item.get("instructions", ""))), 240),
                "answer_check": truncate_text(normalize_whitespace(str(raw_item.get("answer_check", ""))), 200),
                "preferred_kinds": [kind for kind in raw_item.get("preferred_kinds", []) if kind in SLIDE_KINDS],
                "piece_limit": clamp_int(raw_item.get("piece_limit", 4), 2, 6),
                "piece_labels": unique_nonempty([str(item) for item in raw_item.get("piece_labels", [])], limit=6),
                "slot_labels": unique_nonempty([str(item) for item in raw_item.get("slot_labels", [])], limit=4),
                "interaction_style": normalize_whitespace(str(raw_item.get("interaction_style", ""))),
                "keywords": infer_activity_keywords(
                    name,
                    family=family,
                    preferred_kinds=[kind for kind in raw_item.get("preferred_kinds", []) if kind in SLIDE_KINDS],
                    extra_keywords=raw_item.get("keywords", []) or raw_item.get("topic_tags", []) or raw_item.get("concepts", []),
                ),
            }
        )
    return records


def infer_activity_keywords(
    name: str,
    *,
    family: str = "",
    preferred_kinds: Iterable[str] = (),
    extra_keywords: Iterable[str] = (),
) -> list[str]:
    lowered = normalize_whitespace(name).lower()
    keywords: set[str] = {normalize_whitespace(keyword).lower() for keyword in extra_keywords if normalize_whitespace(keyword)}
    keywords.update(FAMILY_KEYWORD_HINTS.get(family, set()))
    for phrase in (
        "graph",
        "table",
        "equation",
        "expression",
        "variable",
        "coordinate",
        "number line",
        "fraction",
        "decimal",
        "percent",
        "ratio",
        "proportional",
        "slope",
        "data",
        "dot plot",
        "box plot",
        "vocabulary",
        "error",
        "claim",
        "evidence",
        "compare",
        "justify",
        "sort",
        "match",
        "build",
        "sequence",
        "choice",
        "teach",
        "discussion",
        "reflection",
        "goal",
    ):
        if phrase in lowered:
            keywords.add(phrase)
    for kind in preferred_kinds:
        keywords.add(kind.replace("_", " "))
    return sorted(keyword for keyword in keywords if keyword)


def load_activity_library(path: Path = ACTIVITY_LIBRARY_PATH) -> list[dict[str, Any]]:
    activities_by_name: dict[str, dict[str, Any]] = {}
    for record in load_activity_database():
        activities_by_name[record["name"].lower()] = record

    if path.exists():
        for raw_line in path.read_text(encoding="utf-8").splitlines():
            name = normalize_whitespace(raw_line)
            if not name or name.startswith("#"):
                continue
            lowered = name.lower()
            if lowered in activities_by_name:
                continue
            family = infer_activity_family(name)
            activities_by_name[lowered] = {
                "name": name,
                "family": family,
                "family_label": ACTIVITY_FAMILY_LABELS[family],
                "family_prompt": ACTIVITY_FAMILY_PROMPTS[family],
                "instructions": "",
                "answer_check": "",
                "preferred_kinds": [],
                "piece_limit": 4,
                "piece_labels": [],
                "slot_labels": [],
                "interaction_style": "",
                "keywords": infer_activity_keywords(name, family=family),
            }
    return sorted(activities_by_name.values(), key=lambda item: (item["family"], item["name"]))


def deck_full_text(deck: dict[str, Any]) -> str:
    parts = [deck.get("lesson_title", ""), deck.get("summary", "")]
    for slide in deck.get("slides", []):
        parts.extend([slide.get("title", ""), slide.get("text", ""), slide.get("notes", "")])
    return normalize_whitespace(" ".join(parts)).lower()


def choose_activity_candidates(deck: dict[str, Any], library: list[dict[str, str]], limit: int = 24) -> list[dict[str, str]]:
    text = deck_full_text(deck)
    topic_groups = {
        "graphing": any(term in text for term in ("graph", "slope", "intercept", "coordinate", "quadrant", "point", "grid", "increasing", "decreasing")),
        "tables": any(term in text for term in ("table", "input", "output", "trend", "data", "rate of change")),
        "algebra": any(term in text for term in ("equation", "expression", "variable", "solve", "coefficient", "constant", "substitution", "term", "algebra")),
        "proportional": any(term in text for term in ("proportional", "ratio", "rate", "unit rate", "scale factor", "percent", "decimal", "fraction")),
        "reasoning": any(term in text for term in ("justify", "explain", "evidence", "claim", "reasoning", "error", "argument", "why")),
    }

    family_scores = {family: 0 for family in ACTIVITY_FAMILIES}
    if topic_groups["graphing"]:
        family_scores["plot_place"] += 5
        family_scores["match_pair"] += 2
        family_scores["compare_rank"] += 2
    if topic_groups["tables"]:
        family_scores["sort_classify"] += 3
        family_scores["match_pair"] += 4
        family_scores["build_construct"] += 3
    if topic_groups["algebra"]:
        family_scores["build_construct"] += 5
        family_scores["sequence_order"] += 4
        family_scores["detect_justify"] += 3
    if topic_groups["proportional"]:
        family_scores["compare_rank"] += 4
        family_scores["sort_classify"] += 3
        family_scores["match_pair"] += 2
    if topic_groups["reasoning"]:
        family_scores["detect_justify"] += 5
        family_scores["reveal_discuss"] += 2
        family_scores["sequence_order"] += 2

    scored: list[tuple[int, str, dict[str, str]]] = []
    for item in library:
        score = family_scores[item["family"]]
        lower_name = item["name"].lower()
        keyword_hits = 0
        for keyword in item.get("keywords", []):
            if keyword and keyword in text:
                keyword_hits += 1
        score += min(keyword_hits * 2, 8)
        for token in ("graph", "table", "equation", "variable", "rate", "proportional", "reasoning", "evidence", "coordinate", "number line", "vocabulary"):
            if token in lower_name and token in text:
                score += 3
        if any(keyword in lower_name for keyword in ("sorting", "matching", "builder", "detective", "discussion", "challenge")):
            score += 1
        scored.append((score, item["name"], item))

    scored.sort(key=lambda row: (-row[0], row[1]))
    selected: list[dict[str, str]] = []
    seen_names: set[str] = set()
    family_counts = {family: 0 for family in ACTIVITY_FAMILIES}
    for score, _name, item in scored:
        if score <= 0 and len(selected) >= min(10, limit):
            break
        if item["name"] in seen_names:
            continue
        if family_counts[item["family"]] >= 5:
            continue
        selected.append(item)
        seen_names.add(item["name"])
        family_counts[item["family"]] += 1
        if len(selected) >= limit:
            break
    return selected


def activity_library_digest(deck: dict[str, Any]) -> str:
    library = load_activity_library()
    if not library:
        return ""
    recommended = choose_activity_candidates(deck, library, limit=24)
    grouped: dict[str, list[str]] = {family: [] for family in ACTIVITY_FAMILIES}
    for item in recommended:
        grouped[item["family"]].append(item["name"])
    lines = [
        "Activity library guidance:",
        "Use the source problems and examples, but wrap them in premium interactive notebook structures when that deepens the lesson.",
        "Choose activities that fit the actual lesson content. Do not add random games that break fidelity.",
        f"Use the broader activity bank intentionally: {len(library)} activity patterns are available, so do not recycle the same few drag-and-drop moves when stronger topic-fit options exist.",
        "Use the database-backed activity records when available so instructions and answer checks stay consistent.",
        "Aim for about 5 high-agency interactive slides per session, using at least 3 distinct engagement modes such as build, compare, justify, choice, or teach-back when the source lesson supports them.",
        "",
        "High-confidence activity matches for this lesson:",
    ]
    for family in ACTIVITY_FAMILIES:
        names = grouped[family]
        if not names:
            continue
        lines.append(
            f"- {ACTIVITY_FAMILY_LABELS[family]}: {', '.join(names[:8])}"
        )
    lines.append("")
    lines.append("Full activity families available:")
    full_grouped: dict[str, list[str]] = {family: [] for family in ACTIVITY_FAMILIES}
    for item in library:
        full_grouped[item["family"]].append(item["name"])
    for family in ACTIVITY_FAMILIES:
        names = full_grouped[family]
        if not names:
            continue
        lines.append(f"- {ACTIVITY_FAMILY_LABELS[family]} ({ACTIVITY_FAMILY_PROMPTS[family]}):")
        lines.append(", ".join(names))
    return "\n".join(lines)


def load_dotenv(dotenv_path: Path = ROOT / ".env") -> None:
    if not dotenv_path.exists():
        return
    for line in dotenv_path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or "=" not in stripped:
            continue
        key, raw_value = stripped.split("=", 1)
        key = key.strip()
        value = raw_value.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = value


def resolve_ssl_cafile() -> str | None:
    env_path = os.getenv("SSL_CERT_FILE", "").strip()
    if env_path and Path(env_path).exists():
        return env_path
    for candidate in COMMON_CA_BUNDLE_PATHS:
        if Path(candidate).exists():
            return candidate
    return None


def validate_source_pptx_path(source_pptx: Path | str) -> Path:
    source_path = Path(source_pptx).expanduser().resolve()
    if not source_path.exists():
        raise RuntimeError(f"Source PPTX was not found: {source_path}")
    if not source_path.is_file():
        raise RuntimeError(f"Source PPTX must be a file: {source_path}")
    if source_path.suffix.lower() != ".pptx":
        raise RuntimeError(f"Source file must be a .pptx deck: {source_path}")
    return source_path


def validate_output_dir_path(output_dir: Path | str) -> Path:
    resolved_output_dir = Path(output_dir).expanduser().resolve()
    if resolved_output_dir.exists() and not resolved_output_dir.is_dir():
        raise RuntimeError(f"Output directory path points to a file: {resolved_output_dir}")
    return resolved_output_dir


def validate_json_artifact_path(path: Path | str, *, label: str) -> Path:
    resolved_path = Path(path).expanduser().resolve()
    if not resolved_path.exists():
        raise RuntimeError(f"{label} was not found: {resolved_path}")
    if not resolved_path.is_file():
        raise RuntimeError(f"{label} must be a file: {resolved_path}")
    if resolved_path.suffix.lower() != ".json":
        raise RuntimeError(f"{label} must be a .json file: {resolved_path}")
    return resolved_path


def notebook_runtime_preflight(
    source_pptx: Path | str,
    *,
    output_dir: Path | str | None = None,
) -> tuple[Path, Path]:
    source_path = validate_source_pptx_path(source_pptx)
    resolved_output_dir = resolve_output_dir(
        source_path,
        Path(output_dir).expanduser() if output_dir else None,
    )
    return source_path, validate_output_dir_path(resolved_output_dir)


def write_json(path: Path, payload: Any) -> None:
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def iter_leaf_shapes(shapes: Any) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def extract_table_text(table: Any) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [normalize_whitespace(cell.text) for cell in row.cells]
        cells = [cell for cell in cells if cell]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_text_items(slide: Any) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for shape in iter_leaf_shapes(slide.shapes):
        text = ""
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text
        elif getattr(shape, "has_table", False):
            text = extract_table_text(shape.table)
        text = normalize_whitespace(text)
        if not text:
            continue
        is_title_placeholder = False
        if getattr(shape, "is_placeholder", False):
            try:
                is_title_placeholder = shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                )
            except Exception:
                is_title_placeholder = False
        items.append(
            {
                "text": text,
                "top": int(getattr(shape, "top", 0)),
                "left": int(getattr(shape, "left", 0)),
                "is_title_placeholder": is_title_placeholder,
            }
        )
    items.sort(key=lambda item: (item["top"], item["left"], len(item["text"])))
    return items


def extract_slide_notes(slide: Any) -> str:
    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        return ""
    notes_text = notes_text.replace("\x0b", "\n")
    return normalize_whitespace(notes_text)


def pick_slide_title(text_items: list[dict[str, Any]]) -> str:
    for item in text_items:
        if item["is_title_placeholder"]:
            return truncate_text(item["text"], 140)
    for item in text_items:
        if len(item["text"]) <= 140:
            return truncate_text(item["text"], 140)
    if text_items:
        return truncate_text(text_items[0]["text"], 140)
    return ""


def build_slide_excerpt(text_items: list[dict[str, Any]], limit: int = 900) -> str:
    pieces = unique_nonempty((item["text"] for item in text_items))
    excerpt = " ".join(pieces)
    return truncate_text(excerpt, limit)


def slide_text_strings(text_items: list[dict[str, Any]], limit: int = 12) -> list[str]:
    return unique_nonempty((item["text"] for item in text_items), limit=limit)


def is_generic_slide_text(text: str, title: str = "") -> bool:
    cleaned = clean_source_prompt(text, title).lower()
    if not cleaned:
        return True
    title_clean = clean_source_prompt(title).lower()
    if title_clean and cleaned == title_clean:
        return True
    generic_labels = {
        "session 1",
        "session 2",
        "this or that",
        "is it reasonable",
        "is it reasonable?",
        "be curious",
        "mindset",
        "workspace",
        "collaborate and connect",
        "let's explore more",
        "let’s explore more",
        "learning targets",
        "summarize",
        "warm-up",
    }
    return cleaned in generic_labels


def is_low_value_vocabulary_term(term: str) -> bool:
    cleaned = display_text_key(term)
    if not cleaned:
        return True
    if cleaned in STOPWORDS or cleaned in GENERIC_VOCAB_TERMS or cleaned in LOW_VALUE_VOCAB_TERMS:
        return True
    words = cleaned.split()
    if len(words) >= 2 and words[0] in LESSON_LEAD_VERBS:
        return True
    if words and all(word in LOW_VALUE_VOCAB_TERMS for word in words):
        return True
    if cleaned.isdigit():
        return True
    return False


def is_placeholder_vocab_definition(definition: str) -> bool:
    cleaned = display_text_key(definition)
    return cleaned in {
        "",
        display_text_key(PLACEHOLDER_VOCAB_DEFINITION),
        "student friendly definition",
        "definition",
        "visual clue",
    }


def is_answer_like_text(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return False
    answer_starters = (
        "reveal:",
        "we can ",
        "we need to ",
        "the area is ",
        "the answer is ",
        "this means ",
        "this tells us ",
        "so, ",
        "therefore",
        "it is ",
    )
    return lowered.startswith(answer_starters)


def dedupe_problem_candidates(candidates: Iterable[str], limit: int = 6) -> list[str]:
    kept: list[str] = []
    kept_normalized: list[str] = []
    for candidate in unique_nonempty(candidates):
        normalized = normalize_whitespace(candidate)
        if not normalized:
            continue
        lowered = normalized.lower()
        if any(lowered == prior or lowered in prior or prior in lowered for prior in kept_normalized):
            continue
        kept.append(normalized)
        kept_normalized.append(lowered)
        if len(kept) >= limit:
            break
    return kept


def build_problem_candidates(text_items: list[dict[str, Any]], title: str, limit: int = 6) -> list[str]:
    raw_items = slide_text_strings(text_items, limit=16)
    if not raw_items:
        return []

    title_clean = normalize_whitespace(title)
    filtered_items = [item for item in raw_items if normalize_whitespace(item) != title_clean] or raw_items
    candidates: list[str] = []

    for index, raw_item in enumerate(filtered_items):
        cleaned_item = clean_source_prompt(raw_item, title_clean)
        if not cleaned_item or is_generic_slide_text(cleaned_item, title_clean):
            continue
        if not is_problem_like_text(cleaned_item):
            continue

        context_parts: list[str] = []
        scan = index - 1
        while scan >= 0 and len(context_parts) < 2:
            previous = clean_source_prompt(filtered_items[scan], title_clean)
            scan -= 1
            if not previous or is_generic_slide_text(previous, title_clean) or is_answer_like_text(previous):
                continue
            context_parts.insert(0, previous)

        combined = clean_source_prompt(" ".join(context_parts + [cleaned_item]), title_clean)
        if combined and not is_answer_like_text(combined):
            candidates.append(combined)
        candidates.append(cleaned_item)

    whole = clean_source_prompt(" ".join(filtered_items), title_clean)
    if whole and is_problem_like_text(whole) and not is_answer_like_text(whole):
        candidates.append(whole)

    return dedupe_problem_candidates((truncate_text(item, 420) for item in candidates), limit=limit)


def extract_slide_images(
    slide: Any,
    slide_number: int,
    output_dir: Path,
    assets_dir: Path,
) -> list[dict[str, Any]]:
    assets: list[dict[str, Any]] = []
    image_index = 0
    for shape in iter_leaf_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        image_index += 1
        try:
            image = shape.image
        except ValueError:
            # Some PPTX files contain picture placeholders or broken handles that
            # python-pptx reports as "no embedded image". Skip those so one bad
            # image object does not kill the entire notebook run.
            continue
        extension = (getattr(image, "ext", None) or "png").lower()
        asset_path = assets_dir / f"slide_{slide_number:03d}_img_{image_index}.{extension}"
        asset_path.write_bytes(image.blob)
        size = getattr(image, "size", (0, 0))
        assets.append(
            {
                "image_index": image_index,
                "slide_number": slide_number,
                "path": str(asset_path),
                "relative_path": str(asset_path.relative_to(output_dir)),
                "content_type": mimetypes.guess_type(asset_path.name)[0] or f"image/{extension}",
                "pixel_width": int(size[0]) if size else 0,
                "pixel_height": int(size[1]) if size else 0,
                "display_width": int(getattr(shape, "width", 0)),
                "display_height": int(getattr(shape, "height", 0)),
            }
        )
    return assets


def detect_lesson_title(prs: Presentation, slides: list[dict[str, Any]], source_pptx: Path) -> str:
    core_title = normalize_whitespace(getattr(prs.core_properties, "title", "") or "")
    if core_title and "delete this slide" not in core_title.lower() and "template" not in core_title.lower():
        return core_title
    for slide in slides[:6]:
        text = clean_source_prompt(slide["text"])
        match = re.search(r"Session\s+[12]\s+(.+)", text, re.IGNORECASE)
        if match:
            lesson_title = normalize_whitespace(match.group(1))
            if lesson_title:
                return lesson_title
    for slide in slides:
        title = normalize_whitespace(slide["title"])
        if not title:
            continue
        if title.lower().startswith("session "):
            continue
        if "delete this slide" in title.lower() or "template" in title.lower():
            continue
        if title:
            return slide["title"]
    return source_pptx.stem.replace("_", " ")


def extract_keyword_candidates(slides: list[dict[str, Any]], limit: int = 6) -> list[str]:
    counter: Counter[str] = Counter()
    for slide in slides:
        combined = f"{slide['title']} {slide['text']}"
        for word in re.findall(r"[A-Za-z][A-Za-z'-]{4,}", combined):
            lowered = word.lower()
            if lowered in STOPWORDS or lowered in GENERIC_VOCAB_TERMS:
                continue
            counter[lowered] += 1
    keywords: list[str] = []
    for word, _count in counter.most_common(limit * 4):
        title_word = word.title()
        if title_word not in keywords:
            keywords.append(title_word)
        if len(keywords) >= limit:
            break
    return keywords


def combined_source_text(source_records: list[dict[str, Any]]) -> str:
    parts: list[str] = []
    for slide in source_records:
        parts.extend(slide.get("problem_texts", []))
        parts.extend(slide.get("text_items", []))
        parts.extend([slide.get("title", ""), slide.get("text", ""), slide.get("notes", "")])
    return normalize_whitespace(" ".join(parts))


def display_term_label(term: str) -> str:
    cleaned = normalize_whitespace(term)
    if not cleaned:
        return ""
    return " ".join(part.capitalize() if part.lower() != "of" else "of" for part in cleaned.split())


def extract_source_deck(source_pptx: Path, output_dir: Path) -> tuple[dict[str, Any], Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    assets_dir = output_dir / "assets"
    assets_dir.mkdir(exist_ok=True)

    prs = Presentation(str(source_pptx))
    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        text_items = extract_text_items(slide)
        title = pick_slide_title(text_items)
        images = extract_slide_images(slide, slide_number, output_dir, assets_dir)
        slides.append(
            {
                "slide_number": slide_number,
                "title": title,
                "text": build_slide_excerpt(text_items),
                "text_items": slide_text_strings(text_items),
                "problem_texts": build_problem_candidates(text_items, title),
                "notes": extract_slide_notes(slide),
                "text_blocks": len(text_items),
                "image_count": len(images),
                "images": images,
            }
        )

    lesson_title = detect_lesson_title(prs, slides, source_pptx)
    deck = {
        "source_file": str(source_pptx.resolve()),
        "source_filename": source_pptx.name,
        "lesson_title": lesson_title,
        "slide_count": len(slides),
        "slides": slides,
        "keyword_candidates": extract_keyword_candidates(slides),
        "summary": truncate_text(
            " | ".join(
                unique_nonempty(
                    [lesson_title] + [slide["title"] for slide in slides if slide["title"]],
                    limit=10,
                )
            ),
            600,
        ),
    }
    output_path = output_dir / "source_deck.json"
    write_json(output_path, deck)
    return deck, output_path


def build_deck_digest(deck: dict[str, Any], max_slides: int = 48) -> str:
    lines = [
        f"Lesson title: {deck['lesson_title']}",
        f"Source filename: {deck['source_filename']}",
        f"Slide count: {deck['slide_count']}",
        f"Keyword candidates: {', '.join(deck['keyword_candidates']) or 'None'}",
        "",
        "Slide-by-slide extracted content:",
    ]
    total_chars = 0
    for slide in deck["slides"][:max_slides]:
        chunk_lines = [
            f"[Slide {slide['slide_number']}] Title: {slide['title'] or '(untitled)'}",
            f"Text: {truncate_text(slide['text'], 700) or '(no extracted text)'}",
        ]
        if slide.get("problem_texts"):
            chunk_lines.append(
                "Exact problem candidates: "
                + " | ".join(truncate_text(item, 260) for item in slide["problem_texts"][:2])
            )
        if slide["notes"]:
            chunk_lines.append(f"Notes: {truncate_text(slide['notes'], 320)}")
        if slide["image_count"]:
            chunk_lines.append(f"Embedded pictures: {slide['image_count']}")
        chunk = "\n".join(chunk_lines)
        total_chars += len(chunk)
        if total_chars > 30000:
            lines.append("... additional slides omitted to stay within prompt budget ...")
            break
        lines.append(chunk)
        lines.append("")
    omitted = deck["slide_count"] - min(deck["slide_count"], max_slides)
    if omitted > 0:
        lines.append(f"Omitted slide count from digest: {omitted}")
    return "\n".join(lines).strip()


def select_prompt_images(deck: dict[str, Any], limit: int = 6) -> list[dict[str, Any]]:
    by_slide: list[dict[str, Any]] = []
    for slide in deck["slides"]:
        if not slide["images"]:
            continue
        best = max(
            slide["images"],
            key=lambda asset: asset["pixel_width"] * asset["pixel_height"],
        )
        by_slide.append(best)
    if not by_slide:
        return []

    midpoint = max(1, deck["slide_count"] // 2)
    first_half = [asset for asset in by_slide if asset["slide_number"] <= midpoint]
    second_half = [asset for asset in by_slide if asset["slide_number"] > midpoint]
    first_half.sort(key=lambda asset: asset["pixel_width"] * asset["pixel_height"], reverse=True)
    second_half.sort(key=lambda asset: asset["pixel_width"] * asset["pixel_height"], reverse=True)
    selected = first_half[: limit // 2] + second_half[: limit - len(first_half[: limit // 2])]
    selected_numbers = {asset["slide_number"] for asset in selected}
    if len(selected) < limit:
        for asset in sorted(
            by_slide,
            key=lambda item: item["pixel_width"] * item["pixel_height"],
            reverse=True,
        ):
            if asset["slide_number"] in selected_numbers:
                continue
            selected.append(asset)
            selected_numbers.add(asset["slide_number"])
            if len(selected) >= limit:
                break
    return selected[:limit]


def image_to_data_url(asset: dict[str, Any]) -> str | None:
    image_path = Path(asset["path"])
    if not image_path.exists():
        return None
    if image_path.stat().st_size > 2_000_000:
        return None
    encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
    return f"data:{asset['content_type']};base64,{encoded}"


def notebook_plan_schema() -> dict[str, Any]:
    vocab_item = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "word": {"type": "string"},
            "definition": {"type": "string"},
            "example": {"type": "string"},
            "visual_cue": {"type": "string"},
        },
        "required": ["word", "definition", "example", "visual_cue"],
    }
    slide_item = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "kind": {"type": "string", "enum": SLIDE_KINDS},
            "section": {"type": "string"},
            "title": {"type": "string"},
            "subtitle": {"type": "string"},
            "primary_text": {"type": "string"},
            "secondary_text": {"type": "string"},
            "bullets": {"type": "array", "items": {"type": "string"}},
            "tasks": {"type": "array", "items": {"type": "string"}},
            "response_prompt": {"type": "string"},
            "sentence_starters": {"type": "array", "items": {"type": "string"}},
            "vocabulary": {"type": "array", "items": vocab_item},
            "activity_name": {"type": "string"},
            "activity_family": {"type": "string", "enum": ACTIVITY_FAMILY_OPTIONS},
            "activity_instructions": {"type": "string"},
            "movable_pieces": {"type": "array", "items": {"type": "string"}},
            "answer_check": {"type": "string"},
            "source_slide_numbers": {"type": "array", "items": {"type": "integer"}},
            "image_source_slide": {"type": "integer"},
            "image_caption": {"type": "string"},
        },
        "required": [
            "kind",
            "section",
            "title",
            "subtitle",
            "primary_text",
            "secondary_text",
            "bullets",
            "tasks",
            "response_prompt",
            "sentence_starters",
            "vocabulary",
            "activity_name",
            "activity_family",
            "activity_instructions",
            "movable_pieces",
            "answer_check",
            "source_slide_numbers",
            "image_source_slide",
            "image_caption",
        ],
    }
    session_item = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "session_title": {"type": "string"},
            "session_subtitle": {"type": "string"},
            "slides": {"type": "array", "items": slide_item},
        },
        "required": ["session_title", "session_subtitle", "slides"],
    }
    return {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "lesson_title": {"type": "string"},
            "subject": {"type": "string"},
            "grade_level": {"type": "string"},
            "standards": {"type": "array", "items": {"type": "string"}},
            "topic_summary": {"type": "string"},
            "session_1": session_item,
        },
        "required": [
            "lesson_title",
            "subject",
            "grade_level",
            "standards",
            "topic_summary",
            "session_1",
        ],
    }


def build_planner_instructions(custom_guidance: str = "") -> str:
    guidance = normalize_whitespace(custom_guidance)
    allowed = ", ".join(SLIDE_KINDS)
    instructions = (
        "You are designing a premium Session 1 student notebook from an extracted lesson slide deck. "
        "Preserve the lesson's real topic, sequence, academic language, visuals, worked examples, and practice when present. "
        "Do not replace the lesson with generic filler. Build only Session 1 and omit session_2 unless the user explicitly asks for it. "
        "Use only these slide kinds: "
        f"{allowed}. "
        "Return JSON only and follow the schema exactly. "
        "When information is missing or unclear, use an empty string, an empty array, or 0 for image_source_slide instead of inventing unsupported details. "
        "source_slide_numbers must reference actual source slides. "
        "If the source includes images that should anchor a slide, set image_source_slide to that source slide number. "
        "Keep the notebook student-friendly, polished, and scaffolded with sentence starters where helpful. "
        "This should feel like a flagship TpT-quality notebook that a teacher would feel confident selling for about $10-$15, not a plain worksheet. "
        f"{REFERENCE_WORKBOOK_BASELINE_GUIDANCE} "
        "Keep the locked notebook architecture intact through this exact six-slide flow: Objectives + Session Map, Be Curious, Vocabulary + Reference Tool, Guided Problem, Interactive Activity, Best-Fit Interactive Review. "
        "Session 1 should contain exactly 6 slides. "
        "Vocabulary slides should include a dedicated visual cue or image column whenever the source deck provides images that can support the terms. "
        "Do not use low-value context words as vocabulary when stronger academic or lesson-specific terms are available. "
        "Preserve the lesson's original problems, but transform only the strongest first problem arc into premium notebook pages so the deck feels polished instead of crowded. "
        "When the extracted deck includes Exact problem candidates, treat those as the highest-fidelity source language for notebook problem pages and lift them directly when used. "
        "Integrate the strongest written source problem directly into the guided problem and interactive pages instead of paraphrasing it away. "
        "If a problem-solving page has clear source problems available, use those source problems as the anchor text instead of generic directions or invented replacement problems. "
        "For a learning_target slide, primary_text must be the content objective and secondary_text must be the language objective. "
        "Each learning_target slide must preserve the actual lesson objective language from the source, include the standard code, and show today's four-step path: Be Curious, Vocabulary, Try It, Discuss. "
        "Any draggable label, checkmark, chip, or movable card must keep its text inside the same shape so it moves as one unit. Never float a separate text box over a draggable shape. "
        "A be_curious slide should combine Notice/Wonder response space, Writing Revolution sentence starters, and compact academic vocabulary support on the same page. "
        "The vocabulary slide must combine a Word | Definition | Example | Visual table with a reference tool below it. "
        "The guided problem slide must keep the source problem wording verbatim when available and include a blank reference table, a student work area, and a TWR frame. "
        "The interactive activity slide must match the lesson type: proportional relationships / rates -> drag-sort, equations / expressions -> error analysis, geometry / measurement -> Partner A/B, statistics / data -> error analysis using the graph or table. "
        "The closing review slide should use the strongest discussion-ready interactive review for the lesson content instead of generic closure filler. "
        "Build in Writing Revolution-style kernel sentence frames, not just loose words, especially on the be curious and first-problem slides. "
        "When a slide uses source problems, worked examples, or practice, keep the exact source problem language whenever possible instead of generic paraphrases. "
        "Keep the last three slides tightly anchored to the same source-problem arc so students can move from guided understanding into interactive practice and review. "
        "After the review page, students should continue the remaining practice in the book. "
        "Before finalizing, perform a detailed professional-publisher copy edit for grammar, punctuation, spacing, capitalization, and clarity while preserving lesson meaning and exact source problems. "
        "Keep text concise enough for large, highly readable student notebook fonts, preserve generous white space, and preserve clear spacing so text lines never feel blended together. "
        "Response areas should feel like premium write-in workspaces: prompt text embedded inside prompt cards, clear boxed writing spaces, and no handwriting lines or worksheet styling. "
        "Use named activities only when they clearly strengthen one of the six slides. Avoid filler interactions. "
        "Give the notebook stronger editorial design energy with image-led anchors when the source deck supports them and clean variation across the six pages. "
        "When source images exist, use them whenever they strengthen modeling or vocabulary understanding. "
        "Only use activity_name values that come from the provided activity library. "
        "activity_name should name the chosen activity, activity_family should match the activity family, activity_instructions should tell students how to use the movable pieces or interactive structure, movable_pieces should list draggable labels/cards/tiles/sticky-notes, and answer_check should explain how the student can verify or justify the result."
    )
    if guidance:
        instructions += f" Additional user guidance: {guidance}"
    return instructions


def build_responses_payload(
    deck: dict[str, Any],
    model: str,
    prompt_images: list[dict[str, Any]],
    custom_guidance: str = "",
) -> dict[str, Any]:
    activity_digest = activity_library_digest(deck)
    user_content: list[dict[str, Any]] = [
        {
            "type": "input_text",
            "text": (
                "Use the extracted source deck below to plan a reusable Session 1 notebook.\n\n"
                + build_deck_digest(deck)
            ),
        }
    ]
    if activity_digest:
        user_content.append({"type": "input_text", "text": activity_digest})
    for asset in prompt_images:
        data_url = image_to_data_url(asset)
        if not data_url:
            continue
        user_content.append(
            {
                "type": "input_text",
                "text": (
                    f"Representative source image from slide {asset['slide_number']} "
                    f"(image {asset['image_index']}). Use it only if it improves fidelity."
                ),
            }
        )
        user_content.append(
            {
                "type": "input_image",
                "image_url": data_url,
                "detail": "low",
            }
        )

    return {
        "model": model,
        "input": [
            {
                "role": "developer",
                "content": [{"type": "input_text", "text": build_planner_instructions(custom_guidance)}],
            },
            {
                "role": "user",
                "content": user_content,
            },
        ],
        "text": {
            "format": {
                "type": "json_schema",
                "name": "session1_notebook_plan",
                "strict": True,
                "schema": notebook_plan_schema(),
            }
        },
        "max_output_tokens": 8000,
    }


def call_openai(payload: dict[str, Any], api_key: str) -> dict[str, Any]:
    req = request.Request(
        f"{OPENAI_API_BASE}/responses",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    cafile = resolve_ssl_cafile()
    context = ssl.create_default_context(cafile=cafile) if cafile else ssl.create_default_context()
    try:
        with request.urlopen(req, timeout=OPENAI_TIMEOUT_SECONDS, context=context) as response:
            return json.loads(response.read().decode("utf-8"))
    except error.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"OpenAI API error {exc.code}: {truncate_text(body, 900)}") from exc
    except error.URLError as exc:
        raise RuntimeError(f"Could not reach the OpenAI API: {exc}") from exc


def extract_response_text(response_payload: dict[str, Any]) -> str:
    top_level = response_payload.get("output_text")
    if isinstance(top_level, str) and top_level.strip():
        return top_level.strip()

    text_parts: list[str] = []
    for item in response_payload.get("output", []):
        if item.get("type") != "message":
            continue
        for content in item.get("content", []):
            if content.get("type") == "output_text" and content.get("text"):
                text_parts.append(content["text"])
            if content.get("type") == "refusal":
                raise RuntimeError(f"Model refusal: {content.get('refusal', 'Unknown refusal')}")

    joined = "".join(text_parts).strip()
    if not joined:
        raise RuntimeError("The Responses API returned no output_text content.")
    return joined


def validate_plan(plan: dict[str, Any], deck: dict[str, Any] | None = None) -> dict[str, Any]:
    for key in ("lesson_title", "subject", "grade_level", "standards", "topic_summary", "session_1"):
        if key not in plan:
            raise RuntimeError(f"Notebook plan is missing required key: {key}")
    for session_key in planned_session_keys(plan):
        session = plan[session_key]
        if not session.get("slides"):
            raise RuntimeError(f"{session_key} contains no slides.")
        slide_count = len(session["slides"])
        min_slides, max_slides = session_slide_count_bounds(session)
        if slide_count < min_slides or slide_count > max_slides:
            raise RuntimeError(
                f"{session_key} must contain between {min_slides} and {max_slides} slides, found {slide_count}."
            )
        required_kinds = required_kinds_for_session(session)
        present_kinds = {slide.get("kind") for slide in session["slides"]}
        missing_kinds = sorted(required_kinds - present_kinds)
        if missing_kinds:
            raise RuntimeError(f"{session_key} is missing required locked-architecture slide kinds: {', '.join(missing_kinds)}.")
        required_roles = required_template_roles_for_session(session)
        if required_roles:
            role_signature = template_role_signature(session)
            if role_signature != required_roles:
                raise RuntimeError(f"{session_key} does not match the exact ESOL workbook template order.")
        objective_slides = [slide for slide in session["slides"] if slide.get("kind") == "learning_target"]
        if not objective_slides:
            raise RuntimeError(f"{session_key} must include a learning_target slide with separate content and language objectives.")
        for slide in objective_slides:
            if not normalize_whitespace(slide.get("primary_text", "")):
                raise RuntimeError(f"{session_key} has a learning_target slide without a content objective.")
            if not normalize_whitespace(slide.get("secondary_text", "")):
                raise RuntimeError(f"{session_key} has a learning_target slide without a language objective.")
            if not is_i_can_objective(slide.get("primary_text", "")):
                raise RuntimeError(f"{session_key} has a learning_target content objective that does not start with 'I can'.")
            if not is_i_can_objective(slide.get("secondary_text", "")):
                raise RuntimeError(f"{session_key} has a learning_target language objective that does not start with 'I can'.")
        practice_like = sum(1 for slide in session["slides"] if slide.get("kind") in {"practice", "challenge"})
        minimum_practice_slides = 1 if uses_exact_esol_template(session) else 2
        if practice_like < minimum_practice_slides:
            raise RuntimeError(f"{session_key} must include the required practice coverage.")
        premium_features = session.get("premium_features", [])
        if premium_features and not 2 <= len(premium_features) <= 4:
            raise RuntimeError(f"{session_key} must select between 2 and 4 premium decision-layer features.")
    run_publisher_copyedit_review(plan)
    run_formal_release_plan_review(plan, deck=deck)
    return plan


def generate_plan_with_openai(
    deck: dict[str, Any],
    output_dir: Path,
    model: str,
    custom_guidance: str = "",
    api_key: str = "",
) -> tuple[dict[str, Any], Path]:
    effective_guidance = enforce_runtime_quality_guidance(custom_guidance)
    load_dotenv()
    resolved_api_key = choose_api_key(api_key, os.getenv("OPENAI_API_KEY", ""))
    if not resolved_api_key:
        raise RuntimeError(
            "OPENAI_API_KEY is missing. Save it in the environment or in a local .env file."
        )

    payload = build_responses_payload(deck, model, select_prompt_images(deck), custom_guidance=effective_guidance)
    response_payload = call_openai(payload, resolved_api_key)
    write_json(output_dir / "openai_response.json", response_payload)
    plan_text = extract_response_text(response_payload)
    plan = validate_plan(
        enforce_plan_requirements(json.loads(plan_text), deck, custom_guidance=effective_guidance),
        deck=deck,
    )
    plan_path = output_dir / "notebook_plan.json"
    write_json(plan_path, plan)
    return plan, plan_path


def build_slide_plan(
    *,
    kind: str,
    title: str,
    subtitle: str = "",
    primary_text: str = "",
    secondary_text: str = "",
    bullets: list[str] | None = None,
    tasks: list[str] | None = None,
    response_prompt: str = "",
    sentence_starters: list[str] | None = None,
    vocabulary: list[dict[str, str]] | None = None,
    activity_name: str = "",
    activity_family: str = "",
    activity_instructions: str = "",
    movable_pieces: list[str] | None = None,
    answer_check: str = "",
    source_slide_numbers: list[int] | None = None,
    image_source_slide: int = 0,
    image_caption: str = "",
    section: str | None = None,
) -> dict[str, Any]:
    return {
        "kind": kind,
        "section": section or KIND_DEFAULT_SECTION[kind],
        "title": truncate_text(title, 110),
        "subtitle": truncate_text(subtitle, 220),
        "primary_text": truncate_text(primary_text, 750),
        "secondary_text": truncate_text(secondary_text, 500),
        "bullets": unique_nonempty(bullets or [], limit=5),
        "tasks": unique_nonempty(tasks or [], limit=4),
        "response_prompt": truncate_text(response_prompt, 300),
        "sentence_starters": unique_nonempty(sentence_starters or [], limit=5),
        "vocabulary": (vocabulary or [])[:6],
        "activity_name": truncate_text(activity_name, 120),
        "activity_family": activity_family if activity_family in ACTIVITY_FAMILY_OPTIONS else "",
        "activity_instructions": truncate_text(activity_instructions, 280),
        "movable_pieces": unique_nonempty(movable_pieces or [], limit=6),
        "answer_check": truncate_text(answer_check, 220),
        "source_slide_numbers": source_slide_numbers or [],
        "image_source_slide": image_source_slide,
        "image_caption": truncate_text(image_caption, 180),
        "premium_layout": "",
        "premium_title": "",
        "premium_text": "",
        "premium_items": [],
        "premium_table": [],
        "partner_prompt": "",
        "discussion_questions": [],
        "context_anchor": "",
        "practice_phase": "",
        "flagship_activity_mode": FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED,
        "flagship_activity": {},
    }


def with_template_role(slide: dict[str, Any], role: str) -> dict[str, Any]:
    slide["template_role"] = role
    return slide


def with_template_metadata(
    slide: dict[str, Any],
    role: str,
    *,
    family: str = EXACT_ESOL_TEMPLATE_FAMILY,
) -> dict[str, Any]:
    slide["template_role"] = role
    slide["template_family"] = family
    return slide


def uses_reference_workbook_template(session: dict[str, Any]) -> bool:
    if session.get("template_family") in {"uploaded_model_session2", EXACT_ESOL_TEMPLATE_FAMILY}:
        return True
    return any(normalize_whitespace(slide.get("template_role", "")) for slide in session.get("slides", []))


def uses_exact_esol_template(session: dict[str, Any]) -> bool:
    return session.get("template_family") == EXACT_ESOL_TEMPLATE_FAMILY


def template_role_signature(session: dict[str, Any]) -> list[tuple[str, str]]:
    signature: list[tuple[str, str]] = []
    for slide in session.get("slides", []):
        signature.append((slide.get("kind", ""), normalize_whitespace(slide.get("template_role", ""))))
    return signature


def required_kinds_for_session(session: dict[str, Any]) -> set[str]:
    if uses_exact_esol_template(session):
        return {"learning_target", "be_curious", "vocabulary", "worked_example", "practice"}
    if uses_reference_workbook_template(session):
        return {"cover", "be_curious", "learning_target", "vocabulary", "worked_example", "practice", "reflection", "exit_ticket"}
    return {"cover", "be_curious", "learning_target", "vocabulary", "guided_notes", "worked_example", "practice", "reflection", "exit_ticket"}


def required_template_roles_for_session(session: dict[str, Any]) -> list[tuple[str, str]]:
    if uses_exact_esol_template(session):
        return EXACT_ESOL_TEMPLATE_SEQUENCE
    return []


def opening_slide_allowed(session: dict[str, Any], plan_slide: dict[str, Any]) -> bool:
    if uses_exact_esol_template(session):
        return (
            plan_slide.get("kind") == "learning_target"
            and normalize_whitespace(plan_slide.get("template_role", "")) == "learning_objectives"
        )
    return plan_slide.get("kind") == "cover"


def closing_slide_allowed(session: dict[str, Any], plan_slide: dict[str, Any]) -> bool:
    if uses_exact_esol_template(session):
        return (
            plan_slide.get("kind") == "practice"
            and normalize_whitespace(plan_slide.get("template_role", "")) == "best_fit_review"
        )
    if plan_slide.get("kind") == "exit_ticket":
        return True
    return plan_slide.get("kind") == "reflection" and normalize_whitespace(plan_slide.get("template_role", "")) == "goal_tracker"


def pick_first_image_slide(slides: list[dict[str, Any]]) -> int:
    for slide in slides:
        if slide["image_count"]:
            return slide["slide_number"]
    return 0


def source_slide_score(slide: dict[str, Any]) -> int:
    title = normalize_whitespace(slide.get("title", "")).lower()
    text = clean_source_prompt(slide.get("text", "")).lower()
    score = max(len(text) // 5, 0)
    if "?" in text or "question:" in text or "think about it:" in text:
        score += 18
    if slide.get("problem_texts"):
        score += 12 * min(len(slide["problem_texts"]), 2)
    if any(term in text for term in ("area", "decompose", "triangle", "polygon", "flag", "flooring", "budget", "congruent")):
        score += 20
    if any(term in title for term in ("stop sign", "flags", "flooring", "explore")):
        score += 34
    if slide.get("image_count"):
        score += 4
    if any(term in title for term in ("session", "this or that", "is it reasonable")):
        score -= 90
    if "be curious" in title:
        score -= 50
    if "mindset" in text:
        score -= 20
    return score


def problem_source_score(slide: dict[str, Any]) -> int:
    title = normalize_whitespace(slide.get("title", "")).lower()
    blob = normalize_whitespace(
        " ".join(
            slide.get("problem_texts", [])
            + slide.get("text_items", [])
            + [slide.get("text", ""), slide.get("notes", "")]
        )
    ).lower()
    prompt_hits = len(source_problem_candidates([slide], limit=4))
    fact_hits = len(source_fact_candidates([slide], limit=3))
    score = max(source_slide_score(slide), 0)
    score += 18 * min(prompt_hits, 3)
    score += 8 * min(fact_hits, 2)
    if "?" in blob or "question:" in blob:
        score += 16
    if any(
        term in blob
        for term in (
            "formula",
            "volume",
            "dimensions",
            "length",
            "width",
            "height",
            "missing",
            "unknown",
            "compare",
            "explain",
            "precision",
            "units",
        )
    ):
        score += 10
    if any(term in blob for term in ("question:", "which size", "how many", "regulations", "what information do you know")):
        score += 8
    if any(
        marker in title or marker in blob
        for marker in ("session 1", "session 2", "this or that", "be curious", "learning target", "learning targets")
    ):
        score -= 90
    if "i can " in blob:
        score -= 120
    if "summarize" in title:
        score -= 28
    if "workspace" in blob and prompt_hits == 0:
        score -= 28
    if "what do you notice" in blob and "what do you wonder" in blob:
        score -= 36
    return score


def problem_window_candidates(session_sources: list[dict[str, Any]], window_size: int = 4) -> list[dict[str, Any]]:
    candidates: list[dict[str, Any]] = []
    seen_signatures: set[tuple[int, ...]] = set()
    max_size = max(1, window_size)
    for start in range(len(session_sources)):
        for size in range(1, max_size + 1):
            window = session_sources[start : start + size]
            if not window:
                continue
            source_numbers = tuple(slide["slide_number"] for slide in window if slide.get("slide_number"))
            if not source_numbers or source_numbers in seen_signatures:
                continue
            prompt_hits = source_problem_candidates(window, limit=4)
            fact_hits = source_fact_candidates(window, limit=4)
            if not prompt_hits and not fact_hits:
                continue
            score = sum(max(problem_source_score(slide), 0) for slide in window)
            score += 16 * min(len(prompt_hits), 3)
            score += 6 * min(len(fact_hits), 2)
            if any(int(slide.get("image_count", 0) or 0) > 0 for slide in window):
                score += 4
            blob = combined_source_text(window).lower()
            if not prompt_hits:
                if any(term in blob for term in ("learning target", "learning targets", "i can ")):
                    continue
                if len(fact_hits) < 2 and not any(
                    term in blob for term in ("formula", "compare", "precision", "regulations", "dimensions", "reveal:")
                ):
                    continue
            if any(term in blob for term in ("would you rather", "organizing yourself", "mindset")):
                score -= 24
            candidates.append(
                {
                    "slides": window,
                    "source_numbers": list(source_numbers),
                    "start_index": start,
                    "score": score,
                    "text_blob": blob,
                }
            )
            seen_signatures.add(source_numbers)
    candidates.sort(key=lambda item: (item["start_index"], -item["score"], len(item["source_numbers"])))
    return candidates


def problem_window_role_score(window: dict[str, Any], *, role: str, ordinal: int, total_windows: int) -> int:
    text_blob = window.get("text_blob", "")
    start_index = int(window.get("start_index", 0) or 0)
    progress = start_index / max(total_windows - 1, 1)
    score = int(window.get("score", 0) or 0)

    if role == "worked_example":
        score += max(0, 16 - start_index * 2)
        if any(term in text_blob for term in ("how can", "how do you use", "what is the volume", "use a formula", "determine")):
            score += 12
        if any(
            term in text_blob
            for term in ("unit cube", "unit cubes", "substitute", "dimensions of the rectangular prism", "find the volume")
        ):
            score += 12
        if "relate to area" in text_blob and not any(
            term in text_blob for term in ("what is the volume", "use a formula", "unit cube", "unit cubes")
        ):
            score -= 14
        if any(term in text_blob for term in ("would you rather", "what do you notice", "what do you wonder")):
            score -= 28
    elif role == "practice":
        score += int(progress * 10) + ordinal * 4
        if any(
            term in text_blob
            for term in (
                "what is the volume",
                "what information do you know",
                "use a formula",
                "divide the volume",
                "determine",
                "how do you use",
                "unit cube",
                "unit cubes",
                "substitute",
            )
        ):
            score += 10
        if "relate to area" in text_blob and not any(
            term in text_blob for term in ("what is the volume", "use a formula", "unit cube", "unit cubes")
        ):
            score -= 12
    elif role == "challenge":
        score += int(progress * 14)
        if any(
            term in text_blob
            for term in ("question:", "which size", "how many", "compare", "regulations", "due to", "apply:", "buy")
        ):
            score += 22
    elif role == "exit_ticket":
        score += int(progress * 16)
        if any(term in text_blob for term in ("question:", "explain", "precision", "compare", "which size", "how many")):
            score += 18
    return score


def select_problem_window_for_slide(
    candidates: list[dict[str, Any]],
    *,
    role: str,
    ordinal: int,
    used_signatures: set[tuple[int, ...]],
    min_start_index: int = 0,
) -> dict[str, Any] | None:
    pool = [item for item in candidates if item["start_index"] >= min_start_index] or candidates
    if not pool:
        return None
    total_windows = max(len(candidates), 1)
    scoring_key = lambda item: (
        problem_window_role_score(item, role=role, ordinal=ordinal, total_windows=total_windows),
        -item["start_index"] if role in {"challenge", "exit_ticket"} else item["start_index"],
    )
    best_any = max(pool, key=scoring_key)
    unused_pool = [item for item in pool if tuple(item["source_numbers"]) not in used_signatures]
    if not unused_pool:
        return best_any
    best_unused = max(unused_pool, key=scoring_key)
    best_any_score = problem_window_role_score(best_any, role=role, ordinal=ordinal, total_windows=total_windows)
    best_unused_score = problem_window_role_score(best_unused, role=role, ordinal=ordinal, total_windows=total_windows)
    if best_unused_score >= best_any_score - 10:
        return best_unused
    return best_any


def core_instructional_slides(slides: list[dict[str, Any]], limit: int = 8) -> list[dict[str, Any]]:
    if not slides:
        return []
    ranked = sorted(slides, key=lambda slide: (-source_slide_score(slide), slide["slide_number"]))
    selected_numbers = {slide["slide_number"] for slide in ranked[:limit] if source_slide_score(slide) > 0}
    if not selected_numbers:
        return slides[:limit]
    ordered = [slide for slide in slides if slide["slide_number"] in selected_numbers]
    return ordered or slides[:limit]


def instructional_focus_numbers(
    deck: dict[str, Any],
    candidate_numbers: list[int],
    fallback_numbers: list[int],
    *,
    limit: int = 3,
) -> list[int]:
    for numbers in (candidate_numbers, fallback_numbers):
        if not numbers:
            continue
        records = source_slides_from_numbers(deck, numbers)
        filtered = [slide["slide_number"] for slide in records if source_slide_score(slide) > 0]
        selected = filtered or [number for number in numbers if isinstance(number, int)]
        if selected:
            return selected[:limit]
    return []


def slide_title_list(slides: list[dict[str, Any]], limit: int = 4) -> list[str]:
    titles = [slide["title"] for slide in slides if slide["title"]]
    return unique_nonempty(titles, limit=limit)


def slide_task_list(slides: list[dict[str, Any]], limit: int = 3) -> list[str]:
    exact_tasks = source_problem_candidates(slides, limit=limit)
    if exact_tasks:
        return [truncate_text(item, 190) for item in exact_tasks[:limit]]
    tasks: list[str] = []
    for slide in slides:
        source = clean_source_prompt(slide["text"] or slide["title"], slide.get("title", ""))
        if not source:
            continue
        if is_generic_slide_text(source, slide.get("title", "")):
            continue
        tasks.append(truncate_text(source, 170))
        if len(tasks) >= limit:
            break
    return tasks


def should_strip_title_prefix(text: str, title: str) -> bool:
    cleaned = normalize_whitespace(text)
    title_clean = normalize_whitespace(title)
    if not cleaned or not title_clean or not cleaned.lower().startswith(title_clean.lower()):
        return False
    suffix = cleaned[len(title_clean) :]
    if not suffix:
        return False
    if suffix[:1] in {":", "-", "|"}:
        return True
    remainder = normalize_whitespace(suffix)
    if not remainder:
        return False
    remainder_lower = remainder.lower()
    title_strip_starters = PROMPT_STARTERS + (
        "step ",
        "precision ",
        "same ",
        "different ",
        "area of ",
        "think about it",
        "let's ",
        "let’s ",
        "analyze ",
        "analyzing ",
        "using symbols appropriately",
        "workspace",
        "mindset",
    )
    if remainder_lower.startswith(title_strip_starters):
        return True
    remainder_first = re.sub(r"^[^a-z0-9]+|[^a-z0-9]+$", "", remainder_lower.split(" ", 1)[0])
    if len(title_clean.split()) <= 4 and (
        remainder_first.isdigit()
        or remainder_first in NUMBER_WORD_STARTERS
    ):
        return True
    return False


def clean_source_prompt(text: str, title: str = "") -> str:
    cleaned = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text or "")
    cleaned = normalize_whitespace(cleaned)
    if not cleaned:
        return ""
    title_clean = normalize_whitespace(title)
    if title_clean and should_strip_title_prefix(cleaned, title_clean):
        cleaned = normalize_whitespace(cleaned[len(title_clean) :].lstrip(" :|-"))
    cleaned = re.sub(r"\s*Reveal:\s*.*$", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+(Mindset|Workspace|Analyzing)\s*$", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"([?!])(?:\s*[?!])+$", r"\1", cleaned)
    cleaned = re.sub(r"\?+\.+$", "?", cleaned)
    cleaned = re.sub(r"!+\.+$", "!", cleaned)
    return cleaned.strip(" -")


def is_problem_like_text(text: str) -> bool:
    lowered = clean_source_prompt(text).lower()
    if not lowered:
        return False
    if "?" in lowered:
        return True
    return lowered.startswith(PROMPT_STARTERS) or any(f" {starter}" in lowered for starter in PROMPT_STARTERS)


def is_dense_numeric_table_text(text: str) -> bool:
    cleaned = normalize_whitespace(text)
    if not cleaned:
        return False
    if cleaned.count("|") >= 5:
        return True
    numeric_count = len(re.findall(r"(?<![A-Za-z])\d[\d,]*(?:\.\d+)?", cleaned))
    alpha_count = len(re.findall(r"[A-Za-z]+", cleaned))
    if numeric_count >= 8 and alpha_count == 0:
        return True
    if numeric_count >= 14 and alpha_count <= 18:
        return True
    if numeric_count >= 10 and alpha_count > 0 and numeric_count > alpha_count * 2:
        return True
    return False


def source_problem_candidates(slides: list[dict[str, Any]], limit: int = 4) -> list[str]:
    question_like: list[str] = []
    statement_like: list[str] = []
    for slide in slides:
        title_candidate = clean_source_prompt(slide.get("title", ""))
        if title_candidate and is_problem_like_text(title_candidate) and not is_generic_slide_text(title_candidate, ""):
            question_like.append(title_candidate)
        if slide.get("problem_texts"):
            for candidate in slide.get("problem_texts", []):
                cleaned_candidate = clean_source_prompt(candidate, slide.get("title", ""))
                if re.search(r"question:\s*", cleaned_candidate, flags=re.IGNORECASE):
                    cleaned_candidate = re.split(r"question:\s*", cleaned_candidate, flags=re.IGNORECASE)[-1]
                lowered = normalize_whitespace(cleaned_candidate).lower()
                if lowered.startswith("i can ") or lowered.count("i can ") >= 2:
                    continue
                if not cleaned_candidate:
                    continue
                if is_dense_numeric_table_text(cleaned_candidate):
                    continue
                if is_generic_slide_text(cleaned_candidate, slide.get("title", "")):
                    continue
                target = question_like if is_problem_like_text(cleaned_candidate) or "question:" in lowered else statement_like
                target.append(cleaned_candidate)
            continue
        whole = clean_source_prompt(slide.get("text", "") or slide.get("title", ""), slide.get("title", ""))
        if not whole:
            continue
        if whole.lower().startswith("i can ") or whole.lower().count("i can ") >= 2:
            continue
        if is_dense_numeric_table_text(whole):
            continue
        if is_generic_slide_text(whole, slide.get("title", "")):
            continue
        if is_problem_like_text(whole):
            question_like.append(whole)
        for sentence in re.split(r"(?<=[?.!])\s+", whole):
            cleaned = clean_source_prompt(sentence)
            if len(cleaned) < 8:
                continue
            if cleaned.lower().startswith("i can ") or cleaned.lower().count("i can ") >= 2:
                continue
            if is_dense_numeric_table_text(cleaned):
                continue
            if is_generic_slide_text(cleaned, slide.get("title", "")):
                continue
            if is_problem_like_text(cleaned):
                question_like.append(cleaned)
            elif cleaned:
                statement_like.append(cleaned)
    ordered_questions = sorted(
        question_like,
        key=lambda text: (-problem_prompt_priority(text), len(normalize_whitespace(text)), normalize_whitespace(text).lower()),
    )
    ordered_statements = sorted(
        statement_like,
        key=lambda text: (-problem_prompt_priority(text), len(normalize_whitespace(text)), normalize_whitespace(text).lower()),
    )
    return dedupe_problem_candidates(ordered_questions + ordered_statements, limit=limit)


def source_problem_text_overlap(text: str, source_candidates: list[str]) -> bool:
    text_key = display_text_key(text)
    if not text_key:
        return False
    text_words = {
        word
        for word in re.findall(r"[a-z0-9]+", text_key)
        if len(word) > 2 and word not in STOPWORDS and word not in GENERIC_VOCAB_TERMS
    }
    for candidate in source_candidates:
        candidate_key = display_text_key(candidate)
        if not candidate_key:
            continue
        if candidate_key in text_key or text_key in candidate_key:
            return True
        candidate_words = {
            word
            for word in re.findall(r"[a-z0-9]+", candidate_key)
            if len(word) > 2 and word not in STOPWORDS and word not in GENERIC_VOCAB_TERMS
        }
        if len(text_words & candidate_words) >= 4:
            return True
    return False


def problem_prompt_priority(text: str) -> int:
    lowered = clean_source_prompt(text).lower()
    if not lowered:
        return 0
    score = 0
    if "?" in lowered or "question:" in lowered:
        score += 3
    if any(
        term in lowered
        for term in (
            "what is the volume",
            "what is the length",
            "which size",
            "how many boxes",
            "how does the length of the box compare",
            "compare to the shipping regulations",
        )
    ):
        score += 16
    if any(
        term in lowered
        for term in (
            "how do you use the formula",
            "use a formula",
            "determine the dimensions",
            "determine the length",
            "how can you determine",
            "how can you use this information to find the volume",
            "unit cube",
            "unit cubes",
            "substitute",
            "what information do you know",
            "find the volume",
        )
    ):
        score += 12
    if any(
        term in lowered
        for term in (
            "median",
            "mean",
            "data set",
            "box plot",
            "box-and-whisker",
            "quartile",
            "distribution",
            "dot plot",
            "line plot",
            "slowest and fastest",
            "compare between",
            "which numbers",
            "between which two values",
        )
    ):
        score += 14
    if any(
        term in lowered
        for term in (
            "to find the volume of a three-dimensional rectangular prism",
            "you can use the volume formulas you know",
            "let's use the dimensions of the unit cube",
            "let’s use the dimensions of the unit cube",
            "let's substitute the dimensions",
            "let’s substitute the dimensions",
        )
    ):
        score -= 6
    if any(term in lowered for term in ("relate to area", "used before", "how is finding area related")):
        score -= 8
    return score


LEARNING_TARGET_MARKERS = (
    "learning target",
    "learning targets",
    "objective",
    "objectives",
    "i can",
    "we will",
    "success criteria",
)


def learning_target_source_slides(deck: dict[str, Any]) -> list[dict[str, Any]]:
    candidates: list[dict[str, Any]] = []
    for slide in deck.get("slides", []):
        blob = normalize_whitespace(f"{slide.get('title', '')} {slide.get('text', '')}").lower()
        if any(marker in blob for marker in LEARNING_TARGET_MARKERS):
            candidates.append(slide)
    return candidates


def extract_learning_target_statements(deck: dict[str, Any]) -> list[str]:
    statements: list[str] = []
    pattern = re.compile(r"(?:I can|We will|Students will)[^.?!]*[.?!]?", re.IGNORECASE)
    for slide in learning_target_source_slides(deck):
        for item in slide.get("text_items", []) or [slide.get("text", "")]:
            cleaned = clean_source_prompt(item, slide.get("title", ""))
            if not cleaned:
                continue
            matches = pattern.findall(cleaned)
            if matches:
                statements.extend(matches)
    return unique_nonempty((normalize_whitespace(text).rstrip(".") + "." for text in statements), limit=6)


def compact_learning_target_statement(text: str) -> str:
    cleaned = normalize_whitespace(text)
    replacements = (
        (
            "decomposing into triangles or composing into a rectangle",
            "decomposing it or composing it into a rectangle",
        ),
        (
            "make use of structure to find the area of a trapezoid by composing it into a rectangle or by using a formula",
            "use structure to find the area of a trapezoid with a rectangle or a formula",
        ),
    )
    lowered = cleaned.lower()
    for source, target in replacements:
        if source in lowered:
            start = lowered.index(source)
            cleaned = cleaned[:start] + target + cleaned[start + len(source):]
            break
    return normalize_i_can_objective(cleaned)


def is_i_can_objective(text: str) -> bool:
    return bool(re.match(r"^\s*i can\b", normalize_whitespace(text), flags=re.IGNORECASE))


def normalize_i_can_objective(text: str, *, fallback: str = "") -> str:
    cleaned = normalize_whitespace(text or fallback).rstrip(".")
    if not cleaned:
        return ""
    cleaned = re.sub(r"^(?:The objective is to|Objective:)\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"^(?:Students will|We will)\s+", "I can ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"^I can\s+I can\s+", "I can ", cleaned, flags=re.IGNORECASE)
    if not is_i_can_objective(cleaned):
        cleaned = re.sub(r"^to\s+", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"^be able to\s+", "", cleaned, flags=re.IGNORECASE)
        cleaned = cleaned[:1].lower() + cleaned[1:] if cleaned else cleaned
        cleaned = f"I can {cleaned}"
    return normalize_whitespace(cleaned).rstrip(".") + "."


def select_session_objective(deck: dict[str, Any], session_key: str) -> str:
    objectives = extract_learning_target_statements(deck)
    if not objectives:
        return f"I can explain and apply the key ideas from {deck['lesson_title']}."
    if session_key == "session_1":
        return normalize_i_can_objective(compact_learning_target_statement(objectives[0]))
    return normalize_i_can_objective(compact_learning_target_statement(objectives[-1]))


def cover_focus_statement(deck: dict[str, Any], session_key: str, reference_text: str = "") -> str:
    focus = lesson_focus_phrase(deck, reference_text)
    lowered = normalize_whitespace(reference_text or deck.get("lesson_title", "")).lower()
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return "Build the volume formula from a visual model and use dimensions and units precisely."
    if "volume" in lowered:
        return "Use the model, dimensions, and units to explain how the volume is found."
    if "surface area" in lowered:
        return "Use nets, dimensions, and units to explain how the surface area is found."
    if "graph" in lowered or "table" in lowered or "equation" in lowered:
        return "Connect the representations and explain how they show the same relationship."
    if "formula" in lowered:
        return "Connect the model, numbers, and formula so the reasoning stays visible."
    if "area" in lowered and "trapezoid" in lowered:
        return "Use structure, decomposition, and clear labels to explain how the area is found."
    if "area" in lowered:
        return "Use the figure, measurements, and labels to explain how the area is found."
    if session_key == "session_1":
        return f"Launch and build understanding of {focus} with carefully sequenced notes, visuals, and practice."
    return f"Review, apply, and extend {focus} with carefully sequenced practice and reflection."


def default_cover_subtitle(deck: dict[str, Any], session_key: str) -> str:
    focus = lesson_focus_phrase(deck, select_session_objective(deck, session_key))
    if session_key == "session_1":
        return f"Launch and concept building for {focus}."
    return f"Review, application, and reflection for {focus}."


def ensure_cover_slide(deck: dict[str, Any], session: dict[str, Any], plan_slide: dict[str, Any], session_key: str) -> None:
    session_label = "Session 1" if session_key == "session_1" else "Session 2"
    session_title = session.get("session_title", "")
    session_subtitle = session.get("session_subtitle", "")
    objective_text = select_session_objective(deck, session_key)
    focus_text = cover_focus_statement(deck, session_key, objective_text)
    standard_text = first_standard_text(deck)

    plan_slide["subtitle"] = first_distinct_text(
        [plan_slide.get("subtitle", ""), session_subtitle, focus_text, default_cover_subtitle(deck, session_key)],
        excluded=[plan_slide.get("title", ""), session_label, session_title],
    ) or default_cover_subtitle(deck, session_key)
    plan_slide["primary_text"] = first_distinct_text(
        [plan_slide.get("primary_text", ""), focus_text, plan_slide.get("secondary_text", ""), objective_text],
        excluded=[plan_slide.get("title", ""), session_label, plan_slide["subtitle"], session_title],
    ) or focus_text
    secondary_default = standard_text or (
        "Use the source problems, visuals, and vocabulary to show clear mathematical thinking."
        if session_key == "session_1"
        else "Use practice, explanation, and reflection to demonstrate what you understand independently."
    )
    plan_slide["secondary_text"] = first_distinct_text(
        [plan_slide.get("secondary_text", ""), standard_text, secondary_default],
        excluded=[plan_slide.get("title", ""), session_label, session_title, plan_slide["subtitle"], plan_slide["primary_text"]],
    ) or secondary_default
    if not plan_slide.get("bullets"):
        plan_slide["bullets"] = source_titles_from_numbers(deck, plan_slide.get("source_slide_numbers", []), limit=3)


def derive_language_objective(deck: dict[str, Any], content_objective: str) -> str:
    lowered = normalize_whitespace(content_objective).lower()
    shape_context = objective_shape_context(content_objective)
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return "I can explain how the dimensions and formula show the volume of a rectangular prism using labels, units, and complete sentences."
    if "volume" in lowered:
        return "I can explain how the model and formula show the volume using labels, units, and complete sentences."
    if "surface area" in lowered:
        return "I can explain how the net, dimensions, and units show the surface area using precise vocabulary and complete sentences."
    if "graph" in lowered or "table" in lowered or "equation" in lowered:
        return "I can explain how the table, graph, or equation show the same relationship using vocabulary and complete sentences."
    if "regular polygon" in lowered or "octagon" in lowered:
        return "I can explain how decomposing the regular polygon into congruent triangles helps me find the total area using labels, measurements, and complete sentences."
    if "decompos" in lowered or "compos" in lowered or "triangle" in lowered:
        if shape_context == "triangle":
            return "I can explain how the base and height show the area of the triangle using labels, vocabulary, and complete sentences."
        if shape_context == "composite figure":
            return "I can explain how I decomposed the composite figure and why the parts show the total area using labels, vocabulary, and complete sentences."
        return f"I can explain how I decomposed or composed the {shape_context} using labels, vocabulary, and complete sentences."
    if "formula" in lowered or "structure" in lowered or "rectangle" in lowered:
        return "I can explain how the model or formula shows the math idea using labels, vocabulary, and complete sentences."
    if "area" in lowered:
        return "I can explain how I found the area using labels, vocabulary, and complete sentences."
    return f"I can explain my strategy for {lesson_focus_phrase(deck, content_objective)} using lesson vocabulary, labels, and complete sentences."


def objective_shape_context(content_objective: str) -> str:
    lowered = normalize_whitespace(content_objective).lower()
    if "regular polygon" in lowered or "octagon" in lowered:
        return "regular polygon"
    if "composite figure" in lowered:
        return "composite figure"
    for shape in ("trapezoid", "triangle", "parallelogram", "rhombus", "rectangle"):
        if shape in lowered:
            return shape
    if "figure" in lowered:
        return "figure"
    return "figure"


def objective_success_criteria(content_objective: str, language_objective: str) -> list[str]:
    lowered = normalize_whitespace(content_objective).lower()
    shape_context = objective_shape_context(content_objective)
    criteria: list[str] = []
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        criteria.append("I can identify the length, width, and height I need.")
        criteria.append("I can connect the model or formula to the prism's volume.")
    elif "volume" in lowered:
        criteria.append("I can identify the measurements I need to find volume.")
        criteria.append("I can connect the model or formula to the volume.")
    elif "surface area" in lowered:
        criteria.append("I can identify the faces or dimensions I need.")
        criteria.append("I can use the net or measurements to find the surface area.")
    elif "graph" in lowered or "table" in lowered or "equation" in lowered:
        criteria.append("I can connect the representations to the same relationship.")
        criteria.append("I can use the correct values, labels, or features to support my explanation.")
    if "regular polygon" in lowered or "octagon" in lowered:
        criteria.append("I can connect the number of sides to the number of congruent triangles.")
        criteria.append("I can find the area of one triangle and use it to find the total area.")
    elif "decompos" in lowered or "compos" in lowered or "triangle" in lowered:
        if shape_context == "triangle":
            criteria.append("I can identify the base and height that show the area of the triangle.")
        elif shape_context == "composite figure":
            criteria.append("I can break apart the composite figure into helpful shapes.")
        else:
            criteria.append(f"I can break apart or rearrange the {shape_context} in a helpful way.")
    if "formula" in lowered or "structure" in lowered:
        criteria.append("I can connect the model to the formula or structure.")
    elif "rectangle" in lowered:
        criteria.append("I can connect the model to a rectangle or rectangular structure.")
    if "area" in lowered:
        criteria.append("I can identify the measurements I need to find area.")
    criteria.append("I can explain my strategy clearly with labels and vocabulary.")
    if "complete sentences" not in language_objective.lower():
        criteria.append("I can explain my thinking clearly.")
    return unique_nonempty(criteria, limit=3)


def objective_sentence_frames(content_objective: str) -> list[str]:
    lowered = normalize_whitespace(content_objective).lower()
    shape_context = objective_shape_context(content_objective)
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return [
            "The dimensions that matter are ___.",
            "I can find the volume by ___.",
            "The units are ___ because ___.",
        ]
    if "volume" in lowered:
        return [
            "The measurements I need are ___.",
            "The model or formula shows volume because ___.",
            "My answer is reasonable because ___.",
        ]
    if "surface area" in lowered:
        return [
            "The faces or parts that matter are ___.",
            "I found the surface area by ___.",
            "The units make sense because ___.",
        ]
    if "graph" in lowered or "table" in lowered or "equation" in lowered:
        return [
            "The table, graph, and equation connect because ___.",
            "One feature that matters is ___.",
            "I can justify my answer by ___.",
        ]
    if "regular polygon" in lowered or "octagon" in lowered:
        return [
            "I decomposed the regular polygon into ___ congruent triangles.",
            "The base and height of one triangle are ___ and ___.",
            "I found the total area by ___.",
        ]
    if "decompos" in lowered or "compos" in lowered or "triangle" in lowered:
        shape_phrase = "figure" if shape_context == "figure" else shape_context
        return [
            f"I decomposed or composed the {shape_phrase} by ___.",
            "The measurements that matter are ___ and ___.",
            "My strategy works because ___.",
        ]
    if "formula" in lowered or "rectangle" in lowered or "structure" in lowered:
        return [
            "I used the rectangle or formula because ___.",
            "The measurements that matter are ___ and ___.",
            "My strategy makes sense because ___.",
        ]
    return [
        "The lesson is asking me to ___.",
        "The measurements that matter are ___.",
        "I can explain my strategy by saying ___.",
    ]


def best_source_prompt_text(slides: list[dict[str, Any]], limit: int = 360, *, prefer_last: bool = False) -> str:
    ordered_slides = list(reversed(slides)) if prefer_last else slides
    candidates = source_problem_candidates(ordered_slides, limit=5)
    if candidates:
        ranked = sorted(
            enumerate(candidates),
            key=lambda item: (-problem_prompt_priority(item[1]), item[0]),
        )
        return trim_dangling_display_text(truncate_display_copy(candidates[ranked[0][0]], limit))
    for slide in ordered_slides:
        cleaned = clean_source_prompt(slide.get("text", "") or slide.get("title", ""), slide.get("title", ""))
        if cleaned:
            return trim_dangling_display_text(truncate_display_copy(cleaned, limit))
    return ""


def wrap_text_segments(text: str, limit: int) -> list[str]:
    words = normalize_whitespace(text).split()
    if not words:
        return []
    segments: list[str] = []
    current: list[str] = []
    for word in words:
        candidate = " ".join(current + [word])
        if current and len(candidate) > limit:
            segments.append(" ".join(current))
            current = [word]
        else:
            current.append(word)
    if current:
        segments.append(" ".join(current))
    return segments


def split_prompt_for_layout(text: str, *, limit: int, max_parts: int = 3) -> list[str]:
    cleaned = clean_source_prompt(text)
    if not cleaned:
        return []

    segments: list[str] = []
    sentence_parts = [normalize_whitespace(part) for part in re.split(r"(?<=[?.!])\s+", cleaned) if normalize_whitespace(part)]
    for sentence in sentence_parts or [cleaned]:
        clause_parts = [sentence]
        if len(sentence) > limit:
            clause_parts = [
                normalize_whitespace(part)
                for part in re.split(r"(?<=[;:])\s+|(?<=,)\s+", sentence)
                if normalize_whitespace(part)
            ] or [sentence]
        for clause in clause_parts:
            if len(clause) <= limit:
                segments.append(clause)
            else:
                segments.extend(wrap_text_segments(clause, limit))
            if len(segments) >= max_parts:
                return unique_nonempty(segments, limit=max_parts)
    return unique_nonempty(segments, limit=max_parts)


def compact_tasks_for_layout(tasks: list[str], *, item_limit: int, max_items: int) -> list[str]:
    compacted: list[str] = []
    for task in tasks:
        if not normalize_whitespace(task):
            continue
        parts = split_prompt_for_layout(task, limit=item_limit, max_parts=max(2, min(4, max_items)))
        if not parts:
            parts = [truncate_text(task, item_limit)]
        compacted.extend(parts)
        if len(compacted) >= max_items:
            break
    return unique_nonempty(compacted, limit=max_items)


def find_term_example_in_records(source_records: list[dict[str, Any]], term: str) -> str:
    if not term:
        return ""
    pattern = re.compile(rf"[^.?!]*\b{re.escape(term)}\b[^.?!]*[.?!]?", re.IGNORECASE)
    for slide in source_records:
        search_texts = slide.get("problem_texts", []) + slide.get("text_items", []) + [slide.get("text", "")]
        for candidate in search_texts:
            match = pattern.search(candidate)
            if not match:
                continue
            example = clean_source_prompt(match.group(0), slide.get("title", ""))
            title_key = display_text_key(slide.get("title", ""))
            example_key = display_text_key(example)
            if (
                example
                and example_key != title_key
                and term.lower() in example.lower()
                and not example.lower().startswith("i can ")
                and not is_generic_slide_text(example, slide.get("title", ""))
            ):
                return truncate_display_copy(example, 96)
    return ""


def normalized_label_term(term: str) -> str:
    cleaned = normalize_whitespace(term)
    if not cleaned:
        return ""
    words = display_text_key(cleaned).split()
    while len(words) > 1 and words[-1] in {"activity", "activities", "move", "moves", "question", "questions", "writing"}:
        words = words[:-1]
    candidate = display_term_label(" ".join(words))
    if not candidate or is_low_value_vocabulary_term(candidate):
        return ""
    if not 1 <= len(display_text_key(candidate).split()) <= 3:
        return ""
    return candidate


def source_label_term_candidates(source_records: list[dict[str, Any]], limit: int = 6) -> list[str]:
    candidates: list[str] = []
    for slide in source_records:
        title = slide.get("title", "")
        raw_texts = [title, *slide.get("text_items", [])]
        for raw_text in raw_texts:
            cleaned = normalize_whitespace(raw_text)
            if not cleaned:
                continue
            if ":" in cleaned:
                suffix = normalize_whitespace(cleaned.split(":", 1)[1]).strip(" -")
                if suffix and "/" not in suffix:
                    candidate = normalized_label_term(suffix)
                    if candidate and not is_generic_slide_text(candidate, title):
                        candidates.append(candidate)
                for part in re.split(r"/", suffix):
                    candidate = normalized_label_term(part)
                    if candidate and not is_generic_slide_text(candidate, title):
                        candidates.append(candidate)
            elif "/" in cleaned:
                for part in re.split(r"/", cleaned):
                    candidate = normalized_label_term(part)
                    if candidate and not is_generic_slide_text(candidate, title):
                        candidates.append(candidate)
    return unique_nonempty(candidates, limit=limit)


CURATED_VOCAB_VISUAL_CUES = {
    "barrier": "Barrier students need removed",
    "claim": "Claim statement students defend",
    "data": "Lab data students can cite",
    "evidence": "Evidence pulled from lab data",
    "lab data": "Lab data students can cite",
    "reasoning": "Reasoning that connects the evidence",
    "sentence frame": "Sentence frame students complete",
    "sentence frames": "Sentence frame students complete",
    "sentence starter": "Sentence starter students finish",
    "sentence starters": "Sentence starter students finish",
}


def curated_vocab_visual_cue(term: str) -> str:
    return CURATED_VOCAB_VISUAL_CUES.get(display_text_key(term), "")


def find_term_example(deck: dict[str, Any], term: str) -> str:
    if not term:
        return ""
    pattern = re.compile(rf"[^.?!]*\b{re.escape(term)}\b[^.?!]*[.?!]?", re.IGNORECASE)
    for slide in deck.get("slides", []):
        search_texts = slide.get("problem_texts", []) + slide.get("text_items", []) + [slide.get("text", "")]
        for candidate in search_texts:
            match = pattern.search(candidate)
            if match:
                example = clean_source_prompt(match.group(0), slide.get("title", ""))
                title_key = display_text_key(slide.get("title", ""))
                example_key = display_text_key(example)
                if example and example_key != title_key and term.lower() in example.lower() and not example.lower().startswith("i can "):
                    return truncate_display_copy(example, 96)
    return truncate_text(f"Use {term.lower()} when you explain the source problem or model.", 110)


def vocabulary_visual_cue(deck: dict[str, Any], term: str) -> str:
    curated = curated_vocab_visual_cue(term)
    if curated:
        return curated
    example = find_term_example(deck, term)
    if example:
        return truncate_display_copy(example.replace(term, term).strip(), 72)
    return truncate_text(f"Look for {term.lower()} in the source figure, model, or problem.", 72)


def student_friendly_definition(term: str) -> str:
    lowered = term.lower()
    if lowered in STUDENT_FRIENDLY_DEFINITIONS:
        return STUDENT_FRIENDLY_DEFINITIONS[lowered]
    if lowered.endswith("s") and lowered[:-1] in STUDENT_FRIENDLY_DEFINITIONS:
        return STUDENT_FRIENDLY_DEFINITIONS[lowered[:-1]]
    if "percent" in lowered:
        return "how many out of every one hundred"
    if "equation" in lowered:
        return "a math sentence that shows two amounts are equal"
    if "ratio" in lowered:
        return "a comparison between two amounts"
    if "score" in lowered:
        return "the number of points earned"
    if "vote" in lowered:
        return "to choose an option"
    if any(char.isdigit() for char in lowered) or lowered in {"d1", "d2", "b1", "b2"}:
        return "a label used to name one measurement in the formula or diagram"
    return f"the meaning of {lowered} in the lesson's model, diagram, or problem"


def build_vocab_placeholders(deck: dict[str, Any]) -> list[dict[str, str]]:
    vocab: list[dict[str, str]] = []
    for term in deck["keyword_candidates"][:5]:
        if is_low_value_vocabulary_term(term):
            continue
        vocab.append(
            {
                "word": term,
                "definition": student_friendly_definition(term),
                "example": find_term_example(deck, term),
                "visual_cue": vocabulary_visual_cue(deck, term),
            }
        )
    return vocab


def source_term_candidates(source_records: list[dict[str, Any]], limit: int = 6) -> list[str]:
    source_text = combined_source_text(source_records).lower()
    if not source_text:
        return []
    terms: list[str] = []
    for term in SOURCE_TERM_PRIORITY:
        pattern = r"\b" + re.escape(term).replace(r"\ ", r"\s+") + r"\b"
        if re.search(pattern, source_text):
            display = display_term_label(term)
            if not is_low_value_vocabulary_term(display):
                terms.append(display)
    terms.extend(source_label_term_candidates(source_records, limit=max(limit + 2, 8)))
    protected_words = {
        word
        for term in terms
        for word in display_text_key(term).split()
        if len(word) > 2
    }
    counter: Counter[str] = Counter()
    for word in re.findall(r"[A-Za-z][A-Za-z'-]{3,}", source_text):
        lowered = word.lower()
        if is_low_value_vocabulary_term(lowered) or lowered in CONTEXT_STOPWORDS:
            continue
        counter[lowered] += 1
    for word, _count in counter.most_common(limit * 3):
        if word in protected_words:
            continue
        if word in STUDENT_FRIENDLY_DEFINITIONS or counter[word] >= 2:
            display = display_term_label(word)
            if not is_low_value_vocabulary_term(display):
                terms.append(display)
    return unique_nonempty(terms, limit=limit)


def source_vocabulary(deck: dict[str, Any], source_numbers: list[int], limit: int = 3) -> list[dict[str, str]]:
    source_records = source_slides_from_numbers(deck, source_numbers)
    source_text = combined_source_text(source_records).lower()
    fallback = build_vocab_placeholders(deck)
    if not source_text:
        return fallback[:limit]
    profile = math_profile_for_text(normalize_whitespace(" ".join([deck.get("lesson_title", ""), combined_source_text(source_records)])))
    profile_terms = [display_term_label(term) for term in profile.get("default_terms", [])]
    source_terms = source_term_candidates(source_records, limit=max(limit + 2, 6))
    ordered_terms = (
        profile_terms + source_terms
        if profile_terms and profile.get("topic") != "generic_math"
        else source_terms + profile_terms
    )
    prioritized_terms = unique_nonempty(
        ordered_terms,
        limit=max(limit + 4, 8),
    )
    prioritized: list[dict[str, str]] = []
    for term in prioritized_terms:
        if is_low_value_vocabulary_term(term):
            continue
        example = find_term_example_in_records(source_records, term) or find_term_example(deck, term)
        visual_cue = curated_vocab_visual_cue(term) or truncate_display_copy(
            example or f"Look for {term.lower()} in the source figure, labels, or model.",
            72,
        )
        prioritized.append(
            {
                "word": term,
                "definition": student_friendly_definition(term),
                "example": example or truncate_display_copy(f"Use {term.lower()} when you explain the source problem or model.", 96),
                "visual_cue": visual_cue,
            }
        )
    keyword_matches = [
        item for item in fallback
        if item["word"].lower() in source_text or any(part in source_text for part in item["word"].lower().split())
    ]
    combined = prioritized + keyword_matches + [item for item in fallback if item not in keyword_matches]
    filtered = [
        item for item in combined
        if not is_low_value_vocabulary_term(item.get("word", ""))
    ]
    results: list[dict[str, str]] = []
    seen_words: set[str] = set()
    seen_roots: set[str] = set()
    for item in filtered:
        word_key = display_text_key(item.get("word", ""))
        word_root = word_key[:-1] if word_key.endswith("s") and len(word_key) > 4 else word_key
        if not word_key or word_key in seen_words or word_root in seen_roots:
            continue
        seen_words.add(word_key)
        seen_roots.add(word_root)
        results.append(item)
        if len(results) >= limit:
            break
    return results


def session_number_value(session_key: str) -> int:
    return 1 if session_key == "session_1" else 2


def session_label_for_key(session_key: str) -> str:
    return f"Session {session_number_value(session_key)}"


def planned_session_keys(plan: dict[str, Any]) -> list[str]:
    return [
        session_key
        for session_key in SESSION_KEY_ORDER
        if isinstance(plan.get(session_key), dict)
    ]


def planned_session_specs(plan: dict[str, Any]) -> list[tuple[str, str, str]]:
    return [
        (session_key, session_label_for_key(session_key), SESSION_OUTPUT_KEYS[session_key])
        for session_key in planned_session_keys(plan)
    ]


def session_slide_count_bounds(session: dict[str, Any]) -> tuple[int, int]:
    if uses_exact_esol_template(session):
        slide_total = len(EXACT_ESOL_TEMPLATE_SEQUENCE)
        return (slide_total, slide_total)
    return (MIN_SESSION_SLIDES, MAX_SESSION_SLIDES)


def infer_unit_label(deck: dict[str, Any]) -> str:
    text = normalize_whitespace(
        " ".join(
            [
                deck.get("lesson_title", ""),
                deck.get("summary", ""),
                " ".join(slide.get("title", "") for slide in deck.get("slides", [])[:10]),
            ]
        )
    )
    match = re.search(r"\bunit\s*(\d+)\b", text, flags=re.IGNORECASE)
    if match:
        return f"Unit {match.group(1)}"
    lowered = text.lower()
    if "volume" in lowered:
        return "Volume Unit"
    if "area" in lowered:
        return "Area Unit"
    return "Math Unit"


def spanish_translation_for_term(term: str) -> str:
    cleaned = display_term_label(term)
    lowered = cleaned.lower()
    if lowered in SPANISH_TRANSLATIONS:
        return SPANISH_TRANSLATIONS[lowered]
    singular = lowered[:-1] if lowered.endswith("s") else lowered
    if singular in SPANISH_TRANSLATIONS:
        translated = SPANISH_TRANSLATIONS[singular]
        if lowered.endswith("s") and not translated.endswith("s"):
            if translated.endswith("z"):
                return translated[:-1] + "ces"
            return translated + "s"
        return translated
    return lowered


def bilingual_term_label(term: str) -> str:
    display = display_term_label(term)
    return f"{display} · {spanish_translation_for_term(display)}"


def vocabulary_symbol(term: str) -> str:
    lowered = display_term_label(term).lower()
    if "formula" in lowered or "equation" in lowered:
        return "📐"
    if any(token in lowered for token in ("median", "mean", "data set", "dot plot", "line plot", "outlier")):
        return "📊"
    if any(token in lowered for token in ("base", "height", "length", "width", "diagonal", "dimension")):
        return "📏"
    if any(token in lowered for token in ("prism", "cube", "volume")):
        return "🧊"
    if "graph" in lowered:
        return "📈"
    if "table" in lowered:
        return "🗂"
    return "🔑"


def math_profile_for_text(text: str) -> dict[str, Any]:
    lowered = normalize_whitespace(text).lower()
    if any(term in lowered for term in ("median", "mean", "data set", "dot plot", "line plot", "outlier", "ordered data")):
        return {
            "topic": "data_analysis",
            "formula": "",
            "vars": ["ordered data", "median", "reasoning"],
            "headers": ["Data Set", "Median", "Reasoning"],
            "answer_label": "Median",
            "shape_label": "Data Display",
            "default_terms": ["median", "data set", "ordered data", "dot plot", "outlier", "mean"],
        }
    if "fraction" in lowered and "decimal" in lowered and ("percent" in lowered or "percentage" in lowered):
        return {
            "topic": "fraction_decimal_percent",
            "formula": "",
            "vars": ["fraction", "decimal", "percent"],
            "headers": ["Fraction", "Decimal", "Percent"],
            "answer_label": "Equivalent Value",
            "shape_label": "Equivalent Representations",
            "default_terms": ["fraction", "decimal", "percent", "equivalent", "notation"],
        }
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return {
            "topic": "volume_prism",
            "formula": "V = l x w x h",
            "vars": ["length", "width", "height"],
            "headers": ["Length", "Width", "Height", "Volume"],
            "answer_label": "Volume",
            "shape_label": "Rectangular Prism",
            "default_terms": ["volume", "rectangular prism", "length", "width", "height", "formula"],
        }
    if "rhombus" in lowered:
        return {
            "topic": "rhombus_area",
            "formula": "A = (d1 x d2) / 2",
            "vars": ["d1", "d2"],
            "headers": ["d1", "d2", "Area"],
            "answer_label": "Area",
            "shape_label": "Rhombus",
            "default_terms": ["area", "rhombus", "diagonal", "formula", "d1", "d2"],
        }
    if "area" in lowered and any(term in lowered for term in ("regular polygon", "octagon")) and any(
        term in lowered for term in ("decompose", "congruent triangle", "congruent triangles", "triangle")
    ):
        return {
            "topic": "regular_polygon_area",
            "formula": "A = n x (b x h) / 2",
            "vars": ["base", "height", "triangles"],
            "headers": ["Base", "Height", "# Triangles", "Area"],
            "answer_label": "Area",
            "shape_label": "Regular Polygon",
            "default_terms": ["area", "regular polygon", "octagon", "decompose", "congruent triangles", "base", "height"],
        }
    if "triangle" in lowered:
        return {
            "topic": "triangle_area",
            "formula": "A = (b x h) / 2",
            "vars": ["base", "height"],
            "headers": ["Base", "Height", "Area"],
            "answer_label": "Area",
            "shape_label": "Triangle",
            "default_terms": ["area", "triangle", "base", "height", "formula", "label"],
        }
    if "trapezoid" in lowered:
        return {
            "topic": "trapezoid_area",
            "formula": "A = ((b1 + b2) x h) / 2",
            "vars": ["b1", "b2", "height"],
            "headers": ["b1", "b2", "Height", "Area"],
            "answer_label": "Area",
            "shape_label": "Trapezoid",
            "default_terms": ["area", "trapezoid", "base", "height", "formula", "compare"],
        }
    if "parallelogram" in lowered:
        return {
            "topic": "parallelogram_area",
            "formula": "A = b x h",
            "vars": ["base", "height"],
            "headers": ["Base", "Height", "Area"],
            "answer_label": "Area",
            "shape_label": "Parallelogram",
            "default_terms": ["area", "parallelogram", "base", "height", "formula", "solve"],
        }
    if "rectangle" in lowered:
        return {
            "topic": "rectangle_area",
            "formula": "A = l x w",
            "vars": ["length", "width"],
            "headers": ["Length", "Width", "Area"],
            "answer_label": "Area",
            "shape_label": "Rectangle",
            "default_terms": ["area", "rectangle", "length", "width", "formula", "solve"],
        }
    return {
        "topic": "generic_math",
        "formula": "",
        "vars": ["value 1", "value 2"],
        "headers": ["Value 1", "Value 2", "Answer"],
        "answer_label": "Answer",
        "shape_label": "Math Model",
        "default_terms": [],
    }


def session_math_profile(deck: dict[str, Any], session_key: str) -> dict[str, Any]:
    source_records = session_source_slides(deck, session_key)
    session_text = normalize_whitespace(
        " ".join([deck.get("lesson_title", ""), combined_source_text(source_records), select_session_objective(deck, session_key)])
    )
    return math_profile_for_text(session_text)


def normalize_formula_text(text: str) -> str:
    return normalize_whitespace(text.replace("×", "x").replace("·", "x"))


def formula_candidates_from_text(text: str) -> list[str]:
    candidates: list[str] = []
    for match in re.finditer(r"\b(?:A|V|SA|y)\s*=\s*[^.;,\n]{3,40}", text, flags=re.IGNORECASE):
        cleaned = normalize_formula_text(match.group(0)).strip(" .")
        if cleaned and cleaned not in candidates:
            candidates.append(cleaned)
    return candidates


def formula_for_session(deck: dict[str, Any], session_key: str) -> str:
    source_records = session_source_slides(deck, session_key)
    session_text = normalize_whitespace(
        " ".join([deck.get("lesson_title", ""), combined_source_text(source_records), select_session_objective(deck, session_key)])
    )
    candidates = formula_candidates_from_text(session_text)
    if candidates:
        return candidates[0]
    return session_math_profile(deck, session_key).get("formula", "")


def prior_formula_for_session(deck: dict[str, Any], session_key: str, current_formula: str) -> str:
    if session_key == "session_2":
        prior_formula = formula_for_session(deck, "session_1")
        if prior_formula and prior_formula != current_formula:
            return prior_formula
    lowered = normalize_whitespace(" ".join([deck.get("lesson_title", ""), current_formula])).lower()
    if "volume" in lowered and ("rectangular prism" in lowered or "prism" in lowered):
        return "A = l x w"
    if any(term in lowered for term in ("triangle", "rhombus", "trapezoid")):
        return "A = b x h"
    return current_formula


def session_esol_vocabulary(deck: dict[str, Any], source_numbers: list[int], *, limit: int = 6) -> list[dict[str, str]]:
    source_records = source_slides_from_numbers(deck, source_numbers)
    profile = math_profile_for_text(normalize_whitespace(" ".join([deck.get("lesson_title", ""), combined_source_text(source_records)])))
    seed_words = [item.get("word", "") for item in source_vocabulary(deck, source_numbers, limit=max(limit + 2, 8))]
    seed_words.extend(profile.get("default_terms", []))
    vocab: list[dict[str, str]] = []
    for term in unique_nonempty(seed_words, limit=limit * 2):
        display = display_term_label(term)
        if is_low_value_vocabulary_term(display):
            continue
        vocab.append(
            {
                "word": display,
                "spanish": spanish_translation_for_term(display),
                "definition": student_friendly_definition(display),
                "example": find_term_example_in_records(source_records, display) or find_term_example(deck, display),
                "visual_cue": vocabulary_visual_cue(deck, display),
            }
        )
    deduped: list[dict[str, str]] = []
    seen: set[str] = set()
    for item in vocab:
        key = display_text_key(item.get("word", ""))
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(item)
        if len(deduped) >= limit:
            break
    return deduped


def vocabulary_activity_terms(vocab_items: list[dict[str, str]], *, limit: int = 4) -> list[str]:
    terms: list[str] = []
    for item in vocab_items:
        word = display_term_label(item.get("word", ""))
        if not word:
            continue
        terms.append(truncate_text(word, 26))
        if len(terms) >= limit:
            break
    return terms


def vocabulary_speaking_frame(vocab_items: list[dict[str, str]]) -> str:
    anchor = display_term_label(vocab_items[0].get("word", "")) if vocab_items else "word"
    anchor = anchor or "word"
    return f"Say: '{anchor}' means ___."


def source_seed_slide(deck: dict[str, Any], source_numbers: list[int], *, kind: str = "practice") -> dict[str, Any]:
    source_records = source_slides_from_numbers(deck, source_numbers)
    slide = build_slide_plan(
        kind=kind,
        title="Source Seed",
        subtitle="",
        primary_text=best_source_prompt_text(source_records, limit=280),
        tasks=slide_task_list(source_records, limit=4),
        source_slide_numbers=source_numbers,
    )
    slide["source_problem_cards"] = source_problem_candidates(source_records, limit=4)
    return slide


def profile_measure_value(profile_var: str, measurements: dict[str, tuple[str, str]], default: str) -> str:
    aliases = {
        "length": ("length", "l"),
        "width": ("width", "w"),
        "height": ("height", "h"),
        "base": ("base", "b"),
        "b1": ("b1",),
        "b2": ("b2",),
        "d1": ("d1",),
        "d2": ("d2",),
    }
    value, _unit = measurement_value(measurements, *(aliases.get(profile_var, (profile_var,))))
    return value or default


def table_seed_values(profile: dict[str, Any], seed_slide: dict[str, Any]) -> list[str]:
    measurements = source_measurements(seed_slide)
    defaults = {
        "volume_prism": ["4", "3", "2"],
        "rhombus_area": ["10", "8"],
        "regular_polygon_area": ["9", "18", "8"],
        "triangle_area": ["10", "6"],
        "trapezoid_area": ["8", "12", "5"],
        "parallelogram_area": ["8", "5"],
        "rectangle_area": ["6", "4"],
        "generic_math": ["6", "4"],
    }
    fallback = defaults.get(profile.get("topic", "generic_math"), defaults["generic_math"])
    values: list[str] = []
    for index, variable in enumerate(profile.get("vars", [])):
        default = fallback[index] if index < len(fallback) else fallback[-1]
        values.append(profile_measure_value(variable, measurements, default))
    return values


def numeric_table_rows(values: list[str], *, seed: int) -> list[list[str]]:
    rows = [values]
    for row_index in range(1, 3):
        rows.append(
            [
                shift_numeric_token(value, seed=seed + row_index, index=index + row_index)
                for index, value in enumerate(values)
            ]
        )
    return rows


def compute_profile_result(profile: dict[str, Any], row: list[str]) -> str:
    try:
        numeric = [float(value.replace(",", "")) for value in row]
    except ValueError:
        return "___"
    topic = profile.get("topic", "")
    if topic == "volume_prism" and len(numeric) >= 3:
        result = numeric[0] * numeric[1] * numeric[2]
    elif topic == "rhombus_area" and len(numeric) >= 2:
        result = (numeric[0] * numeric[1]) / 2
    elif topic == "regular_polygon_area" and len(numeric) >= 3:
        result = ((numeric[0] * numeric[1]) / 2) * numeric[2]
    elif topic == "triangle_area" and len(numeric) >= 2:
        result = (numeric[0] * numeric[1]) / 2
    elif topic == "trapezoid_area" and len(numeric) >= 3:
        result = ((numeric[0] + numeric[1]) * numeric[2]) / 2
    elif len(numeric) >= 2:
        result = numeric[0] * numeric[1]
    else:
        return "___"
    if int(result) == result:
        return str(int(result))
    return f"{result:.1f}".rstrip("0").rstrip(".")


def build_reference_table(profile: dict[str, Any], seed_slide: dict[str, Any], *, completed: bool) -> list[list[str]]:
    base_values = table_seed_values(profile, seed_slide)
    seed = sum(num for num in seed_slide.get("source_slide_numbers", []) if isinstance(num, int)) or 7
    data = [profile.get("headers", ["Value 1", "Value 2", "Answer"])]
    for row in numeric_table_rows(base_values, seed=seed):
        answer = compute_profile_result(profile, row) if completed else "___"
        data.append(row + [answer])
    return data


def exact_content_objectives(deck: dict[str, Any], session_key: str, current_formula: str) -> list[str]:
    profile = session_math_profile(deck, session_key)
    topic = profile.get("topic", "")
    source_objective = select_session_objective(deck, session_key)
    lowered_source = normalize_whitespace(source_objective).lower()
    if topic == "volume_prism":
        primary = source_objective if "volume" in lowered_source else "I can find the volume of a rectangular prism."
        return [primary, "I can use the dimensions and the formula to justify the volume."]
    if topic == "rhombus_area":
        primary = source_objective if "rhombus" in lowered_source else "I can find the area of a rhombus."
        return [primary, "I can use d1 and d2 to justify the area."]
    if topic == "regular_polygon_area":
        primary = (
            source_objective
            if any(term in lowered_source for term in ("regular polygon", "octagon"))
            else "I can find the area of a regular polygon by decomposing it into triangles."
        )
        return [primary, "I can connect each side of the polygon to one congruent triangle."]
    if topic == "triangle_area":
        primary = source_objective if "triangle" in lowered_source else "I can find the area of a triangle."
        return [primary, "I can use the base and height to justify the area."]
    if topic == "trapezoid_area":
        primary = source_objective if "trapezoid" in lowered_source else "I can find the area of a trapezoid."
        return [primary, "I can use both bases and the height to justify the area."]
    if "A =" in current_formula:
        primary = source_objective if "area" in lowered_source else "I can find the area."
        return [primary, "I can use the model or formula to justify the answer."]
    if source_objective:
        return [source_objective, "I can explain my strategy with labels, vocabulary, and evidence."]
    return ["I can solve the problem.", "I can explain my math."]


def spoken_language_frame(profile: dict[str, Any], current_formula: str) -> str:
    topic = profile.get("topic", "")
    if topic == "volume_prism":
        return "I say: 'I use ___ x ___ x ___.'"
    if topic == "rhombus_area":
        return "I say: 'I use ___ x ___, then half.'"
    if topic == "regular_polygon_area":
        return "I say: 'I find one triangle, then multiply by the number of triangles.'"
    if topic == "triangle_area":
        return "I say: 'I use base, height, then half.'"
    if "A =" in current_formula:
        return "I say: 'I use ___ because ___.'"
    return "I say: 'First ___, then ___.'"


def word_help_strip(deck: dict[str, Any], source_numbers: list[int], *, limit: int = 5) -> list[str]:
    return [
        bilingual_term_label(item.get("word", ""))
        for item in session_esol_vocabulary(deck, source_numbers, limit=limit)
        if item.get("word", "")
    ][:limit]


def exact_notice_lines(profile: dict[str, Any]) -> list[str]:
    topic = profile.get("topic", "")
    if topic == "volume_prism":
        return ["I see three dimensions.", "I see a prism model.", "I see the formula."]
    if topic == "regular_polygon_area":
        return ["I see a regular polygon.", "I see congruent triangles.", "I see labeled measures."]
    if "area" in topic:
        return ["I see labeled measures.", "I see a shape model.", "I see the formula."]
    return ["I see a table.", "I see a formula.", "I see labels."]


def exact_wonder_prompt(profile: dict[str, Any]) -> str:
    if profile.get("topic") == "volume_prism":
        return "If one side changes, does volume change?"
    if profile.get("topic") == "regular_polygon_area":
        return "How does one triangle help us find the whole area?"
    if "area" in profile.get("topic", ""):
        return "If one measure changes, does area change?"
    return "If one value changes, what happens?"


def exact_bridge_sentence(deck: dict[str, Any], session_key: str, current_formula: str, prior_formula: str) -> str:
    if prior_formula and prior_formula != current_formula:
        return f"Now connect {prior_formula} to {current_formula}."
    if session_key == "session_2":
        return "Now use that idea alone."
    return "Now use that idea today."


def context_hook_text(context_anchor: str, lesson_title: str) -> str:
    if context_anchor:
        lowered = context_anchor.lower()
        article = "" if lowered.startswith(("the ", "a ", "an ")) else "the "
        return truncate_text(f"Use {article}{lowered} context.", 64)
    return "Use the lesson model and formula."


def error_analysis_steps(profile: dict[str, Any], formula_text: str) -> tuple[list[str], str, str]:
    topic = profile.get("topic", "")
    if topic == "volume_prism":
        return (
            [
                "✅ Step 1: Find l, w, h.",
                f"✅ Step 2: Write {formula_text}.",
                "❌ Step 3: Multiply only two values.",
                "✅ Step 4: Label cubic units.",
            ],
            f"Use {formula_text} with all dimensions.",
            "Why do we multiply, not add?",
        )
    if "area" in topic and topic in {"triangle_area", "rhombus_area", "trapezoid_area"}:
        return (
            [
                "✅ Step 1: Find the measures.",
                f"✅ Step 2: Write {formula_text}.",
                "❌ Step 3: Forget one-half.",
                "✅ Step 4: Label square units.",
            ],
            f"Use {formula_text} and include one-half.",
            "Why do we use one-half, not one?",
        )
    return (
        [
            "✅ Step 1: Read the values.",
            f"✅ Step 2: Write {formula_text or 'the formula'}.",
            "❌ Step 3: Use the wrong operation.",
            "✅ Step 4: Label the answer.",
        ],
        f"Fix it with {formula_text or 'the formula'}.",
        "Why do we use this operation?",
    )


def count_term_hits(text: str, terms: tuple[str, ...]) -> int:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return 0
    return sum(1 for term in terms if term in lowered)


def session_signal_text(deck: dict[str, Any], session_key: str, slides: list[dict[str, Any]]) -> str:
    source_records = session_source_slides(deck, session_key)
    parts = [deck.get("lesson_title", ""), deck.get("summary", "")]
    for slide in source_records:
        parts.extend(slide.get("problem_texts", []))
        parts.extend(slide.get("text_items", []))
        parts.extend([slide.get("title", ""), slide.get("text", ""), slide.get("notes", "")])
    for slide in slides:
        parts.extend(
            [
                slide.get("title", ""),
                slide.get("subtitle", ""),
                slide.get("primary_text", ""),
                slide.get("secondary_text", ""),
                " ".join(slide.get("bullets", [])),
                " ".join(slide.get("tasks", [])),
                slide.get("response_prompt", ""),
            ]
        )
    return normalize_whitespace(" ".join(parts)).lower()


def best_session_context_anchor(deck: dict[str, Any], session_key: str) -> str:
    candidates: list[tuple[int, str]] = []
    lesson_key = display_text_key(deck.get("lesson_title", ""))
    for slide in session_source_slides(deck, session_key):
        title = normalize_whitespace(slide.get("title", ""))
        if title and not is_generic_slide_text(title) and not is_problem_like_text(title) and display_text_key(title) != lesson_key:
            words = [word.lower() for word in re.findall(r"[A-Za-z']+", title)]
            non_math = [word for word in words if len(word) > 2 and word not in CONTEXT_STOPWORDS]
            if non_math:
                score = 8 + min(len(non_math), 3) + (2 if slide.get("image_count") else 0)
                candidates.append((score, title))
        for problem_text in slide.get("problem_texts", [])[:2]:
            cleaned = clean_source_prompt(problem_text, slide.get("title", ""))
            if not cleaned or is_generic_slide_text(cleaned, slide.get("title", "")):
                continue
            if len(cleaned) > 72 or is_problem_like_text(cleaned):
                continue
            words = [word.lower() for word in re.findall(r"[A-Za-z']+", cleaned)]
            non_math = [word for word in words if len(word) > 2 and word not in CONTEXT_STOPWORDS]
            if non_math:
                candidates.append((5 + min(len(non_math), 3), cleaned))
    if not candidates:
        return ""
    candidates.sort(key=lambda item: (-item[0], len(item[1])))
    return truncate_text(candidates[0][1], 56)


def slide_index_by_kind(slides: list[dict[str, Any]], kinds: set[str]) -> int | None:
    for index, slide in enumerate(slides):
        if slide.get("kind") in kinds:
            return index
    return None


def slide_engagement_modes(plan_slide: dict[str, Any]) -> set[str]:
    modes: set[str] = set()
    kind = plan_slide.get("kind", "")
    family = normalize_whitespace(plan_slide.get("activity_family", ""))
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    premium_layout = normalize_whitespace(plan_slide.get("premium_layout", ""))

    modes.update(ENGAGEMENT_MODE_ACTIVITY_FAMILIES.get(family, set()))
    modes.update(ENGAGEMENT_MODE_TEMPLATE_ROLES.get(template_role, set()))
    modes.update(ENGAGEMENT_MODE_PREMIUM_LAYOUTS.get(premium_layout, set()))

    if kind == "be_curious":
        modes.add("notice")
    if normalize_whitespace(plan_slide.get("partner_prompt", "")):
        modes.add("partner")
    if plan_slide.get("discussion_questions"):
        modes.add("discuss")
    if plan_slide.get("choice_paths"):
        modes.add("choice")
    if plan_slide.get("independent_problems"):
        modes.add("apply")
    return modes


def session_engagement_modes(session: dict[str, Any]) -> list[str]:
    modes: set[str] = set()
    for slide in session.get("slides", []):
        modes.update(slide_engagement_modes(slide))
    return sorted(modes)


def session_engagement_slide_count(session: dict[str, Any]) -> int:
    return sum(1 for slide in session.get("slides", []) if slide_engagement_modes(slide))


def engagement_slide_target(session: dict[str, Any]) -> int:
    slide_count = len(session.get("slides", []))
    if slide_count >= PREMIUM_TARGET_SESSION_SLIDES:
        return 5
    if slide_count >= 11:
        return 4
    return 3


def engagement_mode_target(session: dict[str, Any]) -> int:
    slide_count = len(session.get("slides", []))
    return 4 if slide_count >= PREMIUM_TARGET_SESSION_SLIDES else 3


def engagement_rebalance_priority(plan_slide: dict[str, Any]) -> int:
    family = normalize_whitespace(plan_slide.get("activity_family", ""))
    kind = plan_slide.get("kind", "")
    priority = ACTIVITY_KIND_PRIORITY.get(kind, 50)
    score = 0
    if not has_activity(plan_slide):
        score += 6
    elif family not in HIGH_AGENCY_ACTIVITY_FAMILIES:
        score += 4
    if kind in PROBLEM_SOLVING_KINDS:
        score += 3
    if kind in {"guided_notes", "quick_review", "reflection", "exit_ticket"}:
        score += 1
    return score * 10 - priority


def first_slide_by_kind(
    slides: list[dict[str, Any]],
    kinds: set[str],
    *,
    prefer_no_activity: bool | None = None,
    start: int = 0,
) -> dict[str, Any] | None:
    matches = [slide for slide in slides[start:] if slide.get("kind") in kinds]
    if prefer_no_activity is None:
        return matches[0] if matches else None
    preferred = [slide for slide in matches if has_activity(slide) != prefer_no_activity]
    if prefer_no_activity:
        preferred = [slide for slide in matches if not has_activity(slide)]
    else:
        preferred = [slide for slide in matches if has_activity(slide)]
    return (preferred or matches or [None])[0]


def first_unused_premium_slide(
    slides: list[dict[str, Any]],
    kinds: set[str],
    *,
    prefer_no_activity: bool | None = None,
) -> dict[str, Any] | None:
    unused = [slide for slide in slides if slide.get("kind") in kinds and not normalize_whitespace(slide.get("premium_layout", ""))]
    if not unused:
        return first_slide_by_kind(slides, kinds, prefer_no_activity=prefer_no_activity)
    if prefer_no_activity is True:
        return next((slide for slide in unused if not has_activity(slide)), unused[0])
    if prefer_no_activity is False:
        return next((slide for slide in unused if has_activity(slide)), unused[0])
    return unused[0]


def representation_terms_for_text(text: str) -> list[str]:
    lowered = normalize_whitespace(text).lower()
    ordered: list[str] = []
    for term in ("scenario", "table", "graph", "equation", "coordinates", "pattern", "diagram", "model"):
        if term in lowered and term not in ordered:
            ordered.append(term)
    if "coordinate" in lowered and "coordinates" not in ordered:
        ordered.append("coordinates")
    if "verbal" in lowered and "scenario" not in ordered:
        ordered.insert(0, "scenario")
    return ordered


def multi_representation_labels(text: str) -> list[str]:
    terms = representation_terms_for_text(text)
    if {"table", "graph", "equation"} <= set(terms):
        return ["Table", "Graph", "Equation", "Explain the Connection"]
    if "coordinates" in terms:
        return ["Coordinates", "Graph", "Rule / Equation", "Explain the Pattern"]
    if "table" in terms:
        return ["Situation", "Table", "Equation / Rule", "Explain the Match"]
    if "diagram" in terms or "model" in terms:
        return ["Visual Model", "Numbers", "Rule / Formula", "Explain It"]
    return ["Situation", "Representation A", "Representation B", "Explain the Connection"]


def strategy_comparison_rows(text: str) -> list[list[str]]:
    lowered = normalize_whitespace(text).lower()
    if "volume" in lowered and ("prism" in lowered or "rectangular" in lowered):
        return [
            ["Visual model", "When you can picture layers or a base", "How the dimensions build the prism", "Leaving out one dimension"],
            ["Formula", "When all needed dimensions are known", "The calculation efficiently", "Forgetting labels or cubic units"],
        ]
    if "graph" in lowered or "table" in lowered or "equation" in lowered:
        return [
            ["Table", "When you want to track pairs of values", "How the values change", "Mixing up input and output"],
            ["Graph", "When you want to see the trend visually", "The overall relationship", "Plotting points in the wrong place"],
            ["Equation", "When you want a rule you can reuse", "How one quantity depends on the other", "Confusing which variable depends on the other"],
        ]
    return [
        ["Concrete / visual", "When you need to see the math idea", "What each part represents", "Missing a key label or part"],
        ["Symbolic / shortcut", "When the structure is already clear", "A fast calculation path", "Using the shortcut before the meaning is clear"],
    ]


def decision_tree_steps(text: str) -> list[str]:
    lowered = normalize_whitespace(text).lower()
    if "proportional" in lowered:
        return [
            "Do the pairs change by the same constant rate?",
            "If yes, test whether the ratio stays the same.",
            "If the ratio is constant, label it proportional. If not, explain why not.",
        ]
    if "independent variable" in lowered or "dependent variable" in lowered:
        return [
            "Ask which quantity changes first.",
            "If one value depends on the other, that value is the dependent variable.",
            "Use the context sentence to justify the direction of dependence.",
        ]
    if "volume" in lowered:
        return [
            "Do you know the dimensions you need?",
            "If all dimensions are known, multiply them or use the volume formula.",
            "If one dimension is missing, use the known volume and divide by the product of the other dimensions.",
        ]
    return [
        "What is the problem asking you to decide or classify?",
        "What clue tells you which rule or strategy fits?",
        "Use that clue to choose the path and justify the decision.",
    ]


def create_your_own_labels(text: str) -> list[str]:
    lowered = normalize_whitespace(text).lower()
    if "graph" in lowered or "table" in lowered or "equation" in lowered:
        return ["Create a Scenario", "Build a Table or Graph", "Write the Rule", "Explain the Relationship"]
    if "volume" in lowered:
        return ["Create a Prism Story", "Label the Dimensions", "Write the Formula / Solve", "Explain Why It Works"]
    return ["Create the Situation", "Show the Math", "Solve It", "Explain the Reasoning"]


def context_phrase(context_anchor: str, *, with_article: bool = False) -> str:
    cleaned = normalize_whitespace(context_anchor)
    if not cleaned:
        return "this problem" if not with_article else "the source context"
    lowered = cleaned.lower()
    if lowered.startswith(("the ", "a ", "an ")):
        return cleaned
    return f"the {cleaned}" if with_article else cleaned


def contextual_sentence_frames(
    deck: dict[str, Any],
    plan_slide: dict[str, Any],
    *,
    context_anchor: str = "",
) -> list[str]:
    source_records = source_slides_from_numbers(deck, plan_slide.get("source_slide_numbers", []))
    source_text = best_source_prompt_text(source_records, limit=220)
    lowered = normalize_whitespace(f"{source_text} {plan_slide.get('primary_text', '')} {plan_slide.get('secondary_text', '')}").lower()
    context_piece = context_phrase(context_anchor).lower()
    if "independent variable" in lowered or "dependent variable" in lowered:
        return [
            "The independent variable is ___ because ___.",
            "The dependent variable is ___ because ___.",
            "In this relationship, ___ changes when ___ changes.",
        ]
    if "graph" in lowered:
        return [
            "The graph shows ___ because ___.",
            "As ___ increases, ___ ___.",
            "The point or feature that matters is ___ because ___.",
        ]
    if "table" in lowered:
        return [
            "The table shows ___ when ___.",
            "Each time ___ changes by ___, ___ changes by ___.",
            "I can use the table to predict ___ because ___.",
        ]
    if "equation" in lowered:
        return [
            "The equation shows that ___ depends on ___.",
            "The value of ___ changes by ___ when ___.",
            "This rule matches the context because ___.",
        ]
    if "volume" in lowered and ("prism" in lowered or "rectangular" in lowered):
        return [
            "The dimensions that matter are ___ because ___.",
            f"In {context_piece}, the volume is ___ because ___.",
            "The cubic units make sense because ___.",
        ]
    return [
        f"In the {context_piece}, ___.",
        "The math idea works because ___.",
        "I know my strategy fits because ___.",
    ]


def common_mistake_payload(text: str, context_anchor: str = "") -> tuple[str, str, list[str]]:
    lowered = normalize_whitespace(text).lower()
    context_piece = context_anchor or "the source problem"
    if "volume" in lowered and ("prism" in lowered or "rectangular" in lowered):
        return (
            "Error Analysis Block",
            f"A student used the dimensions from {context_piece} but multiplied only two dimensions and forgot to include the height.",
            [
                "Nearly correct: The student identified two useful dimensions.",
                "Mistake sample: They found only the area of one face, not the full volume.",
                "Fix it: Rewrite the formula with all three dimensions and explain why the units must be cubic.",
            ],
        )
    if "graph" in lowered:
        return (
            "Error Analysis Block",
            f"A student read the graph in {context_piece} and matched the wrong coordinate pair to the question.",
            [
                "Nearly correct: The student found a point that looks relevant.",
                "Mistake sample: The x- and y-values were interpreted in the wrong order.",
                "Fix it: Identify what each axis represents and explain how that changes the answer.",
            ],
        )
    if "table" in lowered or "equation" in lowered:
        return (
            "Error Analysis Block",
            f"A student used a rule from {context_piece} but mixed up which quantity depends on the other.",
            [
                "Nearly correct: The student noticed the quantities are related.",
                "Mistake sample: The relationship was written in the wrong direction.",
                "Fix it: Explain which quantity changes based on the other and rewrite the rule.",
            ],
        )
    return (
        "Error Analysis Block",
        f"A student used part of {context_piece} correctly but made one reasoning mistake before the final answer.",
        [
            "Nearly correct: One part of the setup is useful.",
            "Mistake sample: A key label, step, or rule was misread.",
            "Fix it: Name the mistake and explain the corrected reasoning.",
        ],
    )


def select_session_premium_features(
    deck: dict[str, Any],
    session_key: str,
    session: dict[str, Any],
) -> tuple[list[str], str]:
    slides = session.get("slides", [])
    signal_text = session_signal_text(deck, session_key, slides)
    context_anchor = best_session_context_anchor(deck, session_key)
    rep_hits = count_term_hits(signal_text, ("table", "graph", "equation", "coordinate", "pattern", "relationship", "verbal"))
    error_hits = count_term_hits(signal_text, ("error", "mistake", "incorrect", "confuse", "units", "substitute", "variable", "axis"))
    explanation_hits = count_term_hits(signal_text, ("explain", "justify", "because", "interpret", "reason", "evidence"))
    evidence_hits = count_term_hits(signal_text, ("claim", "evidence", "reasoning", "prove", "justify", "because"))
    decision_hits = count_term_hits(signal_text, ("proportional", "dependent variable", "independent variable", "which quantity", "whether", "category"))
    strategy_hits = count_term_hits(signal_text, ("compare", "strategy", "method", "another way", "different way", "model and formula"))
    discussion_hits = count_term_hits(signal_text, ("discuss", "share", "compare", "explain to", "partner", "defend"))
    problem_count = sum(1 for slide in session_source_slides(deck, session_key) if slide.get("problem_texts"))
    image_count = sum(int(slide.get("image_count", 0) or 0) for slide in session_source_slides(deck, session_key))
    progression_ready = bool(first_slide_by_kind(slides, {"worked_example"})) and sum(
        1 for slide in slides if slide.get("kind") in {"practice", "challenge", "exit_ticket"}
    ) >= 2
    collaborative_ready = bool(first_slide_by_kind(slides, {"practice", "challenge", "quick_review"}, prefer_no_activity=False))

    scores: dict[str, float] = {
        "writing_revolution": 5.5 if explanation_hits >= 2 else 3.2,
        "scaffold_fade": 4.8 if progression_ready else 0,
        "scenario_continuity": 4.5 if context_anchor and problem_count >= 2 else 0,
        "mastery_tracker": 4.2 if session_key == "session_2" else 2.4,
        "turn_and_teach": 3.8 if collaborative_ready and (discussion_hits >= 1 or explanation_hits >= 2) else 0,
        "create_your_own": 4.0 if session_key == "session_2" and problem_count >= 3 else 0,
        "error_analysis": 4.6 if error_hits >= 2 else (3.6 if any(term in signal_text for term in ("formula", "units", "graph", "table", "variable")) else 0),
        "multi_representation": 5.2 if rep_hits >= 3 else (3.8 if rep_hits >= 2 else 0),
        "strategy_comparison": 4.0 if strategy_hits >= 2 or ("formula" in signal_text and "model" in signal_text) else 0,
        "evidence_ladder": 4.9 if evidence_hits >= 2 else (3.8 if explanation_hits >= 2 else 0),
        "decision_tree": 4.4 if decision_hits >= 1 else (3.4 if "if" in signal_text and "then" in signal_text else 0),
        "real_world_transfer": 4.8 if session_key == "session_2" and context_anchor and problem_count >= 2 else (3.6 if session_key == "session_2" and problem_count >= 3 else 0),
    }

    selected: list[str] = []
    interactive_priority = [
        ("multi_representation", rep_hits >= 2),
        ("error_analysis", error_hits >= 1 or problem_count >= 2),
        ("strategy_comparison", strategy_hits >= 1 or explanation_hits >= 2),
        ("evidence_ladder", evidence_hits >= 1 or explanation_hits >= 2),
        ("decision_tree", decision_hits >= 1),
        ("turn_and_teach", collaborative_ready and (discussion_hits >= 1 or explanation_hits >= 1)),
        ("create_your_own", session_key == "session_2" and problem_count >= 2),
        ("real_world_transfer", session_key == "session_2" and problem_count >= 2),
    ]
    supporting_priority = [
        ("scenario_continuity", bool(context_anchor and problem_count >= 2)),
        ("writing_revolution", explanation_hits >= 2),
        ("scaffold_fade", progression_ready),
        ("mastery_tracker", session_key == "session_2"),
    ]
    for feature, condition in interactive_priority + supporting_priority:
        if condition and scores.get(feature, 0) > 0 and feature not in selected:
            selected.append(feature)
        if len(selected) >= 4:
            break

    ranked = sorted(scores.items(), key=lambda item: (-item[1], PREMIUM_FEATURE_OPTIONS.index(item[0])))
    for feature, score in ranked:
        if score <= 0 or feature in selected:
            continue
        selected.append(feature)
        if len(selected) >= 4:
            break

    interactive_selected = [feature for feature in selected if feature in HIGH_AGENCY_PREMIUM_FEATURES]
    if len(interactive_selected) < 2:
        for feature, _condition in interactive_priority:
            if feature in interactive_selected:
                continue
            if scores.get(feature, 0) <= 0 and feature not in {"error_analysis", "multi_representation", "turn_and_teach"}:
                continue
            if feature not in selected:
                if len(selected) >= 4:
                    replacement = next(
                        (existing for existing in reversed(selected) if existing not in HIGH_AGENCY_PREMIUM_FEATURES),
                        None,
                    )
                    if not replacement:
                        continue
                    selected[selected.index(replacement)] = feature
                else:
                    selected.append(feature)
            interactive_selected = [item for item in selected if item in HIGH_AGENCY_PREMIUM_FEATURES]
            if len(interactive_selected) >= 2:
                break

    if len(selected) < 2:
        for fallback in ("error_analysis", "multi_representation", "evidence_ladder", "turn_and_teach", "writing_revolution"):
            if fallback not in selected:
                selected.append(fallback)
            if len(selected) >= 2:
                break

    ordered = unique_nonempty(selected, limit=4)
    ordered.sort(key=lambda feature: (0 if feature in HIGH_AGENCY_PREMIUM_FEATURES else 1, PREMIUM_FEATURE_OPTIONS.index(feature)))
    return ordered[:4], context_anchor


def apply_locked_architecture_labels(session: dict[str, Any]) -> None:
    slides = session.get("slides", [])
    collaborative_assigned = False
    independent_assigned = False
    for slide in slides:
        kind = slide.get("kind", "")
        if kind in LOCKED_SECTION_LABELS:
            slide["section"] = LOCKED_SECTION_LABELS[kind]
        if kind == "be_curious" and (not slide.get("title") or "curious" in slide.get("title", "").lower()):
            slide["title"] = "Notice + Wonder"
        if kind == "worked_example":
            slide["section"] = "Guided Practice"
        elif kind == "quick_review":
            slide["section"] = "Review"
        elif kind in {"practice", "challenge"}:
            phase = normalize_whitespace(slide.get("practice_phase", "")).lower()
            if "together" in phase:
                slide["section"] = "Collaborative Practice"
                collaborative_assigned = True
                continue
            if "independent" in phase:
                slide["section"] = "Independent Practice"
                independent_assigned = True
                continue
            if not collaborative_assigned:
                slide["section"] = "Collaborative Practice"
                collaborative_assigned = True
            else:
                slide["section"] = "Independent Practice"
                independent_assigned = True
        elif kind == "reflection":
            slide["section"] = "Reflection"
        elif kind == "exit_ticket":
            slide["section"] = "Exit Ticket"
    if not independent_assigned:
        for slide in reversed(slides):
            if slide.get("kind") in {"practice", "challenge", "quick_review"}:
                slide["section"] = "Independent Practice"
                break


def reorder_locked_architecture_sequence(session: dict[str, Any]) -> None:
    order = {
        "cover": 0,
        "quick_review": 1,
        "be_curious": 2,
        "learning_target": 3,
        "vocabulary": 4,
        "guided_notes": 5,
        "worked_example": 6,
        "practice": 7,
        "challenge": 8,
        "reflection": 9,
        "exit_ticket": 10,
    }
    slides = session.get("slides", [])
    indexed = list(enumerate(slides))
    indexed.sort(key=lambda item: (order.get(item[1].get("kind", ""), 50), item[0]))
    session["slides"] = [slide for _index, slide in indexed]


def apply_writing_revolution_layer(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    for slide in session.get("slides", []):
        if slide.get("kind") not in {"be_curious", "guided_notes", "worked_example", "practice", "reflection", "exit_ticket", "challenge"}:
            continue
        frames = contextual_sentence_frames(deck, slide, context_anchor=context_anchor)
        slide["sentence_starters"] = unique_nonempty(frames + list(slide.get("sentence_starters", [])), limit=5)


def apply_scenario_continuity(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    if not context_anchor:
        return
    context_copy = context_phrase(context_anchor)
    for slide in session.get("slides", []):
        if slide.get("kind") not in {"guided_notes", "worked_example", "practice", "challenge", "reflection", "exit_ticket"}:
            continue
        slide["context_anchor"] = context_anchor
        subtitle = normalize_whitespace(slide.get("subtitle", ""))
        if subtitle and context_copy.lower() not in subtitle.lower() and len(subtitle) <= 145:
            slide["subtitle"] = truncate_text(f"{subtitle} Keep using {context_copy.lower()} as the context.", 220)


def apply_scaffold_fade(session: dict[str, Any]) -> None:
    slides = session.get("slides", [])
    worked = first_slide_by_kind(slides, {"worked_example"})
    practice_slides = [slide for slide in slides if slide.get("kind") in {"practice", "challenge", "exit_ticket"}]
    if worked:
        worked["practice_phase"] = worked.get("practice_phase") or "Fully Guided"
        if not normalize_whitespace(worked.get("subtitle", "")):
            worked["subtitle"] = "Fully guided modeling with the exact lesson idea."
    if practice_slides:
        practice_slides[0]["practice_phase"] = practice_slides[0].get("practice_phase") or "Try It Together"
        practice_slides[0]["section"] = "Collaborative Practice"
        if "independent" not in practice_slides[0].get("subtitle", "").lower():
            practice_slides[0]["subtitle"] = truncate_text(
                practice_slides[0].get("subtitle", "") or "Work with support, compare moves, and explain the strategy.",
                220,
            )
    if len(practice_slides) > 1:
        practice_slides[1]["practice_phase"] = practice_slides[1].get("practice_phase") or "Now Solve Independently"
        practice_slides[1]["section"] = "Independent Practice"
        if not normalize_whitespace(practice_slides[1].get("subtitle", "")) or "together" in practice_slides[1].get("subtitle", "").lower():
            practice_slides[1]["subtitle"] = "Use the same skill focus on your own and show how you know."


def apply_error_analysis(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"worked_example", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    title, text, items = common_mistake_payload(best_source_prompt_text(source_records, limit=240) or target.get("primary_text", ""), context_anchor)
    target["premium_layout"] = "error_analysis"
    target["premium_title"] = title
    target["premium_text"] = text
    target["premium_items"] = items
    clear_activity_fields(target)


def apply_multi_representation(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"guided_notes", "practice"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"guided_notes", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    text = best_source_prompt_text(source_records, limit=260) or target.get("primary_text", "")
    target["premium_layout"] = "multi_representation"
    target["premium_title"] = "Multi-Representation Builder"
    target["premium_text"] = truncate_text(
        f"Use the source task to build connected representations{f' from the {context_anchor.lower()} context' if context_anchor else ''}.",
        220,
    )
    target["premium_items"] = multi_representation_labels(text)
    clear_activity_fields(target)


def apply_strategy_comparison(deck: dict[str, Any], session: dict[str, Any]) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"guided_notes", "challenge", "practice"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"guided_notes", "challenge", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    text = best_source_prompt_text(source_records, limit=240) or target.get("primary_text", "")
    target["premium_layout"] = "strategy_comparison"
    target["premium_title"] = "Strategy Comparison Panel"
    target["premium_text"] = "Compare more than one way to think about the lesson so you can choose the strongest fit."
    target["premium_table"] = [["Strategy", "When to use it", "What it shows best", "Possible mistake"]] + strategy_comparison_rows(text)
    clear_activity_fields(target)


def apply_evidence_ladder(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"worked_example", "practice", "challenge"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"worked_example", "practice", "challenge"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    source_prompt = best_source_prompt_text(source_records, limit=220) or target.get("primary_text", "")
    context_copy = context_phrase(context_anchor) if context_anchor else "the source task"
    target["premium_layout"] = "evidence_ladder"
    target["premium_title"] = "Evidence Ladder"
    target["premium_text"] = truncate_text(
        f"Build a claim you can defend from {context_copy.lower()} by naming the clue, the math evidence, and the reasoning.",
        220,
    )
    target["premium_items"] = [
        truncate_text(source_prompt or f"What is the strongest claim you can make about {context_copy.lower()}?", 88),
        "Clue from the representation or numbers",
        "Math evidence that supports the claim",
        "Reasoning that proves the claim fits",
    ]
    clear_activity_fields(target)


def apply_decision_tree(deck: dict[str, Any], session: dict[str, Any]) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"guided_notes", "quick_review", "practice"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"guided_notes", "quick_review", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    text = best_source_prompt_text(source_records, limit=240) or target.get("primary_text", "")
    target["premium_layout"] = "decision_tree"
    target["premium_title"] = "Math Cheat Code"
    target["premium_text"] = "Follow the decision path before you commit to a strategy."
    target["premium_items"] = decision_tree_steps(text)
    clear_activity_fields(target)


def apply_turn_and_teach(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"practice", "quick_review", "challenge"}, prefer_no_activity=False) or first_unused_premium_slide(session.get("slides", []), {"practice", "challenge"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    source_prompt = best_source_prompt_text(source_records, limit=180) or target.get("primary_text", "")
    prompt = f"Turn and teach: explain how you would solve {context_anchor.lower()} using the exact lesson strategy." if context_anchor else f"Turn and teach: explain how you would solve this source problem and why the strategy fits."
    if "graph" in source_prompt.lower():
        prompt = "Turn and teach: explain what the graph shows and how you know."
    elif "table" in source_prompt.lower():
        prompt = "Turn and teach: explain how the table values connect and what pattern you notice."
    target["partner_prompt"] = truncate_text(prompt, 170)
    target["premium_layout"] = target.get("premium_layout", "") or "turn_and_teach"


def apply_create_your_own(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"exit_ticket", "challenge", "practice"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"exit_ticket", "challenge", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    text = best_source_prompt_text(source_records, limit=220) or target.get("primary_text", "")
    target["premium_layout"] = "create_your_own"
    target["premium_title"] = "Create-Your-Own Task"
    target["premium_text"] = truncate_text(
        f"Create a new example that still matches the lesson idea{f' and stays in the {context_anchor.lower()} context' if context_anchor else ''}.",
        220,
    )
    target["premium_items"] = create_your_own_labels(text)
    clear_activity_fields(target)


def apply_real_world_transfer(deck: dict[str, Any], session: dict[str, Any], context_anchor: str) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"challenge", "exit_ticket", "practice"}, prefer_no_activity=True) or first_unused_premium_slide(session.get("slides", []), {"challenge", "exit_ticket", "practice"})
    if not target:
        return
    source_records = source_slides_from_numbers(deck, target.get("source_slide_numbers", []))
    source_prompt = best_source_prompt_text(source_records, limit=220) or target.get("primary_text", "")
    context_copy = context_phrase(context_anchor) if context_anchor else "a new real-world situation"
    target["premium_layout"] = "real_world_transfer"
    target["premium_title"] = "Real-World Transfer Studio"
    target["premium_text"] = truncate_text(
        f"Use the lesson idea in {context_copy.lower()}, then explain what stays true when the situation changes.",
        220,
    )
    target["premium_items"] = [
        truncate_text(source_prompt or f"Start from the source idea in {context_copy.lower()}.", 84),
        "What stays true from the original lesson?",
        "What changes in the new situation?",
        "How would you justify the transfer?",
    ]
    clear_activity_fields(target)


def apply_mastery_tracker(session: dict[str, Any]) -> None:
    target = first_unused_premium_slide(session.get("slides", []), {"reflection"})
    if not target:
        return
    target["premium_layout"] = "mastery_tracker"
    target["premium_title"] = "Reflection + Mastery Tracker"
    target["premium_text"] = "Pause, rate your confidence, name what you understand, and prove one idea with a mini-check."
    target["premium_items"] = [
        "Confidence: Not yet / Getting there / Ready to teach",
        "What I understand now",
        "What I still need help with",
        "Prove-it mini check",
    ]


def apply_premium_decision_layer(deck: dict[str, Any], session_key: str, session: dict[str, Any]) -> None:
    features, context_anchor = select_session_premium_features(deck, session_key, session)
    session["premium_features"] = features
    session["context_anchor"] = context_anchor
    apply_locked_architecture_labels(session)

    if "scenario_continuity" in features:
        apply_scenario_continuity(deck, session, context_anchor)
    if "scaffold_fade" in features:
        apply_scaffold_fade(session)
    if "writing_revolution" in features:
        apply_writing_revolution_layer(deck, session, context_anchor)
    if "error_analysis" in features:
        apply_error_analysis(deck, session, context_anchor)
    if "multi_representation" in features:
        apply_multi_representation(deck, session, context_anchor)
    if "strategy_comparison" in features:
        apply_strategy_comparison(deck, session)
    if "evidence_ladder" in features:
        apply_evidence_ladder(deck, session, context_anchor)
    if "decision_tree" in features:
        apply_decision_tree(deck, session)
    if "turn_and_teach" in features:
        apply_turn_and_teach(deck, session, context_anchor)
    if "create_your_own" in features:
        apply_create_your_own(deck, session, context_anchor)
    if "real_world_transfer" in features:
        apply_real_world_transfer(deck, session, context_anchor)
    if "mastery_tracker" in features:
        apply_mastery_tracker(session)


def flagship_activity_spec(plan_slide: dict[str, Any]) -> dict[str, Any] | None:
    spec = plan_slide.get("flagship_activity")
    if not isinstance(spec, dict):
        return None
    return spec if normalize_whitespace(str(spec.get("title", "")) or str(spec.get("type", ""))) else None


def flagship_focus_phrase(
    deck: dict[str, Any],
    plan_slide: dict[str, Any],
    *,
    context_anchor: str = "",
) -> str:
    blob = workbook_source_blob(plan_slide)
    rep_terms = representation_terms_for_text(blob)
    if has_data_analysis_context(blob):
        return "Data Evidence"
    if "volume" in blob and ("prism" in blob or "rectangular" in blob):
        return "Volume Reasoning"
    if "area" in blob:
        return "Area Reasoning"
    if {"table", "graph"} <= set(rep_terms):
        return "Table + Graph"
    if {"table", "equation"} <= set(rep_terms):
        return "Table + Rule"
    if "graph" in rep_terms and "equation" in rep_terms:
        return "Graph + Rule"
    vocab_items = source_vocabulary(deck, plan_slide.get("source_slide_numbers", []), limit=2)
    if vocab_items:
        return display_term_label(vocab_items[0].get("word", ""))
    if context_anchor:
        return truncate_text(context_anchor, 28)
    return truncate_text(source_problem_focus_phrase(plan_slide).replace("source ", "").title(), 28) or "Source Strategy"


def build_flagship_session_data(
    deck: dict[str, Any],
    session_key: str,
    session: dict[str, Any],
) -> dict[str, Any]:
    if not deck.get("keyword_candidates"):
        keyword_candidates = source_term_candidates(session_source_slides(deck, session_key) or deck.get("slides", []), limit=8)
        if not keyword_candidates:
            keyword_candidates = unique_nonempty(
                (
                    normalize_whitespace(slide.get("title", ""))
                    for slide in deck.get("slides", [])
                    if normalize_whitespace(slide.get("title", ""))
                ),
                limit=6,
            )
        deck = {**deck, "keyword_candidates": keyword_candidates}
    slides = session.get("slides", [])
    source_records = session_source_slides(deck, session_key)
    source_numbers = [slide["slide_number"] for slide in source_records if slide.get("slide_number")]
    signal_text = session_signal_text(deck, session_key, slides)
    problem_count = sum(1 for slide in source_records if slide.get("problem_texts"))
    vocab_items = session_esol_vocabulary(deck, source_numbers[:6] or source_numbers, limit=6)
    support_score = sum(
        1
        for condition in (
            problem_count >= 2,
            len(vocab_items) >= 3,
            count_term_hits(signal_text, ("table", "graph", "equation", "coordinate", "representation", "diagram", "model")) >= 2,
            count_term_hits(signal_text, ("compare", "justify", "because", "evidence", "strategy", "claim", "explain")) >= 2,
            count_term_hits(signal_text, ("error", "mistake", "incorrect", "fix", "units", "variable", "formula")) >= 1,
            bool(best_session_context_anchor(deck, session_key)),
        )
        if condition
    )
    target_count = 0
    if support_score >= 2 and source_numbers:
        target_count = 2
        if support_score >= 4 and len(slides) >= 11:
            target_count += 1
        if support_score >= 5 and problem_count >= 3 and len(slides) >= 13:
            target_count += 1
    elif support_score == 1 and source_numbers:
        target_count = 1
    return {
        "deck": deck,
        "session_key": session_key,
        "session": session,
        "slides": slides,
        "source_records": source_records,
        "source_numbers": source_numbers,
        "signal_text": signal_text,
        "profile": session_math_profile(deck, session_key),
        "context_anchor": best_session_context_anchor(deck, session_key),
        "vocab_items": vocab_items,
        "problem_count": problem_count,
        "target_count": clamp_int(target_count, 0, 4),
        "support_score": support_score,
    }


def flagship_slide_can_host_activity(plan_slide: dict[str, Any]) -> bool:
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    if role and role not in FLAGSHIP_ACTIVITY_RENDERABLE_TEMPLATE_ROLES:
        return False
    return plan_slide.get("kind") in {
        "be_curious",
        "vocabulary",
        "guided_notes",
        "worked_example",
        "practice",
        "quick_review",
        "challenge",
        "reflection",
        "exit_ticket",
    }


def flagship_activity_candidate_slides(
    session: dict[str, Any],
    placement: str,
    *,
    used_indices: set[int] | None = None,
) -> list[int]:
    slides = session.get("slides", [])
    if not slides:
        return []
    used_indices = used_indices or set()
    midpoint = len(slides) // 2
    candidates: list[tuple[int, int]] = []
    for index, slide in enumerate(slides):
        if index in used_indices or not slide.get("source_slide_numbers") or not flagship_slide_can_host_activity(slide):
            continue
        kind = slide.get("kind", "")
        role = normalize_whitespace(slide.get("template_role", ""))
        score = 0
        if placement in {"middle", "middle_bonus"}:
            score -= abs(index - midpoint)
            if kind in {"guided_notes", "worked_example", "practice"}:
                score += 10
            if role in {"guided_practice", "drag_sort"}:
                score += 6
        elif placement in {"end", "end_bonus"}:
            score += index * 2
            if kind in {"challenge", "practice", "reflection", "exit_ticket"}:
                score += 10
            if role in {"error_analysis", "turn_and_teach"}:
                score += 6
        elif placement == "launch_bonus":
            score -= index
            if kind in {"be_curious", "vocabulary"}:
                score += 12
        else:
            score += 4 if kind in {"practice", "challenge", "quick_review"} else 0
            if role in {"drag_sort", "turn_and_teach"}:
                score += 3
        if normalize_whitespace(slide.get("premium_layout", "")):
            score -= 2
        candidates.append((score, index))
    return [index for _score, index in sorted(candidates, key=lambda item: (-item[0], item[1]))]


def choose_flagship_activity_type_for_slide(
    session_data: dict[str, Any],
    plan_slide: dict[str, Any],
    *,
    placement: str,
    used_types: set[str],
) -> tuple[str, int]:
    blob = workbook_source_blob(plan_slide)
    rep_terms = representation_terms_for_text(blob)
    vocab_count = len(source_vocabulary(session_data["deck"], plan_slide.get("source_slide_numbers", []), limit=5))
    problem_ready = len(source_problem_candidates(source_slides_from_numbers(session_data["deck"], plan_slide.get("source_slide_numbers", [])), limit=4)) >= 2
    kind = plan_slide.get("kind", "")
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    compare_hits = count_term_hits(blob, ("compare", "justify", "because", "evidence", "strategy", "claim", "same", "different"))
    mistake_hits = count_term_hits(blob, ("error", "mistake", "incorrect", "fix", "units", "formula", "variable"))
    scores: dict[str, int] = {}

    if mistake_hits or role == "error_analysis":
        scores["Error Analysis"] = 9 + mistake_hits
        scores["Fix the Mistake"] = 8 + mistake_hits
    if {"table", "graph"} <= set(rep_terms) and ("equation" in rep_terms or "rule" in blob):
        scores["Table / Graph / Equation connection task"] = 11
    if len(rep_terms) >= 2 or any(term in blob for term in ("model", "diagram", "representation")):
        scores["Complete the Representation"] = 8 + len(rep_terms)
    if any(term in blob for term in ("pattern", "rule", "table", "equation", "grows by", "changes by")):
        scores["Build the Rule / Find the Pattern"] = 8
    if kind in {"worked_example", "practice"} and problem_ready:
        scores["My Turn / Your Turn"] = 8 if placement == "middle" else 6
        scores["Explain the Strategy"] = 7
    if compare_hits >= 2:
        scores["Compare and Justify"] = 8 + compare_hits
        scores["Mini Math Debate"] = 7 + compare_hits
    if kind == "be_curious" or placement == "launch_bonus":
        scores["Notice / Wonder extension"] = 8
    if vocab_count >= 3 or kind == "vocabulary":
        scores["Vocabulary in Action"] = 7 + vocab_count
    if session_data.get("context_anchor") and placement in {"end", "end_bonus"}:
        scores["Real-World Transfer"] = 7
    if kind in {"practice", "challenge"} and compare_hits >= 1:
        scores["Sort / Match / Classify"] = 6 + compare_hits
        scores["Which One Doesn't Belong"] = 6
        scores["Example vs. Non-Example"] = 5

    if placement == "middle":
        for activity_type in (
            "Table / Graph / Equation connection task",
            "Complete the Representation",
            "Build the Rule / Find the Pattern",
            "My Turn / Your Turn",
            "Vocabulary in Action",
        ):
            if activity_type in scores:
                scores[activity_type] += 2
    if placement in {"end", "end_bonus"}:
        for activity_type in (
            "Error Analysis",
            "Fix the Mistake",
            "Compare and Justify",
            "Mini Math Debate",
            "Real-World Transfer",
            "Explain the Strategy",
        ):
            if activity_type in scores:
                scores[activity_type] += 2

    ranked = sorted(
        (
            (activity_type, score)
            for activity_type, score in scores.items()
            if activity_type in FLAGSHIP_ACTIVITY_TYPES and activity_type not in used_types
        ),
        key=lambda item: (-item[1], FLAGSHIP_ACTIVITY_TYPES.index(item[0])),
    )
    if not ranked:
        return "", 0
    return ranked[0]


def choose_flagship_activity_mix(session_data: dict[str, Any]) -> list[dict[str, Any]]:
    target_count = int(session_data.get("target_count", 0) or 0)
    if target_count <= 0:
        return []
    placements = ["middle", "end"]
    if target_count >= 3:
        placements.append("launch_bonus" if session_data.get("support_score", 0) >= 4 else "middle_bonus")
    if target_count >= 4:
        placements.append("end_bonus")
    selected: list[dict[str, Any]] = []
    used_indices: set[int] = set()
    used_types: set[str] = set()
    for placement in placements[:target_count]:
        for slide_index in flagship_activity_candidate_slides(session_data["session"], placement, used_indices=used_indices):
            plan_slide = session_data["slides"][slide_index]
            activity_type, score = choose_flagship_activity_type_for_slide(
                session_data,
                plan_slide,
                placement=placement,
                used_types=used_types,
            )
            if not activity_type or score < 6:
                continue
            footprint = FLAGSHIP_ACTIVITY_FOOTPRINT_FULL if (
                placement in {"middle", "end"}
                and (
                    role := normalize_whitespace(plan_slide.get("template_role", ""))
                ) in {"drag_sort", "error_analysis", "turn_and_teach", "guided_practice"}
            ) else FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED
            selected.append(
                {
                    "placement": placement,
                    "slide_index": slide_index,
                    "type": activity_type,
                    "estimatedSlideFootprint": footprint,
                }
            )
            used_indices.add(slide_index)
            used_types.add(activity_type)
            break
    return selected


def flagship_support_lines(
    activity_type: str,
    *,
    deck: dict[str, Any],
    plan_slide: dict[str, Any],
    context_anchor: str,
    vocabulary_used: list[str],
) -> list[str]:
    supports: list[str] = []
    if vocabulary_used:
        supports.append("Word bank: " + ", ".join(vocabulary_used[:3]))
    frames = contextual_sentence_frames(deck, plan_slide, context_anchor=context_anchor)
    if frames:
        supports.append("Sentence starter: " + truncate_text(frames[0], 92))
    if activity_type in {"Error Analysis", "Fix the Mistake", "My Turn / Your Turn", "Explain the Strategy", "Real-World Transfer"}:
        tips = problem_tip_lines(plan_slide, limit=2)
        if tips:
            supports.append("Worked-example reminder: " + truncate_text(tips[0], 92))
    return unique_nonempty(supports, limit=3)


def build_flagship_activity_spec(session_data: dict[str, Any], mix_item: dict[str, Any]) -> dict[str, Any]:
    slide_index = int(mix_item["slide_index"])
    plan_slide = session_data["slides"][slide_index]
    source_numbers = [int(num) for num in plan_slide.get("source_slide_numbers", []) if isinstance(num, int)]
    source_records = source_slides_from_numbers(session_data["deck"], source_numbers)
    seeded = dict(plan_slide)
    seeded["source_problem_cards"] = source_problem_candidates(source_records, limit=4)
    seeded["vocabulary"] = source_vocabulary(session_data["deck"], source_numbers, limit=4)
    focus_phrase = flagship_focus_phrase(session_data["deck"], seeded, context_anchor=session_data["context_anchor"])
    activity_type = mix_item["type"]
    vocabulary_used = [
        display_term_label(item.get("word", ""))
        for item in source_vocabulary(session_data["deck"], source_numbers, limit=4)
        if display_term_label(item.get("word", ""))
    ]
    supports = flagship_support_lines(
        activity_type,
        deck=session_data["deck"],
        plan_slide=seeded,
        context_anchor=session_data["context_anchor"],
        vocabulary_used=vocabulary_used,
    )
    prompts: list[str] = []
    title = ""
    purpose = ""
    directions = ""

    if activity_type in {"Error Analysis", "Fix the Mistake"}:
        steps, fix_text, why_prompt = source_error_analysis_content(seeded)
        title = f"{focus_phrase} {'Error Analysis' if activity_type == 'Error Analysis' else 'Fix the Mistake'}"
        purpose = "Use the exact source details to spot the break in reasoning and repair it with evidence."
        directions = "Read the source-aligned steps, identify the move that does not fit, fix it, and explain why the corrected reasoning matches the lesson."
        prompts = steps[:3] + [fix_text, why_prompt]
    elif activity_type == "Table / Graph / Equation connection task":
        rep_labels = multi_representation_labels(workbook_source_blob(seeded))
        title = f"{focus_phrase} Connection Task"
        purpose = "Connect the source representations so students explain how the same relationship appears in more than one form."
        directions = "Match or build the table, graph, and rule pieces that belong together, then explain how the source relationship stays consistent."
        prompts = rep_labels + source_fact_candidates(source_records, limit=2)
    elif activity_type == "Complete the Representation":
        rep_labels = multi_representation_labels(workbook_source_blob(seeded))
        title = f"Complete the {focus_phrase} Representation"
        purpose = "Use the source model, values, or labels to finish the missing representation without losing the lesson meaning."
        directions = "Use the source clues to complete the missing representation, then write one sentence that proves how the finished parts connect."
        prompts = rep_labels + source_fact_candidates(source_records, limit=2)
    elif activity_type == "Build the Rule / Find the Pattern":
        title = f"Build the {focus_phrase} Rule"
        purpose = "Use the source pattern or changing values to build a reusable rule and defend how it fits."
        directions = "Look for the pattern in the source values, build the rule that matches it, test the rule with one more example, and explain why it works."
        prompts = source_fact_candidates(source_records, limit=4) + representation_terms_for_text(workbook_source_blob(seeded))
    elif activity_type == "My Turn / Your Turn":
        title = f"My Turn / Your Turn: {focus_phrase}"
        purpose = "Move from modeled work to independent reasoning while keeping the source strategy visible."
        directions = "Study the modeled source problem, solve the similar follow-up on your own, and compare how the strategy stayed the same."
        prompts = [
            truncate_display_copy(source_problem_statement(seeded), 62),
            truncate_display_copy(similar_problem_statement(seeded), 62),
            truncate_display_copy(source_compare_prompt(seeded), 62),
        ]
    elif activity_type == "Explain the Strategy":
        title = f"Explain the {focus_phrase} Strategy"
        purpose = "Turn the solve path into a student-friendly explanation that uses lesson evidence."
        directions = "Order the strategy moves, then explain which clue, value, or representation proves the strategy is the right fit."
        prompts = problem_tip_lines(seeded, limit=3) + [source_compare_prompt(seeded)]
    elif activity_type == "Mini Math Debate":
        title = f"{focus_phrase} Mini Math Debate"
        purpose = "Compare competing claims or strategies and justify the stronger one with source evidence."
        directions = "Compare the claims, decide which argument is stronger, and justify your decision with a source detail, value, or representation."
        prompts = [
            truncate_display_copy(source_problem_statement(seeded), 58),
            truncate_display_copy(similar_problem_statement(seeded), 58),
            truncate_display_copy(source_compare_prompt(seeded), 58),
        ]
    elif activity_type == "Vocabulary in Action":
        title = f"Vocabulary in Action: {focus_phrase}"
        purpose = "Use the actual lesson terms in a source-aligned context instead of treating vocabulary as a standalone page."
        directions = "Match each lesson word to a source clue, model, or example, then use at least two words in a complete explanation."
        prompts = vocabulary_used[:4] or source_term_candidates(source_records, limit=4)
    elif activity_type == "Real-World Transfer":
        title = f"{focus_phrase} Real-World Transfer"
        purpose = "Carry the source strategy into a nearby context without lowering the rigor or losing the original structure."
        directions = "Use the source strategy in the new situation, solve or build the transfer task, and explain which part of the lesson transfers directly."
        prompts = [
            truncate_display_copy(source_problem_statement(seeded), 60),
            truncate_display_copy(third_problem_statement(seeded), 60),
            truncate_display_copy(context_phrase(session_data["context_anchor"]), 40),
        ]
    elif activity_type in {"Compare and Justify", "Which One Doesn't Belong", "Example vs. Non-Example", "Sort / Match / Classify"}:
        title = f"{focus_phrase} {activity_type}"
        purpose = "Sort, compare, or classify the source ideas so students have to justify the rule they used."
        directions = "Sort or compare the source cards, decide which idea or example belongs where, and justify the rule or decision with evidence from the lesson."
        prompts = source_drag_sort_pieces(seeded) + source_sort_labels(seeded)
    else:
        title = f"{focus_phrase} Notice + Wonder"
        purpose = "Extend the launch so students use source clues to notice, wonder, predict, and explain."
        directions = "Reveal the source clues one at a time, record what you notice and wonder, and explain which clue matters most for solving the lesson problem."
        notice_kernels, wonder_kernels = source_be_curious_kernels(source_records)
        prompts = notice_kernels[:2] + wonder_kernels[:2]

    prompts = unique_nonempty(
        (trim_dangling_display_text(truncate_display_copy(prompt, 64)) for prompt in prompts),
        limit=6,
    )
    if not prompts:
        prompts = unique_nonempty(source_drag_sort_pieces(seeded) + source_fact_candidates(source_records, limit=4), limit=4)
    if not prompts:
        prompts = ["Use a source clue", "Build the move", "Check with evidence", "Explain why"]

    return {
        "id": f"{session_data['session_key']}-flagship-{slide_index + 1}-{slugify(activity_type)[:18]}",
        "session": session_data["session_key"],
        "type": activity_type,
        "placement": mix_item["placement"],
        "title": truncate_text(title, 120),
        "purpose": truncate_text(purpose, 220),
        "sourceAnchors": source_numbers,
        "vocabularyUsed": unique_nonempty(vocabulary_used, limit=4),
        "directions": truncate_text(directions, 280),
        "prompts": prompts,
        "supports": supports,
        "layoutVariant": FLAGSHIP_ACTIVITY_TYPE_TO_LAYOUT.get(activity_type, "strategy_workspace"),
        "estimatedSlideFootprint": mix_item["estimatedSlideFootprint"],
        "slideIndex": slide_index,
    }


def build_flagship_activity_specs(session_data: dict[str, Any], mix: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [build_flagship_activity_spec(session_data, item) for item in mix]


def flagship_activity_has_generic_markers(spec: dict[str, Any]) -> bool:
    core_blob = normalize_whitespace(
        " ".join(
            [
                str(spec.get("title", "")),
                str(spec.get("purpose", "")),
                str(spec.get("directions", "")),
                " ".join(str(item) for item in spec.get("prompts", [])),
            ]
        )
    ).lower()
    if "___" in core_blob:
        return True
    return any(
        re.search(rf"(?<![a-z0-9]){re.escape(marker)}(?![a-z0-9])", core_blob)
        for marker in FLAGSHIP_ACTIVITY_GENERIC_MARKERS
    )


def flagship_answer_check_text(spec: dict[str, Any]) -> str:
    activity_type = spec.get("type", "")
    if activity_type in {"Error Analysis", "Fix the Mistake"}:
        return "Use the exact source values, labels, and reasoning to prove the corrected step."
    if activity_type in {"Table / Graph / Equation connection task", "Complete the Representation"}:
        return "Each piece should connect to the same source relationship and representation."
    if activity_type in {"Compare and Justify", "Mini Math Debate", "Which One Doesn't Belong"}:
        return "Your choice should be supported by a clear source detail, value, or representation."
    if activity_type == "Vocabulary in Action":
        return "Each term should match the source clue or explanation that best shows its meaning."
    return "Use a source detail, lesson word, or worked-example move to explain why your activity choices fit."


def flagship_activity_payload(spec: dict[str, Any]) -> dict[str, Any]:
    return {
        "activity_name": truncate_text(spec.get("title", ""), 120),
        "activity_family": FLAGSHIP_ACTIVITY_TYPE_TO_FAMILY.get(spec.get("type", ""), "build_construct"),
        "activity_instructions": truncate_text(spec.get("directions", ""), 280),
        "movable_pieces": unique_nonempty(spec.get("prompts", []), limit=6),
        "answer_check": truncate_text(flagship_answer_check_text(spec), 220),
    }


def apply_flagship_activity_spec_to_slide(plan_slide: dict[str, Any], spec: dict[str, Any]) -> None:
    plan_slide["flagship_activity"] = dict(spec)
    plan_slide["flagship_activity_mode"] = spec.get("estimatedSlideFootprint", FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED)
    payload = flagship_activity_payload(spec)
    plan_slide["activity_name"] = payload["activity_name"]
    plan_slide["activity_family"] = payload["activity_family"]
    plan_slide["activity_instructions"] = payload["activity_instructions"]
    plan_slide["movable_pieces"] = payload["movable_pieces"]
    plan_slide["answer_check"] = payload["answer_check"]
    if spec.get("placement") in {"end", "end_bonus"} and plan_slide.get("kind") in {"challenge", "reflection", "exit_ticket"}:
        plan_slide["practice_phase"] = plan_slide.get("practice_phase") or "Flagship reasoning move"


def plan_flagship_activities_for_session(session_data: dict[str, Any]) -> tuple[list[dict[str, Any]], list[str]]:
    if int(session_data.get("target_count", 0) or 0) <= 0:
        return [], ["source content was too thin to require a flagship activity add-on"]
    mix = choose_flagship_activity_mix(session_data)
    specs = build_flagship_activity_specs(session_data, mix)
    fallback_notes: list[str] = []
    if len(specs) < session_data["target_count"]:
        fallback_notes.append("reduced flagship activity count because the strongest slide anchors were limited")
    if not specs:
        fallback_notes.append("kept the original notebook activity set because no lesson-anchored flagship activity was stronger")
    return specs, fallback_notes


def apply_flagship_activity_layer(deck: dict[str, Any], session_key: str, session: dict[str, Any]) -> None:
    session_data = build_flagship_session_data(deck, session_key, session)
    session["flagship_activity_requested_target"] = session_data["target_count"]
    session["flagship_activity_target"] = session_data["target_count"]
    session["flagship_activity_support_score"] = session_data["support_score"]
    session["flagship_activity_fallbacks"] = []
    session["flagship_activities"] = []
    if session_data["target_count"] <= 0:
        session["flagship_activity_fallbacks"] = ["kept the original notebook output because source support for flagship activities was limited"]
        return
    try:
        specs, fallback_notes = plan_flagship_activities_for_session(session_data)
        applied_specs: list[dict[str, Any]] = []
        for spec in specs:
            slide_index = int(spec.get("slideIndex", -1))
            if slide_index < 0 or slide_index >= len(session.get("slides", [])):
                continue
            slide = session["slides"][slide_index]
            apply_flagship_activity_spec_to_slide(slide, spec)
            applied_specs.append(dict(spec))
        session["flagship_activities"] = applied_specs
        session["flagship_activity_target"] = len(applied_specs)
        session["flagship_activity_fallbacks"] = fallback_notes
    except Exception as exc:
        session["flagship_activity_target"] = 0
        session["flagship_activity_fallbacks"] = [f"kept the original notebook output because flagship activity planning fell back safely ({truncate_text(str(exc), 120)})"]


def normalize_flagship_activity_spec(spec: dict[str, Any], *, preserve_source_prompt: bool = True) -> dict[str, Any]:
    normalized = dict(spec)
    normalized["title"] = truncate_text(publisher_copyedit_text(str(spec.get("title", ""))), 120)
    normalized["purpose"] = truncate_text(publisher_copyedit_text(str(spec.get("purpose", "")), preserve_source_prompt=preserve_source_prompt), 220)
    normalized["directions"] = truncate_text(publisher_copyedit_text(str(spec.get("directions", "")), preserve_source_prompt=preserve_source_prompt), 280)
    normalized["sourceAnchors"] = [int(num) for num in spec.get("sourceAnchors", []) if isinstance(num, int)]
    normalized["vocabularyUsed"] = [
        display_term_label(publisher_copyedit_text(str(item)))
        for item in unique_nonempty(spec.get("vocabularyUsed", []), limit=4)
        if display_term_label(publisher_copyedit_text(str(item)))
    ]
    normalized["prompts"] = publisher_copyedit_list(spec.get("prompts", []), preserve_source_prompt=preserve_source_prompt, limit=6)
    normalized["supports"] = publisher_copyedit_list(spec.get("supports", []), preserve_source_prompt=preserve_source_prompt, limit=3)
    normalized["layoutVariant"] = truncate_text(normalize_whitespace(str(spec.get("layoutVariant", ""))), 48)
    normalized["estimatedSlideFootprint"] = (
        spec.get("estimatedSlideFootprint", FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED)
        if spec.get("estimatedSlideFootprint", FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED) in {FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED, FLAGSHIP_ACTIVITY_FOOTPRINT_FULL}
        else FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED
    )
    return normalized


def validate_flagship_activities(session: dict[str, Any], *, deck: dict[str, Any] | None = None) -> list[str]:
    issues: list[str] = []
    target = int(session.get("flagship_activity_target", 0) or 0)
    specs = [spec for spec in session.get("flagship_activities", []) if isinstance(spec, dict)]
    if target >= 2 and len(specs) < 2:
        issues.append("session is missing the required flagship activity count for the available lesson support")
    if target > 0 and len(specs) < min(target, 4):
        issues.append("session fell below the planned flagship activity target")
    for spec in specs:
        anchors = [int(num) for num in spec.get("sourceAnchors", []) if isinstance(num, int)]
        if not anchors:
            issues.append(f"{spec.get('id', 'flagship')}: flagship activity is missing source anchors")
        if deck:
            max_slide_number = len(deck.get("slides", []))
            if any(num < 1 or num > max_slide_number for num in anchors):
                issues.append(f"{spec.get('id', 'flagship')}: flagship activity uses invalid source anchors")
        if flagship_activity_has_generic_markers(spec):
            issues.append(f"{spec.get('id', 'flagship')}: flagship activity still contains generic filler")
        if not normalize_whitespace(spec.get("directions", "")) or len(normalize_whitespace(spec.get("directions", "")).split()) < 7:
            issues.append(f"{spec.get('id', 'flagship')}: flagship activity directions are too thin")
        if spec.get("estimatedSlideFootprint") not in {FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED, FLAGSHIP_ACTIVITY_FOOTPRINT_FULL}:
            issues.append(f"{spec.get('id', 'flagship')}: flagship activity footprint is invalid")
    return issues


def ensure_interactive_engagement_layer(
    deck: dict[str, Any],
    session_key: str,
    session: dict[str, Any],
    *,
    library: list[dict[str, Any]],
) -> None:
    slides = session.get("slides", [])
    if not slides:
        return

    target_count = engagement_slide_target(session)
    target_modes = engagement_mode_target(session)
    candidates = [
        slide
        for slide in slides
        if ACTIVITY_KIND_CAPS.get(slide.get("kind", ""), 0) > 0
        and not normalize_whitespace(slide.get("template_role", ""))
        and problem_render_mode(slide) != PROBLEM_RENDER_MODE_FOCUS
    ]
    if not candidates:
        return

    attempts = 0
    while (
        session_engagement_slide_count(session) < target_count
        or len(session_engagement_modes(session)) < target_modes
    ) and attempts < len(candidates):
        attempts += 1
        before_count = session_engagement_slide_count(session)
        before_modes = set(session_engagement_modes(session))
        used_names = {
            normalize_whitespace(slide.get("activity_name", "")).lower()
            for slide in slides
            if normalize_whitespace(slide.get("activity_name", ""))
        }
        improved = False
        for slide in sorted(candidates, key=engagement_rebalance_priority, reverse=True):
            current_name = normalize_whitespace(slide.get("activity_name", "")).lower()
            current_family = normalize_whitespace(slide.get("activity_family", ""))
            current_modes = slide_engagement_modes(slide)
            should_refresh = (
                not has_activity(slide)
                or current_family not in HIGH_AGENCY_ACTIVITY_FAMILIES
                or current_modes.issubset({"notice", "discuss"})
            )
            if not should_refresh:
                continue

            excluded_names = set(used_names)
            if current_name:
                excluded_names.discard(current_name)
            ensure_activity_from_library(
                slide,
                kind=slide.get("kind", "practice"),
                deck=deck,
                library=library,
                excluded_names=excluded_names,
                force_refresh=True,
            )
            if slide.get("kind") == "be_curious":
                ensure_source_aligned_be_curious_activity(slide, deck=deck)
            if slide.get("kind") in PROBLEM_SOLVING_KINDS:
                ensure_source_anchored_problem_activity(slide)
            ensure_peer_discussion_support(slide)
            refresh_problem_activity_for_context(slide)

            after_count = session_engagement_slide_count(session)
            after_modes = set(session_engagement_modes(session))
            new_name = normalize_whitespace(slide.get("activity_name", "")).lower()
            if new_name:
                used_names.add(new_name)
            if (
                after_count > before_count
                or len(after_modes) > len(before_modes)
                or (new_name and new_name != current_name)
            ):
                improved = True
                break
        if not improved:
            break


def flagship_cover_terms(deck: dict[str, Any], source_numbers: list[int], limit: int = 3) -> list[str]:
    terms: list[str] = []
    for item in source_vocabulary(deck, source_numbers, limit=5):
        word = normalize_whitespace(item.get("word", ""))
        if not word:
            continue
        lowered = word.lower()
        if lowered in GENERIC_VOCAB_TERMS or lowered in STOPWORDS:
            continue
        terms.append(word.title())
    for candidate in deck.get("keyword_candidates", []):
        cleaned = normalize_whitespace(candidate)
        lowered = cleaned.lower()
        if not cleaned or lowered in GENERIC_VOCAB_TERMS or lowered in STOPWORDS:
            continue
        terms.append(cleaned.title())
    for title in source_titles_from_numbers(deck, source_numbers, limit=4):
        cleaned = normalize_whitespace(title)
        if not cleaned or len(cleaned) > 28:
            continue
        terms.append(cleaned)
    return unique_nonempty(terms, limit=limit)


def find_activity_record(name: str, library: list[dict[str, Any]]) -> dict[str, Any] | None:
    lowered = name.lower()
    for item in library:
        if item["name"].lower() == lowered:
            return item
    return None


def source_fact_candidates(source_records: list[dict[str, Any]], limit: int = 6) -> list[str]:
    candidates: list[str] = []
    for slide in source_records:
        title = slide.get("title", "")
        raw_items = slide.get("text_items", []) + [slide.get("text", "")]
        for raw_item in raw_items:
            cleaned = clean_source_prompt(raw_item, title)
            if not cleaned or is_generic_slide_text(cleaned, title):
                continue
            if len(cleaned) < 12:
                continue
            if is_problem_like_text(cleaned) or is_answer_like_text(cleaned):
                candidates.append(cleaned)
                continue
            for sentence in re.split(r"(?<=[?.!])\s+", cleaned):
                sentence = clean_source_prompt(sentence, title)
                if len(sentence) < 12 or is_generic_slide_text(sentence, title):
                    continue
                if is_problem_like_text(sentence) or is_answer_like_text(sentence):
                    candidates.append(sentence)
    return unique_nonempty(
        (
            trim_dangling_display_text(truncate_display_copy(item, 88))
            for item in candidates
        ),
        limit=limit,
    )


def source_question_prompts(source_records: list[dict[str, Any]], limit: int = 4) -> list[str]:
    prompts: list[str] = []
    for slide in source_records:
        raw_texts = slide.get("problem_texts", []) + slide.get("text_items", []) + [slide.get("text", "")]
        for raw_text in raw_texts:
            for piece in re.split(r"\s*•\s*", normalize_whitespace(raw_text)):
                cleaned = clean_source_prompt(piece, slide.get("title", ""))
                if not cleaned or not is_problem_like_text(cleaned):
                    continue
                prompts.append(trim_dangling_display_text(truncate_display_copy(cleaned, 54)))
    return unique_nonempty(prompts, limit=limit)


def source_annotation_terms(source_records: list[dict[str, Any]], limit: int = 6) -> list[str]:
    combined = combined_source_text(source_records).lower()
    terms: list[str] = []
    for term in SOURCE_TERM_PRIORITY:
        pattern = r"\b" + re.escape(term).replace(r"\ ", r"\s+") + r"\b"
        if re.search(pattern, combined):
            terms.append(display_term_label(term))
    if "precise" in combined or "precision" in combined:
        terms.append("Units")
    return unique_nonempty(terms, limit=limit)


def reference_sort_activity_content(
    deck: dict[str, Any],
    source_numbers: list[int],
    *,
    default_tasks: list[str],
    default_pieces: list[str],
    default_instructions: str,
    default_answer_check: str,
) -> dict[str, Any]:
    source_records = source_slides_from_numbers(deck, source_numbers)
    vocab_terms = [item.get("word", "") for item in source_vocabulary(deck, source_numbers, limit=4)]
    prompts = source_question_prompts(source_records, limit=3)
    if vocab_terms and prompts:
        return {
            "tasks": ["Key Lesson Terms", "Discussion Prompts"],
            "movable_pieces": unique_nonempty(vocab_terms[:3] + prompts[:3], limit=6),
            "instructions": "Sort the cards into key lesson terms and discussion prompts, then explain one placement.",
            "answer_check": "Each card should belong with either a key lesson term or a discussion prompt from the source lesson.",
        }
    fact_cards = [
        trim_dangling_display_text(truncate_display_copy(item, 48))
        for item in source_fact_candidates(source_records, limit=4)
    ]
    if vocab_terms and fact_cards:
        return {
            "tasks": ["Key Vocabulary", "Source Details"],
            "movable_pieces": unique_nonempty(vocab_terms[:3] + fact_cards[:3], limit=6),
            "instructions": "Sort the cards into key vocabulary and source details, then explain one placement.",
            "answer_check": "Each card should fit either the key vocabulary from the lesson or a detail taken from the source slide.",
        }
    return {
        "tasks": default_tasks,
        "movable_pieces": default_pieces,
        "instructions": default_instructions,
        "answer_check": default_answer_check,
    }


def source_activity_profile(source_records: list[dict[str, Any]], source_texts: list[str]) -> dict[str, bool]:
    source_blob = normalize_whitespace(
        " ".join(source_texts) or combined_source_text(source_records)
    ).lower()
    title_blob = " ".join(clean_source_prompt(slide.get("title", "")) for slide in source_records).lower()
    strongest_signal = max((source_slide_score(slide) for slide in source_records), default=0)
    return {
        "has_visual": any(int(slide.get("image_count", 0) or 0) > 0 for slide in source_records)
        or any(term in source_blob for term in ("diagram", "figure", "model", "prism", "cube", "cubes")),
        "has_formula": any(
            term in source_blob
            for term in ("formula", "equation", "substitute", "product", "dimensions", "length", "width", "height", "volume", "area")
        ),
        "has_sequence": any(
            term in source_blob
            for term in ("first", "next", "then", "step", "how do you use", "determine the", "substitute", "follow")
        ),
        "has_representation": any(
            term in source_blob
            for term in ("table", "graph", "equation", "coordinate", "number line", "representation", "model")
        ),
        "has_compare": any(
            term in source_blob
            for term in ("compare", "rather", "which", "strongest", "weakest", "longer", "shorter", "different-sized", "size and how many")
        ),
        "has_explanation": any(
            term in source_blob for term in ("explain", "why", "justify", "evidence", "reason", "precision", "how do you know")
        ),
        "has_data_analysis": has_data_analysis_context(f"{title_blob} {source_blob}"),
        "has_context": any(
            term in title_blob or term in source_blob
            for term in ("popcorn", "movie theater", "family")
        )
        or has_shipping_box_context(f"{title_blob} {source_blob}"),
        "has_low_signal": strongest_signal < 30,
        "has_generic_launch": any(term in source_blob for term in ("this or that", "session 1", "session 2")),
    }


def minimum_activity_alignment_score(kind: str) -> int:
    return {
        "be_curious": 4,
        "vocabulary": 5,
        "guided_notes": 6,
        "worked_example": 6,
        "practice": 6,
        "quick_review": 8,
        "challenge": 6,
        "reflection": 5,
        "exit_ticket": 5,
    }.get(kind, 5)


def activity_alignment_score(
    kind: str,
    record: dict[str, Any],
    *,
    profile: dict[str, bool],
) -> int:
    name = record.get("name", "").lower()
    family = record.get("family", "")
    preferred_names = ACTIVITY_KIND_DEFAULTS.get(kind, [])
    score = 0
    if kind in record.get("preferred_kinds", []):
        score += 5
    if record.get("name", "") in preferred_names:
        score += max(1, 4 - preferred_names.index(record["name"]))

    if family == "plot_place" and profile["has_visual"]:
        score += 6
    if family == "build_construct" and profile["has_formula"]:
        score += 6
    if family == "sequence_order" and profile["has_sequence"]:
        score += 5
    if family == "match_pair" and profile["has_representation"]:
        score += 5
    if family in {"sort_classify", "compare_rank"} and profile["has_compare"]:
        score += 4
    if family == "detect_justify" and profile["has_explanation"]:
        score += 4
    if family == "reveal_discuss" and (profile["has_visual"] or profile["has_context"]):
        score += 3
    if family in HIGH_AGENCY_ACTIVITY_FAMILIES and kind not in {"be_curious", "vocabulary"}:
        score += 2
    if family == "reveal_discuss" and kind not in {"be_curious", "vocabulary"}:
        score -= 2
    if profile.get("has_data_analysis") and any(term in name for term in ("equation", "formula", "substitute", "volume", "area", "algebra")):
        score -= 12
    elif not profile["has_formula"] and any(term in name for term in ("equation", "formula", "substitute")):
        score -= 7

    if kind == "vocabulary":
        if "visual vocabulary matching" in name:
            score += 8 if profile["has_visual"] else 5
        if "definition matching" in name and not profile["has_visual"]:
            score += 5
        if "algebra" in name and not profile["has_formula"]:
            score -= 6
        if profile.get("has_data_analysis") and any(term in name for term in ("equation", "formula", "area", "volume")):
            score -= 6
    elif kind in {"guided_notes", "worked_example"}:
        if "label-the-diagram" in name and profile["has_visual"]:
            score += 3
        if "equation" in name and profile["has_formula"]:
            score += 3
        if "claim-evidence" in name:
            score -= 4
        if profile.get("has_data_analysis") and any(term in name for term in ("equation", "formula", "area model", "volume")):
            score -= 8
    elif kind == "practice":
        if "equation" in name and profile["has_formula"]:
            score += 5
        if "build-a-model" in name and profile["has_visual"]:
            score += 5
        if "claim-evidence" in name and not profile["has_explanation"]:
            score -= 5
        if "table" in name and not profile["has_representation"]:
            score -= 6
        if profile.get("has_data_analysis") and any(term in name for term in ("sort", "claim", "evidence", "compare", "match")):
            score += 2
    elif kind == "quick_review":
        if "four-corner" in name and profile["has_compare"]:
            score += 5
        elif "four-corner" in name:
            score -= 4
        if profile["has_low_signal"]:
            score -= 6
        if profile["has_generic_launch"]:
            score -= 8
    elif kind in {"challenge", "exit_ticket"}:
        if family == "detect_justify" and profile["has_explanation"]:
            score += 5
        if family == "build_construct" and profile["has_formula"]:
            score += 3

    if family == "plot_place" and not profile["has_visual"]:
        score -= 3
    if family == "compare_rank" and not profile["has_compare"]:
        score -= 4
    if family == "build_construct" and not profile["has_formula"] and kind not in {"vocabulary", "reflection"}:
        score -= 2
    return score


def pick_activity_for_kind(
    kind: str,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    *,
    source_records: list[dict[str, Any]] | None = None,
    source_texts: list[str] | None = None,
    excluded_names: set[str] | None = None,
) -> dict[str, Any]:
    excluded_names = {item.lower() for item in (excluded_names or set())}
    preferred_names = ACTIVITY_KIND_DEFAULTS.get(kind, [])
    recommended = choose_activity_candidates(deck, library, limit=min(36, max(20, len(library) // 4 or 20)))
    recommended_ranks = {item["name"]: index for index, item in enumerate(recommended)}
    source_records = source_records or []
    source_texts = source_texts or []
    profile = source_activity_profile(source_records, source_texts)
    lesson_text = normalize_whitespace(" ".join(source_texts + [deck_full_text(deck)])).lower()

    candidate_pool: list[dict[str, Any]] = []
    seen_names: set[str] = set()
    for record in library:
        lowered = record["name"].lower()
        if lowered in excluded_names:
            continue
        if kind in PROBLEM_SOLVING_KINDS and record.get("family", "") not in SOURCE_ANCHORED_PROBLEM_ACTIVITY_FAMILIES:
            continue
        if lowered not in seen_names:
            candidate_pool.append(record)
            seen_names.add(lowered)

    scored: list[tuple[int, int, int, str, dict[str, Any]]] = []
    for record in candidate_pool:
        name = record["name"]
        preferred_rank = preferred_names.index(name) if name in preferred_names else 99
        recommended_rank = recommended_ranks.get(name, 999)
        score = activity_alignment_score(kind, record, profile=profile)
        keyword_hits = sum(1 for keyword in record.get("keywords", []) if keyword and keyword in lesson_text)
        score += min(keyword_hits, 4)
        if name in recommended_ranks:
            score += 2
        scored.append((score, preferred_rank, recommended_rank, name, record))

    scored.sort(key=lambda item: (-item[0], item[1], item[2], item[3]))
    if scored and scored[0][0] >= minimum_activity_alignment_score(kind):
        return scored[0][4]
    if scored and kind in PROBLEM_SOLVING_KINDS:
        return scored[0][4]

    return {
        "name": "",
        "family": "",
        "family_label": "",
        "family_prompt": "",
    }


def extract_activity_snippets(source_texts: list[str], limit: int = 6) -> list[str]:
    snippets: list[str] = []
    for source in source_texts:
        normalized = normalize_whitespace(source)
        if not normalized:
            continue
        for piece in re.split(r"[|:;]", normalized):
            cleaned = trim_dangling_display_text(truncate_display_copy(piece.strip(), 58))
            if len(cleaned) < 6:
                continue
            snippets.append(cleaned)
    return unique_nonempty(snippets, limit=limit)


def source_activity_pieces(
    *,
    kind: str,
    record: dict[str, Any],
    deck: dict[str, Any],
    source_records: list[dict[str, Any]],
    source_texts: list[str],
    piece_limit: int,
) -> list[str]:
    family = record.get("family", "")
    solver_blob = normalize_whitespace(" ".join(source_texts) or combined_source_text(source_records)).lower()
    solver_labels: list[str] = []
    if kind == "be_curious":
        launch_terms = source_be_curious_terms(source_records, limit=piece_limit)
        if launch_terms:
            return launch_terms
    if kind in PROBLEM_SOLVING_KINDS and family in {"build_construct", "detect_justify", "compare_rank", "sequence_order", "sort_classify", "match_pair"}:
        if any(term in solver_blob for term in ("which size", "how many boxes", "would you buy", "popcorn")):
            solver_labels = ["Choose a size", "Choose how many", "Use evidence", "Explain why"]
        elif any(term in solver_blob for term in ("percent", "percentage", "students voted", "test scores", "score")):
            if family == "match_pair":
                solver_labels = ["Percent", "Whole", "Part", "Equation"]
            elif family == "sort_classify":
                solver_labels = ["Known percent", "Whole amount", "Math equation", "Answer check"]
            else:
                solver_labels = ["Read the problem", "Find the percent", "Find the whole", "Solve + explain"]
        elif any(term in solver_blob for term in ("regulations", "compare to", "shorter than", "longer than")):
            solver_labels = ["Find the result", "Compare it", "State the conclusion", "Support it"]
        elif any(term in solver_blob for term in ("missing", "unknown", "length of the box", "relate the volume and dimensions")):
            solver_labels = ["Known information", "Volume formula", "Missing value", "Check the answer"]
        elif any(term in solver_blob for term in ("formula", "volume", "dimensions", "length", "width", "height", "substitute")):
            solver_labels = ["Known dimensions", "Write the formula", "Substitute values", "Solve + units"]
        elif "what information do you know" in solver_blob:
            solver_labels = ["What I know", "What I need", "Solve it", "Explain it"]
        if solver_labels:
            return unique_nonempty(solver_labels, limit=piece_limit)
    if kind == "vocabulary":
        vocab_numbers = [slide["slide_number"] for slide in source_records]
        vocab_words = [item.get("word", "") for item in source_vocabulary(deck, vocab_numbers, limit=piece_limit)]
        return unique_nonempty(vocab_words, limit=piece_limit)
    if family == "reveal_discuss":
        prompt_terms = source_term_candidates(source_records, limit=piece_limit)
        if prompt_terms:
            return prompt_terms
    if family in {"detect_justify", "compare_rank"}:
        focus_terms = source_term_candidates(source_records, limit=piece_limit)
        if focus_terms:
            return focus_terms
    if family in {"plot_place", "build_construct"}:
        terms = source_annotation_terms(source_records, limit=piece_limit)
        if terms:
            return terms
    prompts = source_problem_candidates(source_records, limit=piece_limit + 2)
    facts = source_fact_candidates(source_records, limit=piece_limit + 2)
    snippets = extract_activity_snippets(source_texts, limit=piece_limit + 2)
    pieces = prompts + facts + snippets
    cleaned = [
        trim_dangling_display_text(truncate_display_copy(piece, 54))
        for piece in pieces
        if piece and not is_generic_slide_text(piece)
    ]
    return unique_nonempty(cleaned, limit=piece_limit)


def build_activity_payload(
    *,
    kind: str,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    source_texts: list[str],
    source_records: list[dict[str, Any]] | None = None,
    excluded_names: set[str] | None = None,
) -> dict[str, Any]:
    source_records = source_records or []
    record = pick_activity_for_kind(
        kind,
        deck,
        library,
        source_records=source_records,
        source_texts=source_texts,
        excluded_names=excluded_names,
    )
    activity_name = record.get("name", "")
    activity_family = record.get("family", "")
    if not activity_name:
        return {
            "activity_name": "",
            "activity_family": "",
            "activity_instructions": "",
            "movable_pieces": [],
            "answer_check": "",
        }

    piece_limit = clamp_int(record.get("piece_limit", 4), 2, 6)
    pieces = source_activity_pieces(
        kind=kind,
        record=record,
        deck=deck,
        source_records=source_records,
        source_texts=source_texts,
        piece_limit=piece_limit,
    )
    if not pieces:
        pieces = unique_nonempty(record.get("piece_labels", []) + deck["keyword_candidates"], limit=piece_limit)
    instructions_map = {
        "sort_classify": "Drag each card into the category where it belongs, then explain the rule that guided your sort.",
        "match_pair": "Match each source card to the representation, idea, or explanation that fits it best.",
        "sequence_order": "Move the cards into the correct order and defend why that sequence makes sense.",
        "build_construct": "Use the movable pieces to build the model, equation, table, or explanation from the lesson.",
        "plot_place": "Place each marker, label, or value in the correct location before checking your reasoning.",
        "detect_justify": "Identify the error or strongest evidence, then justify your choice using the lesson model.",
        "compare_rank": "Compare the options and rank or group them from strongest to weakest using the source evidence.",
        "reveal_discuss": "Reveal or move one piece at a time, then pause to discuss what the new information tells you.",
    }
    answer_map = {
        "sort_classify": "Check that every card fits the category rule and use a full sentence to justify one placement.",
        "match_pair": "Verify that each pair shows the same idea in two different forms.",
        "sequence_order": "Check that each step leads logically to the next with no gaps or reversals.",
        "build_construct": "Compare your built model to the original source example and revise any mismatched pieces.",
        "plot_place": "Use the source labels, values, or graph features to confirm every placement.",
        "detect_justify": "Use evidence from the source problem or worked example to defend your correction.",
        "compare_rank": "Explain the ordering rule and test whether each item still fits that rule.",
        "reveal_discuss": "After each reveal, revise your prediction and cite a clue from the source material.",
    }
    instruction_text = record.get("instructions") or instructions_map.get(activity_family, "")
    answer_text = record.get("answer_check") or answer_map.get(activity_family, "")
    if kind in PROBLEM_SOLVING_KINDS and activity_family == "build_construct":
        instruction_text = "Use the movable pieces to build the solve path for the exact source problem before you write the full solution."
        answer_text = "Check that your setup, calculations, and final statement all match the source problem and the lesson model."
    elif kind in PROBLEM_SOLVING_KINDS and activity_family == "detect_justify":
        instruction_text = "Use the source problem to decide which move, claim, or conclusion is correct, then justify it with evidence."
        answer_text = "Verify that your correction or claim matches the problem details, formula setup, and final reasoning."
    return {
        "activity_name": activity_name,
        "activity_family": activity_family,
        "activity_instructions": instruction_text,
        "movable_pieces": pieces,
        "answer_check": answer_text,
    }


def slide_lookup(deck: dict[str, Any]) -> dict[int, dict[str, Any]]:
    return {slide["slide_number"]: slide for slide in deck.get("slides", [])}


def source_slides_from_numbers(deck: dict[str, Any], numbers: list[int]) -> list[dict[str, Any]]:
    lookup = slide_lookup(deck)
    return [lookup[number] for number in numbers if number in lookup]


def source_texts_from_numbers(deck: dict[str, Any], numbers: list[int], limit: int = 6) -> list[str]:
    texts: list[str] = []
    for slide in source_slides_from_numbers(deck, numbers):
        texts.extend(slide.get("problem_texts", []))
        texts.extend(slide.get("text_items", []))
        texts.extend([slide.get("title", ""), slide.get("text", ""), slide.get("notes", "")])
    return unique_nonempty(texts, limit=limit)


def source_titles_from_numbers(deck: dict[str, Any], numbers: list[int], limit: int = 4) -> list[str]:
    titles = [
        clean_source_prompt(slide.get("title", ""), slide.get("title", ""))
        for slide in source_slides_from_numbers(deck, numbers)
    ]
    filtered = [title for title in titles if title and not is_generic_slide_text(title)]
    return unique_nonempty(filtered, limit=limit)


def apply_publisher_copyedit(plan: dict[str, Any], deck: dict[str, Any]) -> dict[str, Any]:
    problem_kinds = {"guided_notes", "worked_example", "practice", "quick_review", "challenge", "exit_ticket"}
    plan["lesson_title"] = truncate_text(publisher_copyedit_text(plan.get("lesson_title", "")), 160)
    plan["subject"] = truncate_text(publisher_copyedit_text(plan.get("subject", "")), 80)
    plan["grade_level"] = truncate_text(publisher_copyedit_text(plan.get("grade_level", "")), 80)
    plan["standards"] = publisher_copyedit_list(plan.get("standards", []), limit=8)
    plan["topic_summary"] = truncate_text(publisher_copyedit_text(plan.get("topic_summary", "")), 600)

    for session_key in planned_session_keys(plan):
        session = plan.get(session_key, {})
        session["session_title"] = truncate_text(publisher_copyedit_text(session.get("session_title", "")), 140)
        session["session_subtitle"] = truncate_text(publisher_copyedit_text(session.get("session_subtitle", "")), 220)
        if not normalize_whitespace(session.get("session_subtitle", "")) or "lesson structure" in session["session_subtitle"].lower():
            session["session_subtitle"] = (
                "Launch the lesson with structured notes, vocabulary, and guided practice."
                if session_key == "session_1"
                else "Continue the lesson with collaborative practice, independent work, and reflection."
            )
        session["premium_features"] = [
            feature for feature in unique_nonempty(session.get("premium_features", []), limit=4)
            if feature in PREMIUM_FEATURE_OPTIONS
        ]
        session["context_anchor"] = truncate_text(publisher_copyedit_text(session.get("context_anchor", "")), 80)
        session["flagship_activity_target"] = clamp_int(int(session.get("flagship_activity_target", 0) or 0), 0, 4)
        session["flagship_activity_support_score"] = clamp_int(int(session.get("flagship_activity_support_score", 0) or 0), 0, 8)
        session["flagship_activity_fallbacks"] = publisher_copyedit_list(
            session.get("flagship_activity_fallbacks", []),
            preserve_source_prompt=True,
            limit=4,
        )
        session["flagship_activities"] = [
            normalize_flagship_activity_spec(spec)
            for spec in session.get("flagship_activities", [])
            if isinstance(spec, dict)
        ][:4]
        for slide in session.get("slides", []):
            source_records = source_slides_from_numbers(deck, slide.get("source_slide_numbers", []))
            source_title = source_records[0].get("title", "") if source_records else ""
            preserve_prompt = slide.get("kind") in problem_kinds
            slide["title"] = truncate_text(publisher_copyedit_text(slide.get("title", "")), 110)
            slide["subtitle"] = truncate_text(publisher_copyedit_text(slide.get("subtitle", "")), 220)
            slide["primary_text"] = truncate_text(
                publisher_copyedit_text(
                    slide.get("primary_text", ""),
                    title=source_title,
                    preserve_source_prompt=preserve_prompt,
                ),
                750,
            )
            slide["secondary_text"] = truncate_text(
                publisher_copyedit_text(
                    slide.get("secondary_text", ""),
                    title=source_title,
                    preserve_source_prompt=preserve_prompt and "?" in slide.get("secondary_text", ""),
                ),
                500,
            )
            slide["bullets"] = publisher_copyedit_list(
                slide.get("bullets", []),
                title=source_title,
                preserve_source_prompt=preserve_prompt,
                limit=5,
            )
            slide["tasks"] = publisher_copyedit_list(
                slide.get("tasks", []),
                title=source_title,
                preserve_source_prompt=True,
                limit=4,
            )
            slide["response_prompt"] = truncate_text(publisher_copyedit_text(slide.get("response_prompt", "")), 300)
            slide["sentence_starters"] = publisher_copyedit_list(slide.get("sentence_starters", []), limit=5)
            slide["activity_name"] = truncate_text(publisher_copyedit_text(slide.get("activity_name", "")), 120)
            slide["activity_instructions"] = truncate_text(
                publisher_copyedit_text(slide.get("activity_instructions", "")),
                280,
            )
            slide["movable_pieces"] = publisher_copyedit_list(slide.get("movable_pieces", []), limit=6)
            slide["answer_check"] = truncate_text(publisher_copyedit_text(slide.get("answer_check", "")), 220)
            slide["premium_layout"] = slide.get("premium_layout", "") if slide.get("premium_layout", "") in {
                "",
                "error_analysis",
                "multi_representation",
                "strategy_comparison",
                "evidence_ladder",
                "decision_tree",
                "create_your_own",
                "real_world_transfer",
                "mastery_tracker",
                "turn_and_teach",
            } else ""
            slide["premium_title"] = truncate_text(publisher_copyedit_text(slide.get("premium_title", "")), 120)
            slide["premium_text"] = truncate_text(
                publisher_copyedit_text(slide.get("premium_text", ""), preserve_source_prompt=True),
                260,
            )
            slide["premium_items"] = publisher_copyedit_list(
                slide.get("premium_items", []),
                preserve_source_prompt=True,
                limit=6,
            )
            premium_table: list[list[str]] = []
            for row in slide.get("premium_table", [])[:5]:
                if not isinstance(row, list):
                    continue
                premium_table.append(
                    [
                        truncate_text(
                            publisher_copyedit_text(str(cell), preserve_source_prompt=True),
                            90,
                        )
                        for cell in row[:4]
                    ]
                )
            slide["premium_table"] = premium_table
            slide["partner_prompt"] = truncate_text(
                publisher_copyedit_text(slide.get("partner_prompt", ""), preserve_source_prompt=True),
                180,
            )
            slide["discussion_questions"] = publisher_copyedit_list(
                [clean_discussion_question(item) for item in slide.get("discussion_questions", [])],
                preserve_source_prompt=True,
                limit=3,
            )
            slide["context_anchor"] = truncate_text(publisher_copyedit_text(slide.get("context_anchor", "")), 80)
            slide["practice_phase"] = truncate_text(publisher_copyedit_text(slide.get("practice_phase", "")), 40)
            slide["flagship_activity_mode"] = (
                slide.get("flagship_activity_mode", FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED)
                if slide.get("flagship_activity_mode", FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED) in {FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED, FLAGSHIP_ACTIVITY_FOOTPRINT_FULL}
                else FLAGSHIP_ACTIVITY_FOOTPRINT_EMBEDDED
            )
            if flagship_activity_spec(slide):
                slide["flagship_activity"] = normalize_flagship_activity_spec(slide["flagship_activity"])
            edited_vocab: list[dict[str, str]] = []
            for item in slide.get("vocabulary", [])[:5]:
                if not isinstance(item, dict):
                    continue
                word = truncate_text(publisher_copyedit_text(str(item.get("word", ""))), 50)
                word = display_term_label(word)
                if not word or is_low_value_vocabulary_term(word):
                    continue
                definition = truncate_text(publisher_copyedit_text(str(item.get("definition", ""))), 140)
                if is_placeholder_vocab_definition(definition):
                    definition = student_friendly_definition(word)
                example = truncate_text(
                    publisher_copyedit_text(
                        str(item.get("example", "")),
                        title=source_title,
                        preserve_source_prompt=True,
                    ),
                    110,
                )
                if not example or is_generic_slide_text(example, source_title):
                    example = find_term_example_in_records(source_records, word) or find_term_example(deck, word)
                visual_cue = truncate_text(publisher_copyedit_text(str(item.get("visual_cue", ""))), 110)
                if not visual_cue or is_generic_slide_text(visual_cue, source_title):
                    visual_cue = vocabulary_visual_cue(deck, word)
                edited_vocab.append(
                    {
                        "word": word,
                        "definition": truncate_text(definition, 140),
                        "example": truncate_text(example, 110),
                        "visual_cue": truncate_text(visual_cue, 110),
                    }
                )
            slide["vocabulary"] = [item for item in edited_vocab if item["word"]]
            if source_records and slide.get("kind") in problem_kinds:
                slide["source_problem_cards"] = source_problem_candidates(source_records, limit=4)
                enrich_problem_fidelity(slide, deck)
            strengthen_slide_copy_for_publisher_rigor(slide, session_key=session_key, deck=deck)
            optimize_slide_copy_for_layout(slide)
            if slide.get("kind") in problem_kinds:
                normalized_problem_cards: list[str] = []
                for card in source_problem_cards(slide, limit=4, max_len=220):
                    cleaned_card = trim_dangling_display_text(card)
                    if not cleaned_card or has_dangling_display_text(cleaned_card):
                        continue
                    normalized_problem_cards.append(cleaned_card)

                if len(unique_nonempty(normalized_problem_cards, limit=4)) < 2 and source_records:
                    for candidate in source_problem_candidates(source_records, limit=4):
                        cleaned_candidate = trim_dangling_display_text(truncate_display_copy(candidate, 220))
                        if not cleaned_candidate or has_dangling_display_text(cleaned_candidate):
                            continue
                        normalized_problem_cards.append(cleaned_candidate)
                        if len(unique_nonempty(normalized_problem_cards, limit=4)) >= 2:
                            break

                if len(unique_nonempty(normalized_problem_cards, limit=4)) < 2:
                    primary_fallback = trim_dangling_display_text(truncate_display_copy(slide.get("primary_text", ""), 220))
                    if primary_fallback and not has_dangling_display_text(primary_fallback):
                        normalized_problem_cards.append(primary_fallback)

                slide["source_problem_cards"] = unique_nonempty(normalized_problem_cards, limit=4)
    return plan


def run_publisher_copyedit_review(plan: dict[str, Any]) -> None:
    issues: list[str] = []
    top_level_fields = [
        ("lesson_title", plan.get("lesson_title", "")),
        ("subject", plan.get("subject", "")),
        ("grade_level", plan.get("grade_level", "")),
        ("topic_summary", plan.get("topic_summary", "")),
    ]
    for field_name, field_value in top_level_fields:
        for issue in publisher_copyedit_issues(field_value):
            issues.append(f"{field_name}: {issue}")

    field_names = (
        "title",
        "subtitle",
        "primary_text",
        "secondary_text",
        "response_prompt",
        "activity_name",
        "activity_instructions",
        "answer_check",
        "premium_title",
        "premium_text",
        "partner_prompt",
        "context_anchor",
        "practice_phase",
    )
    list_fields = ("bullets", "tasks", "sentence_starters", "movable_pieces", "premium_items", "source_problem_cards", "discussion_questions")
    for session_key in planned_session_keys(plan):
        session = plan.get(session_key, {})
        for field_name in ("session_title", "session_subtitle", "context_anchor"):
            for issue in publisher_copyedit_issues(str(session.get(field_name, ""))):
                issues.append(f"{session_key} {field_name}: {issue}")
        for note_index, item in enumerate(session.get("flagship_activity_fallbacks", []), start=1):
            for issue in publisher_copyedit_issues(str(item)):
                issues.append(f"{session_key} flagship_activity_fallbacks[{note_index}]: {issue}")
        for activity_index, spec in enumerate(session.get("flagship_activities", []), start=1):
            if not isinstance(spec, dict):
                continue
            for field_name in ("title", "purpose", "directions", "layoutVariant"):
                for issue in publisher_copyedit_issues(str(spec.get(field_name, ""))):
                    issues.append(f"{session_key} flagship_activities[{activity_index}].{field_name}: {issue}")
            for list_name in ("prompts", "supports", "vocabularyUsed"):
                for item_index, item in enumerate(spec.get(list_name, []), start=1):
                    for issue in publisher_copyedit_issues(str(item)):
                        issues.append(f"{session_key} flagship_activities[{activity_index}].{list_name}[{item_index}]: {issue}")
        for page_index, slide in enumerate(plan.get(session_key, {}).get("slides", []), start=1):
            for field_name in field_names:
                for issue in publisher_copyedit_issues(str(slide.get(field_name, ""))):
                    issues.append(f"{session_key} page {page_index} {field_name}: {issue}")
            for list_name in list_fields:
                for item_index, item in enumerate(slide.get(list_name, []), start=1):
                    for issue in publisher_copyedit_issues(str(item)):
                        issues.append(f"{session_key} page {page_index} {list_name}[{item_index}]: {issue}")
            for vocab_index, vocab_item in enumerate(slide.get("vocabulary", []), start=1):
                if not isinstance(vocab_item, dict):
                    continue
                for field_name in ("word", "definition", "example", "visual_cue"):
                    for issue in publisher_copyedit_issues(str(vocab_item.get(field_name, ""))):
                        issues.append(f"{session_key} page {page_index} vocabulary[{vocab_index}].{field_name}: {issue}")
                if is_low_value_vocabulary_term(vocab_item.get("word", "")):
                    issues.append(f"{session_key} page {page_index} vocabulary[{vocab_index}].word: low-value vocabulary term")
                if is_placeholder_vocab_definition(vocab_item.get("definition", "")):
                    issues.append(f"{session_key} page {page_index} vocabulary[{vocab_index}].definition: placeholder definition")
            if flagship_activity_spec(slide):
                spec = slide["flagship_activity"]
                for field_name in ("title", "purpose", "directions", "layoutVariant"):
                    for issue in publisher_copyedit_issues(str(spec.get(field_name, ""))):
                        issues.append(f"{session_key} page {page_index} flagship_activity.{field_name}: {issue}")
                for list_name in ("prompts", "supports", "vocabularyUsed"):
                    for item_index, item in enumerate(spec.get(list_name, []), start=1):
                        for issue in publisher_copyedit_issues(str(item)):
                            issues.append(f"{session_key} page {page_index} flagship_activity.{list_name}[{item_index}]: {issue}")

    if issues:
        raise RuntimeError("Publisher copyedit review failed: " + " | ".join(issues[:8]))


def run_formal_release_plan_review(plan: dict[str, Any], *, deck: dict[str, Any] | None = None) -> None:
    issues: list[str] = []
    if not normalize_whitespace(plan.get("lesson_title", "")):
        issues.append("lesson_title: missing")
    if not normalize_whitespace(plan.get("topic_summary", "")):
        issues.append("topic_summary: missing")

    for session_key in planned_session_keys(plan):
        session = plan.get(session_key, {})
        slides = session.get("slides", [])
        if not slides:
            issues.append(f"{session_key}: missing slides")
            continue
        required_roles = required_template_roles_for_session(session)
        if required_roles and template_role_signature(session) != required_roles:
            issues.append(f"{session_key}: exact ESOL workbook sequence is not locked")
        if not opening_slide_allowed(session, slides[0]):
            issues.append(f"{session_key}: first slide does not match the locked opening slide")
        if not closing_slide_allowed(session, slides[-1]):
            issues.append(f"{session_key}: last slide does not match the locked closing slide")

        first_index = lambda target: next((index for index, slide in enumerate(slides) if slide.get("kind") == target), None)
        learning_target_index = first_index("learning_target")
        worked_example_index = first_index("worked_example")
        practice_index = first_index("practice")
        reflection_index = first_index("reflection")
        exit_ticket_index = first_index("exit_ticket")
        if learning_target_index is not None and worked_example_index is not None and learning_target_index > worked_example_index:
            issues.append(f"{session_key}: learning target should come before the worked example")
        if worked_example_index is not None and practice_index is not None and worked_example_index > practice_index:
            issues.append(f"{session_key}: worked example should come before practice")
        if reflection_index is not None and exit_ticket_index is not None and reflection_index > exit_ticket_index:
            issues.append(f"{session_key}: reflection should come before the exit ticket")

        exact_template = uses_exact_esol_template(session)
        problem_slides = [slide for slide in slides if slide.get("kind") in PROBLEM_SOLVING_KINDS]
        guided_present = any(
            slide.get("kind") == "worked_example"
            or "fully guided" in normalize_whitespace(slide.get("practice_phase", "")).lower()
            for slide in problem_slides
        )
        together_present = any(
            "together" in normalize_whitespace(slide.get("practice_phase", "")).lower()
            or "collaborative" in normalize_whitespace(slide.get("title", "")).lower()
            or "partner" in normalize_whitespace(slide.get("partner_prompt", "")).lower()
            for slide in problem_slides
        ) or (
            exact_template
            and any(normalize_whitespace(slide.get("template_role", "")) == "interactive_activity" for slide in slides)
        )
        independent_present = any(
            "independ" in normalize_whitespace(slide.get("practice_phase", "")).lower()
            or slide.get("kind") in {"challenge", "exit_ticket"}
            for slide in problem_slides
        ) or (
            exact_template
            and any(normalize_whitespace(slide.get("template_role", "")) == "best_fit_review" for slide in slides)
        )
        if problem_slides and not guided_present:
            issues.append(f"{session_key}: workbook baseline is missing a modeled source-problem page")
        if problem_slides and not together_present:
            if exact_template:
                issues.append(f"{session_key}: workbook baseline is missing the interactive activity page")
            else:
                issues.append(f"{session_key}: workbook baseline is missing a collaborative or solve-together page")
        if problem_slides and not independent_present:
            if exact_template:
                issues.append(f"{session_key}: workbook baseline is missing the best-fit review page")
            else:
                issues.append(f"{session_key}: workbook baseline is missing an independent solve page")
        engagement_count = session_engagement_slide_count(session)
        engagement_modes = session_engagement_modes(session)
        if engagement_count < engagement_slide_target(session):
            issues.append(f"{session_key}: session does not include enough high-agency engagement pages")
        if len(engagement_modes) < engagement_mode_target(session):
            issues.append(f"{session_key}: session does not include enough distinct engagement moves")
        interactive_features = [
            feature
            for feature in session.get("premium_features", [])
            if feature in HIGH_AGENCY_PREMIUM_FEATURES
        ]
        if session.get("premium_features") and len(interactive_features) < min(2, len(session.get("premium_features", []))):
            issues.append(f"{session_key}: premium feature mix leans too heavily toward support instead of interaction")
        for issue in validate_flagship_activities(session, deck=deck):
            issues.append(f"{session_key}: {issue}")

        paired_problem_entries = [
            (page_index, slide)
            for page_index, slide in enumerate(slides, start=1)
            if slide.get("kind") in PROBLEM_INTERACTIVE_PAIR_KINDS
            and not normalize_whitespace(slide.get("template_role", ""))
        ]
        for pair_index, (page_index, pair_slide) in enumerate(paired_problem_entries):
            mode = problem_render_mode(pair_slide)
            if pair_index % 2 == 0:
                if mode != PROBLEM_RENDER_MODE_FOCUS:
                    issues.append(f"{session_key} page {page_index}: problem slide must be problem-first before interactive apply")
                if has_activity(pair_slide):
                    issues.append(f"{session_key} page {page_index}: problem-first slide should not contain an embedded activity")
            else:
                if mode != PROBLEM_RENDER_MODE_INTERACTIVE:
                    issues.append(f"{session_key} page {page_index}: follow-up slide must be flagged as interactive apply")
                if not has_activity(pair_slide):
                    issues.append(f"{session_key} page {page_index}: interactive apply slide is missing draggable activity content")

        for page_index, slide in enumerate(slides, start=1):
            kind = slide.get("kind", "")
            if not normalize_whitespace(slide.get("title", "")):
                issues.append(f"{session_key} page {page_index}: missing title")
            if not normalize_whitespace(slide.get("section", "")):
                issues.append(f"{session_key} page {page_index}: missing section label")
            for vocab_index, vocab_item in enumerate(slide.get("vocabulary", []), start=1):
                if not isinstance(vocab_item, dict):
                    continue
                if is_low_value_vocabulary_term(vocab_item.get("word", "")):
                    issues.append(f"{session_key} page {page_index}: vocabulary item {vocab_index} is a low-value term")
                if is_placeholder_vocab_definition(vocab_item.get("definition", "")):
                    issues.append(f"{session_key} page {page_index}: vocabulary item {vocab_index} still has a placeholder definition")
            if kind == "learning_target" and display_text_key(slide.get("primary_text", "")) == display_text_key(slide.get("secondary_text", "")):
                issues.append(f"{session_key} page {page_index}: content and language objectives should not duplicate each other")
            if kind == "learning_target" and not is_i_can_objective(slide.get("primary_text", "")):
                issues.append(f"{session_key} page {page_index}: content objective must start with 'I can'")
            if kind == "learning_target" and not is_i_can_objective(slide.get("secondary_text", "")):
                issues.append(f"{session_key} page {page_index}: language objective must start with 'I can'")
            if flagship_activity_spec(slide) and kind in {"cover", "learning_target"}:
                issues.append(f"{session_key} page {page_index}: flagship activities should not replace core launch or objective sections")
            if slide_needs_peer_discussion_support(slide):
                if len([question for question in slide.get("discussion_questions", []) if is_specific_discussion_question(question)]) < 2:
                    issues.append(f"{session_key} page {page_index}: slide is missing two specific peer-discussion questions")
                if not normalize_whitespace(slide.get("partner_prompt", "")):
                    issues.append(f"{session_key} page {page_index}: slide is missing a partner discussion direction")
            if kind in PROBLEM_SOLVING_KINDS:
                source_numbers = [int(num) for num in slide.get("source_slide_numbers", []) if isinstance(num, int)]
                if not source_numbers:
                    issues.append(f"{session_key} page {page_index}: problem-solving slide is missing source anchors")
                problem_copy_blob = normalize_whitespace(
                    " ".join(
                        [
                            slide.get("response_prompt", ""),
                            slide.get("activity_name", ""),
                            slide.get("activity_instructions", ""),
                            slide.get("answer_check", ""),
                            " ".join(slide.get("movable_pieces", [])),
                        ]
                    )
                ).lower()
                if has_data_analysis_context(workbook_source_blob(slide)) and any(
                    term in problem_copy_blob
                    for term in ("formula", "equation", "substitute", "volume", "area", "prism", "missing value")
                ):
                    issues.append(f"{session_key} page {page_index}: data-analysis slide still uses formula-specific activity copy")
                if has_activity(slide) and problem_activity_has_generic_markers(slide):
                    issues.append(f"{session_key} page {page_index}: problem-solving activity is still too generic instead of source-anchored")
                template_role = normalize_whitespace(slide.get("template_role", ""))
                if template_role == "choice_board":
                    choice_labels = [
                        normalize_whitespace(item.get("label", ""))
                        for item in slide.get("choice_paths", [])
                        if isinstance(item, dict)
                    ]
                    if any(label.lower().startswith("path ") for label in choice_labels):
                        issues.append(f"{session_key} page {page_index}: choice board still uses generic path labels")
                if template_role == "independent_practice":
                    independent_prompts = [
                        normalize_whitespace(item.get("prompt", ""))
                        for item in slide.get("independent_problems", [])
                        if isinstance(item, dict)
                    ]
                    if any("compare s" in prompt.lower() for prompt in independent_prompts):
                        issues.append(f"{session_key} page {page_index}: independent practice still uses generic session-to-session prompts")
                if deck:
                    max_slide_number = len(deck.get("slides", []))
                    if any(num < 1 or num > max_slide_number for num in source_numbers):
                        issues.append(f"{session_key} page {page_index}: problem-solving slide has invalid source slide numbers")
                    source_records = source_slides_from_numbers(deck, source_numbers)
                    source_candidates = source_problem_candidates(source_records, limit=4)
                    if source_candidates:
                        source_cards = slide.get("source_problem_cards", [])
                        if not source_cards:
                            issues.append(f"{session_key} page {page_index}: problem-solving slide is missing exact source problem cards")
                        overlap_texts = [slide.get("primary_text", ""), *slide.get("tasks", []), *source_cards]
                        overlap_hits = sum(
                            1
                            for candidate in source_candidates
                            if any(source_problem_text_overlap(text, [candidate]) for text in overlap_texts)
                        )
                        if overlap_hits < min(2, len(source_candidates)):
                            issues.append(f"{session_key} page {page_index}: problem-solving slide drifted away from exact source problem language")
                prompts = workbook_problem_prompts(slide, limit=3)
                if len(unique_nonempty(prompts, limit=3)) < 2:
                    issues.append(f"{session_key} page {page_index}: problem-solving slide needs at least two clear solve prompts")
                if action_prompt_count(prompts) < 2:
                    issues.append(f"{session_key} page {page_index}: problem-solving slide is not action-oriented enough")
                if not normalize_whitespace(slide.get("response_prompt", "")):
                    issues.append(f"{session_key} page {page_index}: problem-solving slide is missing explanation/check guidance")
                if len(normalize_whitespace(slide.get("subtitle", ""))) > 76:
                    issues.append(f"{session_key} page {page_index}: problem-solving subtitle is too long")

    if issues:
        raise RuntimeError("Formal release plan review failed: " + " | ".join(issues[:8]))


def session_source_slides(deck: dict[str, Any], session_key: str) -> list[dict[str, Any]]:
    slides = deck.get("slides", [])
    def session_marker_index(target_number: int) -> int | None:
        pattern = re.compile(rf"\bsession\s*{target_number}\b", re.IGNORECASE)
        for index, slide in enumerate(slides):
            candidates = [
                slide.get("title", ""),
                *slide.get("text_items", [])[:4],
                slide.get("text", ""),
            ]
            if any(pattern.search(normalize_whitespace(candidate)) for candidate in candidates if normalize_whitespace(candidate)):
                return index
        return None

    session1_index = session_marker_index(1) or 0
    session2_index = session_marker_index(2)
    if session2_index is not None and session2_index > session1_index:
        if session_key == "session_1":
            return slides[session1_index:session2_index]
        return slides[session2_index:] or slides[session2_index - 1 :]
    midpoint = max(1, len(slides) // 2)
    if session_key == "session_1":
        return slides[:midpoint]
    return slides[midpoint:] or slides[midpoint - 1 :]


def is_community_building_prompt(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    return any(
        marker in lowered
        for marker in (
            "work productively with your classmates",
            "connecting with your classmates",
            "during math class",
            "your classmates",
            "math class",
        )
    )


def launch_source_score(slide: dict[str, Any]) -> int:
    title = normalize_whitespace(slide.get("title", "")).lower()
    blob = normalize_whitespace(
        " ".join(
            [
                title,
                *slide.get("problem_texts", [])[:2],
                *slide.get("text_items", [])[:4],
                slide.get("text", ""),
            ]
        )
    ).lower()
    score = 0
    if "be curious" in title:
        score += 8
    if any(
        marker in blob
        for marker in (
            "which doesn",
            "which doesn't",
            "what math do you see",
            "notice",
            "wonder",
            "estimate the shaded",
            "about, between, or exact",
            "it is about",
            "it is between",
            "it is exactly",
        )
    ):
        score += 6
    if slide.get("image_count"):
        score += 4
    if slide.get("problem_texts"):
        score += 2
    if is_community_building_prompt(blob):
        score -= 5
    return score


def session_launch_slides(deck: dict[str, Any], session_key: str, limit: int = 3) -> list[dict[str, Any]]:
    session_sources = session_source_slides(deck, session_key)
    if not session_sources:
        return []
    early_pool = session_sources[: max(limit + 2, 6)] or session_sources
    window_size = max(1, min(limit, len(early_pool)))
    best_window = early_pool[:window_size]
    best_score = -999
    for start in range(0, len(early_pool) - window_size + 1):
        window = early_pool[start : start + window_size]
        score = sum(launch_source_score(slide) for slide in window)
        if any("be curious" in normalize_whitespace(slide.get("title", "")).lower() for slide in window):
            score += 2
        if any(slide.get("image_count") for slide in window):
            score += 2
        if score > best_score:
            best_window = window
            best_score = score
    if best_score > 0:
        return best_window
    fallback = [slide for slide in early_pool if slide.get("image_count") or slide.get("problem_texts")] or early_pool
    return fallback[:window_size]


def pick_launch_image_slide(source_records: list[dict[str, Any]]) -> int:
    for slide in source_records:
        if slide.get("image_count") and "be curious" in normalize_whitespace(slide.get("title", "")).lower():
            return int(slide.get("slide_number", 0) or 0)
    for slide in source_records:
        blob = normalize_whitespace(
            " ".join([slide.get("title", ""), *slide.get("text_items", [])[:3], slide.get("text", "")])
        ).lower()
        if slide.get("image_count") and launch_source_score(slide) > 0 and not is_community_building_prompt(blob):
            return int(slide.get("slide_number", 0) or 0)
    return int(pick_first_image_slide(source_records) or 0)


def be_curious_prompt_candidates(source_records: list[dict[str, Any]], limit: int = 5) -> list[str]:
    scored: list[tuple[int, int, str]] = []
    seen: set[str] = set()
    raw_candidates = source_problem_candidates(source_records, limit=max(limit + 3, 8))
    for index, candidate in enumerate(raw_candidates):
        cleaned = normalize_whitespace(candidate)
        if not cleaned:
            continue
        key = display_text_key(cleaned)
        if not key or key in seen or is_community_building_prompt(cleaned):
            continue
        seen.add(key)
        lowered = cleaned.lower()
        score = 0
        if "?" in lowered:
            score += 4
        if any(
            marker in lowered
            for marker in (
                "what math do you see",
                "estimate the shaded",
                "which doesn",
                "which doesn't",
                "what percent",
                "part of the goal",
                "least to greatest",
                "different notations",
            )
        ):
            score += 8
        if any(
            marker in lowered
            for marker in (
                "blue whales",
                "body weight",
                "figure",
                "shaded",
                "goal",
                "cans",
                "fraction",
                "decimal",
                "percent",
            )
        ):
            score += 4
        if len(cleaned) >= 36:
            score += 2
        if len(cleaned) >= 90:
            score += 1
        if len(cleaned) > 170:
            score -= 2
        scored.append((score, index, cleaned))
    ordered = [text for _score, _index, text in sorted(scored, key=lambda item: (-item[0], item[1]))]
    return unique_nonempty(ordered, limit=limit)


def polish_be_curious_prompt(text: str) -> str:
    cleaned = normalize_whitespace(text)
    lowered = cleaned.lower()
    if lowered in {"which doesn’t belong?", "which doesn't belong?"}:
        return "Which figure doesn't belong, and what visual clue helped you decide?"
    if "blue whales" in lowered and "what percent" in lowered:
        return "What percent of a blue whale's body weight is taken in during one gulp?"
    return cleaned


def be_curious_prompt_pair(source_records: list[dict[str, Any]]) -> tuple[str, str]:
    candidates = [polish_be_curious_prompt(text) for text in be_curious_prompt_candidates(source_records, limit=5)]
    primary = next(
        (
            text
            for text in candidates
            if any(
                marker in text.lower()
                for marker in ("estimate the shaded", "what math do you see", "which figure doesn't belong")
            )
        ),
        candidates[0] if candidates else "",
    )
    secondary = next(
        (
            text
            for text in candidates
            if display_text_key(text) != display_text_key(primary)
            and len(normalize_whitespace(text)) >= 20
        ),
        "",
    )
    source_blob = combined_source_text(source_records).lower()
    if not secondary:
        if any(term in source_blob for term in ("blue whales", "body weight", "what percent")):
            secondary = "What percent of the whale's body weight is being described in the problem?"
        elif any(term in source_blob for term in ("shaded", "about, between, or exact", "which doesn", "which doesn't")):
            secondary = "Which figure doesn't belong, and what visual clue helped you decide?"
        elif any(term in source_blob for term in ("food drive", "goal that they have achieved")):
            secondary = "How could the students represent the part of the goal they have reached?"
        elif any(term in source_blob for term in ("least to greatest", "different notations")):
            secondary = "Which number form would make the comparison easier, and why?"
    return (
        truncate_display_copy(primary, 140),
        truncate_display_copy(secondary, 180),
    )


def source_be_curious_terms(source_records: list[dict[str, Any]], *, limit: int = 4) -> list[str]:
    source_blob = combined_source_text(source_records).lower()
    terms: list[str] = []
    if any(term in source_blob for term in ("shaded", "it is about", "it is between", "it is exactly", "which doesn", "which doesn't")):
        terms.extend(["Estimate", "Exact", "Part", "Whole"])
    if any(term in source_blob for term in ("food drive", "cans", "part of the goal", "goal that they have achieved")):
        terms.extend(["Part", "Whole", "Percent", "Represent"])
    if any(term in source_blob for term in ("body weight", "blue whales", "what percent")):
        terms.extend(["Percent", "Part", "Whole", "Body Weight"])
    if any(term in source_blob for term in ("least to greatest", "order", "notation", "different notations")):
        terms.extend(["Fraction", "Decimal", "Percent", "Notation"])
    if not terms:
        terms.extend(source_term_candidates(source_records, limit=limit + 2))
    return unique_nonempty(terms, limit=limit)


def source_be_curious_vocabulary(
    deck: dict[str, Any],
    source_numbers: list[int],
    *,
    limit: int = 3,
) -> list[dict[str, str]]:
    source_records = source_slides_from_numbers(deck, source_numbers)
    terms = source_be_curious_terms(source_records, limit=max(limit + 1, 4))
    vocab: list[dict[str, str]] = []
    for term in terms:
        example = find_term_example_in_records(source_records, term) or find_term_example(deck, term)
        visual_cue = curated_vocab_visual_cue(term) or truncate_display_copy(
            example or f"Look for {term.lower()} in the source image, prompt, or model.",
            72,
        )
        vocab.append(
            {
                "word": display_term_label(term),
                "definition": student_friendly_definition(term),
                "example": example or truncate_display_copy(f"Use {term.lower()} when you explain what you notice.", 96),
                "visual_cue": visual_cue,
            }
        )
        if len(vocab) >= limit:
            break
    return vocab


def source_be_curious_response_prompt(source_records: list[dict[str, Any]]) -> str:
    source_blob = combined_source_text(source_records).lower()
    if any(term in source_blob for term in ("shaded", "it is about", "it is between", "it is exactly", "which doesn", "which doesn't")):
        return "What visual clue tells you whether the amount is about, between, or exact?"
    if any(term in source_blob for term in ("food drive", "cans", "part of the goal", "goal that they have achieved")):
        return "What clue shows the part, the whole, and how the goal could be represented?"
    if any(term in source_blob for term in ("body weight", "blue whales", "what percent")):
        return "What clue in the image helps you identify the part, the whole, and the percent?"
    if any(term in source_blob for term in ("least to greatest", "order", "different notations")):
        return "What clue tells you these numbers need to be rewritten before you compare them?"
    return "What clue in the source image or prompt should shape your first move?"


def source_be_curious_kernels(source_records: list[dict[str, Any]]) -> tuple[list[str], list[str]]:
    source_blob = combined_source_text(source_records).lower()
    notice: list[str] = []
    wonder: list[str] = []
    if any(term in source_blob for term in ("shaded", "it is about", "it is between", "it is exactly", "which doesn", "which doesn't")):
        notice.extend(
            [
                "I notice one figure is about ___.",
                "I notice one figure is exact because ___.",
            ]
        )
        wonder.extend(
            [
                "I wonder which figure doesn't belong because ___.",
                "I predict the model shows a part of a whole.",
            ]
        )
    if any(term in source_blob for term in ("food drive", "cans", "part of the goal", "goal that they have achieved")):
        notice.extend(
            [
                "I notice the problem gives a part and a whole.",
                "I notice 75 out of 100 can be shown more than one way.",
            ]
        )
        wonder.extend(
            [
                "I wonder how to show the same amount as a fraction, decimal, and percent.",
                "I predict the model will help me compare part to whole.",
            ]
        )
    if any(term in source_blob for term in ("body weight", "blue whales", "what percent")):
        notice.extend(
            [
                "I notice the story gives a part and a whole.",
                "I notice the image helps me picture the whole amount.",
            ]
        )
        wonder.extend(
            [
                "I wonder what whole this percent uses.",
                "I predict I will compare a part to a whole.",
            ]
        )
    if any(term in source_blob for term in ("least to greatest", "order", "notation", "written in different notations")):
        notice.extend(
            [
                "I notice the numbers are written in different forms.",
                "I notice I may need one notation to compare them fairly.",
            ]
        )
        wonder.extend(
            [
                "I wonder which notation will make the comparison easier.",
                "I predict rewriting the numbers will help me order them.",
            ]
        )
    if not notice:
        notice = [
            "I notice ___ because ___.",
            "I notice the image or model shows ___.",
        ]
    if not wonder:
        wonder = [
            "I wonder ___ because ___.",
            "I predict ___ because ___.",
        ]
    return unique_nonempty(notice, limit=2), unique_nonempty(wonder, limit=2)


def is_generic_be_curious_notice(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    return not lowered or any(
        marker in lowered
        for marker in (
            "record what stands out",
            "what do you notice first when you look back",
            "notice: ___",
            "notice: ___.",
        )
    )


def is_generic_be_curious_wonder(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    return not lowered or any(
        marker in lowered
        for marker in (
            "ask one question",
            "what do you wonder or predict before you solve",
            "wonder: ___",
            "wonder: ___.",
        )
    )


def ensure_be_curious_slide(deck: dict[str, Any], plan_slide: dict[str, Any]) -> None:
    source_numbers = plan_slide.get("source_slide_numbers", [])
    source_records = source_slides_from_numbers(deck, source_numbers)
    primary_prompt, secondary_prompt = be_curious_prompt_pair(source_records)
    notice_kernels, wonder_kernels = source_be_curious_kernels(source_records)
    plan_slide["image_source_slide"] = pick_launch_image_slide(source_records) or int(plan_slide.get("image_source_slide", 0) or 0)
    notice_text = plan_slide.get("primary_text", "")
    wonder_text = plan_slide.get("secondary_text", "")
    if is_generic_be_curious_notice(notice_text) or has_dangling_display_text(notice_text):
        plan_slide["primary_text"] = primary_prompt or "Record two details you notice in the source image, figure, or problem setup."
    if (
        is_generic_be_curious_wonder(wonder_text)
        or has_dangling_display_text(wonder_text)
        or wonder_text.lower().startswith("use the source prompt")
        or is_community_building_prompt(wonder_text)
    ):
        plan_slide["secondary_text"] = secondary_prompt or "Write one question or prediction that could help you solve the source problem."
    plan_slide["notice_kernels"] = notice_kernels
    plan_slide["wonder_kernels"] = wonder_kernels
    plan_slide["sentence_starters"] = unique_nonempty(
        notice_kernels
        + wonder_kernels
        + plan_slide.get("sentence_starters", [])
        + [
            "I notice ___ because ___.",
            "I wonder ___ because ___.",
            "One lesson word that may matter is ___.",
            "I predict the problem will ask me to ___.",
        ],
        limit=5,
    )
    plan_slide["vocabulary"] = source_be_curious_vocabulary(deck, source_numbers, limit=3) or plan_slide.get("vocabulary", [])
    response_prompt = normalize_whitespace(plan_slide.get("response_prompt", ""))
    if (
        not response_prompt
        or "predict the lesson focus" in response_prompt.lower()
        or "session 1" in response_prompt.lower()
        or "help you most today" in response_prompt.lower()
    ):
        plan_slide["response_prompt"] = source_be_curious_response_prompt(source_records)


def build_extension_slide(
    *,
    session_key: str,
    deck: dict[str, Any],
    slides: list[dict[str, Any]],
    library: list[dict[str, Any]],
) -> dict[str, Any]:
    source_pool = session_source_slides(deck, session_key)
    used_numbers = {
        number
        for slide in slides
        for number in slide.get("source_slide_numbers", [])
        if isinstance(number, int)
    }
    unused_pool = [slide for slide in source_pool if slide["slide_number"] not in used_numbers] or source_pool
    ranked_unused = sorted(
        unused_pool,
        key=lambda slide: (
            -source_slide_score(slide),
            -len(slide.get("problem_texts", [])),
            -int(slide.get("image_count", 0)),
            slide["slide_number"],
        ),
    )
    selected_records = sorted(ranked_unused[:4], key=lambda slide: slide["slide_number"])
    source_numbers = [slide["slide_number"] for slide in selected_records]
    source_records = source_slides_from_numbers(deck, source_numbers)
    source_texts = source_texts_from_numbers(deck, source_numbers, limit=6)
    exact_prompt = best_source_prompt_text(source_records, limit=320)
    tasks = slide_task_list(source_records, limit=3)
    counts = Counter(slide.get("kind", "") for slide in slides)
    problem_rich = sum(1 for slide in source_pool if slide.get("problem_texts"))
    if problem_rich >= 5 and counts.get("practice", 0) < (5 if session_key == "session_1" else 4):
        kind = "practice"
    elif session_key == "session_2" and problem_rich >= 6 and counts.get("challenge", 0) < 1:
        kind = "challenge"
    elif counts.get("worked_example", 0) == 0:
        kind = "worked_example"
    elif session_key == "session_1" and counts.get("guided_notes", 0) == 0:
        kind = "guided_notes"
    else:
        kind = "practice"

    title_map = {
        "guided_notes": "Source-Aligned Guided Notes",
        "worked_example": "Worked Example from the Slides",
        "practice": "Source Problem Practice",
        "challenge": "Source Application Challenge",
        "reflection": "Reflection + Wrap-Up",
    }
    subtitle_map = {
        "guided_notes": "Keep the exact lesson language visible as you take notes.",
        "worked_example": "Walk through the exact source reasoning before you solve independently.",
        "practice": "Use the exact source problem language and show your work clearly.",
        "challenge": "Apply the source idea to a later problem from the deck.",
        "reflection": "Write about what you noticed, solved, and still want to revisit.",
    }
    return build_slide_plan(
        kind=kind,
        title=title_map[kind],
        subtitle=subtitle_map[kind],
        primary_text=exact_prompt or deck.get("summary", ""),
        bullets=source_titles_from_numbers(deck, source_numbers, limit=4),
        tasks=tasks,
        response_prompt=(
            "Use the exact source problem language and explain how you know your strategy fits."
            if kind in {"practice", "challenge", "worked_example"}
            else "Use complete sentences and lesson vocabulary in your reflection."
        ),
        sentence_starters=kernel_sentence_frames(kind, deck, {"sentence_starters": []}),
        source_slide_numbers=source_numbers,
        image_source_slide=source_numbers[0] if source_numbers else 0,
        **build_activity_payload(
            kind=kind,
            deck=deck,
            library=library,
            source_texts=source_texts,
            source_records=source_records,
        ),
    )


def trim_session_slides(slides: list[dict[str, Any]]) -> None:
    protected = {"cover", "be_curious", "learning_target", "worked_example", "exit_ticket"}
    while len(slides) > MAX_SESSION_SLIDES:
        counts = Counter(slide.get("kind", "") for slide in slides)
        removable_index = None
        removable_score = -1
        for index, slide in enumerate(slides):
            kind = slide.get("kind", "")
            if kind in protected:
                continue
            score = {
                "vocabulary": 5,
                "quick_review": 4,
                "reflection": 3,
                "practice": 2,
                "guided_notes": 2,
                "challenge": 1,
            }.get(kind, 1)
            if counts.get(kind, 0) > 1:
                score += 2
            if index >= len(slides) - 2:
                score += 1
            if score > removable_score:
                removable_score = score
                removable_index = index
        if removable_index is None:
            slides.pop()
        else:
            slides.pop(removable_index)


def premium_session_slide_target(deck: dict[str, Any]) -> int:
    slide_count = int(deck.get("slide_count", 0) or 0)
    problem_rich = sum(1 for slide in deck.get("slides", []) if slide.get("problem_texts"))
    image_rich = sum(1 for slide in deck.get("slides", []) if slide.get("image_count"))
    target = PREMIUM_TARGET_SESSION_SLIDES if slide_count >= 8 else MIN_SESSION_SLIDES
    if slide_count >= 12:
        target += 1
    if problem_rich >= 6:
        target += 1
    if problem_rich >= 10:
        target += 1
    if image_rich >= 4:
        target += 1
    return clamp_int(target, MIN_SESSION_SLIDES, MAX_SESSION_SLIDES)


def premium_activity_target(deck: dict[str, Any]) -> int:
    target = MIN_ACTIVITY_SLIDES_PER_SESSION
    problem_rich = sum(1 for slide in deck.get("slides", []) if slide.get("problem_texts"))
    image_rich = sum(1 for slide in deck.get("slides", []) if slide.get("image_count"))
    if premium_session_slide_target(deck) >= 14 or problem_rich >= 6:
        target += 1
    if problem_rich >= 9:
        target += 1
    if problem_rich >= 12 and image_rich >= 4:
        target += 1
    return clamp_int(target, MIN_ACTIVITY_SLIDES_PER_SESSION, MAX_ACTIVITY_SLIDES_PER_SESSION)


def premium_activity_cap(deck: dict[str, Any]) -> int:
    return clamp_int(
        premium_activity_target(deck) + 1,
        MIN_ACTIVITY_SLIDES_PER_SESSION,
        MAX_ACTIVITY_SLIDES_PER_SESSION,
    )


def kernel_sentence_frames(kind: str, deck: dict[str, Any], plan_slide: dict[str, Any]) -> list[str]:
    keyword = deck.get("keyword_candidates", ["the lesson idea"])[0]
    lesson = deck.get("lesson_title", "this lesson")
    frames_map = {
        "be_curious": [
            "I notice ___ because ___.",
            "I wonder ___ because ___.",
            "A lesson word that might matter is ___.",
            "I predict the problem will ask me to ___.",
        ],
        "learning_target": [
            "The content of today's lesson is ___.",
            "The objective is to ___.",
            "Before the lesson, I could ___.",
            "After the lesson, I can ___ because ___.",
        ],
        "guided_notes": [
            f"One important idea about {keyword.lower()} is ___.",
            "The source slide shows ___.",
            "This matters because ___.",
            "I can explain this by saying ___.",
        ],
        "worked_example": [
            "First, ___.",
            "Next, ___.",
            "Then, ___.",
            "Therefore, ___.",
        ],
        "practice": [
            "The problem is asking me to ___.",
            "I started by ___.",
            "I used ___ because ___.",
            "My answer makes sense because ___.",
        ],
        "quick_review": [
            "One thing I remember is ___.",
            "A strategy that still matters is ___.",
            "This connects to ___ because ___.",
        ],
        "challenge": [
            f"This challenge extends {lesson} by ___.",
            "A stronger strategy here is ___.",
            "The evidence for my choice is ___.",
        ],
        "reflection": [
            "At first, ___.",
            "Now I understand ___.",
            "The evidence that helped me was ___.",
            "I can use this when ___.",
        ],
        "exit_ticket": [
            "My final answer is ___.",
            "I know this because ___.",
            "The most important idea from the lesson is ___.",
        ],
    }
    frames = frames_map.get(kind, [
        "I noticed ___.",
        "I learned ___.",
        "I can explain ___ because ___.",
    ])
    if plan_slide.get("sentence_starters"):
        frames = plan_slide["sentence_starters"] + frames
    return unique_nonempty(frames, limit=5)


def be_curious_sentence_kernels(plan_slide: dict[str, Any]) -> tuple[list[str], list[str]]:
    direct_notice = unique_nonempty(plan_slide.get("notice_kernels", []), limit=2)
    direct_wonder = unique_nonempty(plan_slide.get("wonder_kernels", []), limit=2)
    starters = unique_nonempty(plan_slide.get("sentence_starters", []), limit=6)
    notice = unique_nonempty(direct_notice + [starter for starter in starters if "notice" in starter.lower()], limit=2)
    wonder = unique_nonempty(
        direct_wonder
        + [
            starter
            for starter in starters
            if any(token in starter.lower() for token in ("wonder", "predict", "question"))
        ],
        limit=2,
    )
    if not notice:
        notice = [
            "I notice ___ because ___.",
            "I notice the image shows ___.",
        ]
    if not wonder:
        wonder = [
            "I wonder ___ because ___.",
            "I predict ___ because ___.",
        ]
    return notice[:2], wonder[:2]


def be_curious_panel_prompt(prompt: str, kernels: list[str], *, fallback: str) -> str:
    lead = normalize_whitespace(prompt) or fallback
    return lead


def compact_vocab_snapshot_definition(word: str, definition: str) -> str:
    curated = {
        "about": "close amount",
        "between": "in the middle",
        "body weight": "full weight",
        "decimal": "place-value form",
        "estimate": "close answer",
        "exact": "precise amount",
        "fraction": "part of a whole",
        "part": "piece of the whole",
        "percent": "out of 100",
        "whole": "complete amount",
    }
    key = display_text_key(word)
    if key in curated:
        return curated[key]
    return trim_dangling_display_text(truncate_display_copy(definition or "lesson word", 30))


def be_curious_vocabulary_items(plan_slide: dict[str, Any], session_plan: dict[str, Any]) -> list[dict[str, str]]:
    # First, try to get vocabulary from the be_curious slide itself
    structured = [
        item
        for item in plan_slide.get("vocabulary", [])
        if isinstance(item, dict) and normalize_whitespace(item.get("word", ""))
    ]
    # If empty, look through all other slides in the session plan for vocabulary
    if not structured:
        for slide in session_plan.get("slides", []):
            if slide.get("vocabulary"):
                structured = [
                    item
                    for item in slide.get("vocabulary", [])
                    if isinstance(item, dict) and normalize_whitespace(item.get("word", ""))
                ]
                if structured:
                    break

    if structured:
        return [
            {
                **item,
                "definition": compact_vocab_snapshot_definition(item.get("word", ""), item.get("definition", "")),
            }
            for item in structured[:2]
        ]
    fallback_terms = unique_nonempty(plan_slide.get("word_help", []), limit=2)
    return [
        {
            "word": term,
            "definition": compact_vocab_snapshot_definition(term, "Use this lesson word in a complete sentence while you notice or wonder."),
        }
        for term in fallback_terms
    ]


def content_objective_payload(
    deck: dict[str, Any],
    plan_slide: dict[str, Any],
    library: list[dict[str, Any]],
    session_key: str,
) -> dict[str, Any]:
    objective_source_slides = learning_target_source_slides(deck)
    objective_source_numbers = [slide["slide_number"] for slide in objective_source_slides]
    content_objective = normalize_whitespace(plan_slide.get("primary_text", ""))
    generic_content_markers = (
        "explain and apply the key ideas from",
        "use the source lesson",
        "main content focus",
        "track your understanding",
    )
    if not content_objective or any(marker in content_objective.lower() for marker in generic_content_markers):
        content_objective = select_session_objective(deck, session_key)
    content_objective = normalize_i_can_objective(
        content_objective,
        fallback=f"I can explain and apply the key ideas from {deck['lesson_title']}.",
    )
    language_objective = normalize_whitespace(plan_slide.get("secondary_text", ""))
    generic_language_markers = (
        "drag checkmarks",
        "before and after",
        "track your understanding",
        "use the source lesson",
        "show what you can do",
    )
    if (
        not language_objective
        or language_objective.lower() == content_objective.lower()
        or any(marker in language_objective.lower() for marker in generic_language_markers)
    ):
        language_objective = derive_language_objective(deck, content_objective)
    language_objective = normalize_i_can_objective(
        language_objective,
        fallback="I can explain my strategy with labels, vocabulary, and complete sentences.",
    )

    source_numbers = objective_source_numbers or plan_slide.get("source_slide_numbers", [])
    criteria = objective_success_criteria(content_objective, language_objective)
    activity_name = "peer answer check placement"
    record = find_activity_record(activity_name, library)
    return {
        "title": "Content + Language Objectives",
        "subtitle": "Read each objective and check your progress before and after the lesson.",
        "primary_text": content_objective,
        "secondary_text": language_objective,
        "bullets": criteria,
        "sentence_starters": objective_sentence_frames(content_objective),
        "source_slide_numbers": source_numbers,
        "activity_name": record["name"] if record else activity_name,
        "activity_family": record["family"] if record else infer_activity_family(activity_name),
        "activity_instructions": "Place checkmarks in the Before and After boxes for the content and language objectives to track your growth.",
        "movable_pieces": [CHECKMARK_CHIP] * 8,
        "answer_check": "Move a check to After only when your notes, speaking, or practice show evidence that you met the objective.",
    }


def ensure_learning_target_slide(deck: dict[str, Any], plan_slide: dict[str, Any], library: list[dict[str, Any]], session_key: str) -> None:
    payload = content_objective_payload(deck, plan_slide, library, session_key)
    plan_slide["title"] = payload["title"]
    plan_slide["subtitle"] = payload["subtitle"]
    plan_slide["primary_text"] = payload["primary_text"]
    plan_slide["secondary_text"] = payload["secondary_text"]
    plan_slide["bullets"] = payload["bullets"]
    plan_slide["sentence_starters"] = payload["sentence_starters"]
    plan_slide["source_slide_numbers"] = payload["source_slide_numbers"]
    plan_slide["activity_name"] = payload["activity_name"]
    plan_slide["activity_family"] = payload["activity_family"]
    plan_slide["activity_instructions"] = payload["activity_instructions"]
    plan_slide["movable_pieces"] = payload["movable_pieces"]
    plan_slide["answer_check"] = payload["answer_check"]


def ensure_activity_from_library(
    plan_slide: dict[str, Any],
    *,
    kind: str,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    excluded_names: set[str] | None = None,
    force_refresh: bool = False,
) -> None:
    source_numbers = plan_slide.get("source_slide_numbers", [])
    source_records = source_slides_from_numbers(deck, source_numbers)
    source_texts = source_texts_from_numbers(deck, source_numbers, limit=6)
    activity_name = normalize_whitespace(plan_slide.get("activity_name", ""))
    excluded_names = {item.lower() for item in (excluded_names or set())}
    record = find_activity_record(activity_name, library) if activity_name else None
    profile = source_activity_profile(source_records, source_texts)
    record_score = activity_alignment_score(kind, record, profile=profile) if record else -999
    if (
        record
        and not force_refresh
        and record["name"].lower() not in excluded_names
        and record_score >= minimum_activity_alignment_score(kind)
    ):
        plan_slide["activity_name"] = record["name"]
        plan_slide["activity_family"] = record["family"]
        payload = build_activity_payload(
            kind=kind,
            deck=deck,
            library=library,
            source_texts=source_texts,
            source_records=source_records,
            excluded_names=excluded_names,
        )
        plan_slide["activity_instructions"] = payload["activity_instructions"] or plan_slide.get("activity_instructions", "")
        plan_slide["movable_pieces"] = payload["movable_pieces"]
        plan_slide["answer_check"] = payload["answer_check"] or plan_slide.get("answer_check", "")
        return

    payload = build_activity_payload(
        kind=kind,
        deck=deck,
        library=library,
        source_texts=source_texts,
        source_records=source_records,
        excluded_names=excluded_names,
    )
    plan_slide["activity_name"] = payload["activity_name"]
    plan_slide["activity_family"] = payload["activity_family"]
    plan_slide["activity_instructions"] = payload["activity_instructions"]
    plan_slide["movable_pieces"] = payload["movable_pieces"]
    plan_slide["answer_check"] = payload["answer_check"]


def ensure_source_aligned_be_curious_activity(plan_slide: dict[str, Any], *, deck: dict[str, Any]) -> None:
    if plan_slide.get("kind") != "be_curious":
        return
    source_records = source_slides_from_numbers(deck, plan_slide.get("source_slide_numbers", []))
    terms = source_be_curious_terms(source_records, limit=4) if source_records else []
    if not terms:
        terms = [
            item.get("word", "")
            for item in plan_slide.get("vocabulary", [])
            if isinstance(item, dict) and normalize_whitespace(item.get("word", ""))
        ][:4]
    plan_slide["activity_name"] = f"source clue reveal {source_activity_signature(plan_slide)}"
    plan_slide["activity_family"] = "reveal_discuss"
    plan_slide["activity_instructions"] = "Reveal or sort the source clues, terms, or details and discuss what you notice, wonder, or predict."
    plan_slide["movable_pieces"] = unique_nonempty([display_term_label(term) for term in terms], limit=4)
    plan_slide["answer_check"] = "Each clue should connect to something you can notice, wonder, or predict from the source prompt or image."


def clear_activity_fields(plan_slide: dict[str, Any]) -> None:
    plan_slide["activity_name"] = ""
    plan_slide["activity_family"] = ""
    plan_slide["activity_instructions"] = ""
    plan_slide["movable_pieces"] = []
    plan_slide["answer_check"] = ""


def problem_render_mode(plan_slide: dict[str, Any]) -> str:
    mode = normalize_whitespace(plan_slide.get("render_mode", "")).lower()
    if mode in {PROBLEM_RENDER_MODE_FOCUS, PROBLEM_RENDER_MODE_INTERACTIVE}:
        return mode
    return ""


def enforce_problem_then_interactive_sequence(
    session: dict[str, Any],
    *,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
) -> None:
    slides = session.get("slides", [])
    candidates: list[tuple[int, dict[str, Any]]] = []
    for index, slide in enumerate(slides):
        kind = slide.get("kind", "")
        if kind not in PROBLEM_INTERACTIVE_PAIR_KINDS:
            continue
        if normalize_whitespace(slide.get("template_role", "")):
            continue
        candidates.append((index, slide))

    if not candidates:
        return

    for position, (_index, slide) in enumerate(candidates):
        has_follow_up = position + 1 < len(candidates)
        if position % 2 == 0:
            slide["render_mode"] = PROBLEM_RENDER_MODE_FOCUS
            clear_activity_fields(slide)
            subtitle = normalize_whitespace(slide.get("subtitle", "")).lower()
            if not subtitle or "interactive" in subtitle:
                if has_follow_up:
                    slide["subtitle"] = "Solve the full source problem on this page, then use the next page to apply the strategy interactively."
                else:
                    slide["subtitle"] = "Solve the full source problem on this page and explain each step clearly."
            continue

        slide["render_mode"] = PROBLEM_RENDER_MODE_INTERACTIVE
        if not has_activity(slide):
            ensure_activity_from_library(
                slide,
                kind=slide.get("kind", "practice"),
                deck=deck,
                library=library,
            )
        slide["section"] = "Practice Extension"
        title = normalize_whitespace(slide.get("title", ""))
        if not title:
            slide["title"] = "Practice Extension"
        elif "interactive apply" in title.lower():
            slide["title"] = replace_case_insensitive(title, r"\binteractive apply:\s*", "Practice Extension: ")
            slide["title"] = replace_case_insensitive(slide["title"], r"\binteractive apply\b", "Practice Extension")
        elif "practice extension" not in title.lower():
            slide["title"] = truncate_text(f"Practice Extension: {title}", 84)
        subtitle = normalize_whitespace(slide.get("subtitle", "")).lower()
        if not subtitle or "solve the full source problem on this page" in subtitle:
            slide["subtitle"] = "Use draggable pieces to plan, solve, and justify the full source problem."
        if not normalize_whitespace(slide.get("response_prompt", "")):
            slide["response_prompt"] = "Use the interactive pieces to prove your solve path and justify your final answer."
        previous = candidates[position - 1][1] if position > 0 else None
        if not slide.get("source_problem_cards") and previous and previous.get("source_problem_cards"):
            slide["source_problem_cards"] = list(previous.get("source_problem_cards", []))
        if not slide.get("source_slide_numbers") and previous:
            slide["source_slide_numbers"] = list(previous.get("source_slide_numbers", []))
        if not int(slide.get("image_source_slide", 0) or 0) and previous:
            slide["image_source_slide"] = int(previous.get("image_source_slide", 0) or 0)

    ensure_unique_activity_names(slides, deck=deck, library=library)


def ensure_unique_activity_names(
    slides: list[dict[str, Any]],
    *,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
) -> None:
    # Keep activity names unique so design review does not fail.
    used_activity_names: set[str] = set()
    for slide in slides:
        if problem_render_mode(slide) == PROBLEM_RENDER_MODE_FOCUS:
            clear_activity_fields(slide)
            continue
        activity_name = normalize_whitespace(slide.get("activity_name", "")).lower()
        if not activity_name:
            continue
        if activity_name in used_activity_names and not normalize_whitespace(slide.get("template_role", "")):
            ensure_activity_from_library(
                slide,
                kind=slide.get("kind", "practice"),
                deck=deck,
                library=library,
                excluded_names=used_activity_names,
                force_refresh=True,
            )
            activity_name = normalize_whitespace(slide.get("activity_name", "")).lower()
        if not activity_name or activity_name in used_activity_names:
            if problem_render_mode(slide) == PROBLEM_RENDER_MODE_INTERACTIVE:
                ensure_activity_from_library(
                    slide,
                    kind=slide.get("kind", "practice"),
                    deck=deck,
                    library=library,
                    excluded_names=used_activity_names,
                )
                activity_name = normalize_whitespace(slide.get("activity_name", "")).lower()
        if not activity_name or activity_name in used_activity_names:
            clear_activity_fields(slide)
            continue
        used_activity_names.add(activity_name)


def activity_priority(plan_slide: dict[str, Any]) -> int:
    return ACTIVITY_KIND_PRIORITY.get(plan_slide.get("kind", ""), 50)


def assign_visual_anchors(slides: list[dict[str, Any]], *, deck: dict[str, Any], session_key: str) -> None:
    session_images = [slide["slide_number"] for slide in session_source_slides(deck, session_key) if slide.get("image_count")]
    if not session_images:
        return
    primary = session_images[0]
    secondary = session_images[1] if len(session_images) > 1 else primary
    tertiary = session_images[2] if len(session_images) > 2 else secondary
    quaternary = session_images[3] if len(session_images) > 3 else tertiary
    for slide in slides:
        if slide.get("kind") == "cover" and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = primary
        elif slide.get("kind") == "be_curious" and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = primary
        elif slide.get("kind") == "vocabulary" and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = secondary
        elif slide.get("kind") in {"guided_notes", "worked_example"} and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = secondary
        elif slide.get("kind") in {"practice", "quick_review"} and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = tertiary
        elif slide.get("kind") in {"challenge", "reflection"} and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = tertiary
        elif slide.get("kind") == "exit_ticket" and not int(slide.get("image_source_slide", 0) or 0):
            slide["image_source_slide"] = quaternary


def normalize_activity_distribution(
    slides: list[dict[str, Any]],
    *,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
) -> None:
    target = premium_activity_target(deck)
    cap = premium_activity_cap(deck)

    for slide in slides:
        if ACTIVITY_KIND_CAPS.get(slide.get("kind", ""), 0) == 0:
            clear_activity_fields(slide)

    active_entries = [
        (index, slide)
        for index, slide in enumerate(slides)
        if normalize_whitespace(slide.get("activity_name", ""))
        and ACTIVITY_KIND_CAPS.get(slide.get("kind", ""), 0) > 0
    ]
    active_entries.sort(key=lambda item: (-activity_priority(item[1]), item[0]))

    kept_indices: set[int] = set()
    kind_counts: Counter[str] = Counter()
    used_names: set[str] = set()
    for index, slide in active_entries:
        kind = slide.get("kind", "")
        name = normalize_whitespace(slide.get("activity_name", "")).lower()
        if kind_counts[kind] >= ACTIVITY_KIND_CAPS.get(kind, 1) or len(kept_indices) >= cap:
            clear_activity_fields(slide)
            continue
        if name and name in used_names:
            ensure_activity_from_library(
                slide,
                kind=kind,
                deck=deck,
                library=library,
                excluded_names=used_names,
                force_refresh=True,
            )
            name = normalize_whitespace(slide.get("activity_name", "")).lower()
            if not name or name in used_names:
                clear_activity_fields(slide)
                continue
        kept_indices.add(index)
        kind_counts[kind] += 1
        if name:
            used_names.add(name)

    if len(kept_indices) >= target:
        return

    candidate_entries = [
        (index, slide)
        for index, slide in enumerate(slides)
        if ACTIVITY_KIND_CAPS.get(slide.get("kind", ""), 0) > 0 and index not in kept_indices
    ]
    candidate_entries.sort(key=lambda item: (-activity_priority(item[1]), item[0]))

    for index, slide in candidate_entries:
        kind = slide.get("kind", "")
        if kind_counts[kind] >= ACTIVITY_KIND_CAPS.get(kind, 1):
            continue
        ensure_activity_from_library(slide, kind=kind, deck=deck, library=library, excluded_names=used_names)
        name = normalize_whitespace(slide.get("activity_name", "")).lower()
        if not name or name in used_names:
            continue
        kept_indices.add(index)
        kind_counts[kind] += 1
        used_names.add(name)
        if len(kept_indices) >= target:
            break


def source_problem_response_prompt(source_records: list[dict[str, Any]], kind: str, existing_prompt: str = "") -> str:
    blob = combined_source_text(source_records).lower()
    existing = normalize_whitespace(existing_prompt)
    if any(term in blob for term in ("which size", "how many boxes", "would you buy", "popcorn")):
        return "Show the math, compare the options, and defend your choice with source evidence."
    if has_data_analysis_context(blob):
        return "Use the data display, key values, and comparison evidence to explain your answer."
    if any(term in blob for term in ("regulations", "compare to", "shorter than", "longer than")):
        return "Solve the source problem, compare your result to the condition, and explain your conclusion."
    if any(term in blob for term in ("precision", "units", "unit cube", "cubic")):
        return "Show the setup, include precise units, and explain why the units matter."
    if any(term in blob for term in ("missing", "unknown", "length of the box", "relate the volume and dimensions")):
        return "Use the formula to find the missing value and explain why your setup works."
    if "what information do you know" in blob:
        return "Organize the known information, solve the problem, and explain how you know."
    if kind == "worked_example":
        return "Record each step from the model and explain why it fits the source problem."
    if kind == "exit_ticket":
        return "Show the setup, the answer, and one sentence that proves it makes sense."
    return existing or "Show your work and explain how you know your strategy fits the source problem."


def apply_problem_window_to_slide(plan_slide: dict[str, Any], window: dict[str, Any], deck: dict[str, Any]) -> None:
    source_records = window.get("slides", [])
    source_numbers = window.get("source_numbers", [])
    if not source_records or not source_numbers:
        return
    kind = plan_slide.get("kind", "")
    prefer_last = kind in {"challenge", "exit_ticket"}
    raw_prompts = source_problem_candidates(source_records, limit=6)
    ranked_prompts = sorted(
        enumerate(raw_prompts),
        key=lambda item: (-problem_prompt_priority(item[1]), item[0]),
    )
    exact_prompts = [raw_prompts[index] for index, _text in ranked_prompts[:4]]
    support_facts = [
        item
        for item in source_fact_candidates(source_records, limit=4)
        if display_text_key(item) not in {display_text_key(prompt) for prompt in exact_prompts}
    ]
    primary_text = best_source_prompt_text(source_records, limit=320, prefer_last=prefer_last)

    plan_slide["source_slide_numbers"] = list(source_numbers)
    plan_slide["image_source_slide"] = pick_first_image_slide(source_records) or int(plan_slide.get("image_source_slide", 0) or 0)
    if primary_text:
        plan_slide["primary_text"] = primary_text
    if exact_prompts:
        selected_prompt_cards = [
            trim_dangling_display_text(truncate_display_copy(item, 190))
            for item in exact_prompts[:4]
        ]
        plan_slide["tasks"] = selected_prompt_cards
        plan_slide["source_problem_cards"] = selected_prompt_cards
    elif not plan_slide.get("tasks"):
        plan_slide["tasks"] = slide_task_list(source_records, limit=4)
    if support_facts:
        plan_slide["bullets"] = [
            trim_dangling_display_text(truncate_display_copy(item, 180))
            for item in support_facts[:3]
        ]
    elif not plan_slide.get("bullets"):
        plan_slide["bullets"] = source_titles_from_numbers(deck, list(source_numbers), limit=3)
    plan_slide["response_prompt"] = source_problem_response_prompt(source_records, kind, plan_slide.get("response_prompt", ""))
    if kind == "worked_example" and (
        not normalize_whitespace(plan_slide.get("subtitle", ""))
        or "walk through a model" in plan_slide.get("subtitle", "").lower()
    ):
        plan_slide["subtitle"] = "Model the exact source problem step by step before students solve on their own."
    if kind == "practice" and (
        not normalize_whitespace(plan_slide.get("subtitle", ""))
        or "use the source lesson tasks" in plan_slide.get("subtitle", "").lower()
        or "source problem practice" in plan_slide.get("title", "").lower()
    ):
        plan_slide["subtitle"] = "Use the exact source problem language and work through the solve path in the notebook."


def is_exploratory_problem_source(slide: dict[str, Any]) -> bool:
    title = normalize_whitespace(slide.get("title", "")).lower()
    blob = combined_source_text([slide]).lower()
    return any(
        marker in title or marker in blob
        for marker in (
            "be curious",
            "what do you notice",
            "what do you wonder",
            "notice and wonder",
            "turn and talk",
        )
    )


def single_slide_problem_window(slide: dict[str, Any], index: int) -> dict[str, Any]:
    source_number = int(slide.get("slide_number", 0) or 0)
    return {
        "slides": [slide],
        "source_numbers": [source_number],
        "score": problem_source_score(slide),
        "text_blob": combined_source_text([slide]),
        "start_index": index,
    }


def reference_problem_windows(deck: dict[str, Any]) -> list[dict[str, Any]]:
    source_slides = deck.get("slides", [])
    windows: list[tuple[int, bool, dict[str, Any]]] = []
    for index, slide in enumerate(source_slides):
        if not slide.get("slide_number"):
            continue
        prompts = source_problem_candidates([slide], limit=4)
        if not prompts:
            continue
        best_priority = max(problem_prompt_priority(prompt) for prompt in prompts)
        windows.append(
            (
                best_priority,
                is_exploratory_problem_source(slide),
                single_slide_problem_window(slide, index),
            )
        )

    preferred = [window for priority, exploratory, window in windows if priority >= 3 and not exploratory]
    if preferred:
        return preferred
    fallback = [window for priority, _exploratory, window in windows if priority >= 3]
    if fallback:
        return fallback
    if windows:
        return [window for _priority, _exploratory, window in windows]
    scored = problem_window_candidates(source_slides, window_size=2)
    return sorted(
        scored,
        key=lambda item: (-int(item.get("score", 0) or 0), int(item.get("start_index", 0) or 0)),
    )


def align_reference_problem_sources(session: dict[str, Any], deck: dict[str, Any]) -> None:
    windows = reference_problem_windows(deck)
    if not windows:
        return

    problem_index = 0
    for plan_slide in session.get("slides", []):
        kind = plan_slide.get("kind", "")
        template_role = normalize_whitespace(plan_slide.get("template_role", ""))
        if kind not in PROBLEM_SOLVING_KINDS or template_role not in REFERENCE_PROBLEM_TEMPLATE_ROLES:
            continue
        window = windows[min(problem_index, len(windows) - 1)]
        apply_problem_window_to_slide(plan_slide, window, deck)
        problem_index += 1


def align_problem_solving_sources(session_key: str, session: dict[str, Any], deck: dict[str, Any]) -> None:
    session_sources = session_source_slides(deck, session_key)
    candidates = problem_window_candidates(session_sources)
    if not candidates:
        return

    used_signatures: set[tuple[int, ...]] = set()
    practice_ordinal = 0
    rolling_start = 0
    later_start = max(0, len(session_sources) // 2 - 2)
    for plan_slide in session.get("slides", []):
        kind = plan_slide.get("kind", "")
        if kind not in PROBLEM_SOLVING_KINDS:
            continue
        role = kind if kind in {"worked_example", "challenge", "exit_ticket"} else "practice"
        min_start = rolling_start if role in {"worked_example", "practice"} else max(rolling_start, later_start)
        ordinal = practice_ordinal if role == "practice" else 0
        window = select_problem_window_for_slide(
            candidates,
            role=role,
            ordinal=ordinal,
            used_signatures=used_signatures,
            min_start_index=min_start,
        )
        if not window:
            continue
        apply_problem_window_to_slide(plan_slide, window, deck)
        used_signatures.add(tuple(window["source_numbers"]))
        if role == "practice":
            practice_ordinal += 1
            rolling_start = min(window["start_index"] + 1, max(len(session_sources) - 1, 0))
        elif role == "worked_example":
            rolling_start = min(window["start_index"] + 1, max(len(session_sources) - 1, 0))


def enrich_problem_fidelity(plan_slide: dict[str, Any], deck: dict[str, Any]) -> None:
    source_numbers = plan_slide.get("source_slide_numbers", [])
    source_records = source_slides_from_numbers(deck, source_numbers)
    if not source_records:
        return
    prompt_records = list(reversed(source_records)) if plan_slide["kind"] in {"challenge", "exit_ticket"} else source_records
    exact_prompts = source_problem_candidates(prompt_records, limit=4)
    generic_markers = (
        "work through the practice problems",
        "show your work",
        "create a short response",
        "use one final prompt",
        "stretch the lesson",
        "use a later source example",
        "review the lesson",
        "apply the idea",
        "deeper practice",
    )
    if plan_slide["kind"] in {"worked_example", "practice", "quick_review", "challenge", "exit_ticket"}:
        if exact_prompts:
            selected_prompt_cards = [
                trim_dangling_display_text(truncate_display_copy(item, 190))
                for item in exact_prompts[:4]
            ]
            plan_slide["tasks"] = selected_prompt_cards
            plan_slide["source_problem_cards"] = selected_prompt_cards
        elif not plan_slide.get("tasks", []):
            plan_slide["tasks"] = slide_task_list(source_records, limit=4)
    if not plan_slide.get("bullets") and plan_slide["kind"] in {"worked_example", "learning_target"}:
        plan_slide["bullets"] = source_titles_from_numbers(deck, source_numbers, limit=4)
    primary_text = normalize_whitespace(plan_slide.get("primary_text", ""))
    exact_primary = best_source_prompt_text(
        source_records,
        limit=360,
        prefer_last=plan_slide["kind"] in {"challenge", "exit_ticket"},
    )
    if plan_slide["kind"] in {"worked_example", "practice", "quick_review", "challenge", "exit_ticket"} and exact_primary:
        plan_slide["primary_text"] = exact_primary
    elif (
        not primary_text
        or any(marker in primary_text.lower() for marker in generic_markers)
        or "reveal:" in primary_text.lower()
    ) and plan_slide["kind"] in {"guided_notes", "worked_example", "practice", "challenge", "exit_ticket"}:
        plan_slide["primary_text"] = exact_primary or truncate_text(
            source_records[0].get("text") or source_records[0].get("title") or "",
            360,
        )


def ensure_locked_architecture_coverage(
    *,
    session_key: str,
    session: dict[str, Any],
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    objective_source_numbers: list[int],
) -> None:
    slides = session.get("slides", [])
    session_sources = session_source_slides(deck, session_key)
    seed_numbers = [slide["slide_number"] for slide in session_sources[:4] if slide.get("slide_number")]
    launch_numbers = [slide["slide_number"] for slide in session_sources[:3] if slide.get("slide_number")]
    closing_numbers = [slide["slide_number"] for slide in session_sources[-4:] if slide.get("slide_number")]

    def insert_after(target_kinds: set[str], new_slide: dict[str, Any]) -> None:
        insert_at = len(slides)
        for index, slide in enumerate(slides):
            if slide.get("kind") in target_kinds:
                insert_at = index + 1
        slides.insert(insert_at, new_slide)

    if not any(slide.get("kind") == "cover" for slide in slides):
        slides.insert(
            0,
            build_slide_plan(
                kind="cover",
                title=deck.get("lesson_title", "Student Notebook"),
                subtitle=session.get("session_title", "") or ("Session 1 Student Notebook" if session_key == "session_1" else "Session 2 Student Notebook"),
                primary_text=cover_focus_statement(deck, session_key, select_session_objective(deck, session_key)),
                secondary_text=default_cover_subtitle(deck, session_key),
                source_slide_numbers=seed_numbers,
                image_source_slide=pick_first_image_slide(session_sources),
            ),
        )
    if not any(slide.get("kind") == "vocabulary" for slide in slides):
        insert_after(
            {"learning_target", "be_curious"},
            build_slide_plan(
                kind="vocabulary",
                title="Academic Vocabulary",
                subtitle="Keep lesson words visible while you read, talk, and solve.",
                primary_text="Use the source lesson words when you annotate, explain, and justify.",
                vocabulary=source_vocabulary(deck, seed_numbers, limit=4),
                source_slide_numbers=seed_numbers,
                **build_activity_payload(kind="vocabulary", deck=deck, library=library, source_texts=source_texts_from_numbers(deck, seed_numbers, limit=6)),
            ),
        )
    if not any(slide.get("kind") == "guided_notes" for slide in slides):
        insert_after(
            {"vocabulary"},
            build_slide_plan(
                kind="guided_notes",
                title="Guided Notes",
                subtitle="Capture the most important ideas from the source lesson.",
                primary_text=best_source_prompt_text(session_sources[:3], limit=320),
                bullets=slide_title_list(session_sources[:4], limit=4),
                response_prompt="Summarize the main idea in your own words.",
                source_slide_numbers=seed_numbers,
                image_source_slide=pick_first_image_slide(session_sources),
            ),
        )
    if not any(slide.get("kind") == "worked_example" for slide in slides):
        insert_after(
            {"guided_notes"},
            build_slide_plan(
                kind="worked_example",
                title="Guided Practice",
                subtitle="Work through the source example before you solve independently.",
                primary_text=best_source_prompt_text(session_sources, limit=320),
                tasks=slide_task_list(session_sources, limit=3),
                response_prompt="Explain each step and why it fits the source problem.",
                source_slide_numbers=seed_numbers,
                image_source_slide=pick_first_image_slide(session_sources),
                **build_activity_payload(kind="worked_example", deck=deck, library=library, source_texts=source_texts_from_numbers(deck, seed_numbers, limit=6)),
            ),
        )
    practice_like = [slide for slide in slides if slide.get("kind") in {"practice", "challenge"}]
    while len(practice_like) < 2:
        insert_after(
            {"worked_example", "practice", "challenge"},
            build_slide_plan(
                kind="practice",
                title="Source Problem Practice",
                subtitle="Use the exact source problem language and show your work clearly.",
                primary_text=best_source_prompt_text(session_sources, limit=320),
                tasks=slide_task_list(session_sources, limit=3),
                response_prompt="Show your work and explain how you know your strategy fits.",
                source_slide_numbers=closing_numbers or seed_numbers,
                image_source_slide=pick_first_image_slide(session_sources),
                **build_activity_payload(kind="practice", deck=deck, library=library, source_texts=source_texts_from_numbers(deck, closing_numbers or seed_numbers, limit=6)),
            ),
        )
        practice_like = [slide for slide in slides if slide.get("kind") in {"practice", "challenge"}]
    if not any(slide.get("kind") == "reflection" for slide in slides):
        insert_after(
            {"practice", "challenge"},
            build_slide_plan(
                kind="reflection",
                title="Reflection",
                subtitle="Pause, summarize the lesson, and name the idea you want to strengthen.",
                primary_text="What idea feels strongest to you right now?",
                secondary_text="What do you still want to review or apply next?",
                response_prompt="Write a complete-sentence reflection using lesson vocabulary.",
                source_slide_numbers=closing_numbers or seed_numbers,
            ),
        )
    if not any(slide.get("kind") == "exit_ticket" for slide in slides):
        slides.append(
            build_slide_plan(
                kind="exit_ticket",
                title="Exit Ticket",
                subtitle="Show what you can do independently.",
                primary_text=select_session_objective(deck, session_key),
                tasks=slide_task_list(session_sources[-4:] or session_sources, limit=2),
                response_prompt="What is the most important evidence you want your teacher to see?",
                source_slide_numbers=closing_numbers or seed_numbers or objective_source_numbers,
            )
        )


def enforce_plan_requirements(
    plan: dict[str, Any],
    deck: dict[str, Any],
    *,
    custom_guidance: str = "",
) -> dict[str, Any]:
    effective_guidance = enforce_runtime_quality_guidance(custom_guidance)
    library = load_activity_library()
    interactive_kinds = {"be_curious", "vocabulary", "guided_notes", "worked_example", "practice", "quick_review", "challenge", "reflection", "exit_ticket"}
    objective_source_numbers = [slide["slide_number"] for slide in learning_target_source_slides(deck)]
    use_exact_workbook = guidance_has_exact_workbook_contract(effective_guidance)
    if use_exact_workbook:
        plan.pop("session_2", None)

    for session_key in planned_session_keys(plan):
        session = plan.get(session_key, {})
        slides = session.get("slides", [])
        session_sources = session_source_slides(deck, session_key)
        session_seed_numbers = [slide["slide_number"] for slide in session_sources[:4] if slide.get("slide_number")]
        launch_source_records = session_launch_slides(deck, session_key, limit=3)
        launch_source_numbers = [slide["slide_number"] for slide in launch_source_records if slide.get("slide_number")]
        if slides and not any(slide.get("kind") == "learning_target" for slide in slides):
            insert_at = 1 if len(slides) > 1 else len(slides)
            slides.insert(
                insert_at,
                build_slide_plan(
                    kind="learning_target",
                    title="Content + Language Objectives",
                    subtitle="Read each objective and check your progress before and after the lesson.",
                    primary_text=select_session_objective(deck, session_key),
                    source_slide_numbers=objective_source_numbers or session_seed_numbers,
                ),
            )

        if slides and not any(slide.get("kind") == "be_curious" for slide in slides):
            converted = False
            for slide in slides:
                search_blob = " ".join(
                    [
                        slide.get("title", ""),
                        slide.get("subtitle", ""),
                        " ".join(source_titles_from_numbers(deck, slide.get("source_slide_numbers", []), limit=3)),
                    ]
                ).lower()
                if any(marker in search_blob for marker in ("be curious", "notice", "wonder", "same", "different")):
                    slide["kind"] = "be_curious"
                    slide["section"] = KIND_DEFAULT_SECTION["be_curious"]
                    converted = True
                    break
            if not converted:
                slides.insert(
                    1 if len(slides) > 1 else len(slides),
                    build_slide_plan(
                        kind="be_curious",
                        title="Be Curious + Notice / Wonder",
                        subtitle="Look closely before solving and use sentence supports to talk about what you see.",
                        primary_text="Record two details you notice in the source figure, image, or problem setup.",
                        secondary_text="Write one question or prediction that could help you solve the source problem.",
                        sentence_starters=kernel_sentence_frames("be_curious", deck, {"sentence_starters": []}),
                        vocabulary=source_be_curious_vocabulary(deck, launch_source_numbers, limit=3),
                        source_slide_numbers=launch_source_numbers,
                        image_source_slide=pick_launch_image_slide(launch_source_records),
                        **build_activity_payload(
                            kind="be_curious",
                            deck=deck,
                            library=library,
                            source_texts=source_texts_from_numbers(deck, launch_source_numbers, limit=6),
                        ),
                    ),
                )

        if use_exact_workbook:
            rebuilt_session = build_exact_esol_workbook_session(
                deck,
                library,
                session_key,
                objective_source_numbers or session_seed_numbers,
            )
            session["session_title"] = rebuilt_session["session_title"]
            session["session_subtitle"] = rebuilt_session["session_subtitle"]
            session["template_family"] = rebuilt_session["template_family"]
            session["context_anchor"] = rebuilt_session.get("context_anchor", "")
            session["slides"] = rebuilt_session["slides"]
            slides = session.get("slides", [])

        if session_key == "session_2" and not uses_exact_esol_template(session):
            rebuilt_session = build_reference_session2_workbook(deck, library, objective_source_numbers or session_seed_numbers)
            session["session_title"] = rebuilt_session["session_title"]
            session["session_subtitle"] = rebuilt_session["session_subtitle"]
            session["template_family"] = rebuilt_session["template_family"]
            session["slides"] = rebuilt_session["slides"]
            slides = session.get("slides", [])

        if not uses_reference_workbook_template(session):
            ensure_locked_architecture_coverage(
                session_key=session_key,
                session=session,
                deck=deck,
                library=library,
                objective_source_numbers=objective_source_numbers or session_seed_numbers,
            )
            reorder_locked_architecture_sequence(session)
            slides = session.get("slides", [])
        if uses_reference_workbook_template(session):
            align_reference_problem_sources(session, deck)
        else:
            align_problem_solving_sources(session_key, session, deck)
        slides = session.get("slides", [])

        for slide in slides:
            kind = slide.get("kind", "")
            template_role = normalize_whitespace(slide.get("template_role", ""))
            if not template_role:
                enrich_problem_fidelity(slide, deck)
            else:
                source_records = source_slides_from_numbers(deck, slide.get("source_slide_numbers", []))
                if source_records and not slide.get("source_problem_cards"):
                    slide["source_problem_cards"] = source_problem_candidates(source_records, limit=4)
            if kind == "cover":
                if not template_role:
                    ensure_cover_slide(deck, session, slide, session_key)
            elif kind == "learning_target":
                if not template_role:
                    ensure_learning_target_slide(deck, slide, library, session_key)
            elif kind == "be_curious":
                current_sources = source_slides_from_numbers(deck, slide.get("source_slide_numbers", []))
                current_launch_score = sum(launch_source_score(record) for record in current_sources)
                preferred_launch_score = sum(launch_source_score(record) for record in launch_source_records)
                if (
                    launch_source_numbers
                    and preferred_launch_score > current_launch_score
                    and slide.get("template_family") != EXACT_ESOL_TEMPLATE_FAMILY
                ):
                    slide["source_slide_numbers"] = launch_source_numbers
                    slide["image_source_slide"] = pick_launch_image_slide(launch_source_records)
                    slide["vocabulary"] = source_be_curious_vocabulary(deck, launch_source_numbers, limit=3)
                if not template_role:
                    ensure_be_curious_slide(deck, slide)
                    ensure_activity_from_library(slide, kind=kind, deck=deck, library=library)
            elif kind in interactive_kinds and not template_role:
                ensure_activity_from_library(slide, kind=kind, deck=deck, library=library)
            if kind == "vocabulary":
                stronger_vocab = session_esol_vocabulary(deck, slide.get("source_slide_numbers", []), limit=6 if template_role else 4)
                current_vocab = slide.get("vocabulary", [])
                weak_vocab = (
                    not current_vocab
                    or len(current_vocab) < (4 if template_role else 3)
                    or any(
                        is_low_value_vocabulary_term(item.get("word", ""))
                        or item.get("definition", "") == PLACEHOLDER_VOCAB_DEFINITION
                        for item in current_vocab
                        if isinstance(item, dict)
                    )
                )
                if weak_vocab and stronger_vocab:
                    slide["vocabulary"] = stronger_vocab
                if not template_role and (
                    not normalize_whitespace(slide.get("subtitle", "")) or "composite figures" in slide.get("subtitle", "").lower()
                ):
                    slide["subtitle"] = "Keep lesson words visible while you model, solve, and explain the source math."
                if not template_role and (
                    not normalize_whitespace(slide.get("primary_text", "")) or "talk, and justify" in slide.get("primary_text", "").lower()
                ):
                    slide["primary_text"] = "Use the source lesson words when you label models, explain your strategy, and justify your answer."
            if kind not in {"cover", "vocabulary", "learning_target"} and not template_role:
                slide["sentence_starters"] = kernel_sentence_frames(kind, deck, slide)
        assign_visual_anchors(slides, deck=deck, session_key=session_key)
        if uses_reference_workbook_template(session):
            session["context_anchor"] = best_session_context_anchor(deck, session_key)
            for slide in slides:
                if slide.get("kind") in {"worked_example", "practice", "reflection", "exit_ticket"}:
                    slide["context_anchor"] = session["context_anchor"]
        else:
            normalize_activity_distribution(slides, deck=deck, library=library)

            target_slide_count = premium_session_slide_target(deck)
            while len(slides) < target_slide_count:
                slides.append(build_extension_slide(session_key=session_key, deck=deck, slides=slides, library=library))
            trim_session_slides(slides)
            reorder_locked_architecture_sequence(session)
            slides = session.get("slides", [])
            align_problem_solving_sources(session_key, session, deck)
            slides = session.get("slides", [])
            normalize_activity_distribution(slides, deck=deck, library=library)
            apply_premium_decision_layer(deck, session_key, session)
            apply_locked_architecture_labels(session)

        enforce_problem_then_interactive_sequence(session, deck=deck, library=library)
        for slide in session.get("slides", []):
            ensure_source_aligned_be_curious_activity(slide, deck=deck)
            ensure_source_anchored_problem_activity(slide)
            ensure_peer_discussion_support(slide)
            refresh_problem_activity_for_context(slide)
        ensure_unique_activity_names(session.get("slides", []), deck=deck, library=library)
        if uses_exact_esol_template(session):
            session["flagship_activity_requested_target"] = 0
            session["flagship_activity_target"] = 0
            session["flagship_activity_support_score"] = 0
            session["flagship_activity_fallbacks"] = [
                "kept the notebook intentionally compact so students move into the book after the first problem arc"
            ]
            session["flagship_activities"] = []
            continue
        ensure_interactive_engagement_layer(deck, session_key, session, library=library)
        apply_flagship_activity_layer(deck, session_key, session)
        ensure_unique_activity_names(session.get("slides", []), deck=deck, library=library)

    plan = apply_publisher_copyedit(plan, deck)
    for session_key in planned_session_keys(plan):
        session = plan.get(session_key, {})
        for slide in session.get("slides", []):
            ensure_peer_discussion_support(slide)
    return validate_plan(plan, deck=deck)


def optimize_slide_copy_for_layout(slide: dict[str, Any]) -> None:
    kind = slide.get("kind", "")
    template_role = normalize_whitespace(slide.get("template_role", ""))
    primary_limit = PRIMARY_TEXT_LIMITS.get(kind, 150)
    secondary_limit = SECONDARY_TEXT_LIMITS.get(kind, 100)
    task_limit = TASK_TEXT_LIMITS.get(kind, 104)
    task_count_limit = TASK_COUNT_LIMITS.get(kind, 3)
    bullet_limit = BULLET_TEXT_LIMITS.get(kind, 88)
    response_limit = RESPONSE_PROMPT_LIMITS.get(kind, 96)
    sentence_limit = 78 if kind in {"learning_target", "be_curious"} else 84

    if kind in PROBLEM_SOLVING_KINDS:
        primary_segments: list[str] = []
        secondary_segments: list[str] = []
        short_subtitle = "" if template_role else short_problem_solving_subtitle(slide)
        if short_subtitle:
            slide["subtitle"] = short_subtitle
        slide["primary_text"] = truncate_display_copy(slide.get("primary_text", ""), primary_limit)
        slide["secondary_text"] = truncate_display_copy(slide.get("secondary_text", ""), secondary_limit)
    else:
        primary_segments = split_prompt_for_layout(slide.get("primary_text", ""), limit=primary_limit, max_parts=3)
        secondary_segments = split_prompt_for_layout(slide.get("secondary_text", ""), limit=secondary_limit, max_parts=2)
        slide["primary_text"] = primary_segments[0] if primary_segments else truncate_text(normalize_whitespace(slide.get("primary_text", "")), primary_limit)
        slide["secondary_text"] = secondary_segments[0] if secondary_segments else truncate_text(normalize_whitespace(slide.get("secondary_text", "")), secondary_limit)
    slide["title"] = trim_dangling_display_text(slide.get("title", ""))
    subtitle_limit = 74 if kind in PROBLEM_SOLVING_KINDS else 132
    slide["subtitle"] = trim_dangling_display_text(truncate_display_copy(slide.get("subtitle", ""), subtitle_limit))
    slide["primary_text"] = trim_dangling_display_text(slide.get("primary_text", ""))
    slide["secondary_text"] = trim_dangling_display_text(slide.get("secondary_text", ""))

    overflow_segments = primary_segments[1:]
    if overflow_segments:
        if kind in {"guided_notes", "worked_example"}:
            slide["bullets"] = unique_nonempty(overflow_segments + list(slide.get("bullets", [])), limit=4)
        elif kind in {"practice", "quick_review", "challenge", "exit_ticket"}:
            slide["tasks"] = overflow_segments + list(slide.get("tasks", []))

    if kind in PROBLEM_SOLVING_KINDS:
        merged_tasks = merge_fragmented_items(list(slide.get("tasks", [])), limit=task_count_limit * 2)
        slide["tasks"] = [
            truncate_display_copy(item, task_limit)
            for item in unique_nonempty(merged_tasks, limit=task_count_limit)
            if item and not has_dangling_display_text(item)
        ]
    else:
        slide["tasks"] = compact_tasks_for_layout(list(slide.get("tasks", [])), item_limit=task_limit, max_items=task_count_limit)
    bullet_count_limit = 3 if kind == "learning_target" else 4
    if kind in PROBLEM_SOLVING_KINDS:
        slide["bullets"] = [truncate_display_copy(item, bullet_limit) for item in unique_nonempty(slide.get("bullets", []), limit=bullet_count_limit)]
    else:
        slide["bullets"] = [truncate_text(item, bullet_limit) for item in unique_nonempty(slide.get("bullets", []), limit=bullet_count_limit)]
    slide["response_prompt"] = trim_dangling_display_text(truncate_text(normalize_whitespace(slide.get("response_prompt", "")), response_limit))
    slide["tasks"] = [item for item in slide["tasks"] if item and not has_dangling_display_text(item)]
    slide["bullets"] = [item for item in slide["bullets"] if item and not has_dangling_display_text(item)]
    slide["sentence_starters"] = [truncate_text(item, sentence_limit) for item in unique_nonempty(slide.get("sentence_starters", []), limit=4)]
    slide["activity_instructions"] = truncate_display_copy(slide.get("activity_instructions", ""), 104)
    slide["answer_check"] = truncate_display_copy(slide.get("answer_check", ""), 92)
    slide["premium_title"] = trim_dangling_display_text(truncate_text(normalize_whitespace(slide.get("premium_title", "")), 96))
    slide["premium_text"] = truncate_display_copy(slide.get("premium_text", ""), 156)
    slide["premium_items"] = [
        trim_dangling_display_text(truncate_text(item, 72)) for item in unique_nonempty(slide.get("premium_items", []), limit=6)
        if item and not has_dangling_display_text(item)
    ]
    premium_table: list[list[str]] = []
    for row in slide.get("premium_table", [])[:5]:
        if not isinstance(row, list):
            continue
        premium_table.append([truncate_text(normalize_whitespace(str(cell)), 64) for cell in row[:4]])
    slide["premium_table"] = premium_table
    slide["partner_prompt"] = trim_dangling_display_text(truncate_text(normalize_whitespace(slide.get("partner_prompt", "")), 116))
    slide["discussion_questions"] = [
        clean_discussion_question(truncate_text(item, 104))
        for item in unique_nonempty(slide.get("discussion_questions", []), limit=3)
        if clean_discussion_question(item)
    ]
    slide["context_anchor"] = trim_dangling_display_text(truncate_text(normalize_whitespace(slide.get("context_anchor", "")), 52))
    slide["practice_phase"] = trim_dangling_display_text(truncate_text(normalize_whitespace(slide.get("practice_phase", "")), 32))


def named_template_activity_payload(
    *,
    name: str,
    family: str,
    instructions: str,
    movable_pieces: list[str],
    answer_check: str,
    library: list[dict[str, Any]],
) -> dict[str, Any]:
    record = find_activity_record(name, library)
    return {
        "activity_name": record["name"] if record else name,
        "activity_family": record["family"] if record else family,
        "activity_instructions": instructions,
        "movable_pieces": unique_nonempty(movable_pieces, limit=6),
        "answer_check": answer_check,
    }


def compressed_notebook_lesson_mode(profile: dict[str, Any], signal_text: str) -> str:
    lowered = normalize_whitespace(signal_text).lower()
    if profile.get("topic") == "data_analysis" or any(
        term in lowered
        for term in ("data set", "ordered data", "line plot", "dot plot", "graph", "table", "median", "mean", "outlier")
    ):
        return "statistics_data"
    if any(
        term in lowered
        for term in (
            "proportional",
            "ratio",
            "rate",
            "unit rate",
            "constant of proportionality",
            "independent variable",
            "dependent variable",
            "relationship",
        )
    ):
        return "proportional_relationships"
    if any(
        term in lowered
        for term in ("equation", "expression", "variable", "coefficient", "constant", "simplify", "solve for")
    ):
        return "equations_expressions"
    if profile.get("topic") not in {"generic_math", "fraction_decimal_percent"} or any(
        term in lowered
        for term in (
            "area",
            "volume",
            "measurement",
            "perimeter",
            "length",
            "width",
            "height",
            "base",
            "triangle",
            "trapezoid",
            "parallelogram",
            "rhombus",
            "prism",
        )
    ):
        return "geometry_measurement"
    return "equations_expressions"


def exact_reference_flow_payload(mode: str) -> tuple[str, list[str]]:
    if mode == "proportional_relationships":
        return (
            "IV -> DV Reference Tool",
            [
                "IV: identify the quantity you control or choose.",
                "Rule: explain how the relationship changes from one value to the next.",
                "DV: name the quantity that changes because of the IV.",
            ],
        )
    if mode == "equations_expressions":
        return (
            "Expression / Equation Tool",
            [
                "Name the variable or missing quantity.",
                "Build the expression or equation that matches the source pattern.",
                "Substitute carefully and check that the relationship still makes sense.",
            ],
        )
    if mode == "statistics_data":
        return (
            "Graph / Table Reference Tool",
            [
                "Read the graph or table before you make a claim.",
                "Pull the strongest values or comparisons as evidence.",
                "Use the evidence to justify your conclusion in words and numbers.",
            ],
        )
    return (
        "Solve Path Reference Tool",
        [
            "Name the measurements or labels you know.",
            "Choose the model or formula that fits the figure.",
            "Substitute, solve, label the answer, and explain what it means.",
        ],
    )


def exact_guided_twr_frames(mode: str) -> list[str]:
    if mode == "proportional_relationships":
        return [
            "Think: The independent variable is ___.",
            "Write: The dependent variable is ___.",
            "Reason: The rule works because ___.",
        ]
    if mode == "equations_expressions":
        return [
            "Think: The variable represents ___.",
            "Write: The matching equation or expression is ___.",
            "Reason: This step works because ___.",
        ]
    if mode == "statistics_data":
        return [
            "Think: The graph or table shows ___.",
            "Write: The strongest evidence is ___.",
            "Reason: That evidence supports the claim because ___.",
        ]
    return [
        "Think: The important measurement is ___.",
        "Write: The formula or model is ___.",
        "Reason: The answer makes sense because ___.",
    ]


def build_exact_interactive_activity_slide(
    *,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    profile: dict[str, Any],
    mode: str,
    guided_problem: str,
    context_anchor: str,
    source_numbers: list[int],
    source_records: list[dict[str, Any]],
    source_problem_cards: list[str],
    vocab_items: list[dict[str, str]],
) -> dict[str, Any]:
    slide = with_template_metadata(
        build_slide_plan(
            kind="practice",
            section="Interactive Activity",
            title="Interactive Activity",
            subtitle="Use the lesson language, visuals, and source problem to try the idea with support.",
            primary_text=guided_problem,
            source_slide_numbers=source_numbers,
            image_source_slide=pick_first_image_slide(source_records),
        ),
        "interactive_activity",
    )
    slide["source_problem_cards"] = source_problem_cards

    if mode == "proportional_relationships":
        slide.update(
            {
                "subtitle": "Sort the lesson quantities into IV or DV, then defend one choice.",
                "interactive_render_role": "drag_sort",
                "tasks": ["Independent Variable", "Dependent Variable"],
                "response_prompt": "Which quantity depends on the other, and how do you know?",
                **named_template_activity_payload(
                    name="strategy sort mat",
                    family="sort_classify",
                    instructions="Drag each card into IV or DV. Then explain one sort with source evidence.",
                    movable_pieces=[item.get("word", "") for item in vocab_items[:4]] + ["x-value", "y-value"],
                    answer_check="Each card should land with the quantity it describes in the source relationship.",
                    library=library,
                ),
            }
        )
        return slide

    if mode == "equations_expressions":
        _title, text, items = common_mistake_payload(guided_problem, context_anchor)
        slide.update(
            {
                "subtitle": "Find the mistake, fix it, and explain why the corrected equation or expression works.",
                "interactive_render_role": "error_analysis",
                "primary_text": text,
                "response_prompt": "The mistake is ___. The correct step is ___.",
                "sentence_starters": ["The mistake is ___.", "The correct step is ___.", "This works because ___."],
                "error_steps": items,
                "fix_it_text": "Fix it: ___.",
                "why_prompt": "Why does this fix work?",
                **named_template_activity_payload(
                    name="error analysis repair cards",
                    family="detect_justify",
                    instructions="Find the incorrect move, repair it, and justify the corrected reasoning.",
                    movable_pieces=["What is correct?", "What is incorrect?", "Fix it", "Explain it"],
                    answer_check="Your correction should match the source relationship, labels, and reasoning.",
                    library=library,
                ),
            }
        )
        return slide

    if mode == "statistics_data":
        _title, text, items = common_mistake_payload(
            best_source_prompt_text(source_records, limit=220) or guided_problem,
            context_anchor,
        )
        slide.update(
            {
                "subtitle": "Use the graph or table to find the mistake, fix it, and defend your correction.",
                "interactive_render_role": "error_analysis",
                "primary_text": text,
                "response_prompt": "Which graph or table clue shows the mistake?",
                "sentence_starters": ["The data shows ___.", "The mistake is ___.", "The fix works because ___."],
                "error_steps": items,
                "fix_it_text": "Fix it with the correct graph or table evidence: ___.",
                "why_prompt": "How does the graph or table prove your fix?",
                **named_template_activity_payload(
                    name="error analysis repair cards",
                    family="detect_justify",
                    instructions="Use the graph or table evidence to repair the mistake and justify the corrected claim.",
                    movable_pieces=["Read the data", "Find the mistake", "Fix the claim", "Justify with evidence"],
                    answer_check="A strong correction cites the graph or table evidence directly.",
                    library=library,
                ),
            }
        )
        return slide

    comparison_slide = {**slide}
    comparison_slide.update(
        {
            "subtitle": "Partner A and Partner B solve related problems, then compare the setup and reasoning.",
            "interactive_render_role": "two_column_compare",
            "response_prompt": "How are the two solve paths alike or different?",
            "sentence_starters": ["Both setups use ___.", "One difference is ___.", "The labeled diagram shows ___."],
            **named_template_activity_payload(
                name="compare-and-rank cards",
                family="compare_rank",
                instructions="Solve the paired problems, then compare the setup, labels, and reasoning.",
                movable_pieces=["Partner A", "Partner B", "Same", "Different", "Best evidence"],
                answer_check="A strong comparison explains what stays the same, what changes, and why.",
                library=library,
            ),
        }
    )
    comparison_slide["partner_a_problem"] = source_problem_statement(comparison_slide)
    comparison_slide["partner_b_problem"] = similar_problem_statement(comparison_slide)
    comparison_slide["compare_frame"] = "The setup is alike because ___. It changes because ___."
    return comparison_slide


def build_exact_review_slide(
    *,
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    profile: dict[str, Any],
    mode: str,
    guided_problem: str,
    context_anchor: str,
    source_numbers: list[int],
    source_records: list[dict[str, Any]],
    source_problem_cards: list[str],
) -> dict[str, Any]:
    slide = with_template_metadata(
        build_slide_plan(
            kind="practice",
            section="Interactive Review",
            title="Best-Fit Interactive Review",
            subtitle="Close with the strongest discussion-ready review for this lesson.",
            primary_text=guided_problem,
            source_slide_numbers=source_numbers,
            image_source_slide=pick_first_image_slide(source_records),
        ),
        "best_fit_review",
    )
    slide["source_problem_cards"] = source_problem_cards

    if mode == "statistics_data":
        slide.update(
            {
                "review_render_role": "two_column_compare",
                "subtitle": "Compare two evidence-based explanations and decide which one is stronger.",
                "response_prompt": "Which explanation uses stronger evidence from the graph or table?",
                "sentence_starters": ["Both explanations use ___.", "The stronger evidence is ___.", "I would revise ___ because ___."],
                **named_template_activity_payload(
                    name="compare-and-rank cards",
                    family="compare_rank",
                    instructions="Compare two evidence-based explanations, then defend the stronger one.",
                    movable_pieces=["Claim A", "Claim B", "Same", "Different", "Best evidence"],
                    answer_check="A strong review names the better evidence and explains why it is stronger.",
                    library=library,
                ),
            }
        )
        slide["partner_a_problem"] = source_problem_statement(slide)
        slide["partner_b_problem"] = similar_problem_statement(slide)
        slide["compare_frame"] = "Explanation ___ is stronger because the graph or table shows ___."
        return slide

    if mode == "geometry_measurement":
        slide.update(
            {
                "review_render_role": "collaborative_practice",
                "subtitle": "Teach the setup with a partner, then compare which labels and measurements mattered most.",
                "response_prompt": "What should your partner notice first in the labeled diagram?",
                "sentence_starters": ["Partner A noticed ___.", "Partner B solved ___.", "The diagram proves ___."],
                **named_template_activity_payload(
                    name="reasoning ladder build",
                    family="build_construct",
                    instructions="Work with a partner to plan the solve path before you explain it aloud.",
                    movable_pieces=["What do we know?", "Choose the formula", "Show the labels", "Check the answer"],
                    answer_check="A strong review keeps the labels, formula, and explanation aligned.",
                    library=library,
                ),
            }
        )
        slide["partner_a_problem"] = source_problem_statement(slide)
        slide["partner_b_problem"] = similar_problem_statement(slide)
        slide["discussion_questions"] = [
            "What stays the same in both problems?",
            "Which label or measurement changes your setup?",
            "Which explanation is easiest to follow, and why?",
        ]
        return slide

    slide.update(
        {
            "review_render_role": "turn_and_teach",
            "subtitle": "Teach the strategy to a partner, listen closely, and revise one part of your explanation.",
            "response_prompt": "Which clue should your partner notice first, and why?",
            "sentence_starters": ["My first move is ___.", "The key clue is ___.", "After the discussion, I revised ___."],
            "partner_prompt": "Solve first. Then teach the strategy to a partner and revise one sentence after the discussion.",
            **named_template_activity_payload(
                name="turn-and-teach flow",
                family="sequence_order",
                instructions="Build the teach-back flow in order, then coach a partner through the source strategy.",
                movable_pieces=["Name the clue", "Choose the strategy", "Show the reasoning", "Teach the check"],
                answer_check="A strong teach-back uses source evidence, vocabulary, and clear reasoning.",
                library=library,
            ),
        }
    )
    slide["discussion_questions"] = [
        f"Which clue in the {source_problem_focus_phrase(slide)} should your partner notice first?",
        "Which part of the explanation makes the strategy easiest to follow?",
        "What would you revise after listening to your partner?",
    ]
    return slide


def build_exact_esol_workbook_session(
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    session_key: str,
    objective_source_numbers: list[int],
) -> dict[str, Any]:
    session_sources = session_source_slides(deck, session_key)
    core = core_instructional_slides(session_sources, limit=6)
    launch = session_launch_slides(deck, session_key, limit=3) or session_sources[:3] or session_sources
    seed_numbers = [slide["slide_number"] for slide in (core or session_sources)[:4] if slide.get("slide_number")]
    launch_numbers = [slide["slide_number"] for slide in launch[:3] if slide.get("slide_number")]
    problem_windows = problem_window_candidates(session_sources, window_size=3)
    first_problem_window = (
        select_problem_window_for_slide(problem_windows, role="worked_example", ordinal=0, used_signatures=set())
        if problem_windows
        else None
    )
    problem_numbers = list(first_problem_window.get("source_numbers", [])) if first_problem_window else []
    if not problem_numbers:
        problem_numbers = [
            slide["slide_number"]
            for slide in (core[1:4] or session_sources[1:4] or session_sources[:3])
            if slide.get("slide_number")
        ]
    problem_numbers = instructional_focus_numbers(deck, problem_numbers, seed_numbers)
    problem_records = source_slides_from_numbers(deck, problem_numbers or seed_numbers)
    session_label = session_label_for_key(session_key)
    session_number = session_number_value(session_key)
    standard_text = first_standard_text(deck) or "Standard"
    context_anchor = best_session_context_anchor(deck, session_key)
    profile = session_math_profile(deck, session_key)
    current_formula = formula_for_session(deck, session_key) or profile.get("formula", "")
    source_numbers = instructional_focus_numbers(deck, problem_numbers or seed_numbers, seed_numbers)
    focus_records = source_slides_from_numbers(deck, source_numbers)
    focus_image_source = pick_first_image_slide(focus_records) or pick_first_image_slide(
        source_slides_from_numbers(deck, launch_numbers or seed_numbers)
    )
    objective_rows = exact_content_objectives(deck, session_key, current_formula)
    language_objective = derive_language_objective(deck, objective_rows[0])
    spoken_frame = spoken_language_frame(profile, current_formula)
    vocab_items = session_esol_vocabulary(deck, source_numbers, limit=6)
    guided_seed = source_seed_slide(deck, problem_numbers or seed_numbers, kind="worked_example")
    guided_problem = source_problem_statement(guided_seed)
    first_problem_cards = source_problem_candidates(problem_records, limit=4)
    if first_problem_cards:
        guided_problem = first_problem_cards[0]
    guided_problem = guided_problem or best_source_prompt_text(problem_records, limit=320)
    data_analysis_focus = has_data_analysis_context(
        combined_source_text(problem_records)
    )
    guided_table = build_reference_table(profile, guided_seed, completed=False)
    lesson_mode = compressed_notebook_lesson_mode(
        profile,
        " ".join(
            [
                deck.get("lesson_title", ""),
                select_session_objective(deck, session_key),
                combined_source_text(session_sources),
            ]
        ),
    )
    reference_flow_title, reference_flow_lines = exact_reference_flow_payload(lesson_mode)

    learning_slide = with_template_metadata(
        build_slide_plan(
            kind="learning_target",
            section="Objectives + Session Map",
            title="Objectives + Session Map",
            subtitle="Start with the goal, the standard, and today's four-step path.",
            primary_text=objective_rows[0],
            secondary_text=language_objective,
            source_slide_numbers=objective_source_numbers or seed_numbers,
        ),
        "learning_objectives",
    )
    learning_slide["content_objective_b"] = objective_rows[1]
    learning_slide["standard_label"] = standard_text
    learning_slide["language_frame_support"] = spoken_frame
    learning_slide["session_map_steps"] = ["Be Curious", "Vocabulary", "Try It", "Discuss"]

    curious_slide = with_template_metadata(
        build_slide_plan(
            kind="be_curious",
            section="Be Curious",
            title="Be Curious",
            subtitle="Look closely at the lesson setup before you solve.",
            primary_text="👁 Notice: ___.",
            secondary_text="🤔 Wonder: ___.",
            response_prompt="What clue might help with the first problem?",
            source_slide_numbers=source_numbers,
            image_source_slide=focus_image_source,
            **named_template_activity_payload(
                name="notice and wonder reveal",
                family="reveal_discuss",
                instructions="Read one clue at a time. Then share a notice or wonder.",
                movable_pieces=["Look", "Notice", "Wonder", "Share"],
                answer_check="Each notice or wonder should connect to the image, model, or first problem.",
                library=library,
            ),
        ),
        "prior_session_review",
    )
    curious_slide["reference_panel_title"] = "First Problem"
    curious_slide["bridge_sentence"] = "Use one notice to predict how the first problem might be solved."
    curious_slide["prior_formula"] = current_formula
    curious_slide["notice_lines"] = exact_notice_lines(profile)
    curious_slide["wonder_prompt"] = exact_wonder_prompt(profile)
    curious_slide["word_help"] = word_help_strip(deck, source_numbers, limit=5)

    vocabulary_slide = with_template_metadata(
        build_slide_plan(
            kind="vocabulary",
            section="Vocabulary + Reference Tool",
            title="Vocabulary + Reference Tool",
            subtitle="Keep the lesson words and the solve-reference tool visible while you work.",
            primary_text="Match each word to a definition, example, and visual clue from the lesson.",
            source_slide_numbers=source_numbers,
            image_source_slide=focus_image_source,
        ),
        "vocabulary_table",
    )
    vocabulary_slide["vocabulary"] = vocab_items
    vocabulary_slide["secondary_text"] = vocabulary_speaking_frame(vocab_items)
    vocabulary_slide["reference_flow_title"] = reference_flow_title
    vocabulary_slide["reference_flow_lines"] = reference_flow_lines

    if data_analysis_focus:
        worked_example_activity = named_template_activity_payload(
            name="claim and evidence organizer",
            family="compare_rank",
            instructions="Organize the data evidence before writing your conclusion.",
            movable_pieces=["Read the data", "Find evidence", "Make a claim", "Check the claim"],
            answer_check="A strong explanation connects a clear claim to source evidence.",
            library=library,
        )
        worked_example_prompt = "Explain how your evidence supports the claim."
        worked_example_steps = [
            "1. Read the source data or chart: ___.",
            "2. Collect two evidence points: ___.",
            "3. Write and justify your claim: ___.",
        ]
        worked_example_label = f"Evidence Focus (S{session_number})"
        worked_example_text = "Use source data, trend language, and evidence to justify your claim."
        guided_practice_activity = named_template_activity_payload(
            name="reasoning pathway",
            family="sequence_order",
            instructions="Sequence the reasoning moves, then write a supported conclusion.",
            movable_pieces=["Read the data", "Identify trend", "State claim", "Support with evidence"],
            answer_check="A strong path includes a clear claim and source evidence.",
            library=library,
        )
        follow_up_tasks = [
            "Re-read the first problem and underline the key data or display clue.",
            "Solve it on your own and defend the claim with evidence.",
            "Continue the remaining practice in your book.",
        ]
    else:
        worked_example_activity = named_template_activity_payload(
            name="equation builder",
            family="build_construct",
            instructions="Build the formula path before you solve.",
            movable_pieces=["Read the problem", "Write the formula", "Use the values", "Label the answer"],
            answer_check="A strong solve path matches the formula and the final answer.",
            library=library,
        )
        worked_example_prompt = "Explain how the formula fits the problem."
        worked_example_steps = [
            "1. Write the formula: ___.",
            "2. Add the values: ___.",
            "3. Solve and label: ___.",
        ]
        worked_example_label = f"Formula (S{session_number})"
        worked_example_text = current_formula or profile.get("formula", "")
        guided_practice_activity = named_template_activity_payload(
            name="equation builder",
            family="build_construct",
            instructions="Build the solve path, then write the final labeled answer.",
            movable_pieces=["Read", "Write formula", "Substitute", "Solve + check"],
            answer_check="A strong path matches the source strategy and includes units.",
            library=library,
        )
        follow_up_tasks = [
            "Re-read the first problem and label the givens.",
            "Solve it on your own and check that your work matches the model.",
            "Continue the remaining practice in your book.",
        ]

    guided_slide = with_template_metadata(
        build_slide_plan(
            kind="worked_example",
            section="Guided Problem",
            title="Guided Problem",
            subtitle="Keep the source problem visible and work through the setup together.",
            primary_text=guided_problem,
            response_prompt=worked_example_prompt,
            source_slide_numbers=source_numbers,
            image_source_slide=pick_first_image_slide(problem_records),
            **worked_example_activity,
        ),
        "guided_practice",
    )
    guided_slide["source_problem_cards"] = first_problem_cards or source_problem_candidates(problem_records, limit=4)
    guided_slide["context_hook"] = context_hook_text(context_anchor, deck.get("lesson_title", "Student Notebook"))
    guided_slide["formula_label"] = worked_example_label
    guided_slide["formula_text"] = worked_example_text
    guided_slide["guided_table"] = guided_table
    guided_slide["guided_steps"] = worked_example_steps
    guided_slide["twr_frames"] = exact_guided_twr_frames(lesson_mode)

    interactive_activity_slide = build_exact_interactive_activity_slide(
        deck=deck,
        library=library,
        profile=profile,
        mode=lesson_mode,
        guided_problem=guided_problem,
        context_anchor=context_anchor,
        source_numbers=source_numbers,
        source_records=problem_records,
        source_problem_cards=guided_slide["source_problem_cards"],
        vocab_items=vocab_items,
    )
    review_slide = build_exact_review_slide(
        deck=deck,
        library=library,
        profile=profile,
        mode=lesson_mode,
        guided_problem=guided_problem,
        context_anchor=context_anchor,
        source_numbers=source_numbers,
        source_records=problem_records,
        source_problem_cards=guided_slide["source_problem_cards"],
    )

    slides = [
        learning_slide,
        curious_slide,
        vocabulary_slide,
        guided_slide,
        interactive_activity_slide,
        review_slide,
    ]
    return {
        "session_title": f"{session_label} Student Notebook",
        "session_subtitle": "A six-slide launch notebook that moves from source-faithful setup into interactive practice and review.",
        "template_family": EXACT_ESOL_TEMPLATE_FAMILY,
        "context_anchor": context_anchor,
        "slides": slides,
    }


def build_reference_session2_workbook(
    deck: dict[str, Any],
    library: list[dict[str, Any]],
    objective_source_numbers: list[int],
) -> dict[str, Any]:
    session_sources = session_source_slides(deck, "session_2")
    core = core_instructional_slides(session_sources, limit=10)
    launch = session_launch_slides(deck, "session_2", limit=3) or session_sources[:3] or session_sources
    seed_numbers = [slide["slide_number"] for slide in (core or session_sources)[:4] if slide.get("slide_number")]
    launch_numbers = [slide["slide_number"] for slide in launch[:3] if slide.get("slide_number")]
    middle_numbers = [slide["slide_number"] for slide in (core[1:5] or session_sources[1:5] or session_sources[:4]) if slide.get("slide_number")]
    later_numbers = [slide["slide_number"] for slide in (session_sources[-5:] or core[-5:] or session_sources) if slide.get("slide_number")]
    closing_numbers = [slide["slide_number"] for slide in (session_sources[-4:] or core[-4:] or session_sources[:4]) if slide.get("slide_number")]
    context_anchor = best_session_context_anchor(deck, "session_2")
    lesson_title = deck.get("lesson_title", "Student Notebook")
    content_objective = select_session_objective(deck, "session_2")
    language_objective = derive_language_objective(deck, content_objective)

    sort_activity = reference_sort_activity_content(
        deck,
        middle_numbers or seed_numbers,
        default_tasks=["Fits the source strategy", "Needs a second look"],
        default_pieces=[
            "Shipping rule",
            "Write V = l x w x h",
            "Divide by width x height",
            "Compare to 12 inches",
            "Dimensions",
            "Explain the result",
        ],
        default_instructions="Move each card into the category that best matches the source lesson thinking.",
        default_answer_check="Every card should land in the category that best matches the source idea or strategy.",
    )
    sort_slide = with_template_role(
        build_slide_plan(
            kind="practice",
            section="Drag-and-Sort",
            title="Drag-and-Sort Activity",
            subtitle="Sort the lesson ideas or properties into the category that fits best.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, middle_numbers), limit=240),
            tasks=sort_activity["tasks"],
            response_prompt="Sort each card, then explain one choice using the lesson vocabulary.",
            source_slide_numbers=middle_numbers or seed_numbers,
            image_source_slide=pick_first_image_slide(source_slides_from_numbers(deck, middle_numbers or seed_numbers)),
            **named_template_activity_payload(
                name="strategy sort mat",
                family="sort_classify",
                instructions=sort_activity["instructions"],
                movable_pieces=sort_activity["movable_pieces"],
                answer_check=sort_activity["answer_check"],
                library=library,
            ),
        ),
        "drag_sort",
    )

    error_title, error_text, error_items = common_mistake_payload(
        best_source_prompt_text(source_slides_from_numbers(deck, later_numbers or closing_numbers), limit=240),
        context_anchor,
    )
    error_slide = with_template_role(
        build_slide_plan(
            kind="practice",
            section="Error Analysis",
            title="Error Analysis Activity",
            subtitle="Find the mistake, fix it, and explain why the correction works.",
            primary_text=error_text,
            tasks=error_items,
            response_prompt="Name the mistake, rewrite the step correctly, and explain why your fix works.",
            sentence_starters=["The mistake is ___.", "The correct step is ___.", "This works because ___."],
            source_slide_numbers=later_numbers or closing_numbers or seed_numbers,
            **named_template_activity_payload(
                name="error analysis repair cards",
                family="detect_justify",
                instructions="Find the incorrect reasoning move, fix it, and justify the corrected version.",
                movable_pieces=["What is correct?", "What is incorrect?", "Fix it", "Explain it"],
                answer_check="Your correction should match the source numbers, labels, and reasoning.",
                library=library,
            ),
        ),
        "error_analysis",
    )

    compare_slide = with_template_role(
        build_slide_plan(
            kind="practice",
            section="Two-Column Compare",
            title="Two-Column Compare",
            subtitle="Solve two related problems, then compare the strategies or results.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, later_numbers or closing_numbers), limit=220),
            response_prompt="How are the two solutions alike or different? Explain with evidence from the lesson.",
            sentence_starters=["Both problems use ___.", "One difference is ___.", "I know this because ___."],
            source_slide_numbers=later_numbers or closing_numbers or seed_numbers,
            **named_template_activity_payload(
                name="compare-and-rank cards",
                family="compare_rank",
                instructions="Compare the two solution paths, rank the stronger evidence, and explain what makes it stronger.",
                movable_pieces=["Strategy A", "Strategy B", "Same", "Different", "Best evidence"],
                answer_check="Your comparison should name what stays the same, what changes, and which evidence is strongest.",
                library=library,
            ),
        ),
        "two_column_compare",
    )

    choice_slide = with_template_role(
        build_slide_plan(
            kind="practice",
            section="Choice Board",
            title="Choice Board",
            subtitle="Choose a path, solve it, and explain why you picked it.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, closing_numbers or later_numbers), limit=220),
            response_prompt="I chose Path ___ because ___.",
            source_slide_numbers=closing_numbers or later_numbers or seed_numbers,
            image_source_slide=pick_first_image_slide(source_slides_from_numbers(deck, closing_numbers or later_numbers)),
            **named_template_activity_payload(
                name="strategy ranking activity",
                family="compare_rank",
                instructions="Choose the path that feels like the best fit, then defend your choice with evidence.",
                movable_pieces=["Path A", "Path B", "Path C", "Best fit"],
                answer_check="Your chosen path should match the type of practice or challenge you are ready for.",
                library=library,
            ),
        ),
        "choice_board",
    )

    collaborative_slide = with_template_role(
        {
            **build_slide_plan(
            kind="practice",
            section="Collaborative Practice",
            title="Collaborative Practice",
            subtitle="Partner up and solve the source problem together before comparing your thinking.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, later_numbers or closing_numbers), limit=220),
            response_prompt="How is this problem like the guided example, and what did your partner notice first?",
            sentence_starters=["Partner A solved ___.", "Partner B noticed ___.", "Together we found ___."],
            source_slide_numbers=later_numbers or closing_numbers or seed_numbers,
            **named_template_activity_payload(
                name="reasoning ladder build",
                family="build_construct",
                instructions="Build the solve plan together before you write the final answer.",
                movable_pieces=["What do we know?", "What do we need?", "Write the equation", "Check the answer"],
                answer_check="A strong partner plan should lead to a correct setup and a clear explanation.",
                library=library,
            ),
            ),
            "practice_phase": "Try It Together",
            "partner_prompt": "Partner A solves the first problem. Partner B solves the second problem. Then compare.",
        },
        "collaborative_practice",
    )

    independent_slide = with_template_role(
        {
            **build_slide_plan(
            kind="practice",
            section="Independent Practice",
            title="Independent Practice",
            subtitle="Show what you know by solving a small set of lesson-aligned problems.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, closing_numbers or later_numbers), limit=220),
            response_prompt="Solve each problem clearly and check your work.",
            source_slide_numbers=closing_numbers or later_numbers or seed_numbers,
            **named_template_activity_payload(
                name="independent solve path builder",
                family="build_construct",
                instructions="Use the move cards to plan your solve path first, then complete the work on your own.",
                movable_pieces=["Read the question", "Choose a strategy", "Show the math", "Check the answer"],
                answer_check="A strong independent solve path should show the strategy, the calculations, and the final check.",
                library=library,
            ),
            ),
            "practice_phase": "Now Solve Independently",
        },
        "independent_practice",
    )

    turn_and_teach_slide = with_template_role(
        build_slide_plan(
            kind="practice",
            section="Turn + Teach",
            title="Turn + Teach Challenge",
            subtitle="Solve a similar lesson-aligned problem and coach a partner through your thinking.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, closing_numbers or later_numbers), limit=220),
            response_prompt="Teach your partner the strategy you used and record one revision after the discussion.",
            sentence_starters=[
                "My first move was ___.",
                "The key clue was ___.",
                "After the discussion, I revised ___.",
            ],
            source_slide_numbers=closing_numbers or later_numbers or seed_numbers,
            **named_template_activity_payload(
                name="turn-and-teach flow",
                family="sequence_order",
                instructions="Build the solve-and-teach flow in order before you explain the problem to a partner.",
                movable_pieces=["Name the problem", "Choose the strategy", "Show the reasoning", "Teach the check"],
                answer_check="A strong teaching flow should help a partner understand both the steps and the reason behind them.",
                library=library,
            ),
        ),
        "turn_and_teach",
    )

    twr_slide = with_template_role(
        build_slide_plan(
            kind="reflection",
            section="TWR Frame",
            title="TWR Frame",
            subtitle="Think, write, and explain how the lesson strategy works.",
            primary_text="Use the frames to explain the strategy in complete sentences.",
            response_prompt="Complete each frame using lesson words, numbers, and labels.",
            sentence_starters=[
                "The problem is asking me to ___.",
                "First, I ___.",
                "My explanation shows ___.",
                "This connects to the lesson because ___.",
            ],
            source_slide_numbers=closing_numbers or later_numbers or seed_numbers,
        ),
        "twr_frame",
    )

    exit_slide = with_template_role(
        build_slide_plan(
            kind="exit_ticket",
            section="Exit Ticket",
            title="Exit Ticket",
            subtitle="Choose a tier and show what you can do on your own.",
            primary_text=best_source_prompt_text(source_slides_from_numbers(deck, closing_numbers or later_numbers), limit=220),
            response_prompt="Choose the tier that lets you show your strongest understanding.",
            sentence_starters=["I solved it by ___.", "My evidence is ___.", "This shows I can ___."],
            source_slide_numbers=closing_numbers or later_numbers or seed_numbers,
        ),
        "tiered_exit",
    )

    goal_slide = with_template_role(
        build_slide_plan(
            kind="reflection",
            section="Goal Tracker",
            title="Goal Tracker",
            subtitle="Rate your confidence and name one thing you learned today.",
            primary_text=content_objective,
            secondary_text="Write one new idea or connection from today's lesson.",
            source_slide_numbers=objective_source_numbers or closing_numbers or seed_numbers,
        ),
        "goal_tracker",
    )

    session = {
        "session_title": "Session 2 Student Notebook",
        "session_subtitle": "Continue the lesson with collaborative practice, independent work, and reflection.",
        "template_family": "uploaded_model_session2",
        "slides": [
            build_slide_plan(
                kind="cover",
                section="Cover",
                title=lesson_title,
                subtitle="Session 2 Student Notebook",
                primary_text=cover_focus_statement(deck, "session_2", lesson_title),
                secondary_text=default_cover_subtitle(deck, "session_2"),
                bullets=slide_title_list(core or session_sources, limit=4),
                source_slide_numbers=seed_numbers,
                image_source_slide=pick_first_image_slide(core or session_sources),
            ),
            build_slide_plan(
                kind="learning_target",
                section="Learning Target",
                title="Content + Language Objectives",
                subtitle="Read each objective and check your progress before and after the lesson.",
                primary_text=content_objective,
                secondary_text=language_objective,
                source_slide_numbers=objective_source_numbers or seed_numbers,
            ),
            build_slide_plan(
                kind="be_curious",
                section="Be Curious",
                title="Be Curious",
                subtitle="Review from Session 1, notice key details, and get ready to solve.",
                primary_text="What do you notice first when you look back at the lesson model or diagram?",
                secondary_text="What do you wonder or predict before you solve today's problems?",
                sentence_starters=["I notice ___.", "I wonder ___.", "This might help because ___."],
                vocabulary=source_vocabulary(deck, launch_numbers or seed_numbers, limit=3),
                response_prompt="What clue from Session 1 will help you most today?",
                source_slide_numbers=launch_numbers or seed_numbers,
                image_source_slide=pick_first_image_slide(source_slides_from_numbers(deck, launch_numbers or seed_numbers)),
                **named_template_activity_payload(
                    name="cover-and-reveal",
                    family="reveal_discuss",
                    instructions="Use the reveal cards to discuss what you notice before solving.",
                    movable_pieces=[item.get("word", "") for item in source_vocabulary(deck, launch_numbers or seed_numbers, limit=3)[:3]],
                    answer_check="A strong notice or wonder should connect to the source image, figure, or prompt.",
                    library=library,
                ),
            ),
            build_slide_plan(
                kind="vocabulary",
                section="Vocabulary",
                title="Vocabulary",
                subtitle="Keep the key lesson terms visible while you model, solve, and explain.",
                primary_text="Use these lesson words when you talk, write, and justify your reasoning.",
                vocabulary=source_vocabulary(deck, seed_numbers, limit=4),
                source_slide_numbers=seed_numbers,
                **named_template_activity_payload(
                    name="visual vocabulary matching",
                    family="match_pair",
                    instructions="Match each lesson word to the visual, model, or example that helps explain it.",
                    movable_pieces=[item.get("word", "") for item in source_vocabulary(deck, seed_numbers, limit=4)],
                    answer_check="Each vocabulary word should connect to the model, picture, or example that best explains it.",
                    library=library,
                ),
            ),
            with_template_role(
                build_slide_plan(
                    kind="worked_example",
                    section="Guided Practice",
                    title="Guided Practice",
                    subtitle="Model the source problem before students try one on their own.",
                    primary_text=best_source_prompt_text(source_slides_from_numbers(deck, middle_numbers or seed_numbers), limit=320),
                    tasks=slide_task_list(source_slides_from_numbers(deck, middle_numbers or seed_numbers), limit=3),
                    response_prompt="Explain each step and why it fits the source problem.",
                    sentence_starters=["First, I ___.", "Next, I ___.", "So, the answer is ___."],
                    source_slide_numbers=middle_numbers or seed_numbers,
                    image_source_slide=pick_first_image_slide(source_slides_from_numbers(deck, middle_numbers or seed_numbers)),
                    **named_template_activity_payload(
                        name="equation builder",
                        family="build_construct",
                        instructions="Build the solve path before you write the full solution.",
                        movable_pieces=["Known information", "Write the formula", "Substitute values", "Solve + units"],
                        answer_check="A complete solve path should move from given information to a labeled answer.",
                        library=library,
                    ),
                ),
                "guided_practice",
            ),
            sort_slide,
            error_slide,
            compare_slide,
            choice_slide,
            collaborative_slide,
            independent_slide,
            turn_and_teach_slide,
            twr_slide,
            exit_slide,
            goal_slide,
        ],
    }
    return session


def generate_heuristic_plan(
    deck: dict[str, Any],
    output_dir: Path,
    *,
    custom_guidance: str = "",
) -> tuple[dict[str, Any], Path]:
    effective_guidance = enforce_runtime_quality_guidance(custom_guidance)
    slides = deck["slides"]
    midpoint = max(1, len(slides) // 2)
    first_half = slides[:midpoint]
    second_half = slides[midpoint:] or slides[midpoint - 1 :]
    first_core = core_instructional_slides(first_half, limit=8)
    second_core = core_instructional_slides(second_half, limit=8)
    first_launch = [slide for slide in first_half if "be curious" in f"{slide['title']} {slide['text']}".lower()] or first_half[:2]
    second_launch = [slide for slide in second_half if "be curious" in f"{slide['title']} {slide['text']}".lower()] or second_half[:3]
    lesson_title = deck["lesson_title"]
    session1_image = pick_first_image_slide(first_half)
    session2_image = pick_first_image_slide(second_half)
    standards: list[str] = []
    vocabulary = build_vocab_placeholders(deck)
    library = load_activity_library()

    session1_slides = [
        build_slide_plan(
            kind="cover",
            title=lesson_title,
            subtitle="Session 1 Student Notebook",
            primary_text="Preview the topic, activate background knowledge, and get ready to take structured notes.",
            secondary_text="Use the source deck to preserve the original lesson sequence and key language.",
            bullets=slide_title_list(first_core or first_half, limit=3),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core or first_half)[:4]],
            image_source_slide=session1_image,
        ),
        build_slide_plan(
            kind="quick_review",
            title="Warm-Up",
            subtitle="Start with a choice, an estimate, and a quick justification.",
            primary_text="Use the opening warm-ups to make a decision, estimate, and defend your thinking.",
            tasks=[
                "Would you rather have THIS or THAT tokens for prizes? Defend your choice.",
                "Tell whether an estimate is reasonable and explain why.",
            ],
            response_prompt="What kind of thinking do both warm-ups ask you to use?",
            **build_activity_payload(
                kind="quick_review",
                deck=deck,
                library=library,
                source_texts=[slide["title"] for slide in first_half[:3]] + [slide["text"] for slide in first_half[:3]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in first_half[:3]],
        ),
        build_slide_plan(
            kind="be_curious",
            title="Be Curious + Notice / Wonder",
            subtitle="Look closely, use sentence supports, and preview key vocabulary before solving.",
            primary_text="Record what stands out in the source image or opening scenario.",
            secondary_text="Ask one question that would help you understand the lesson more deeply.",
            sentence_starters=["I notice that ...", "I wonder why ...", "This might matter because ..."],
            vocabulary=source_vocabulary(deck, [slide["slide_number"] for slide in first_launch[:3]], limit=3),
            response_prompt="What clues in the opening slide help you predict the lesson focus?",
            **build_activity_payload(
                kind="be_curious",
                deck=deck,
                library=library,
                source_texts=[slide["title"] for slide in first_launch[:2]] + [slide["text"] for slide in first_launch[:2]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in first_launch[:2]],
            image_source_slide=pick_first_image_slide(first_launch) or session1_image,
        ),
        build_slide_plan(
            kind="learning_target",
            title="Learning Target",
            subtitle="Turn the lesson into a student-friendly goal.",
            primary_text=f"I can explain and apply the key ideas from {lesson_title}.",
            bullets=slide_title_list(first_core[1:] or first_core or first_half[1:], limit=4),
            response_prompt="What do you already know that will help you meet this goal?",
            source_slide_numbers=[slide["slide_number"] for slide in (first_core or first_half)[:4]],
        ),
        build_slide_plan(
            kind="vocabulary",
            title="Academic Vocabulary",
            subtitle="Keep important lesson words visible while you work.",
            primary_text="Use these key terms as you annotate, talk, and write during the lesson.",
            vocabulary=vocabulary,
            **build_activity_payload(
                kind="vocabulary",
                deck=deck,
                library=library,
                source_texts=deck["keyword_candidates"],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core or first_half)[:5]],
        ),
        build_slide_plan(
            kind="guided_notes",
            title="Guided Notes",
            subtitle="Capture the big ideas in your own words.",
            primary_text=truncate_text((first_core[0]["text"] if first_core else first_half[1]["text"] if len(first_half) > 1 else deck["summary"]), 320),
            bullets=slide_title_list(first_core or first_half[1:], limit=4),
            response_prompt="Summarize the main concept from today's lesson in one or two sentences.",
            sentence_starters=["The main idea is ...", "One important detail is ...", "This connects to ..."],
            **build_activity_payload(
                kind="guided_notes",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (first_core or first_half[1:5])[:4]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core or first_half[1:5])[:4]],
            image_source_slide=pick_first_image_slide(first_core or first_half[1:]),
        ),
        build_slide_plan(
            kind="worked_example",
            title="Worked Example",
            subtitle="Walk through a model before you try it on your own.",
            primary_text=truncate_text((first_core[1]["text"] if len(first_core) > 1 else first_core[0]["text"] if first_core else first_half[2]["text"] if len(first_half) > 2 else deck["summary"]), 360),
            bullets=slide_task_list(first_core[1:] or first_core or first_half[2:], limit=4),
            response_prompt="Explain the steps in the example using complete sentences.",
            sentence_starters=["First, I ...", "Next, I ...", "This shows that ...", "So, the answer is ..."],
            **build_activity_payload(
                kind="worked_example",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (first_core[1:] or first_core or first_half[2:6])[:4]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core[1:] or first_core or first_half[2:6])[:4]],
        ),
        build_slide_plan(
            kind="practice",
            title="Practice",
            subtitle="Use the source lesson tasks to build confidence.",
            primary_text="Work through the practice problems and show your thinking.",
            tasks=slide_task_list(first_core[2:] or first_core or first_half[3:], limit=3),
            response_prompt="Which strategy from the model will you reuse first?",
            **build_activity_payload(
                kind="practice",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (first_core[2:] or first_core or first_half[3:7])[:4]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core[2:] or first_core or first_half[3:7])[:4]],
        ),
        build_slide_plan(
            kind="reflection",
            title="Reflection",
            subtitle="Pause and make sense of the lesson so far.",
            primary_text="What part of the lesson feels clear right now?",
            secondary_text="What do you still need to review before Session 2?",
            sentence_starters=["I understand that ...", "I still need help with ...", "A useful strategy was ..."],
            response_prompt="Write a short reflection about your current understanding.",
            **build_activity_payload(
                kind="reflection",
                deck=deck,
                library=library,
                source_texts=deck["keyword_candidates"] + [deck["summary"]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core[-3:] or first_half[-3:])],
        ),
        build_slide_plan(
            kind="exit_ticket",
            title="Exit Ticket",
            subtitle="Show what you can do with the main idea from Session 1.",
            primary_text="Explain how decomposing a regular polygon into congruent triangles can help you find area.",
            tasks=slide_task_list(first_core[-4:] or first_half[-4:], limit=2),
            response_prompt="What evidence would convince someone that your area strategy works?",
            **build_activity_payload(
                kind="exit_ticket",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (first_core[-4:] or first_half[-4:])],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (first_core[-4:] or first_half[-4:])],
        ),
    ]

    session2_slides = [
        build_slide_plan(
            kind="cover",
            title=lesson_title,
            subtitle="Session 2 Student Notebook",
            primary_text="Review the lesson, practice more deeply, and apply the ideas in a new way.",
            secondary_text="Use evidence from the source deck to justify your thinking.",
            bullets=slide_title_list(second_core or second_half, limit=3),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core or second_half)[:4]],
            image_source_slide=session2_image,
        ),
        build_slide_plan(
            kind="be_curious",
            title="Be Curious + Notice / Wonder",
            subtitle="Compare the figures, use sentence supports, and preview lesson words before solving.",
            primary_text="What do you notice first when you compare the figures or situations?",
            tasks=slide_task_list(second_launch[:3], limit=3),
            secondary_text="What do you wonder, predict, or want to compare before you solve?",
            vocabulary=source_vocabulary(deck, [slide["slide_number"] for slide in second_launch[:3]], limit=3),
            response_prompt="Which clue from the source comparison feels most important before you solve?",
            **build_activity_payload(
                kind="be_curious",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in second_launch[:3]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in second_launch[:3]],
            image_source_slide=pick_first_image_slide(second_launch) or session2_image,
        ),
        build_slide_plan(
            kind="learning_target",
            title="Learning Target",
            subtitle="Turn the next part of the lesson into clear goals.",
            primary_text=f"I can explain and apply the key ideas from {lesson_title}.",
            bullets=slide_title_list(second_core or second_half[:4], limit=4),
            response_prompt="What skill from Session 1 will help you most today?",
            source_slide_numbers=[slide["slide_number"] for slide in (second_core or second_half[:4])[:4]],
        ),
        build_slide_plan(
            kind="vocabulary",
            title="Academic Vocabulary",
            subtitle="Keep important words visible while you solve composite figures.",
            primary_text="Use these lesson words when you talk, write, and justify your reasoning.",
            vocabulary=vocabulary,
            **build_activity_payload(
                kind="vocabulary",
                deck=deck,
                library=library,
                source_texts=deck["keyword_candidates"],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core or second_half[:5])[:5]],
        ),
        build_slide_plan(
            kind="guided_notes",
            title="Concept Review",
            subtitle="Reconnect the new lesson segment to the earlier learning.",
            primary_text=truncate_text((second_core[0]["text"] if second_core else second_half[0]["text"] if second_half else deck["summary"]), 320),
            bullets=slide_title_list(second_core or second_half, limit=4),
            response_prompt="How does this new part of the lesson connect to what you learned before?",
            sentence_starters=["This is similar to ...", "This is different because ...", "The connection is ..."],
            **build_activity_payload(
                kind="guided_notes",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (second_core or second_half[:4])[:4]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core or second_half[:4])[:4]],
            image_source_slide=pick_first_image_slide(second_core or second_half),
        ),
        build_slide_plan(
            kind="worked_example",
            title="Apply the Idea",
            subtitle="Model the next step before independent work.",
            primary_text=truncate_text((second_core[1]["text"] if len(second_core) > 1 else second_core[0]["text"] if second_core else second_half[1]["text"] if len(second_half) > 1 else deck["summary"]), 360),
            bullets=slide_task_list(second_core[1:] or second_core or second_half[1:], limit=4),
            response_prompt="What step should come next, and why?",
            sentence_starters=["First, I ...", "Then, I ...", "I know this because ..."],
            **build_activity_payload(
                kind="worked_example",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (second_core[1:] or second_core or second_half[1:5])[:4]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core[1:] or second_core or second_half[1:5])[:4]],
        ),
        build_slide_plan(
            kind="practice",
            title="Deeper Practice",
            subtitle="Use the later source slides for independent work.",
            primary_text="Show your work and use the lesson language in your explanations.",
            tasks=slide_task_list(second_core[2:] or second_core or second_half[2:], limit=3),
            response_prompt="Circle one task that feels like a strong check for understanding.",
            **build_activity_payload(
                kind="practice",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (second_core[2:] or second_core or second_half[2:7])[:5]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core[2:] or second_core or second_half[2:7])[:5]],
        ),
        build_slide_plan(
            kind="challenge",
            title="Challenge",
            subtitle="Stretch the lesson into a new application or comparison.",
            primary_text="Use a later source example, scenario, or image to explain how the concept transfers.",
            tasks=slide_task_list(second_core[-3:] or second_half[-3:], limit=2),
            response_prompt="What makes this challenge more demanding than the earlier practice?",
            **build_activity_payload(
                kind="challenge",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (second_core[-3:] or second_half[-3:])],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core[-3:] or second_half[-3:])],
            image_source_slide=session2_image,
        ),
        build_slide_plan(
            kind="reflection",
            title="Reflection",
            subtitle="Summarize how your understanding changed during Session 2.",
            primary_text="What strategy, idea, or pattern helped you most today?",
            secondary_text="Where could you use this idea again?",
            sentence_starters=["At first, I ...", "Now I can ...", "I could use this when ..."],
            response_prompt="Write a complete-sentence reflection about what you learned.",
            **build_activity_payload(
                kind="reflection",
                deck=deck,
                library=library,
                source_texts=deck["keyword_candidates"] + [deck["summary"]],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core[-4:] or second_half[-4:])],
        ),
        build_slide_plan(
            kind="exit_ticket",
            title="Exit Ticket",
            subtitle="Use one final prompt to show what you can do on your own.",
            primary_text="Create a short response, solve a final item, or explain the central lesson idea.",
            tasks=slide_task_list(second_core[-2:] or second_half[-2:], limit=2),
            response_prompt="What is the most important thing you want your teacher to see in your answer?",
            **build_activity_payload(
                kind="exit_ticket",
                deck=deck,
                library=library,
                source_texts=[slide["text"] for slide in (second_core[-3:] or second_half[-3:])],
            ),
            source_slide_numbers=[slide["slide_number"] for slide in (second_core[-3:] or second_half[-3:])],
        ),
    ]

    plan = {
        "lesson_title": lesson_title,
        "subject": "Source Lesson",
        "grade_level": "Unspecified",
        "standards": standards,
        "topic_summary": deck["summary"],
        "session_1": {
            "session_title": "Session 1 Student Notebook",
            "session_subtitle": "A six-slide launch notebook before students continue in the book.",
            "slides": [],
        },
    }
    plan = enforce_plan_requirements(plan, deck, custom_guidance=effective_guidance)
    plan_path = output_dir / "notebook_plan.json"
    write_json(plan_path, plan)
    return plan, plan_path


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_background(slide: Any, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_label_text_color(fill_color: RGBColor) -> RGBColor:
    if fill_color in {GOLD, CORAL, PALE_GOLD, PALE_CORAL, PAPER_WARM}:
        return NAVY
    return PAPER


def add_soft_circle(slide: Any, x: int, y: int, size: int, fill_color: RGBColor) -> Any:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    return circle


def add_cover_background_decor(slide: Any) -> None:
    edge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    edge.fill.solid()
    edge.fill.fore_color.rgb = NAVY
    edge.line.fill.background()
    add_soft_circle(slide, Inches(8.75), Inches(-0.55), Inches(4.7), PALE_GOLD)
    add_soft_circle(slide, Inches(10.75), Inches(4.65), Inches(2.75), RGBColor(252, 207, 110))


def add_rect(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    fill_color: RGBColor,
    *,
    line_color: RGBColor = SOFT_LINE,
    rounded: bool = True,
    line_width: float = 0.95,
) -> Any:
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if rounded
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def emu_to_inches(value: int) -> float:
    return float(value) / 914400.0


def estimate_text_lines(text: str, width_in: float, font_size: float) -> int:
    cleaned_lines = (text or "").split("\n")
    if not cleaned_lines:
        return 1
    chars_per_line = max(10, int(width_in * max(6.5, 18.0 - min(font_size, 22.0) * 0.46)))
    total_lines = 0
    for raw_line in cleaned_lines:
        cleaned = normalize_whitespace(raw_line)
        if not cleaned:
            total_lines += 1
            continue
        line_capacity = chars_per_line - 2 if cleaned.startswith("-") else chars_per_line
        total_lines += max(1, (len(cleaned) + max(line_capacity, 1) - 1) // max(line_capacity, 1))
    return total_lines


def fit_text_style(
    text: str,
    width: int,
    height: int,
    *,
    base_size: float,
    min_size: float | None = None,
    line_spacing: float = DEFAULT_TEXT_LINE_SPACING,
) -> tuple[float, float]:
    width_in = max(emu_to_inches(max(int(width), 1)), 0.7)
    height_in = max(emu_to_inches(max(int(height), 1)), 0.22)
    current_size = float(base_size)
    minimum = float(min_size if min_size is not None else max(10.8, base_size - 3.1))
    if current_size >= AUTO_FIT_MIN_FONT_PT:
        minimum = max(minimum, AUTO_FIT_MIN_FONT_PT)
    minimum = min(minimum, current_size)
    current_spacing = max(1.12, float(line_spacing))

    for _ in range(18):
        estimated_lines = estimate_text_lines(text, width_in, current_size)
        line_height_in = max((current_size / 72.0) * current_spacing * 1.34, 0.16)
        capacity = max(1, int(height_in / line_height_in))
        if estimated_lines <= capacity or current_size <= minimum:
            break
        current_size = max(minimum, current_size - 0.45)
        if current_size <= base_size - 1.0:
            current_spacing = min(1.22, current_spacing + 0.02)

    if estimate_text_lines(text, width_in, current_size) <= 2 and height_in >= 0.85:
        current_size = min(base_size + 0.45, current_size + 0.3)

    return round(current_size, 2), round(current_spacing, 2)


def set_textbox_text(
    box: Any,
    text: str,
    *,
    size: float = 17.0,
    color: RGBColor = INK,
    bold: bool = False,
    font: str = FONT_BODY,
    align: Any = PP_ALIGN.LEFT,
    margin: float = 0.1,
    line_spacing: float = DEFAULT_TEXT_LINE_SPACING,
    min_size: float | None = None,
    auto_fit: bool = True,
    vertical_anchor: Any = MSO_VERTICAL_ANCHOR.TOP,
) -> Any:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    resolved_size = size
    resolved_spacing = line_spacing
    if auto_fit:
        resolved_size, resolved_spacing = fit_text_style(
            text,
            int(getattr(box, "width", 1)),
            int(getattr(box, "height", 1)),
            base_size=size,
            min_size=min_size,
            line_spacing=line_spacing,
        )
    lines = (text or "").split("\n")
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(resolved_size)
        paragraph.font.name = font
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.alignment = align
        paragraph.line_spacing = resolved_spacing
        paragraph.space_after = Pt(2.6 if index < len(lines) - 1 else 0.8)
    return box


def set_shape_text(shape: Any, text: str, **kwargs: Any) -> Any:
    return set_textbox_text(shape, text, **kwargs)


def set_shape_heading_body(
    shape: Any,
    heading: str,
    body: str = "",
    *,
    heading_size: float = 12.0,
    body_size: float = 11.0,
    heading_color: RGBColor = NAVY,
    body_color: RGBColor = INK,
    heading_font: str = FONT_HEAD,
    body_font: str = FONT_BODY,
    margin: float = 0.08,
    body_min_size: float | None = None,
    line_spacing: float = DEFAULT_TEXT_LINE_SPACING,
    align: Any = PP_ALIGN.LEFT,
    vertical_anchor: Any = MSO_VERTICAL_ANCHOR.TOP,
) -> Any:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)

    heading_text = normalize_whitespace(heading)
    body_text = normalize_whitespace(body)
    resolved_body_size = body_size
    resolved_spacing = line_spacing
    if body_text:
        usable_width = max(int(getattr(shape, "width", 1) - Inches(margin * 2)), 1)
        heading_line_height_in = max((heading_size / 72.0) * 1.1 * 1.28, 0.18)
        reserved_height = heading_line_height_in + 0.16 + margin * 2
        usable_height = max(int(getattr(shape, "height", 1) - Inches(reserved_height)), 1)
        resolved_body_size, resolved_spacing = fit_text_style(
            body_text,
            usable_width,
            usable_height,
            base_size=body_size,
            min_size=body_min_size,
            line_spacing=line_spacing,
        )

    heading_paragraph = tf.paragraphs[0]
    heading_paragraph.text = heading_text
    heading_paragraph.font.size = Pt(heading_size)
    heading_paragraph.font.name = heading_font
    heading_paragraph.font.color.rgb = heading_color
    heading_paragraph.font.bold = True
    heading_paragraph.alignment = align
    heading_paragraph.line_spacing = 1.1
    heading_paragraph.space_after = Pt(1.6 if body_text else 0.8)

    if body_text:
        body_paragraph = tf.add_paragraph()
        body_paragraph.text = body_text
        body_paragraph.font.size = Pt(resolved_body_size)
        body_paragraph.font.name = body_font
        body_paragraph.font.color.rgb = body_color
        body_paragraph.font.bold = False
        body_paragraph.alignment = align
        body_paragraph.line_spacing = max(1.1, resolved_spacing)
        body_paragraph.space_after = Pt(0.8)
    return shape


def add_text(slide: Any, x: int, y: int, w: int, h: int, text: str, **kwargs: Any) -> Any:
    safe_w = max(int(w), 1)
    safe_h = max(int(h), 1)
    box = slide.shapes.add_textbox(x, y, safe_w, safe_h)
    return set_textbox_text(box, text, **kwargs)


def add_card(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    body: str = "",
    *,
    fill: RGBColor = PAPER,
    accent: RGBColor = NAVY,
    title_size: float = 18.0,
    body_size: float = 15.2,
    body_color: RGBColor = INK,
) -> None:
    line_color = accent if fill == PAPER else SOFT_LINE
    card = add_rect(slide, x, y, w, h, fill, line_color=line_color, line_width=1.05)
    card.adjustments[0] = 0.18
    compact = float(h) <= float(Inches(0.72))
    label_margin_x = Inches(0.18 if compact else 0.22)
    label_y = y + Inches(0.14 if compact else 0.18)
    label_h = Inches(0.22 if compact else 0.28)
    label_width = min(
        int(w - label_margin_x * 2),
        max(int(Inches(1.34)), int(float(Inches(0.78 + max(len(normalize_whitespace(title)), 10) * 0.05)))),
    )
    label = add_rect(
        slide,
        x + label_margin_x,
        label_y,
        label_width,
        label_h,
        accent,
        line_color=accent,
        line_width=0.7,
    )
    label.adjustments[0] = 0.32
    add_text(
        slide,
        x + label_margin_x,
        label_y + Inches(0.01),
        label_width,
        label_h - Inches(0.02),
        truncate_text(title, 40 if compact else 56),
        size=10.8 if compact else min(title_size, 12.4),
        min_size=10.6 if compact else 10.8,
        color=accent_label_text_color(accent),
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )
    if body:
        body_top = y + (Inches(0.42) if compact else Inches(0.62))
        body_height = max(int(h - (Inches(0.50) if compact else Inches(0.82))), 1)
        add_text(
            slide,
            x + label_margin_x,
            body_top,
            w - label_margin_x * 2,
            body_height,
            body,
            size=min(body_size, 10.8) if compact else body_size,
            min_size=10.4 if compact else 11.2,
            color=body_color,
            font=FONT_BODY,
            margin=0.02 if compact else 0.1,
        )


def add_tag(slide: Any, x: int, y: int, text: str, fill_color: RGBColor) -> None:
    width = max(float(Inches(1.55)), min(float(Inches(2.85)), float(Inches(0.94 + max(len(normalize_whitespace(text)), 8) * 0.07))))
    tag = add_rect(slide, x, y, int(width), Inches(0.34), fill_color, line_color=fill_color)
    tag.adjustments[0] = 0.35
    set_shape_text(
        tag,
        text,
        size=10.8,
        color=accent_label_text_color(fill_color),
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def soft_panel_fill_for_accent(accent: RGBColor) -> RGBColor:
    if accent == TEAL:
        return PALE_BLUE
    if accent == GOLD:
        return PALE_GOLD
    if accent == CORAL:
        return PALE_CORAL
    if accent == SAGE:
        return PALE_SAGE
    return PALE_NAVY


def add_footer_bar(slide: Any, footer_text: str) -> None:
    if not normalize_whitespace(footer_text):
        return
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.34), Inches(7.01), Inches(12.65), Inches(0.04))
    footer.fill.solid()
    footer.fill.fore_color.rgb = TEAL
    footer.line.fill.background()
    add_text(
        slide,
        Inches(0.48),
        Inches(6.84),
        Inches(12.0),
        Inches(0.18),
        truncate_text(footer_text, 170),
        size=10.4,
        color=SOFT_NAVY,
        font=FONT_BODY,
        margin=0.01,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_header(
    slide: Any,
    *,
    section: str,
    title: str,
    subtitle: str,
    page: int,
    accent: RGBColor,
    footer_text: str,
) -> None:
    set_background(slide)
    add_tag(slide, Inches(0.48), Inches(0.38), truncate_text(section, 28), GOLD)
    title_width = min(int(Inches(7.0)), max(int(Inches(3.6)), int(float(Inches(1.88 + max(len(normalize_whitespace(title)), 18) * 0.08)))))
    title_banner = add_rect(
        slide,
        Inches(0.48),
        Inches(0.84),
        title_width,
        Inches(0.52),
        NAVY,
        line_color=NAVY,
        line_width=0.7,
    )
    title_banner.adjustments[0] = 0.24
    add_text(
        slide,
        Inches(0.48),
        Inches(0.92),
        title_width,
        Inches(0.30),
        truncate_text(title, 110),
        size=18.2,
        color=PAPER,
        bold=True,
        font=FONT_HEAD,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    subtitle_fill = soft_panel_fill_for_accent(accent)
    subtitle_plate = add_rect(
        slide,
        Inches(0.48),
        Inches(1.39),
        Inches(10.78),
        Inches(0.30),
        subtitle_fill,
        line_color=subtitle_fill,
        line_width=0.6,
    )
    subtitle_plate.adjustments[0] = 0.18
    add_text(
        slide,
        Inches(0.62),
        Inches(1.45),
        Inches(10.36),
        Inches(0.20),
        truncate_text(subtitle, 180),
        size=12.0,
        color=accent,
        font=FONT_BODY,
        margin=0.01,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    page_pill = add_rect(
        slide,
        Inches(11.92),
        Inches(0.58),
        Inches(0.62),
        Inches(0.62),
        GOLD,
        line_color=GOLD,
        line_width=0.7,
    )
    page_pill.adjustments[0] = 0.3
    add_text(
        slide,
        Inches(11.92),
        Inches(0.73),
        Inches(0.62),
        Inches(0.24),
        str(page),
        size=12.2,
        color=PAPER,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.02,
    )
    accent_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.48), Inches(1.78), Inches(12.06), Inches(0.03))
    accent_rule.fill.solid()
    accent_rule.fill.fore_color.rgb = accent
    accent_rule.line.fill.background()
    add_footer_bar(slide, footer_text)


def add_name_bar(slide: Any) -> None:
    top_rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.45),
        Inches(6.40),
        Inches(12.32),
        Inches(0.06),
    )
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = NAVY
    top_rule.line.fill.background()
    panel = add_rect(
        slide,
        Inches(0.45),
        Inches(6.48),
        Inches(12.32),
        Inches(0.82),
        PAPER,
        line_color=LINE,
        line_width=0.9,
    )
    panel.adjustments[0] = 0.16
    add_text(
        slide,
        Inches(0.66),
        Inches(6.58),
        Inches(2.0),
        Inches(0.16),
        "Student Record",
        size=10.8,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    field_specs = [
        ("Name", Inches(0.66), Inches(2.95)),
        ("Date", Inches(4.02), Inches(2.10)),
        ("Class", Inches(6.52), Inches(1.70)),
    ]
    for label, field_x, field_w in field_specs:
        add_text(
            slide,
            field_x,
            Inches(6.78),
            field_w,
            Inches(0.14),
            label,
            size=10.4,
            color=SOFT_NAVY,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )
        slot = add_rect(
            slide,
            field_x,
            Inches(6.98),
            field_w,
            Inches(0.18),
            PAPER_WARM,
            line_color=SOFT_LINE,
        )
        slot.adjustments[0] = 0.12


def add_picture_contain(slide: Any, asset: dict[str, Any], x: int, y: int, w: int, h: int) -> None:
    image_path = Path(asset["path"])
    if not image_path.exists():
        return
    pixel_width = asset.get("pixel_width") or asset.get("display_width") or 0
    pixel_height = asset.get("pixel_height") or asset.get("display_height") or 0
    if not pixel_width or not pixel_height:
        slide.shapes.add_picture(BytesIO(image_path.read_bytes()), x, y, width=w, height=h)
        return
    box_ratio = float(w) / float(h)
    image_ratio = pixel_width / pixel_height
    if image_ratio >= box_ratio:
        width = w
        height = int(float(w) / image_ratio)
        left = x
        top = int(float(y) + (float(h) - height) / 2)
    else:
        height = h
        width = int(float(h) * image_ratio)
        top = y
        left = int(float(x) + (float(w) - width) / 2)
    slide.shapes.add_picture(BytesIO(image_path.read_bytes()), left, top, width=width, height=height)


def add_image_panel(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    asset: dict[str, Any] | None,
    *,
    fill: RGBColor = PALE_NAVY,
    label: str = "Source image",
    picture_padding: int | None = None,
) -> None:
    frame = add_rect(slide, x, y, w, h, fill)
    frame.adjustments[0] = 0.1
    padding = picture_padding if picture_padding is not None else Inches(0.14)
    if asset:
        add_picture_contain(slide, asset, x + padding, y + padding, w - padding * 2, h - padding * 2)
        return
    add_text(
        slide,
        x + Inches(0.26),
        y + Inches(0.26),
        w - Inches(0.52),
        h - Inches(0.52),
        label,
        size=16,
        color=MUTED,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )


def add_response_lines(slide: Any, x: int, y: int, w: int, count: int, spacing: float = 0.31) -> None:
    for index in range(count):
        y_pos = y + Inches(index * spacing)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y_pos, x + w, y_pos)
        line.line.color.rgb = LINE
        line.line.width = Pt(1)


def response_box_accent(fill: RGBColor) -> RGBColor:
    if fill == PALE_BLUE:
        return TEAL
    if fill == PALE_GOLD:
        return GOLD
    if fill == PALE_CORAL:
        return CORAL
    if fill == PALE_SAGE:
        return SAGE
    return NAVY


def add_lined_area(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    prompt: str,
    *,
    lines: int = 4,
    fill: RGBColor = PAPER,
) -> None:
    prompt_text = trim_dangling_display_text(truncate_display_copy(prompt or "", 160))
    title_text = truncate_text(title, 60)
    accent = response_box_accent(fill)
    total_h_in = emu_to_inches(int(h))
    frame = add_rect(slide, x, y, w, h, fill, line_color=accent if fill == PAPER else SOFT_LINE, line_width=1.0)
    frame.adjustments[0] = 0.16
    inner_x = x + Inches(0.16)
    inner_w = max(int(w - Inches(0.32)), 1)
    inner_bottom = y + h - Inches(0.16)

    if prompt_text and total_h_in <= 0.9:
        prompt_box = add_rect(
            slide,
            inner_x,
            y + Inches(0.16),
            inner_w,
            max(int(h - Inches(0.32)), 1),
            PAPER,
            line_color=SOFT_LINE,
        )
        prompt_box.adjustments[0] = 0.18
        set_shape_heading_body(
            prompt_box,
            title_text,
            prompt_text,
            heading_size=11.0,
            body_size=10.6,
            heading_color=accent,
            body_color=INK,
            margin=0.05,
            body_min_size=10.4,
            line_spacing=1.14,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )
        return

    response_top = y + Inches(0.44)
    if prompt_text:
        prompt_lines = max(
            1,
            estimate_text_lines(
                prompt_text,
                max(emu_to_inches(inner_w - Inches(0.18)), 0.8),
                11.8,
            ),
        )
        prompt_h_in = min(
            max(0.58, 0.24 * prompt_lines + 0.24),
            max(0.72, total_h_in * 0.46),
        )
        prompt_box = add_rect(
            slide,
            inner_x,
            y + Inches(0.16),
            inner_w,
            int(Inches(prompt_h_in)),
            PAPER,
            line_color=SOFT_LINE,
        )
        prompt_box.adjustments[0] = 0.18
        set_shape_heading_body(
            prompt_box,
            title_text,
            prompt_text,
            heading_size=12.8 if total_h_in >= 1.4 else 12.0,
            body_size=12.0 if total_h_in >= 1.4 else 11.2,
            heading_color=accent,
            body_color=INK,
            margin=0.06,
            body_min_size=10.4,
            line_spacing=1.14,
        )
        response_top = prompt_box.top + prompt_box.height + Inches(0.12)
    else:
        chip_w = min(
            int(Inches(2.15)),
            max(int(Inches(1.36)), int(float(Inches(0.78 + max(len(title_text), 12) * 0.052)))),
        )
        add_chip(
            slide,
            inner_x,
            y + Inches(0.16),
            chip_w,
            Inches(0.18),
            title_text,
            fill=PALE_NAVY if fill == PAPER else PAPER,
            accent=accent,
        )

    response_h = max(int(inner_bottom - response_top), int(Inches(0.30 if lines <= 1 else 0.44)))
    response_box = add_rect(
        slide,
        inner_x,
        response_top,
        inner_w,
        response_h,
        PAPER,
        line_color=SOFT_LINE,
    )
    response_box.adjustments[0] = 0.18


def add_small_math_visual(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    profile: dict[str, Any],
    accent: RGBColor = TEAL,
) -> None:
    add_rect(slide, x, y, w, h, PAPER, line_color=accent)
    add_text(
        slide,
        x + Inches(0.10),
        y + Inches(0.08),
        w - Inches(0.20),
        Inches(0.16),
        truncate_text(profile.get("shape_label", "Math Model"), 24),
        size=10.6,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
        align=PP_ALIGN.CENTER,
    )
    inset_x = x + Inches(0.24)
    inset_y = y + Inches(0.32)
    inset_w = w - Inches(0.48)
    inset_h = h - Inches(0.56)
    if profile.get("topic") == "volume_prism":
        front = add_rect(slide, inset_x + Inches(0.20), inset_y + Inches(0.18), inset_w - Inches(0.48), inset_h - Inches(0.36), PALE_BLUE, line_color=accent)
        back = add_rect(slide, inset_x, inset_y, inset_w - Inches(0.48), inset_h - Inches(0.36), PAPER, line_color=accent)
        for start_x, start_y, end_x, end_y in (
            (inset_x, inset_y, inset_x + Inches(0.20), inset_y + Inches(0.18)),
            (inset_x + inset_w - Inches(0.48), inset_y, inset_x + inset_w - Inches(0.28), inset_y + Inches(0.18)),
            (inset_x, inset_y + inset_h - Inches(0.36), inset_x + Inches(0.20), inset_y + inset_h - Inches(0.18)),
            (inset_x + inset_w - Inches(0.48), inset_y + inset_h - Inches(0.36), inset_x + inset_w - Inches(0.28), inset_y + inset_h - Inches(0.18)),
        ):
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
            line.line.color.rgb = accent
            line.line.width = Pt(1.2)
        add_text(slide, x + Inches(0.12), y + h - Inches(0.26), Inches(0.44), Inches(0.12), "l", size=10.4, color=accent, bold=True, font=FONT_HEAD, margin=0.01)
        add_text(slide, x + w - Inches(0.46), y + Inches(0.30), Inches(0.28), Inches(0.12), "h", size=10.4, color=accent, bold=True, font=FONT_HEAD, margin=0.01)
        add_text(slide, x + Inches(0.16), y + Inches(0.20), Inches(0.28), Inches(0.12), "w", size=10.4, color=accent, bold=True, font=FONT_HEAD, margin=0.01)
        _ = front
        _ = back
        return
    if profile.get("topic") == "triangle_area":
        mid_x = int(inset_x + inset_w / 2)
        line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inset_x, inset_y + inset_h, inset_x + inset_w, inset_y + inset_h)
        line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inset_x, inset_y + inset_h, mid_x, inset_y)
        line3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mid_x, inset_y, inset_x + inset_w, inset_y + inset_h)
        for line in (line1, line2, line3):
            line.line.color.rgb = accent
            line.line.width = Pt(1.2)
    else:
        add_rect(slide, inset_x, inset_y, inset_w, inset_h, PALE_BLUE, line_color=accent)
    labels = profile.get("vars", [])[:3]
    for index, label in enumerate(labels):
        add_text(
            slide,
            x + Inches(0.10) + index * Inches(0.42),
            y + h - Inches(0.24),
            Inches(0.38),
            Inches(0.12),
            truncate_text(label[:1] if len(label) > 2 else label, 4),
            size=10.4,
            color=accent,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )


def workbook_source_blob(plan_slide: dict[str, Any]) -> str:
    parts = [
        plan_slide.get("title", ""),
        plan_slide.get("subtitle", ""),
        plan_slide.get("primary_text", ""),
        plan_slide.get("secondary_text", ""),
        plan_slide.get("response_prompt", ""),
        " ".join(plan_slide.get("tasks", [])),
        " ".join(plan_slide.get("bullets", [])),
        " ".join(plan_slide.get("movable_pieces", [])),
    ]
    return normalize_whitespace(" ".join(parts)).lower()


def has_box_plot_context(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    return any(
        term in lowered
        for term in ("box plot", "box-and-whisker", "box and whisker", "whisker", "quartile", "quartiles")
    )


def has_data_analysis_context(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    return any(
        term in lowered
        for term in (
            "data analysis",
            "data-analysis",
            "median",
            "mean",
            "data set",
            "dot plot",
            "line plot",
            "ordered data",
            "outlier",
            "distribution",
            "finish times",
            "slowest and fastest",
            "lower extreme",
            "upper extreme",
            "box plot",
            "box-and-whisker",
            "quartile",
            "quartiles",
        )
    )


def has_shipping_box_context(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    if has_box_plot_context(lowered):
        return False
    return any(
        term in lowered
        for term in ("shipping", "shipping rule", "shipping box", "length of the box", "box length", "regulations")
    )


def problem_workbook_mode(blob: str) -> str:
    if any(term in blob for term in ("which size", "how many boxes", "would you buy", "popcorn")):
        return "decision"
    if has_shipping_box_context(blob):
        return "shipping"
    if any(term in blob for term in ("how can you determine the length", "unit cube", "unit cubes")):
        return "unit_model"
    if any(term in blob for term in ("how do you use the formula", "volume of the prism", "what is the volume", "rectangular prism")):
        return "volume_formula"
    return "general"


def problem_workbook_topic_phrase(mode: str) -> str:
    if mode == "shipping":
        return "shipping-box problem"
    if mode == "decision":
        return "popcorn decision problem"
    return "lesson problem"


def problem_fact_statements(plan_slide: dict[str, Any], *, limit: int = 3) -> list[str]:
    facts: list[str] = []
    for text in source_problem_cards(plan_slide, limit=6, max_len=200):
        lowered = text.lower()
        if lowered.startswith(("the width is", "the height is", "the length is")):
            facts.append(text)
        elif "12-inch shipping rule" in lowered or "shipping rule" in lowered:
            facts.append(text)
        elif "popcorn box information" in lowered:
            facts.append(text)
    return unique_nonempty(facts, limit=limit)


def inferred_problem_focus_phrase(plan_slide: dict[str, Any]) -> str:
    text = " ".join(
        [
            plan_slide.get("title", ""),
            plan_slide.get("subtitle", ""),
            plan_slide.get("primary_text", ""),
            " ".join(plan_slide.get("tasks", [])),
            " ".join(plan_slide.get("bullets", [])),
        ]
    ).lower()
    if any(
        term in text
        for term in ("median", "mean", "data set", "dot plot", "line plot", "ordered data", "outlier", "box plot", "quartile", "quartiles")
    ):
        return "data analysis problem"
    if has_shipping_box_context(text):
        return "shipping-box problem"
    if any(term in text for term in ("choice", "which size", "how many boxes", "would you buy", "best option", "popcorn")):
        return "decision problem"
    if any(term in text for term in ("volume", "cubic")):
        return "volume problem"
    if "area" in text:
        return "area problem"
    if any(term in text for term in ("graph", "coordinate", "plot")):
        return "graph problem"
    if "table" in text:
        return "table problem"
    if any(term in text for term in ("equation", "expression", "variable")):
        return "equation problem"
    if re.search(r"\b(ratio|rate|rates|proportion|proportional)\b", text):
        return "ratio problem"
    return "lesson problem"


def display_problem_focus_phrase(plan_slide: dict[str, Any]) -> str:
    focus_phrase = inferred_problem_focus_phrase(plan_slide)
    return "problem" if focus_phrase == "lesson problem" else focus_phrase


def source_problem_focus_phrase(plan_slide: dict[str, Any]) -> str:
    focus_phrase = display_problem_focus_phrase(plan_slide)
    return "source problem" if focus_phrase == "problem" else f"source {focus_phrase}"


def similar_problem_focus_phrase(plan_slide: dict[str, Any]) -> str:
    focus_phrase = display_problem_focus_phrase(plan_slide)
    return "similar problem" if focus_phrase == "problem" else f"similar {focus_phrase}"


SUBSCRIPT_DIGIT_TRANSLATION = str.maketrans("₀₁₂₃₄₅₆₇₈₉", "0123456789")


def source_measure_blob(plan_slide: dict[str, Any]) -> str:
    parts = (
        list(plan_slide.get("source_problem_cards", []))
        + list(plan_slide.get("tasks", []))
        + list(plan_slide.get("bullets", []))
        + [plan_slide.get("primary_text", ""), plan_slide.get("secondary_text", ""), plan_slide.get("response_prompt", "")]
    )
    return normalize_whitespace(" ".join(part for part in parts if part)).translate(SUBSCRIPT_DIGIT_TRANSLATION)


def source_measurements(plan_slide: dict[str, Any]) -> dict[str, tuple[str, str]]:
    blob = source_measure_blob(plan_slide)
    measurements: dict[str, tuple[str, str]] = {}

    def remember(label: str, value: str, unit: str = "") -> None:
        normalized_label = label.lower().replace(" ", "")
        if normalized_label not in measurements and value:
            measurements[normalized_label] = (value, unit.strip())

    for match in re.finditer(
        r"(?:there are|is)\s+(\d+(?:\.\d+)?)\s+unit cubes?\s+along the\s+(width|height|length)",
        blob,
        flags=re.IGNORECASE,
    ):
        remember(match.group(2), match.group(1), "units")

    for match in re.finditer(
        r"\b(d1|d2|b1|b2|length|width|height|base|b|h)\s*=\s*(\d+(?:\.\d+)?)\s*([A-Za-z]+(?:\^?2|²|3)?)?",
        blob,
        flags=re.IGNORECASE,
    ):
        remember(match.group(1), match.group(2), match.group(3) or "")

    for match in re.finditer(
        r"\b(\d+(?:\.\d+)?)\s*(cm|mm|m|in|ft|units|inches)\s*[x×]\s*(\d+(?:\.\d+)?)\s*(cm|mm|m|in|ft|units|inches)",
        blob,
        flags=re.IGNORECASE,
    ):
        remember("width", match.group(1), match.group(2) or "")
        remember("height", match.group(3), match.group(4) or match.group(2) or "")

    return measurements


def measurement_value(measurements: dict[str, tuple[str, str]], *labels: str) -> tuple[str, str]:
    for label in labels:
        key = label.lower().replace(" ", "")
        if key in measurements:
            return measurements[key]
    return ("", "")


def measurement_unit(measurements: dict[str, tuple[str, str]], default: str = "units") -> str:
    for _label, (_value, unit) in measurements.items():
        if unit:
            return unit
    return default


def concrete_problem_text(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return False
    if re.search(r"\d", lowered):
        return True
    return any(term in lowered for term in ("source model", "source prism", "shipping", "popcorn", "table", "graph", "coordinate"))


def shifted_measure_value(value: str, *, seed: int, index: int) -> str:
    return shift_numeric_token(value, seed=seed, index=index)


def synthesized_problem_statement(plan_slide: dict[str, Any], *, similar: bool = False) -> str:
    blob = workbook_source_blob(plan_slide)
    lowered = source_measure_blob(plan_slide).lower()
    mode = problem_workbook_mode(blob)
    measurements = source_measurements(plan_slide)
    unit = measurement_unit(measurements, default="units")
    seed = sum(num for num in plan_slide.get("source_slide_numbers", []) if isinstance(num, int)) + (1 if similar else 0)

    if mode == "decision":
        if similar:
            small_volume = "24"
            large_volume = "36"
            target_volume = "72"
            return (
                f"Box A holds {shifted_measure_value(small_volume, seed=seed, index=0)} cubic units and "
                f"Box B holds {shifted_measure_value(large_volume, seed=seed, index=1)} cubic units. "
                f"You need at least {shifted_measure_value(target_volume, seed=seed, index=2)} cubic units. "
                "Which option would you choose, and how many boxes would you buy?"
            )
        return "Use the source popcorn information to compare the choices, decide how many boxes are needed, and justify the best option."

    if mode == "shipping":
        if similar:
            volume = shifted_measure_value("54", seed=seed, index=0)
            width = shifted_measure_value("3", seed=seed, index=1)
            height = shifted_measure_value("2", seed=seed, index=2)
            return (
                f"A box has volume {volume} cubic inches, width {width} inches, and height {height} inches. "
                "Find the missing length. Does it meet the 12-inch shipping rule?"
            )
        return "Use the source dimensions and volume to find the missing box length, then compare it to the 12-inch shipping rule."

    if "volume" in lowered or mode in {"unit_model", "volume_formula"}:
        length_value, length_unit = measurement_value(measurements, "length")
        width_value, width_unit = measurement_value(measurements, "width")
        height_value, height_unit = measurement_value(measurements, "height")
        active_unit = length_unit or width_unit or height_unit or unit
        if width_value and height_value and not length_value:
            if similar:
                sim_width = shifted_measure_value(width_value, seed=seed, index=0)
                sim_height = shifted_measure_value(height_value, seed=seed, index=1)
                sim_length = shifted_measure_value("4", seed=seed, index=2)
                return (
                    f"A rectangular prism has length {sim_length} {active_unit}, width {sim_width} {active_unit}, "
                    f"and height {sim_height} {active_unit}. Find the volume."
                )
            return (
                f"Use the source prism model. The width is {width_value} {active_unit} and the height is {height_value} {active_unit}. "
                "Determine the length from the model and find the volume."
            )
        if length_value and width_value and height_value:
            if similar:
                return (
                    f"A rectangular prism has length {shifted_measure_value(length_value, seed=seed, index=0)} {active_unit}, "
                    f"width {shifted_measure_value(width_value, seed=seed, index=1)} {active_unit}, "
                    f"and height {shifted_measure_value(height_value, seed=seed, index=2)} {active_unit}. Find the volume."
                )
            return (
                f"A rectangular prism has length {length_value} {active_unit}, width {width_value} {active_unit}, "
                f"and height {height_value} {active_unit}. Find the volume."
            )
        if similar:
            return "A rectangular prism has length 5 units, width 3 units, and height 2 units. Find the volume."
        return "Use the source prism model to determine the dimensions and find the volume."

    if "area" in lowered and "rhombus" in lowered:
        d1_value, d1_unit = measurement_value(measurements, "d1")
        d2_value, d2_unit = measurement_value(measurements, "d2")
        base_value, base_unit = measurement_value(measurements, "base", "b")
        height_value, height_unit = measurement_value(measurements, "height", "h")
        active_unit = d1_unit or d2_unit or base_unit or height_unit or measurement_unit(measurements, default="cm")
        if d1_value and d2_value:
            if similar:
                return (
                    f"A rhombus has diagonals d1 = {shifted_measure_value(d1_value, seed=seed, index=0)} {active_unit} "
                    f"and d2 = {shifted_measure_value(d2_value, seed=seed, index=1)} {active_unit}. Find the area."
                )
            return f"A rhombus has diagonals d1 = {d1_value} {active_unit} and d2 = {d2_value} {active_unit}. Find the area."
        if base_value and height_value:
            if similar:
                return (
                    f"A rhombus has base {shifted_measure_value(base_value, seed=seed, index=0)} {active_unit} "
                    f"and height {shifted_measure_value(height_value, seed=seed, index=1)} {active_unit}. Find the area."
                )
            return f"A rhombus has base {base_value} {active_unit} and height {height_value} {active_unit}. Find the area."

    if "area" in lowered and "parallelogram" in lowered:
        base_value, base_unit = measurement_value(measurements, "base", "b")
        height_value, height_unit = measurement_value(measurements, "height", "h")
        active_unit = base_unit or height_unit or measurement_unit(measurements, default="units")
        if base_value and height_value:
            if similar:
                return (
                    f"A parallelogram has base {shifted_measure_value(base_value, seed=seed, index=0)} {active_unit} "
                    f"and height {shifted_measure_value(height_value, seed=seed, index=1)} {active_unit}. Find the area."
                )
            return f"A parallelogram has base {base_value} {active_unit} and height {height_value} {active_unit}. Find the area."

    return ""


def is_solve_ready_problem_text(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return False
    if lowered.startswith(("the width is", "the height is", "the length is")):
        return False
    if lowered in {"volume of a rectangular prism.", "volume of a rectangular prism"}:
        return False
    if len(lowered.split()) <= 3:
        return False
    if "?" in lowered:
        return True
    return any(
        lowered.startswith(starter)
        for starter in ("solve", "find", "determine", "write", "graph", "plot", "compare", "choose", "explain", "show", "use")
    )


def full_problem_context_score(text: str) -> int:
    cleaned = normalize_whitespace(text)
    if not cleaned:
        return 0
    score = 0
    word_count = len(cleaned.split())
    if concrete_problem_text(cleaned):
        score += 4
    if word_count >= 12:
        score += 2
    if word_count >= 18:
        score += 3
    if len(cleaned) >= 95:
        score += 4
    if len(cleaned) >= 140:
        score += 3
    punctuation_count = len(re.findall(r"[?.!]", cleaned))
    if punctuation_count >= 2:
        score += 4
    lowered = cleaned.lower()
    if lowered.startswith(("think about it:", "what is", "how can", "in general, how")) and word_count <= 12:
        score -= 4
    return score


def problem_statement_candidates(plan_slide: dict[str, Any], *, limit: int = 4) -> list[str]:
    candidates = source_problem_cards(plan_slide, limit=max(limit + 2, 6), max_len=220)
    scored: list[tuple[int, int, str]] = []
    for index, candidate in enumerate(candidates):
        cleaned = normalize_whitespace(candidate)
        if not cleaned:
            continue
        default_end = "?" if is_problem_like_text(cleaned) or "?" in cleaned else "."
        statement = ensure_terminal_punctuation(cleaned, default=default_end)
        lowered = statement.lower()
        score = 0
        if is_solve_ready_problem_text(statement):
            score += 12
        if is_problem_like_text(statement):
            score += 4
        if re.search(r"\d", statement):
            score += 2
        score += full_problem_context_score(statement)
        if lowered.startswith(("the width is", "the height is", "the length is")):
            score -= 8
        if lowered in {"volume of a rectangular prism.", "volume of a rectangular prism"}:
            score -= 12
        scored.append((score, -index, statement))
    ordered = [statement for _score, _neg_index, statement in sorted(scored, key=lambda item: (-item[0], -item[1]))]
    return unique_nonempty(ordered, limit=limit)


def primary_source_problem_targets(plan_slide: dict[str, Any], *, limit: int = 2) -> list[str]:
    statements = problem_statement_candidates(plan_slide, limit=max(limit + 2, 4))
    targets: list[str] = []
    if statements and is_solve_ready_problem_text(statements[0]):
        targets.append(normalize_whitespace(statements[0]))
    for statement in statements:
        cleaned = normalize_whitespace(statement)
        if not cleaned:
            continue
        if (
            full_problem_context_score(cleaned) >= 8
            or len(cleaned) >= 90
            or ("?" in cleaned and len(cleaned) >= 60)
        ):
            targets.append(cleaned)
    if not targets and statements:
        targets.append(statements[0])
    return unique_nonempty(targets, limit=limit)


def rendered_source_problem_statement(plan_slide: dict[str, Any]) -> str:
    targets = primary_source_problem_targets(plan_slide, limit=1)
    if targets:
        return targets[0]
    return source_problem_statement(plan_slide)


def shift_numeric_token(token: str, *, seed: int, index: int) -> str:
    raw_value = token.replace(",", "")
    try:
        numeric = float(raw_value)
    except ValueError:
        return token
    if numeric == 0:
        return token
    if "." in raw_value:
        step = 0.5 if numeric < 10 else 1.0
        adjusted = numeric + step * (1 + ((seed + index) % 2))
        return f"{adjusted:.1f}".rstrip("0").rstrip(".")

    integer = int(round(numeric))
    step = 1 if integer < 6 else 2 if integer < 20 else 5 if integer < 100 else 10
    adjusted = integer + step * (1 + ((seed + index) % 2))
    return f"{adjusted:,}" if "," in token else str(adjusted)


def numeric_variant_problem(text: str, *, seed: int) -> str:
    matches = list(re.finditer(r"(?<![A-Za-z])\d[\d,]*(?:\.\d+)?", text))
    if not matches:
        return ""
    pieces: list[str] = []
    last_end = 0
    changed = 0
    for match_index, match in enumerate(matches):
        pieces.append(text[last_end:match.start()])
        token = match.group(0)
        if changed < 4:
            pieces.append(shift_numeric_token(token, seed=seed, index=match_index))
            changed += 1
        else:
            pieces.append(token)
        last_end = match.end()
    pieces.append(text[last_end:])
    variant = normalize_whitespace("".join(pieces))
    if display_text_key(variant) == display_text_key(text):
        return ""
    return ensure_terminal_punctuation(variant, default="?" if "?" in text else ".")


def source_problem_statement(plan_slide: dict[str, Any]) -> str:
    candidates = problem_statement_candidates(plan_slide, limit=4)
    synthesized = synthesized_problem_statement(plan_slide, similar=False)
    if candidates:
        lead = candidates[0]
        if synthesized and (not concrete_problem_text(lead) or not is_solve_ready_problem_text(lead)):
            return synthesized
        return lead
    if synthesized:
        return synthesized
    facts = problem_fact_statements(plan_slide, limit=2)
    detail_text = f" {' '.join(facts)}" if facts else ""
    focus_phrase = source_problem_focus_phrase(plan_slide)
    return f"Solve the {focus_phrase}.{detail_text}".strip()


def exact_source_followup_problem_statement(plan_slide: dict[str, Any], *, ordinal: int = 1) -> str:
    candidates = problem_statement_candidates(plan_slide, limit=max(ordinal + 3, 6))
    solve_ready = [candidate for candidate in candidates if is_solve_ready_problem_text(candidate)]
    if len(solve_ready) > ordinal:
        return solve_ready[ordinal]
    return ""


def alternate_problem_statement(plan_slide: dict[str, Any], *, ordinal: int = 1) -> str:
    candidates = problem_statement_candidates(plan_slide, limit=4)
    solve_ready = [candidate for candidate in candidates if is_solve_ready_problem_text(candidate)]
    if len(solve_ready) > ordinal:
        return solve_ready[ordinal]
    synthesized = synthesized_problem_statement(plan_slide, similar=True)
    if synthesized:
        return synthesized
    source_numbers = [num for num in plan_slide.get("source_slide_numbers", []) if isinstance(num, int)]
    seed = sum(source_numbers) + ordinal
    variant = numeric_variant_problem(source_problem_statement(plan_slide), seed=seed)
    if variant:
        return variant
    return ""


def similar_problem_statement(plan_slide: dict[str, Any]) -> str:
    exact_followup = exact_source_followup_problem_statement(plan_slide, ordinal=1)
    if exact_followup:
        return exact_followup
    alternate = alternate_problem_statement(plan_slide, ordinal=1)
    synthesized = synthesized_problem_statement(plan_slide, similar=True)
    if synthesized and (not alternate or not concrete_problem_text(alternate) or not re.search(r"\d", alternate)):
        return synthesized
    if alternate:
        return alternate
    focus_phrase = similar_problem_focus_phrase(plan_slide)
    return f"Write and solve one {focus_phrase} that uses the same strategy."


def build_problem_creation_prompt(plan_slide: dict[str, Any]) -> str:
    source_phrase = source_problem_focus_phrase(plan_slide)
    similar_phrase = similar_problem_focus_phrase(plan_slide)
    return (
        f"Change one number, measurement, or condition from the {source_phrase}. "
        f"Write a {similar_phrase} and solve it."
    )


def third_problem_statement(plan_slide: dict[str, Any]) -> str:
    similar = similar_problem_statement(plan_slide)
    alternate = alternate_problem_statement(plan_slide, ordinal=2)
    if (
        alternate
        and concrete_problem_text(alternate)
        and display_text_key(alternate) != display_text_key(similar)
    ):
        return alternate
    seed = sum(num for num in plan_slide.get("source_slide_numbers", []) if isinstance(num, int)) + 7
    if re.search(r"\d", similar):
        shifted = numeric_variant_problem(similar, seed=seed)
        if shifted and display_text_key(shifted) != display_text_key(similar):
            return shifted
    synthesized = synthesized_problem_statement(plan_slide, similar=True)
    if synthesized:
        shifted = numeric_variant_problem(synthesized, seed=seed)
        if shifted and display_text_key(shifted) != display_text_key(synthesized):
            return shifted
        return synthesized
    return build_problem_creation_prompt(plan_slide)


def problem_tip_lines(plan_slide: dict[str, Any], *, include_visual_tip: bool = False, limit: int = 3) -> list[str]:
    blob = workbook_source_blob(plan_slide)
    mode = problem_workbook_mode(blob)
    source_blob = source_measure_blob(plan_slide).lower()
    tips: list[str] = []
    if mode == "decision":
        tips.extend(
            [
                "Find the volume of each option before you decide.",
                "Compare each choice to the amount you need.",
                "Use numbers from the problem to justify your decision.",
            ]
        )
    elif mode == "shipping":
        tips.extend(
            [
                "Write V = l x w x h before solving.",
                "Divide the volume by width x height to find the missing length.",
                "Compare your length to the 12-inch rule at the end.",
            ]
        )
    elif "area" in source_blob and "rhombus" in source_blob:
        tips.extend(
            [
                "Identify the diagonals or the base and height first.",
                "Use the matching area formula before you calculate.",
                "Label the final answer in square units.",
            ]
        )
    elif "area" in source_blob:
        tips.extend(
            [
                "Find the base and height before you multiply.",
                "Check that you are using the right area formula.",
                "Label the final answer in square units.",
            ]
        )
    elif "volume" in source_blob or mode in {"unit_model", "volume_formula"}:
        tips.extend(
            [
                "Find or count each dimension before you multiply.",
                "Use V = l x w x h to set up the problem.",
                "Label the final answer in cubic units.",
            ]
        )
    elif has_data_analysis_context(source_blob):
        tips.extend(
            [
                "Identify the key values or sections of the data display first.",
                "Compare the spread, center, or extremes that answer the question.",
                "Use the box plot or ordered data as evidence in your explanation.",
            ]
        )
    else:
        tips.extend(
            [
                "Underline what you know and what you need to find.",
                "Choose the formula or strategy that matches the lesson.",
                "Check whether your answer makes sense in the context.",
            ]
        )
    if include_visual_tip:
        tips.insert(0, "Use the source model, image, or diagram to help set up the problem.")
    return unique_nonempty(tips, limit=limit)


def problem_tip_body(plan_slide: dict[str, Any], *, include_visual_tip: bool = False, limit: int = 3) -> str:
    lines = problem_tip_lines(plan_slide, include_visual_tip=include_visual_tip, limit=limit)
    return "\n".join(f"- {truncate_text(line, 74)}" for line in lines)


def workbook_prompt_stack_subtitle(variant: str) -> str:
    if variant == "guided":
        return "Study the model, then solve a similar problem with the same strategy."
    if variant == "exit":
        return "Solve the final check clearly and explain how the evidence supports your answer."
    return "Solve the source problem, then solve a similar problem and explain your reasoning."


def workbook_prompt_stack_kicker(variant: str) -> str:
    if variant == "guided":
        return "Guided Example"
    if variant == "exit":
        return "Quick Check"
    return ""


def prompt_stack_layout_mode(problems: list[str]) -> str:
    cleaned = [normalize_whitespace(problem) for problem in problems if normalize_whitespace(problem)]
    if not cleaned:
        return "standard"
    if any(len(problem) >= 140 for problem in cleaned):
        return "focus"
    if len(cleaned) >= 2 and sum(len(problem) for problem in cleaned[:2]) >= 250 and any(len(problem) >= 115 for problem in cleaned[:2]):
        return "focus"
    if len(cleaned) >= 3 and sum(len(problem) for problem in cleaned[:3]) >= 300:
        return "focus"
    return "standard"


def problem_display_cards(plan_slide: dict[str, Any], *, variant: str = "practice") -> list[str]:
    source_cards = unique_nonempty(
        primary_source_problem_targets(plan_slide, limit=2) + problem_statement_candidates(plan_slide, limit=4),
        limit=4,
    )
    source_problem = source_cards[0] if source_cards else rendered_source_problem_statement(plan_slide)
    if variant == "guided":
        cards = [f"Model This: {source_problem}"]
        partner_prompt = (
            source_cards[1]
            if len(source_cards) > 1 and is_solve_ready_problem_text(source_cards[1])
            else similar_problem_statement(plan_slide)
        )
        cards.append(f"Your Turn: {partner_prompt}")
        return unique_nonempty(cards, limit=2)
    if variant == "exit":
        cards = [f"Solve This: {source_problem}"]
        cards.append(
            "Prove It: "
            + truncate_display_copy(
                plan_slide.get("response_prompt", "") or "Show why your answer makes sense.",
                118,
            )
        )
        return unique_nonempty(cards, limit=2)
    cards: list[str] = []
    full_source_prompts = [prompt for prompt in source_cards if is_solve_ready_problem_text(prompt)]
    if source_problem and source_problem not in full_source_prompts:
        full_source_prompts.insert(0, source_problem)
    for index, prompt in enumerate(full_source_prompts[:2], start=1):
        cards.append(f"P{index}: {prompt}")
    if not cards:
        cards.append(f"P1: {source_problem}")
    if len(cards) == 1:
        cards.append(f"P2: {similar_problem_statement(plan_slide)}")
    long_prompt_stack = any(len(prompt) >= 150 for prompt in cards)
    if len(cards) < 3 and not long_prompt_stack:
        cards.append(f"P3: {third_problem_statement(plan_slide)}")
    return unique_nonempty(cards, limit=3 if not long_prompt_stack else 2)


def fallback_problem_activity(plan_slide: dict[str, Any]) -> dict[str, Any]:
    blob = workbook_source_blob(plan_slide)
    mode = problem_workbook_mode(blob)
    source_phrase = source_problem_focus_phrase(plan_slide)
    similar_phrase = similar_problem_focus_phrase(plan_slide)
    if mode == "decision":
        return {
            "activity_name": "decision move sort",
            "activity_family": "compare_rank",
            "activity_instructions": "Move the cards into the order you will use to compare the popcorn choices, then explain your choice.",
            "movable_pieces": ["Find each volume", "Compare the choices", "Choose the best fit", "Justify the choice"],
            "answer_check": "Your final choice should match the option that best meets the need with clear volume evidence.",
        }
    if mode == "shipping":
        return {
            "activity_name": "missing-length move sort",
            "activity_family": "sequence_order",
            "activity_instructions": "Place the solve moves in order before you find the missing length and compare it to the shipping rule.",
            "movable_pieces": ["Read the condition", "Write V = l x w x h", "Find the missing length", "Compare to 12 inches"],
            "answer_check": "Check that you solved for length first and then stated whether the box meets the rule.",
        }
    if has_data_analysis_context(blob):
        return {
            "activity_name": "data analysis solve builder",
            "activity_family": "build_construct",
            "activity_instructions": "Build the analysis path before you write the full answer and explanation.",
            "movable_pieces": ["Read the data display", "Find the key values", "Compare the data", "Explain the evidence"],
            "answer_check": "A strong solve path identifies the correct values and uses the data display as evidence.",
        }
    return {
        "activity_name": f"{display_problem_focus_phrase(plan_slide)} solve builder",
        "activity_family": "build_construct",
        "activity_instructions": f"Move the cards into a strong solve plan, then use that plan to solve the {source_phrase} or a {similar_phrase}.",
        "movable_pieces": ["What do I know?", "What do I need?", "Set up the math", "Check the answer"],
        "answer_check": "Your completed plan should help you solve the problem and explain why the answer makes sense.",
    }


def refresh_problem_activity_for_context(plan_slide: dict[str, Any]) -> None:
    family = normalize_whitespace(plan_slide.get("activity_family", ""))
    if not family or plan_slide.get("kind") not in PROBLEM_SOLVING_KINDS:
        return
    source_blob = workbook_source_blob(plan_slide)
    if not has_data_analysis_context(source_blob):
        return
    original_name = normalize_whitespace(plan_slide.get("activity_name", "")).lower()
    uniqueness_seed = normalize_whitespace(
        " ".join(
            [
                original_name,
                plan_slide.get("template_role", ""),
                plan_slide.get("title", ""),
                plan_slide.get("section", ""),
            ]
        )
    ).lower()
    name_suffix = str(sum((index + 1) * ord(char) for index, char in enumerate(uniqueness_seed)) % 10000 or 1)
    source_signature = "-".join(str(num) for num in plan_slide.get("source_slide_numbers", [])[:3] if isinstance(num, int)) or "analysis"
    stale_formula_pieces = any(
        term in " ".join(plan_slide.get("movable_pieces", [])).lower()
        for term in ("formula", "equation", "volume", "missing value", "prism")
    )
    if family == "build_construct":
        plan_slide["activity_name"] = f"data analysis solve builder {name_suffix}-{source_signature}"
        plan_slide["activity_instructions"] = "Build the analysis path before you write the full answer and explanation."
        plan_slide["movable_pieces"] = [
            "Read the data display",
            "Find the key values",
            "Compare the data",
            "Explain the evidence",
        ]
        plan_slide["answer_check"] = "A strong solve path identifies the correct values and uses the data display as evidence."
    elif family == "match_pair":
        plan_slide["activity_name"] = f"data evidence matching {name_suffix}-{source_signature}"
        plan_slide["activity_instructions"] = "Match each data clue, comparison idea, or evidence sentence to the part of the display that supports it."
        plan_slide["movable_pieces"] = [
            "Key values",
            "Median / center",
            "Spread clue",
            "Evidence sentence",
        ]
        plan_slide["answer_check"] = "Each match should connect a claim or comparison to the data evidence that supports it."
    elif family == "sequence_order":
        plan_slide["activity_name"] = f"data analysis solve flow {name_suffix}-{source_signature}"
        plan_slide["activity_instructions"] = "Order the analysis moves before you teach the strategy to a partner."
        plan_slide["movable_pieces"] = [
            "Read the question",
            "Find the key values",
            "Compare the data",
            "Explain the evidence",
        ]
        plan_slide["answer_check"] = "A strong solve-and-teach flow uses the correct values and cites the data display as evidence."
    elif family == "sort_classify":
        plan_slide["activity_name"] = f"data clue sort {name_suffix}-{source_signature}"
        plan_slide["activity_instructions"] = "Sort the clues by whether they name a key value, a comparison idea, or evidence from the data display."
        if stale_formula_pieces or not plan_slide.get("movable_pieces"):
            plan_slide["movable_pieces"] = [
                "Key values",
                "Spread clue",
                "Comparison idea",
                "Evidence",
            ]
        plan_slide["answer_check"] = "Each card should help you identify a value, compare the data, or support your explanation."
    elif family == "detect_justify":
        plan_slide["activity_name"] = f"data claim and evidence sort {name_suffix}-{source_signature}"
        plan_slide["activity_instructions"] = "Sort the claim, evidence, and explanation pieces so they match what the data display shows."
        plan_slide["movable_pieces"] = [
            "Best claim",
            "Key values",
            "Spread evidence",
            "Explain the reasoning",
        ]
        plan_slide["answer_check"] = "A strong final check matches the claim to the correct values and explains how the data supports it."


def activity_slide_for_render(plan_slide: dict[str, Any]) -> dict[str, Any]:
    if has_activity(plan_slide):
        rendered = dict(plan_slide)
        refresh_problem_activity_for_context(rendered)
        return rendered
    rendered = dict(plan_slide)
    rendered.update(fallback_problem_activity(plan_slide))
    refresh_problem_activity_for_context(rendered)
    return rendered


def short_problem_solving_subtitle(plan_slide: dict[str, Any]) -> str:
    kind = plan_slide.get("kind", "")
    phase = normalize_whitespace(plan_slide.get("practice_phase", "")).lower()
    focus_phrase = display_problem_focus_phrase(plan_slide).replace("data analysis", "data-analysis")
    lead_phrase = "source problem" if focus_phrase == "problem" else focus_phrase
    stretch_phrase = "stretch problem" if focus_phrase == "problem" else f"stretch {focus_phrase}"
    if kind == "worked_example":
        return f"Model the {lead_phrase}, then solve a similar one."
    if kind == "practice":
        if "together" in phase:
            return f"Solve the {lead_phrase} together, then solve a similar one."
        if "independently" in phase:
            return f"Solve the {lead_phrase} and a similar one independently."
        return f"Solve the {lead_phrase}, then solve a similar one."
    if kind == "challenge":
        return f"Solve the {stretch_phrase}, then compare or extend it."
    if kind == "exit_ticket":
        final_phrase = "problem" if focus_phrase == "problem" else focus_phrase
        return f"Solve one final {final_phrase} and explain your reasoning."
    return ""


def publisher_problem_response_prompt(plan_slide: dict[str, Any]) -> str:
    source_blob = source_measure_blob(plan_slide).lower()
    if any(term in source_blob for term in ("median", "mean", "data set", "dot plot", "line plot", "ordered data", "outlier")):
        return "Explain how the ordered data or data display support your answer."
    if "volume" in source_blob or any(term in source_blob for term in ("prism", "unit cube", "cubic")):
        return "Explain how the dimensions, formula, and units support your answer."
    if "area" in source_blob:
        return "Explain how the measurements, labels, and formula support your answer."
    if any(term in source_blob for term in ("graph", "table", "equation", "expression", "coordinate")):
        return "Explain how the representation and values support your answer."
    return "Explain how your strategy and evidence support your answer."


def strengthen_slide_copy_for_publisher_rigor(
    slide: dict[str, Any],
    *,
    session_key: str,
    deck: dict[str, Any],
) -> None:
    kind = slide.get("kind", "")
    lower_subtitle = normalize_whitespace(slide.get("subtitle", "")).lower()
    lower_response = normalize_whitespace(slide.get("response_prompt", "")).lower()

    if kind in PROBLEM_SOLVING_KINDS:
        slide["subtitle"] = short_problem_solving_subtitle(slide)
        if (
            not normalize_whitespace(slide.get("response_prompt", ""))
            or any(
                token in lower_response
                for token in (
                    "what step should come next",
                    "which strategy from the model",
                    "show your work",
                    "show the work",
                )
            )
        ):
            slide["response_prompt"] = publisher_problem_response_prompt(slide)
    elif kind == "reflection":
        if not lower_subtitle or "next step" in lower_subtitle or "lesson so far" in lower_subtitle:
            slide["subtitle"] = "Pause, summarize the lesson, and name the idea you want to strengthen."
        if not normalize_whitespace(slide.get("response_prompt", "")) or "short reflection" in lower_response:
            slide["response_prompt"] = "Write a brief reflection that names your strategy, evidence, and next focus."
    elif kind == "be_curious" and "lesson problem" in lower_subtitle:
        slide["subtitle"] = replace_case_insensitive(slide.get("subtitle", ""), r"\blesson problem\b", "source problem")
    elif kind == "cover":
        focus_text = cover_focus_statement(deck, session_key, slide.get("primary_text", "") or deck.get("lesson_title", ""))
        if not normalize_whitespace(slide.get("primary_text", "")) or publisher_copyedit_issues(slide.get("primary_text", "")):
            slide["primary_text"] = focus_text


GENERIC_DISCUSSION_QUESTION_MARKERS = (
    "q1: ___",
    "q2: ___",
    "q3: ___",
    "the ___ changes because ___.",
    "how is this problem like the guided example, and what did your partner notice first?",
    "explain the idea to a partner using the lesson vocabulary.",
)


def clean_discussion_question(text: str) -> str:
    cleaned = normalize_whitespace(re.sub(r"^\s*q\s*\d+\s*[:.)-]?\s*", "", text, flags=re.IGNORECASE))
    cleaned = trim_dangling_display_text(cleaned)
    if not cleaned:
        return ""
    if "?" not in cleaned and cleaned.split() and cleaned.split()[0].lower() in {"what", "which", "how", "why", "when", "where"}:
        cleaned = ensure_terminal_punctuation(cleaned, default="?")
    return truncate_text(cleaned, 150)


def is_specific_discussion_question(text: str) -> bool:
    cleaned = clean_discussion_question(text)
    lowered = cleaned.lower()
    if not cleaned or "___" in cleaned or len(cleaned.split()) < 6:
        return False
    if any(marker == lowered for marker in GENERIC_DISCUSSION_QUESTION_MARKERS):
        return False
    specificity_markers = (
        "source",
        "problem",
        "strategy",
        "answer",
        "setup",
        "equation",
        "formula",
        "representation",
        "representations",
        "table",
        "graph",
        "dot plot",
        "median",
        "mean",
        "measurement",
        "measurements",
        "dimension",
        "dimensions",
        "area",
        "volume",
        "evidence",
        "model",
        "data",
        "pattern",
        "values",
        "units",
        "step",
        "clue",
        "reasoning",
        "work",
        "notebook",
    )
    return any(marker in lowered for marker in specificity_markers)


def discussion_focus_terms_for_slide(plan_slide: dict[str, Any]) -> tuple[str, str]:
    blob = normalize_whitespace(
        " ".join(
            [
                source_problem_statement(plan_slide),
                similar_problem_statement(plan_slide),
                plan_slide.get("primary_text", ""),
                plan_slide.get("secondary_text", ""),
                " ".join(plan_slide.get("tasks", [])),
                " ".join(plan_slide.get("bullets", [])),
            ]
        )
    )
    lowered = blob.lower()
    profile = math_profile_for_text(blob)
    topic = profile.get("topic", "")

    if topic == "data_analysis":
        return ("ordered values in the data set", "dot plot or ordered list")
    if topic == "fraction_decimal_percent":
        return ("equivalent representations", "fraction-decimal-percent relationship")
    if topic == "volume_prism":
        return ("dimensions from the prism model", "volume formula and units")
    if topic in {"rectangle_area", "parallelogram_area", "triangle_area", "trapezoid_area", "rhombus_area", "regular_polygon_area"}:
        return ("measurements in the figure", "area formula or area model")
    if "graph" in lowered:
        return ("feature of the graph", "graph evidence")
    if "table" in lowered:
        return ("table row or value pair", "table pattern")
    if "equation" in lowered or "rule" in lowered:
        return ("part of the equation", "equation")
    if "ratio" in lowered or "unit rate" in lowered:
        return ("ratio relationship", "unit-rate reasoning")
    if "fraction" in lowered or "percent" in lowered:
        return ("part-whole relationship", "visual model")
    return ("key quantity or step", "setup or evidence")


def generated_partner_prompt(plan_slide: dict[str, Any]) -> str:
    kind = plan_slide.get("kind", "")
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    if role in {"two_column_compare", "collaborative_practice"} or (
        normalize_whitespace(plan_slide.get("partner_a_problem", ""))
        and normalize_whitespace(plan_slide.get("partner_b_problem", ""))
    ):
        return "Solve your assigned problem first. Then compare what stayed the same, what changed, and which explanation is clearest."
    if role == "turn_and_teach" or normalize_whitespace(plan_slide.get("premium_layout", "")) == "turn_and_teach":
        return "Solve first. Then teach the strategy to a partner and revise one sentence after the discussion."
    if kind == "worked_example":
        return "Explain the setup to a partner before either of you writes the final answer."
    if kind in {"practice", "challenge", "exit_ticket"}:
        return "Solve independently first. Then compare one move and one piece of evidence with a partner."
    if kind == "reflection":
        return "Share your strongest strategy with a partner and listen for one idea you want to borrow."
    return ""


def generated_discussion_questions(plan_slide: dict[str, Any]) -> list[str]:
    kind = plan_slide.get("kind", "")
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    focus_phrase, evidence_phrase = discussion_focus_terms_for_slide(plan_slide)
    source_blob = normalize_whitespace(
        " ".join(
            [
                source_problem_statement(plan_slide),
                similar_problem_statement(plan_slide),
                plan_slide.get("primary_text", ""),
                plan_slide.get("subtitle", ""),
            ]
        )
    )
    profile = math_profile_for_text(source_blob)
    topic = profile.get("topic", "")
    source_focus = source_problem_focus_phrase(plan_slide)

    if topic == "data_analysis":
        base_questions = [
            "Which ordered values matter most, and why?",
            "How does the dot plot prove the median?",
        ]
    elif topic == "fraction_decimal_percent":
        base_questions = [
            "Which representation helps you most, and why?",
            "How do the equivalent values prove your answer?",
        ]
    elif topic == "volume_prism":
        base_questions = [
            "Which dimensions matter first, and why?",
            "How do the formula and units prove your answer?",
        ]
    elif topic in {"rectangle_area", "parallelogram_area", "triangle_area", "trapezoid_area", "rhombus_area", "regular_polygon_area"}:
        base_questions = [
            "Which measurements matter first, and why?",
            "How does the area model or formula prove your answer?",
        ]
    elif "graph" in source_blob.lower():
        base_questions = [
            "Which graph feature matters most, and why?",
            "How does the graph prove your answer?",
        ]
    elif "table" in source_blob.lower():
        base_questions = [
            "Which table values matter most, and why?",
            "How does the table pattern prove your answer?",
        ]
    elif "equation" in source_blob.lower() or "rule" in source_blob.lower():
        base_questions = [
            "Which part of the equation matters most, and why?",
            "How does the equation prove your answer?",
        ]
    else:
        base_questions = [
            f"Which {focus_phrase} matters most, and why?",
            f"How does the {evidence_phrase} prove your answer?",
        ]

    if role in {"two_column_compare", "collaborative_practice"} or (
        normalize_whitespace(plan_slide.get("partner_a_problem", ""))
        and normalize_whitespace(plan_slide.get("partner_b_problem", ""))
    ):
        return [
            "What stays the same in both problems?",
            "What changes, and how does that change your setup?",
            "Which explanation is easier to follow, and why?",
        ]
    if role == "turn_and_teach" or normalize_whitespace(plan_slide.get("premium_layout", "")) == "turn_and_teach":
        return [
            f"Which clue in the {source_focus} should your partner notice first?",
            base_questions[1],
            "After you listen to your partner, what sentence would you revise?",
        ]
    if kind == "worked_example":
        return [
            "Which clue tells you the first step?",
            f"How does the {evidence_phrase} show the setup works?",
        ]
    if kind in {"practice", "challenge", "exit_ticket"}:
        return base_questions
    if kind == "reflection":
        return [
            "Which strategy would you teach a partner first?",
            "What evidence from your notebook would you point to?",
        ]
    return []


def slide_needs_peer_discussion_support(plan_slide: dict[str, Any]) -> bool:
    kind = plan_slide.get("kind", "")
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    if role in {"two_column_compare", "collaborative_practice", "turn_and_teach"}:
        return True
    if kind in {"worked_example", "practice", "challenge", "reflection"} and role not in {
        "drag_sort",
        "error_analysis",
        "choice_board",
        "independent_practice",
        "twr_frame",
        "goal_tracker",
        "tiered_exit",
    }:
        return True
    return False


def should_refresh_partner_prompt(text: str) -> bool:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return True
    return any(
        marker in lowered
        for marker in (
            "then compare.",
            "talk with a partner",
            "explain the idea to a partner",
            "what did your partner notice first",
            "partner a solves the first problem",
        )
    )


def ensure_peer_discussion_support(plan_slide: dict[str, Any]) -> None:
    if not slide_needs_peer_discussion_support(plan_slide):
        return
    generated_questions = generated_discussion_questions(plan_slide)
    merged_questions: list[str] = []
    seen: set[str] = set()
    for question in list(plan_slide.get("discussion_questions", [])) + generated_questions:
        cleaned = clean_discussion_question(str(question))
        key = display_text_key(cleaned)
        if not key or key in seen or not is_specific_discussion_question(cleaned):
            continue
        seen.add(key)
        merged_questions.append(cleaned)
        if len(merged_questions) >= 3:
            break
    plan_slide["discussion_questions"] = merged_questions

    partner_prompt = normalize_whitespace(plan_slide.get("partner_prompt", ""))
    if should_refresh_partner_prompt(partner_prompt):
        plan_slide["partner_prompt"] = truncate_text(generated_partner_prompt(plan_slide), 180)

    if normalize_whitespace(plan_slide.get("template_role", "")) == "two_column_compare" and plan_slide["discussion_questions"]:
        third_question = plan_slide["discussion_questions"][2] if len(plan_slide["discussion_questions"]) > 2 else plan_slide["discussion_questions"][-1]
        if not normalize_whitespace(plan_slide.get("compare_frame", "")) or "same because" in plan_slide.get("compare_frame", "").lower():
            plan_slide["compare_frame"] = third_question


def source_activity_signature(plan_slide: dict[str, Any]) -> str:
    source_numbers = [str(num) for num in plan_slide.get("source_slide_numbers", [])[:3] if isinstance(num, int)]
    return "-".join(source_numbers) or "source"


def problem_activity_has_generic_markers(plan_slide: dict[str, Any]) -> bool:
    blob = normalize_whitespace(
        " ".join(
            [
                plan_slide.get("activity_name", ""),
                plan_slide.get("activity_instructions", ""),
                plan_slide.get("answer_check", ""),
                " ".join(plan_slide.get("movable_pieces", [])),
            ]
        )
    ).lower()
    if any(marker in blob for marker in GENERIC_PROBLEM_ACTIVITY_MARKERS):
        return True
    piece_keys = {
        display_text_key(piece)
        for piece in plan_slide.get("movable_pieces", [])
        if display_text_key(piece)
    }
    return any(piece in piece_keys for piece in GENERIC_PROBLEM_ACTIVITY_PIECES)


def source_focused_problem_activity_payload(plan_slide: dict[str, Any]) -> dict[str, Any]:
    refreshed = dict(plan_slide)
    if (
        not normalize_whitespace(refreshed.get("activity_name", ""))
        or normalize_whitespace(refreshed.get("activity_family", "")) not in SOURCE_ANCHORED_PROBLEM_ACTIVITY_FAMILIES
        or problem_activity_has_generic_markers(refreshed)
    ):
        refreshed.update(fallback_problem_activity(refreshed))
    refresh_problem_activity_for_context(refreshed)
    return {
        "activity_name": refreshed.get("activity_name", ""),
        "activity_family": refreshed.get("activity_family", ""),
        "activity_instructions": refreshed.get("activity_instructions", ""),
        "movable_pieces": list(refreshed.get("movable_pieces", [])),
        "answer_check": refreshed.get("answer_check", ""),
    }


def source_sort_labels(plan_slide: dict[str, Any]) -> list[str]:
    blob = workbook_source_blob(plan_slide)
    mode = problem_workbook_mode(blob)
    if has_data_analysis_context(blob):
        return ["Key values or clues", "Claims or evidence"]
    if mode == "decision":
        return ["Choices or quantities", "Evidence or reason"]
    if mode == "shipping":
        return ["Setup or value", "Condition or conclusion"]
    return ["What helps solve it", "What proves it"]


def source_drag_sort_pieces(plan_slide: dict[str, Any]) -> list[str]:
    blob = workbook_source_blob(plan_slide)
    if has_data_analysis_context(blob):
        return ["Key values", "Data clue", "Comparison idea", "Evidence"]
    mode = problem_workbook_mode(blob)
    if mode == "decision":
        return ["Choice", "Quantity needed", "Work shown", "Reason"]
    if mode == "shipping":
        return ["Known dimensions", "Missing length", "12-inch rule", "Conclusion"]
    source_cards = source_problem_cards(plan_slide, limit=4, max_len=72)
    fallback_cards = unique_nonempty(plan_slide.get("tasks", []) + plan_slide.get("bullets", []), limit=4)
    pieces = [
        trim_dangling_display_text(truncate_display_copy(item, 56))
        for item in source_cards + fallback_cards
    ]
    return unique_nonempty(pieces, limit=4) or ["Known information", "Strategy move", "Evidence", "Explain"]


def source_error_analysis_content(plan_slide: dict[str, Any]) -> tuple[list[str], str, str]:
    blob = workbook_source_blob(plan_slide)
    mode = problem_workbook_mode(blob)
    if has_data_analysis_context(blob):
        return (
            [
                "✅ Step 1: Read the data display carefully.",
                "✅ Step 2: Find the key values that answer the question.",
                "❌ Step 3: Make a claim without checking the median, quartiles, or spread.",
                "✅ Step 4: Use the data evidence to explain the answer.",
            ],
            "Fix the claim by citing the correct values or sections of the data display.",
            "Which value or part of the display gives the strongest evidence?",
        )
    if mode == "decision":
        return (
            [
                "✅ Step 1: Read what each choice gives you.",
                "✅ Step 2: Solve for the amount in each option.",
                "❌ Step 3: Choose before comparing the evidence.",
                "✅ Step 4: Justify the final choice with numbers.",
            ],
            "Fix the choice by comparing all of the source quantities first.",
            "Which numbers prove the best choice?",
        )
    if mode == "shipping":
        return (
            [
                "✅ Step 1: Write V = l x w x h.",
                "✅ Step 2: Use the known width and height.",
                "❌ Step 3: Skip the missing length and jump to the conclusion.",
                "✅ Step 4: Compare the result to the rule.",
            ],
            "Fix the work by solving for the missing length before checking the rule.",
            "Which step proves whether the box meets the condition?",
        )
    return (
        [
            "✅ Step 1: Read the exact source problem.",
            "✅ Step 2: Identify the quantities, values, or clues that matter most.",
            "❌ Step 3: Use a step or claim that does not match the source details.",
            "✅ Step 4: Explain the answer with evidence.",
        ],
        "Fix the step or claim so it matches the exact source problem.",
        "Which detail from the source problem proves your correction?",
    )


def source_compare_prompt(plan_slide: dict[str, Any]) -> str:
    questions = generated_discussion_questions(plan_slide)
    if len(questions) >= 3:
        return questions[2]
    if len(questions) >= 2:
        return questions[1]
    blob = workbook_source_blob(plan_slide)
    if has_data_analysis_context(blob):
        return "How do the values, comparisons, or evidence you used stay the same or change across the two problems?"
    return "What stays the same, what changes, and which evidence best supports your answer?"


def source_choice_paths(plan_slide: dict[str, Any]) -> list[dict[str, str]]:
    hints = problem_tip_lines(plan_slide, limit=3)
    prompts = unique_nonempty(
        [
            source_problem_statement(plan_slide),
            similar_problem_statement(plan_slide),
            third_problem_statement(plan_slide),
        ],
        limit=3,
    )
    paths: list[dict[str, str]] = []
    for index, prompt in enumerate(prompts):
        hint = hints[index] if index < len(hints) else hints[-1] if hints else "Use the source strategy and show your evidence."
        paths.append(
            {
                "label": f"Problem {index + 1}",
                "problem": prompt,
                "hint": "Hint: " + trim_dangling_display_text(truncate_display_copy(hint, 72)),
                "answer": "My answer: ___",
            }
        )
    return paths


def source_independent_problem_set(plan_slide: dict[str, Any]) -> list[dict[str, str]]:
    second_source_problem = exact_source_followup_problem_statement(plan_slide, ordinal=1)
    third_source_problem = exact_source_followup_problem_statement(plan_slide, ordinal=2)
    prompts = unique_nonempty(
        [
            source_problem_statement(plan_slide),
            second_source_problem or similar_problem_statement(plan_slide),
            third_source_problem or third_problem_statement(plan_slide),
        ],
        limit=3,
    )
    problems: list[dict[str, str]] = []
    for index, prompt in enumerate(prompts):
        problems.append(
            {
                "label": f"P{index + 1}",
                "prompt": prompt,
                "answer": "Answer: ___",
            }
        )
    return problems


def retarget_template_role_slide_to_source(plan_slide: dict[str, Any]) -> None:
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    if role not in SOURCE_ALIGNED_TEMPLATE_ACTIVITY_ROLES:
        return

    signature = source_activity_signature(plan_slide)
    source_focus = source_problem_focus_phrase(plan_slide)
    source_phrase = "source problem" if source_focus == "problem" else source_focus
    primary_problem = source_problem_statement(plan_slide)
    similar_problem = similar_problem_statement(plan_slide)
    discussion_questions = generated_discussion_questions(plan_slide)

    if role == "drag_sort":
        plan_slide["title"] = "Source Clue Sort"
        plan_slide["subtitle"] = "Sort the exact source clues, values, or evidence that help solve the problem."
        plan_slide["tasks"] = source_sort_labels(plan_slide)
        plan_slide["response_prompt"] = "Sort one card and explain how it helps solve or justify the source problem."
        plan_slide["activity_name"] = f"source clue sort {signature}"
        plan_slide["activity_family"] = "sort_classify"
        plan_slide["activity_instructions"] = "Sort the source clues, values, or evidence into the category that best supports the solve path."
        plan_slide["movable_pieces"] = source_drag_sort_pieces(plan_slide)
        plan_slide["answer_check"] = "Each card should help identify a value, a strategy move, or evidence for the source problem."
        return

    if role == "error_analysis":
        steps, fix_text, why_prompt = source_error_analysis_content(plan_slide)
        plan_slide["title"] = "Check the Reasoning"
        plan_slide["subtitle"] = "Find the step, claim, or conclusion that does not match the source problem."
        plan_slide["response_prompt"] = "Which step or claim does not match the source problem, and how would you fix it?"
        plan_slide["error_steps"] = []
        plan_slide["fix_it_text"] = fix_text
        plan_slide["why_prompt"] = why_prompt
        plan_slide["activity_name"] = f"source reasoning check {signature}"
        plan_slide["activity_family"] = "detect_justify"
        plan_slide["activity_instructions"] = "Use the source problem to decide which statement, step, or conclusion needs to be checked or fixed."
        if has_data_analysis_context(workbook_source_blob(plan_slide)):
            plan_slide["movable_pieces"] = ["Best claim", "Wrong claim", "Data evidence", "Explain the fix"]
        else:
            plan_slide["movable_pieces"] = ["Correct step", "Wrong step", "Best evidence", "Explain the fix"]
        plan_slide["answer_check"] = "A strong correction matches the exact source values, setup, and evidence."
        return

    if role == "two_column_compare":
        plan_slide["title"] = "Compare the Source Problems"
        plan_slide["subtitle"] = "Each partner solves one source-aligned problem, then compares the strategy or evidence."
        plan_slide["partner_a_problem"] = primary_problem
        plan_slide["partner_b_problem"] = similar_problem
        plan_slide["compare_frame"] = source_compare_prompt(plan_slide)
        clear_activity_fields(plan_slide)
        return

    if role == "choice_board":
        plan_slide["title"] = "Choose a Source Problem"
        plan_slide["subtitle"] = "Choose one source-aligned problem, solve it, and explain why it fits the skill you are practicing."
        plan_slide["response_prompt"] = "I chose Problem ___ because it lets me show the source strategy clearly."
        plan_slide["choice_paths"] = source_choice_paths(plan_slide)
        clear_activity_fields(plan_slide)
        return

    if role == "collaborative_practice":
        plan_slide["title"] = "Collaborative Source Practice"
        plan_slide["subtitle"] = "Partners solve source-aligned problems and compare what stayed the same or changed."
        plan_slide["partner_a_problem"] = primary_problem
        plan_slide["partner_b_problem"] = similar_problem
        plan_slide["discussion_questions"] = discussion_questions
        clear_activity_fields(plan_slide)
        return

    if role == "independent_practice":
        plan_slide["title"] = "Independent Source Practice"
        plan_slide["subtitle"] = f"Solve three {source_phrase} prompts and show your reasoning clearly."
        plan_slide["independent_problems"] = source_independent_problem_set(plan_slide)
        plan_slide["helpful_tips"] = problem_tip_lines(plan_slide, limit=3)
        clear_activity_fields(plan_slide)
        return

    if role == "turn_and_teach":
        base_payload = source_focused_problem_activity_payload(plan_slide)
        plan_slide["title"] = "Turn + Teach the Strategy"
        plan_slide["subtitle"] = "Solve a source-aligned problem, then teach the strategy with evidence from the problem."
        plan_slide["activity_name"] = base_payload["activity_name"]
        plan_slide["activity_family"] = "sequence_order"
        plan_slide["activity_instructions"] = "Order the source solve-and-teach moves before you explain the strategy to a partner."
        if has_data_analysis_context(workbook_source_blob(plan_slide)):
            plan_slide["movable_pieces"] = ["Read the data display", "Find the key values", "Compare the data", "Teach the evidence"]
        else:
            plan_slide["movable_pieces"] = ["Read the problem", "Choose the strategy", "Show the math", "Teach the check"]
        plan_slide["answer_check"] = "A strong teaching flow should help a partner understand the strategy and the evidence."


def ensure_source_anchored_problem_activity(plan_slide: dict[str, Any]) -> None:
    if plan_slide.get("kind") not in PROBLEM_SOLVING_KINDS:
        return
    if plan_slide.get("kind") == "exit_ticket" and normalize_whitespace(plan_slide.get("template_role", "")) == "tiered_exit":
        clear_activity_fields(plan_slide)
        return
    if problem_render_mode(plan_slide) == PROBLEM_RENDER_MODE_FOCUS:
        clear_activity_fields(plan_slide)
        return
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    if role in SOURCE_ALIGNED_TEMPLATE_ACTIVITY_ROLES:
        retarget_template_role_slide_to_source(plan_slide)
        return
    if (
        not has_activity(plan_slide)
        or normalize_whitespace(plan_slide.get("activity_family", "")) not in SOURCE_ANCHORED_PROBLEM_ACTIVITY_FAMILIES
        or problem_activity_has_generic_markers(plan_slide)
    ):
        payload = source_focused_problem_activity_payload(plan_slide)
        plan_slide["activity_name"] = payload["activity_name"]
        plan_slide["activity_family"] = payload["activity_family"]
        plan_slide["activity_instructions"] = payload["activity_instructions"]
        plan_slide["movable_pieces"] = payload["movable_pieces"]
        plan_slide["answer_check"] = payload["answer_check"]

def split_problem_card_fragments(text: str) -> list[str]:
    cleaned = clean_source_prompt(text)
    if not cleaned:
        return []
    if re.search(r"question:\s*", cleaned, flags=re.IGNORECASE):
        parts = [normalize_whitespace(part) for part in re.split(r"question:\s*", cleaned, flags=re.IGNORECASE) if normalize_whitespace(part)]
        return parts
    pieces = [normalize_whitespace(part) for part in re.split(r"(?<=[?.!])\s+", cleaned) if normalize_whitespace(part)]
    return pieces or [cleaned]


def rewrite_problem_card_text(text: str) -> str:
    cleaned = normalize_whitespace(text)
    if not cleaned:
        return ""
    lowered = cleaned.lower()
    if "how can you determine the length, width, and height of the prism" in lowered:
        return "Determine the prism's length, width, and height from the source model."
    if lowered.startswith("to find the volume of a three-dimensional rectangular prism"):
        return "Use V = l x w x h to solve the rectangular prism problem."
    if "let's use the dimensions of the unit cube" in lowered or "let’s use the dimensions of the unit cube" in lowered:
        return "Use the unit-cube model to determine the prism dimensions."
    width_match = re.search(r"there are\s+(.+?)\s+unit cubes?\s+along the width", cleaned, flags=re.IGNORECASE)
    if width_match:
        return f"The width is {normalize_whitespace(width_match.group(1))} units in the source model."
    height_match = re.search(r"there are\s+(.+?)\s+unit cubes?\s+along the height", cleaned, flags=re.IGNORECASE)
    if height_match:
        return f"The height is {normalize_whitespace(height_match.group(1))} units in the source model."
    length_match = re.search(r"there are\s+(.+?)\s+unit cubes?\s+along the length", cleaned, flags=re.IGNORECASE)
    if length_match:
        return f"The length is {normalize_whitespace(length_match.group(1))} units in the source model."
    if "how do you use the formula to find the volume of the prism" in lowered:
        return "Use the volume formula to find the prism's volume."
    if "let's substitute the dimensions of the prism into the formula" in lowered or "let’s substitute the dimensions of the prism into the formula" in lowered:
        return "Substitute the prism dimensions into V = l x w x h."
    if lowered.startswith("what is the volume of the prism"):
        return "Find the volume of the prism."
    if lowered.startswith("what is the volume of the rectangular prism shown"):
        return "Find the volume of the rectangular prism and explain your answer."
    if "use a formula for volume to relate the volume and dimensions" in lowered and "what is the length of the box" in lowered:
        return "Use V = l x w x h to find the missing length of the box."
    if "divide the volume by the product of the width and height" in lowered:
        return "Divide the volume by width x height to find the missing length."
    if "due to shipping regulations" in lowered and "compare" in lowered:
        return "Compare the box length to the 12-inch shipping rule."
    if "which size and how many boxes would you buy" in lowered:
        return "Choose the popcorn size and number of boxes that make the most sense."
    if "you go to movie theater with your family" in lowered and "popcorn" in lowered:
        return "Use the popcorn box information from the source to compare your options."
    if lowered.startswith("you can use the volume formulas you know"):
        return ""
    return cleaned


def problem_card_priority(text: str) -> int:
    lowered = normalize_whitespace(text).lower()
    if not lowered:
        return 0
    score = problem_prompt_priority(text)
    if any(
        term in lowered
        for term in (
            "use v = l x w x h",
            "determine the prism's length",
            "find the volume of the prism",
            "find the volume of the rectangular prism",
            "choose the popcorn size",
            "compare the box length",
        )
    ):
        score += 8
    if any(
        term in lowered
        for term in (
            "use the popcorn box information",
            "the width is",
            "the height is",
            "the length is",
            "use the unit-cube model",
        )
    ):
        score += 4
    if lowered.startswith("volume of a rectangular prism") or lowered.startswith("summarize:"):
        score -= 18
    return score


def workbook_problem_prompts(plan_slide: dict[str, Any], *, limit: int = 3) -> list[str]:
    prompts: list[str] = []
    ordered_source_prompts = unique_nonempty(
        primary_source_problem_targets(plan_slide, limit=2) + problem_statement_candidates(plan_slide, limit=max(limit + 1, 4)),
        limit=max(limit + 1, 4),
    )
    if ordered_source_prompts:
        prompts.append(f"Solve this source problem: {ordered_source_prompts[0]}")
        if len(ordered_source_prompts) > 1:
            prompts.append(f"Explain or solve this source follow-up: {ordered_source_prompts[1]}")
    else:
        prompts.append(f"Solve this source problem: {rendered_source_problem_statement(plan_slide)}")
    if len(prompts) < 2:
        prompts.append(f"Solve a similar problem: {similar_problem_statement(plan_slide)}")
    prompts.append(f"Write and solve one more: {build_problem_creation_prompt(plan_slide)}")
    prompts.extend(f"Use the source problem language: {prompt}" for prompt in ordered_source_prompts[2:])
    return unique_nonempty(prompts, limit=limit)


def problem_workbook_content(plan_slide: dict[str, Any], *, variant: str = "practice") -> dict[str, Any]:
    explain_prompt = plan_slide.get("response_prompt", "") or "Explain why your answer makes sense."
    source_problem = rendered_source_problem_statement(plan_slide)
    second_source_problem = exact_source_followup_problem_statement(plan_slide, ordinal=1)
    similar_problem = similar_problem_statement(plan_slide)
    similar_phrase = similar_problem_focus_phrase(plan_slide)
    if variant == "guided":
        return {
            "panel_title": "Model + Try",
            "show_problem_cards": True,
            "top_cards": [
                ("Model This", source_problem),
                ("Your Turn", similar_problem),
            ],
            "left_workspace": (
                "Model the Source Problem",
                "Show the setup, labels, and calculations clearly.",
                4,
                PAPER,
            ),
            "right_workspaces": [
                ("Solve a Similar Problem", f"Use the same strategy on your own with a {similar_phrase}.", 2, PALE_GOLD),
                ("Explain + Check", truncate_display_copy(explain_prompt, 112), 2, PALE_BLUE),
            ],
        }
    if variant == "exit":
        return {
            "panel_title": "Final Solve + Prove",
            "show_problem_cards": True,
            "top_cards": [("Solve This", source_problem)],
            "left_workspace": (
                "Solve the Final Check",
                "Show the setup, the work, and a clearly labeled final answer.",
                5,
                PAPER,
            ),
            "right_workspaces": [
                ("Prove It", truncate_display_copy(explain_prompt, 112), 3, PALE_CORAL),
            ],
        }
    return {
        "panel_title": "Solve, Check, Explain",
        "show_problem_cards": True,
        "top_cards": [
            ("P1 Solve This", source_problem),
            ("P2 From Slides" if second_source_problem else "P2 Similar Problem", second_source_problem or similar_problem),
        ],
        "left_workspace": (
            "Solve the Source Problem",
            "Show the setup, labels, and calculations neatly.",
            5,
            PAPER,
        ),
        "right_workspaces": [
            ("Solve a Similar Problem", f"Use the same strategy on a {similar_phrase} and keep the work organized.", 2, PALE_GOLD),
            ("Explain + Check", truncate_display_copy(explain_prompt, 112), 2, PALE_SAGE),
        ],
    }


def add_problem_prompt_card(
    slide: Any,
    *,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    prompt: str,
    accent: RGBColor,
    fill: RGBColor = PAPER,
) -> None:
    card = add_rect(slide, x, y, w, h, fill, line_color=SOFT_LINE, rounded=True)
    card.adjustments[0] = 0.16
    set_shape_heading_body(
        card,
        truncate_text(title, 34),
        trim_dangling_display_text(truncate_display_copy(prompt, 260 if emu_to_inches(int(w)) >= 3.0 else 210)),
        heading_size=11.2,
        body_size=12.0,
        heading_color=accent,
        body_color=INK,
        margin=0.08,
        body_min_size=10.8,
        line_spacing=1.14,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_problem_prompt_stack(
    slide: Any,
    *,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    problems: list[str],
    accent: RGBColor,
    fill: RGBColor = PAPER,
    subtitle: str = "",
    kicker: str = "",
) -> None:
    add_rect(slide, x, y, w, h, fill, line_color=SOFT_LINE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    title_x = x + Inches(0.16)
    title_w = w - Inches(0.32)
    if kicker:
        kicker_w = max(
            int(Inches(1.32)),
            min(int(Inches(1.82)), int(float(Inches(0.92 + max(len(kicker), 12) * 0.048)))),
        )
        add_chip(
            slide,
            x + w - kicker_w - Inches(0.14),
            y + Inches(0.16),
            kicker_w,
            Inches(0.18),
            kicker,
            fill=PALE_GOLD if accent == GOLD else PAPER_WARM,
            accent=accent,
        )
        title_w -= kicker_w + Inches(0.12)
    add_text(
        slide,
        title_x,
        y + Inches(0.16),
        title_w,
        Inches(0.22),
        truncate_text(title, 42),
        size=15.0,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    if subtitle:
        add_text(
            slide,
            x + Inches(0.16),
            y + Inches(0.38),
            w - Inches(0.32),
            Inches(0.16),
            truncate_text(subtitle, 104),
            size=10.8,
            color=MUTED,
            font=FONT_BODY,
            margin=0.01,
        )
    prompt_items = problems[:3] or ["Use the exact source problem language and solve it in the workbook."]
    layout_mode = prompt_stack_layout_mode(prompt_items)
    if layout_mode == "focus":
        prompt_items = prompt_items[:2]
    card_gap = Inches(0.08 if layout_mode == "focus" else 0.10)
    card_top = y + (Inches(0.56) if subtitle and layout_mode == "focus" else Inches(0.58) if subtitle else Inches(0.46) if layout_mode == "focus" else Inches(0.48))
    available_h = h - (Inches(0.66) if subtitle and layout_mode == "focus" else Inches(0.70) if subtitle else Inches(0.58) if layout_mode == "focus" else Inches(0.62))
    min_card_h = Inches(0.46 if layout_mode == "focus" else 0.24 if len(prompt_items) >= 3 else 0.38)
    card_h = max(int(min_card_h), int((available_h - card_gap * max(len(prompt_items) - 1, 0)) / max(len(prompt_items), 1)))
    if layout_mode == "focus":
        prompt_limit = 250 if len(prompt_items) == 1 else 220
        prompt_size = 13.0 if len(prompt_items) == 1 else 12.0
        prompt_min_size = 10.8 if len(prompt_items) == 1 else 10.4
    else:
        prompt_limit = 136 if len(prompt_items) >= 3 else 250 if len(prompt_items) == 1 else 200
        prompt_size = 12.8 if len(prompt_items) == 1 else 11.6 if len(prompt_items) == 2 else 10.9
        prompt_min_size = 10.8 if len(prompt_items) == 1 else 10.4
    for index, prompt in enumerate(prompt_items):
        item_y = card_top + index * (card_h + card_gap)
        inner_fill = PAPER if index == 0 else BG
        inner = add_rect(
            slide,
            x + Inches(0.14),
            item_y,
            w - Inches(0.28),
            card_h,
            inner_fill,
            line_color=SOFT_LINE,
            rounded=True,
        )
        inner.adjustments[0] = 0.16
        set_shape_text(
            inner,
            trim_dangling_display_text(truncate_display_copy(prompt, prompt_limit)),
            size=prompt_size,
            min_size=prompt_min_size,
            color=INK,
            font=FONT_BODY,
            margin=0.08,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )


def render_problem_workbook_panel(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
    accent: RGBColor,
    variant: str = "practice",
) -> None:
    content = problem_workbook_content(plan_slide, variant=variant)
    add_rect(slide, x, y, w, h, PAPER, line_color=accent)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(
        slide,
        x + Inches(0.16),
        y + Inches(0.16),
        w - Inches(0.32),
        Inches(0.22),
        content["panel_title"],
        size=14.7,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    inner_x = x + Inches(0.14)
    inner_w = w - Inches(0.28)
    card_gap = Inches(0.10)
    card_y = y + Inches(0.58)
    card_h = max(int(Inches(0.62)), min(int(Inches(0.86)), int(float(h) * 0.26)))
    top_cards = content["top_cards"]
    show_problem_cards = content.get("show_problem_cards", True)
    if show_problem_cards:
        if len(top_cards) <= 1:
            add_problem_prompt_card(
                slide,
                x=inner_x,
                y=card_y,
                w=inner_w,
                h=card_h,
                title=top_cards[0][0],
                prompt=top_cards[0][1],
                accent=accent,
                fill=PALE_GOLD if accent == GOLD else PAPER,
            )
        else:
            left_card_w = int(float(inner_w - card_gap) * 0.56)
            right_card_w = inner_w - left_card_w - card_gap
            add_problem_prompt_card(
                slide,
                x=inner_x,
                y=card_y,
                w=left_card_w,
                h=card_h,
                title=top_cards[0][0],
                prompt=top_cards[0][1],
                accent=accent,
                fill=PAPER,
            )
            add_problem_prompt_card(
                slide,
                x=inner_x + left_card_w + card_gap,
                y=card_y,
                w=right_card_w,
                h=card_h,
                title=top_cards[1][0],
                prompt=top_cards[1][1],
                accent=accent,
                fill=PALE_GOLD if accent != GOLD else PAPER,
            )
        workspace_y = card_y + card_h + Inches(0.12)
        workspace_h = max(int(Inches(1.0)), h - card_h - Inches(0.84))
    else:
        workspace_y = y + Inches(0.58)
        workspace_h = max(int(Inches(1.2)), h - Inches(0.76))
    workspace_w = int(float(inner_w) * 0.58) if content["right_workspaces"] else inner_w
    explain_x = inner_x + workspace_w + Inches(0.12)
    explain_w = inner_w - workspace_w - Inches(0.12)
    left_title, left_prompt, left_lines, left_fill = content["left_workspace"]
    add_lined_area(
        slide,
        inner_x,
        workspace_y,
        workspace_w,
        workspace_h,
        left_title,
        left_prompt,
        lines=left_lines,
        fill=left_fill,
    )
    if content["right_workspaces"] and explain_w > int(Inches(1.6)):
        right_gap = Inches(0.10)
        if len(content["right_workspaces"]) == 1:
            title, prompt, lines, fill = content["right_workspaces"][0]
            add_lined_area(slide, explain_x, workspace_y, explain_w, workspace_h, title, prompt, lines=lines, fill=fill)
        else:
            top_title, top_prompt, top_lines, top_fill = content["right_workspaces"][0]
            bottom_title, bottom_prompt, bottom_lines, bottom_fill = content["right_workspaces"][1]
            top_h = max(int(Inches(0.94)), int(float(workspace_h - right_gap) * 0.48))
            bottom_h = workspace_h - top_h - right_gap
            add_lined_area(slide, explain_x, workspace_y, explain_w, top_h, top_title, top_prompt, lines=top_lines, fill=top_fill)
            add_lined_area(
                slide,
                explain_x,
                workspace_y + top_h + right_gap,
                explain_w,
                bottom_h,
                bottom_title,
                bottom_prompt,
                lines=bottom_lines,
                fill=bottom_fill,
            )


def add_sentence_starters(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    starters: list[str],
    *,
    accent: RGBColor = CORAL,
) -> None:
    body = "\n".join(f"- {truncate_text(item, 85)}" for item in starters[:5]) or "- I can explain that ..."
    compact = float(h) <= float(Inches(1.0))
    add_card(
        slide,
        x,
        y,
        w,
        h,
        "Sentence Starters",
        body,
        fill=PALE_CORAL if accent == CORAL else PALE_SAGE,
        accent=accent,
        title_size=15.0 if compact else 16.3,
        body_size=12.3 if compact else 14.0,
    )


def add_vocabulary_snapshot(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    vocabulary: list[dict[str, str]],
) -> None:
    compact = float(h) <= float(Inches(1.95))
    add_rect(slide, x, y, w, h, PAPER, line_color=SAGE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = SAGE
    band.line.fill.background()
    add_text(
        slide,
        x + Inches(0.16),
        y + Inches(0.16),
        w - Inches(0.32),
        Inches(0.22),
        "Academic Vocabulary",
        size=14.3,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        margin=0.02,
    )
    max_rows = 2
    rows = vocabulary[:max_rows] or [
        {
            "word": "Math word",
            "definition": "a word to use when you explain your thinking",
            "example": "Use it in a sentence while you solve.",
        }
    ]
    if compact:
        gutter = Inches(0.14)
        col_w = int((float(w) - float(Inches(0.32)) - float(gutter)) / 2)
        row_top = y + Inches(0.48)
        for index, item in enumerate(rows[:2]):
            left = x + Inches(0.16) + index * (col_w + gutter)
            add_text(
                slide,
                left,
                row_top,
                col_w,
                Inches(0.18),
                truncate_text(item.get("word", "") or f"Word {index + 1}", 24),
                size=11.4,
                color=SAGE,
                bold=True,
                font=FONT_HEAD,
                margin=0.01,
            )
            add_text(
                slide,
                left,
                row_top + Inches(0.16),
                col_w,
                Inches(0.18),
                truncate_text(item.get("definition", "") or "lesson word", 42),
                size=10.4,
                color=INK,
                font=FONT_BODY,
                margin=0.01,
            )
    else:
        row_top = y + Inches(0.48)
        for index, item in enumerate(rows[:2]):
            top = row_top + index * Inches(0.56)
            add_text(
                slide,
                x + Inches(0.16),
                top,
                w - Inches(0.32),
                Inches(0.18),
                truncate_text(item.get("word", "") or f"Word {index + 1}", 34),
                size=12.0,
                color=SAGE,
                bold=True,
                font=FONT_HEAD,
                margin=0.01,
            )
            add_text(
                slide,
                x + Inches(0.16),
                top + Inches(0.16),
                w - Inches(0.32),
                Inches(0.26),
                truncate_text(item.get("definition", "") or "a lesson word to use in your explanation", 84),
                size=10.9,
                color=INK,
                font=FONT_BODY,
                margin=0.01,
            )
    add_slot_box(
        slide,
        x + Inches(0.16),
        y + h - Inches(0.42),
        w - Inches(0.32),
        Inches(0.26),
        "My word: ____________________",
        fill=PALE_SAGE,
        accent=SAGE,
    )


def add_vocabulary_row(
    slide: Any,
    *,
    y: int,
    word: str,
    definition: str,
    example: str,
    visual_cue: str,
    asset: dict[str, Any] | None,
) -> None:
    left = Inches(0.78)
    row_h = Inches(0.84)
    word_w = Inches(2.0)
    definition_w = Inches(3.15)
    example_w = Inches(3.45)
    visual_w = Inches(2.8)
    gap = Inches(0.12)
    add_rect(slide, left, y, word_w, row_h, PAPER, line_color=SAGE)
    add_rect(slide, left + word_w + gap, y, definition_w, row_h, PAPER, line_color=LINE)
    add_rect(slide, left + word_w + definition_w + gap * 2, y, example_w, row_h, PAPER, line_color=LINE)
    visual_x = left + word_w + definition_w + example_w + gap * 3
    add_rect(slide, visual_x, y, visual_w, row_h, PAPER, line_color=TEAL)

    add_text(
        slide,
        left + Inches(0.12),
        y + Inches(0.14),
        word_w - Inches(0.24),
        row_h - Inches(0.24),
        truncate_text(word or "Word", 28),
        size=13.2,
        min_size=11.0,
        color=SAGE,
        bold=True,
        font=FONT_HEAD,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        left + word_w + gap + Inches(0.12),
        y + Inches(0.10),
        definition_w - Inches(0.24),
        row_h - Inches(0.20),
        truncate_text(definition or "Student-friendly definition", 110),
        size=11.3,
        min_size=10.4,
        color=INK,
        font=FONT_BODY,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        left + word_w + definition_w + gap * 2 + Inches(0.12),
        y + Inches(0.10),
        example_w - Inches(0.24),
        row_h - Inches(0.20),
        truncate_text(example or "Use this word in a source-based sentence or explanation.", 120),
        size=11.0,
        min_size=10.4,
        color=INK,
        font=FONT_BODY,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    if asset:
        add_image_panel(
            slide,
            visual_x + Inches(0.08),
            y + Inches(0.10),
            Inches(0.82),
            Inches(0.58),
            asset,
            fill=PALE_BLUE,
        )
        add_text(
            slide,
            visual_x + Inches(0.98),
            y + Inches(0.18),
            visual_w - Inches(1.10),
            Inches(0.46),
            truncate_text(visual_cue or "Visual clue from the source image.", 62),
            size=10.6,
            min_size=10.4,
            color=MUTED,
            font=FONT_BODY,
            margin=0.01,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )
    else:
        add_text(
            slide,
            visual_x + Inches(0.10),
            y + Inches(0.10),
            visual_w - Inches(0.20),
            row_h - Inches(0.20),
            truncate_text(visual_cue or "Sketch or label a visual cue here.", 80),
            size=10.7,
            min_size=10.4,
            color=MUTED,
            font=FONT_BODY,
            margin=0.02,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
        )


def add_vocabulary_feature_card(
    slide: Any,
    *,
    x: int,
    y: int,
    w: int,
    h: int,
    word: str,
    definition: str,
    example: str,
    visual_cue: str,
    asset: dict[str, Any] | None,
    accent: RGBColor,
    fill: RGBColor,
) -> None:
    visual_cue_text = trim_dangling_display_text(truncate_display_copy(visual_cue or "Visual cue", 34))
    example_text = trim_dangling_display_text(
        truncate_display_copy(example or visual_cue or "Use this word in the lesson context.", 72)
    )
    add_rect(slide, x, y, w, h, fill, line_color=accent)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(
        slide,
        x + Inches(0.16),
        y + Inches(0.14),
        w - Inches(0.32),
        Inches(0.36),
        truncate_text(word or "Term", 28),
        size=22.0,
        min_size=18.0,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    image_w = Inches(1.42)
    image_h = Inches(0.88)
    image_y = y + Inches(0.50)
    add_rect(slide, x + Inches(0.16), image_y, image_w, image_h, PAPER, line_color=accent)
    if asset:
        add_picture_contain(slide, asset, x + Inches(0.22), image_y + Inches(0.06), image_w - Inches(0.12), image_h - Inches(0.12))
    else:
        add_text(
            slide,
            x + Inches(0.20),
            image_y + Inches(0.24),
            image_w - Inches(0.08),
            Inches(0.40),
            visual_cue_text,
            size=10.6,
            min_size=10.4,
            color=MUTED,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
            margin=0.01,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )
    text_x = x + Inches(1.68)
    text_w = w - Inches(1.80)
    add_text(
        slide,
        text_x,
        image_y,
        text_w,
        Inches(0.40),
        truncate_text(definition or "Student-friendly definition", 112),
        size=12.1,
        min_size=10.9,
        color=INK,
        font=FONT_BODY,
        margin=0.01,
    )
    add_text(
        slide,
        text_x,
        image_y + Inches(0.46),
        text_w,
        Inches(0.16),
        "Source Example",
        size=10.5,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    add_text(
        slide,
        text_x,
        image_y + Inches(0.62),
        text_w,
        Inches(0.24),
        example_text,
        size=11.1,
        min_size=10.4,
        color=MUTED,
        font=FONT_BODY,
        margin=0.01,
    )


def add_checkmark_tracker(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    content_objective: str,
    language_objective: str,
) -> None:
    add_rect(slide, x, y, w, h, PAPER, line_color=TEAL)
    add_text(
        slide,
        x + Inches(0.14),
        y + Inches(0.10),
        Inches(5.6),
        Inches(0.22),
        "Before / After Objective Check-In",
        size=14.3,
        color=TEAL,
        bold=True,
        font=FONT_HEAD,
        margin=0.02,
    )
    add_text(
        slide,
        x + w - Inches(2.22),
        y + Inches(0.12),
        Inches(0.92),
        Inches(0.18),
        "Before",
        size=11.3,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )
    add_text(
        slide,
        x + w - Inches(1.14),
        y + Inches(0.12),
        Inches(0.92),
        Inches(0.18),
        "After",
        size=11.3,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )
    add_text(
        slide,
        x + Inches(0.14),
        y + Inches(0.34),
        w - Inches(0.28),
        Inches(0.16),
        "List each objective one time, then place the checkmarks in the Before and After boxes as your evidence grows.",
        size=10.6,
        color=MUTED,
        font=FONT_BODY,
        margin=0.01,
    )

    rows = [
        ("Content Objective", content_objective or "I can explain the main lesson idea.", PALE_BLUE, TEAL),
        (
            "Language Objective",
            language_objective or "I can explain my strategy with labels, vocabulary, and complete sentences.",
            PALE_GOLD,
            GOLD,
        ),
    ]
    row_y = y + Inches(0.64)
    row_h = Inches(0.56)
    text_box_w = w - Inches(2.62)
    for index, (label, objective_text, fill_color, accent) in enumerate(rows):
        item_y = row_y + index * Inches(0.70)
        add_rect(slide, x + Inches(0.14), item_y, text_box_w, row_h, fill_color, line_color=accent)
        add_text(
            slide,
            x + Inches(0.26),
            item_y + Inches(0.08),
            Inches(1.55),
            Inches(0.14),
            label,
            size=10.4,
            color=accent,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )
        add_text(
            slide,
            x + Inches(1.68),
            item_y + Inches(0.06),
            text_box_w - Inches(1.66),
            Inches(0.34),
            truncate_text(objective_text, 132),
            size=12.0,
            color=INK,
            font=FONT_BODY,
            margin=0.01,
        )
        add_slot_box(slide, x + w - Inches(2.18), item_y + Inches(0.13), Inches(0.86), Inches(0.26), "", fill=PALE_BLUE, accent=TEAL)
        add_slot_box(slide, x + w - Inches(1.10), item_y + Inches(0.13), Inches(0.86), Inches(0.26), "", fill=PALE_GOLD, accent=GOLD)

    add_text(
        slide,
        x + Inches(0.14),
        y + h - Inches(0.38),
        Inches(1.85),
        Inches(0.16),
        "Checkmarks",
        size=10.4,
        color=TEAL,
        bold=True,
        font=FONT_HEAD,
        margin=0.01,
    )
    chip_y = y + h - Inches(0.36)
    for index in range(4):
        chip_x = x + Inches(1.72) + index * Inches(0.48)
        add_chip(slide, chip_x, chip_y, Inches(0.38), Inches(0.18), CHECKMARK_CHIP, fill=PAPER, accent=TEAL)


def style_table(
    table: Any,
    header_fill: RGBColor = NAVY,
    body_fill: RGBColor = PAPER,
    *,
    header_font_size: float = 12.0,
    body_font_size: float = 11.4,
    header_align: Any = PP_ALIGN.CENTER,
    body_align: Any = PP_ALIGN.CENTER,
    column_alignments: list[Any] | None = None,
) -> None:
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_index == 0 else body_fill
            cell.text_frame.word_wrap = True
            cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            cell.text_frame.margin_left = Inches(0.03)
            cell.text_frame.margin_right = Inches(0.03)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)
            for paragraph in cell.text_frame.paragraphs:
                if row_index == 0:
                    paragraph.alignment = header_align
                elif column_alignments and col_index < len(column_alignments):
                    paragraph.alignment = column_alignments[col_index]
                else:
                    paragraph.alignment = body_align
                paragraph.line_spacing = TABLE_TEXT_LINE_SPACING
                for run in paragraph.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(header_font_size if row_index == 0 else body_font_size)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = PAPER if row_index == 0 else INK


def add_table(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    data: list[list[str]],
    *,
    column_widths: list[int] | None = None,
    row_heights: list[int] | None = None,
    header_font_size: float = 12.0,
    body_font_size: float = 11.4,
    header_align: Any = PP_ALIGN.CENTER,
    body_align: Any = PP_ALIGN.CENTER,
    column_alignments: list[Any] | None = None,
) -> Any:
    normalized_rows = [list(row) for row in data if isinstance(row, list)]
    if not normalized_rows:
        normalized_rows = [[""]]
    max_cols = max((len(row) for row in normalized_rows), default=1)
    max_cols = max(1, max_cols)
    for row in normalized_rows:
        if len(row) < max_cols:
            row.extend([""] * (max_cols - len(row)))
        elif len(row) > max_cols:
            del row[max_cols:]

    table = slide.shapes.add_table(len(normalized_rows), max_cols, x, y, w, h).table
    for row_index, row in enumerate(normalized_rows):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = truncate_text(value, 120)
    if column_widths:
        for index, width in enumerate(column_widths[: len(table.columns)]):
            table.columns[index].width = int(width)
    if row_heights:
        for index, height in enumerate(row_heights[: len(table.rows)]):
            table.rows[index].height = int(height)
    style_table(
        table,
        header_font_size=header_font_size,
        body_font_size=body_font_size,
        header_align=header_align,
        body_align=body_align,
        column_alignments=column_alignments,
    )
    return table


def bullet_text(lines: list[str], limit: int = 4, max_len: int = 160) -> str:
    cleaned = [truncate_text(line, max_len) for line in unique_nonempty(lines, limit=limit)]
    return "\n".join(f"- {line}" for line in cleaned)


def has_activity(plan_slide: dict[str, Any]) -> bool:
    return bool(normalize_whitespace(plan_slide.get("activity_name", "")))


def activity_accent(plan_slide: dict[str, Any]) -> RGBColor:
    family = plan_slide.get("activity_family", "")
    return {
        "sort_classify": TEAL,
        "match_pair": GOLD,
        "sequence_order": CORAL,
        "build_construct": NAVY,
        "plot_place": TEAL,
        "detect_justify": CORAL,
        "compare_rank": GOLD,
        "reveal_discuss": SAGE,
    }.get(family, TEAL)


def activity_fill(plan_slide: dict[str, Any]) -> RGBColor:
    family = plan_slide.get("activity_family", "")
    return {
        "sort_classify": PALE_BLUE,
        "match_pair": PALE_GOLD,
        "sequence_order": PALE_CORAL,
        "build_construct": PAPER,
        "plot_place": PALE_BLUE,
        "detect_justify": PALE_CORAL,
        "compare_rank": PALE_GOLD,
        "reveal_discuss": PALE_SAGE,
    }.get(family, PALE_BLUE)


def activity_record_for_slide(plan_slide: dict[str, Any]) -> dict[str, Any] | None:
    activity_name = normalize_whitespace(plan_slide.get("activity_name", ""))
    if not activity_name:
        return None
    return find_activity_record(activity_name, load_activity_library())


def activity_display_title(plan_slide: dict[str, Any]) -> str:
    family = plan_slide.get("activity_family", "")
    kind = plan_slide.get("kind", "")
    activity_name = publisher_copyedit_text(plan_slide.get("activity_name", ""))
    lowered = activity_name.lower()
    source_blob = workbook_source_blob(plan_slide)
    if family == "build_construct":
        if has_data_analysis_context(source_blob):
            return "Build the Analysis Path"
        if any(term in lowered for term in ("equation", "formula")):
            return "Equation Builder"
        if any(term in lowered for term in ("model", "representation")):
            return "Build the Model"
        return "Build the Solve Path"
    if family == "sequence_order":
        return "Step Sort"
    if family == "compare_rank":
        return "Compare the Choices"
    if family == "detect_justify":
        return "Fix It + Explain"
    if family == "match_pair":
        if has_data_analysis_context(source_blob):
            return "Match the Data Evidence"
        return "Match the Vocabulary" if kind == "vocabulary" else "Match + Connect"
    if family == "sort_classify":
        if has_data_analysis_context(source_blob):
            return "Sort the Data Clues"
        return "Sort the Vocabulary" if kind == "vocabulary" else "Sort + Justify"
    if family == "plot_place":
        return "Place the Labels"
    if family == "reveal_discuss":
        return "Reveal + Discuss"
    if not activity_name:
        return "Interactive Practice"
    cleaned = re.sub(r"[-_]+", " ", activity_name)
    return truncate_text(cleaned.title(), 42)


def activity_zone_labels(plan_slide: dict[str, Any], limit: int = 4) -> list[str]:
    spec = flagship_activity_spec(plan_slide)
    if spec:
        labels = FLAGSHIP_ACTIVITY_LAYOUT_ZONE_LABELS.get(spec.get("layoutVariant", ""), [])
        if labels:
            return labels[:limit]
    family = plan_slide.get("activity_family", "")
    blob = workbook_source_blob(plan_slide)
    if plan_slide.get("kind") == "vocabulary" and family == "match_pair":
        return ["Picture clue", "Meaning clue", "Use in math", "Say it"][:limit]
    if plan_slide.get("kind") in PROBLEM_SOLVING_KINDS:
        if family == "build_construct":
            if has_data_analysis_context(blob):
                return ["Choose the key values", "Use the data evidence"]
            if any(term in blob for term in ("which size", "how many", "would you buy", "question:")):
                return ["Solve the choices", "Decision evidence"]
            if any(term in blob for term in ("missing", "unknown", "length of the box", "relate the volume and dimensions")):
                return ["Set up the equation", "Check the missing value"]
            return ["Build the solve path", "Support pieces"]
        if family == "detect_justify":
            return ["Best claim or fix", "Evidence from the source"]
        if family == "compare_rank":
            return ["Strongest fit", "Middle", "Needs more proof"]
        if family == "sequence_order":
            if has_data_analysis_context(blob):
                return ["Step 1", "Step 2", "Step 3", "Step 4"][:limit]
            return ["Step 1", "Step 2", "Step 3"]
    record = activity_record_for_slide(plan_slide)
    if not record:
        return []
    labels = [trim_dangling_display_text(truncate_text(item, 28)) for item in record.get("piece_labels", [])]
    return unique_nonempty(labels, limit=limit)


def activity_pieces(plan_slide: dict[str, Any], limit: int = 4) -> list[str]:
    candidates = (
        plan_slide.get("movable_pieces", [])
        or plan_slide.get("tasks", [])
        or plan_slide.get("bullets", [])
        or plan_slide.get("sentence_starters", [])
    )
    return [trim_dangling_display_text(truncate_display_copy(item, 60)) for item in unique_nonempty(candidates, limit=limit)]


def flagship_activity_support_panel_height(plan_slide: dict[str, Any], available_h: int) -> int:
    spec = flagship_activity_spec(plan_slide)
    if not spec:
        return 0
    supports = unique_nonempty(spec.get("supports", []), limit=3)
    if not supports or float(available_h) < float(Inches(2.15)):
        return 0
    return int(Inches(0.46 if len(supports) <= 2 else 0.62))


def render_flagship_activity_supports(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
    accent: RGBColor,
) -> None:
    spec = flagship_activity_spec(plan_slide)
    if not spec:
        return
    supports = unique_nonempty(spec.get("supports", []), limit=3)
    if not supports:
        return
    add_card(
        slide,
        x,
        y,
        w,
        h,
        "Activity Support",
        bullet_text([truncate_text(item, 76) for item in supports], limit=3, max_len=76),
        fill=PAPER,
        accent=accent,
        title_size=11.8,
        body_size=10.6,
    )


def add_chip(slide: Any, x: int, y: int, w: int, h: int, text: str, *, fill: RGBColor = PAPER, accent: RGBColor = NAVY) -> None:
    chip = add_rect(slide, x, y, w, h, fill, line_color=accent, rounded=True)
    chip.adjustments[0] = 0.22
    chip_font_size = 11.0 if float(h) <= float(Inches(0.2)) else 12.8
    set_shape_text(
        chip,
        trim_dangling_display_text(truncate_display_copy(text, 52)),
        size=chip_font_size,
        color=INK,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
        margin=0.01,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_draggable_piece_box(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    text: str,
    *,
    fill: RGBColor = PAPER,
    accent: RGBColor = NAVY,
) -> None:
    piece_box = add_rect(slide, x, y, w, h, fill, line_color=accent, rounded=True)
    piece_box.adjustments[0] = 0.12
    set_shape_text(
        piece_box,
        trim_dangling_display_text(truncate_display_copy(text, 72)),
        size=11.7,
        min_size=10.4,
        color=INK,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
        margin=0.02,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_activity_header(slide: Any, x: int, y: int, w: int, title: str, family: str, accent: RGBColor) -> None:
    family_label = ACTIVITY_FAMILY_LABELS.get(family, "Interactive Activity")
    if float(w) <= float(Inches(3.0)):
        add_text(
            slide,
            x + Inches(0.02),
            y,
            w - Inches(0.04),
            Inches(0.22),
            truncate_text(title or ACTIVITY_FAMILY_SHORT_LABELS.get(family, "Activity"), 24),
            size=11.2,
            color=accent,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )
        return
    family_width = max(
        float(Inches(1.68)),
        min(float(Inches(2.16)), float(Inches(0.92 + max(len(family_label), 12) * 0.055))),
    )
    title_width = max(float(Inches(1.9)), float(w) - family_width - float(Inches(0.20)))
    add_text(
        slide,
        x + Inches(0.04),
        y,
        int(title_width),
        Inches(0.24),
        truncate_text(title, 80),
        size=13.4,
        color=accent,
        bold=True,
        font=FONT_HEAD,
        margin=0.02,
    )
    pill_x = x + w - int(family_width)
    pill = add_rect(slide, pill_x, y - Inches(0.01), int(family_width), Inches(0.28), PALE_NAVY, line_color=PALE_NAVY)
    pill.adjustments[0] = 0.25
    add_text(
        slide,
        pill_x,
        y + Inches(0.02),
        int(family_width),
        Inches(0.18),
        truncate_text(family_label, 22),
        size=10.4,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )


def add_slot_box(slide: Any, x: int, y: int, w: int, h: int, label: str, *, fill: RGBColor = PAPER, accent: RGBColor = LINE) -> None:
    slot = add_rect(slide, x, y, w, h, fill, line_color=accent, rounded=True)
    set_shape_text(
        slot,
        trim_dangling_display_text(truncate_display_copy(label, 48)),
        size=11.9,
        color=MUTED,
        font=FONT_BODY,
        align=PP_ALIGN.CENTER,
        margin=0.03,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_simple_grid(slide: Any, x: int, y: int, w: int, h: int, *, accent: RGBColor) -> None:
    add_rect(slide, x, y, w, h, PAPER, line_color=accent, rounded=False)
    for index in range(1, 4):
        x_pos = x + int(float(w) * index / 4)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_pos, y, x_pos, y + h)
        line.line.color.rgb = LINE
        line.line.width = Pt(1)
    for index in range(1, 3):
        y_pos = y + int(float(h) * index / 3)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y_pos, x + w, y_pos)
        line.line.color.rgb = LINE
        line.line.width = Pt(1)


def add_premium_meta_chips(slide: Any, plan_slide: dict[str, Any]) -> None:
    chips: list[tuple[str, RGBColor]] = []
    if normalize_whitespace(plan_slide.get("context_anchor", "")):
        chips.append((f"Context: {plan_slide['context_anchor']}", TEAL))
    if normalize_whitespace(plan_slide.get("practice_phase", "")):
        chips.append((plan_slide["practice_phase"], GOLD))
    if not chips:
        return
    x = Inches(0.5)
    y = Inches(1.82)
    for label, accent in chips[:2]:
        width = Inches(2.9) if len(label) > 22 else Inches(2.15)
        add_chip(slide, x, y, width, Inches(0.18), label, fill=PAPER, accent=accent)
        x += width + Inches(0.12)


def discussion_panel_title(plan_slide: dict[str, Any]) -> str:
    role = normalize_whitespace(plan_slide.get("template_role", ""))
    kind = plan_slide.get("kind", "")
    if role == "two_column_compare":
        return "Compare + Discuss"
    if role == "collaborative_practice":
        return "Talk It Through"
    if role == "turn_and_teach" or normalize_whitespace(plan_slide.get("premium_layout", "")) == "turn_and_teach":
        return "Turn + Teach"
    if kind == "reflection":
        return "Reflect + Share"
    return "Talk It Through"


def discussion_panel_body(plan_slide: dict[str, Any], *, max_questions: int = 2, include_prompt: bool = True) -> str:
    lines: list[str] = []
    prompt = normalize_whitespace(plan_slide.get("partner_prompt", ""))
    if include_prompt and prompt:
        lines.append(truncate_display_copy(prompt, 86 if max_questions <= 2 else 104))
    for index, question in enumerate(plan_slide.get("discussion_questions", [])[:max_questions], start=1):
        cleaned = clean_discussion_question(question)
        if not cleaned:
            continue
        lines.append(f"Q{index}. {truncate_display_copy(cleaned, 58 if max_questions <= 2 else 72)}")
    return "\n".join(lines[: max_questions + (1 if include_prompt and prompt else 0)])


def add_discussion_panel(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
    fill: RGBColor,
    accent: RGBColor,
    max_questions: int = 2,
    include_prompt: bool = True,
) -> None:
    add_card(
        slide,
        x,
        y,
        w,
        h,
        discussion_panel_title(plan_slide),
        discussion_panel_body(plan_slide, max_questions=max_questions, include_prompt=include_prompt),
        fill=fill,
        accent=accent,
        title_size=15.0,
        body_size=11.6 if float(h) <= float(Inches(1.7)) else 12.2,
    )


def render_premium_panel(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
    accent: RGBColor,
) -> bool:
    layout = normalize_whitespace(plan_slide.get("premium_layout", ""))
    if not layout and normalize_whitespace(plan_slide.get("partner_prompt", "")):
        layout = "turn_and_teach"
    if not layout or layout not in PREMIUM_BLOCKING_LAYOUTS | {"turn_and_teach"}:
        return False

    title = plan_slide.get("premium_title", "") or PREMIUM_FEATURE_LABELS.get(layout, "Premium Notebook Move")
    body = plan_slide.get("premium_text", "")
    items = plan_slide.get("premium_items", [])
    table = plan_slide.get("premium_table", [])

    if layout == "error_analysis":
        left_w = int(float(w) * 0.56)
        right_w = w - left_w - Inches(0.16)
        error_body_lines = items[:2] or ([body] if body else [])
        add_card(
            slide,
            x,
            y,
            left_w,
            h,
            title,
            "\n".join(f"- {trim_dangling_display_text(truncate_text(item, 68))}" for item in error_body_lines),
            fill=PALE_CORAL,
            accent=CORAL,
            title_size=15.2,
            body_size=12.6,
        )
        add_lined_area(
            slide,
            x + left_w + Inches(0.16),
            y,
            right_w,
            h,
            "Fix It",
            items[2] if len(items) > 2 else body or "Explain the mistake and write the corrected reasoning.",
            lines=4 if float(h) >= float(Inches(2.0)) else 3,
            fill=PAPER,
        )
        return True

    if layout == "multi_representation":
        add_rect(slide, x, y, w, h, PAPER, line_color=TEAL)
        add_text(
            slide,
            x + Inches(0.16),
            y + Inches(0.12),
            w - Inches(0.32),
            Inches(0.20),
            title,
            size=14.6,
            color=TEAL,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )
        if body:
            add_text(
                slide,
                x + Inches(0.16),
                y + Inches(0.34),
                w - Inches(0.32),
                Inches(0.24),
                body,
                size=11.0,
                color=MUTED,
                font=FONT_BODY,
                margin=0.01,
            )
        box_y = y + Inches(0.68)
        box_h = max(int((h - Inches(0.84)) / 2), int(Inches(0.56)))
        box_w = int((w - Inches(0.24)) / 2)
        labels = items[:4] or ["Representation A", "Representation B", "Reasoning", "Explain"]
        for index, label in enumerate(labels[:4]):
            row = index // 2
            col = index % 2
            add_slot_box(
                slide,
                x + col * (box_w + Inches(0.12)),
                box_y + row * (box_h + Inches(0.12)),
                box_w,
                box_h,
                label,
                fill=PALE_BLUE if index < 2 else PALE_GOLD,
                accent=TEAL if index < 2 else GOLD,
            )
        return True

    if layout == "strategy_comparison":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Compare the available strategies so you can choose the strongest fit for the lesson.",
            fill=PAPER,
            accent=GOLD,
            title_size=15.0,
            body_size=11.5,
        )
        if table:
            add_table(
                slide,
                x + Inches(0.14),
                y + Inches(0.72),
                w - Inches(0.28),
                h - Inches(0.86),
                table,
            )
        return True

    if layout == "evidence_ladder":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Use a claim, evidence, and reasoning ladder so your explanation feels provable.",
            fill=PAPER,
            accent=CORAL,
            title_size=15.0,
            body_size=11.5,
        )
        ladder_y = y + Inches(0.78)
        ladder_h = max(int((h - Inches(1.02)) / 3), int(Inches(0.42)))
        rung_labels = items[1:4] or [
            "Clue from the representation",
            "Math evidence that supports the claim",
            "Reasoning that proves the claim",
        ]
        claim_text = items[0] if items else "What is the strongest claim you can defend?"
        add_slot_box(
            slide,
            x + Inches(0.14),
            ladder_y - Inches(0.34),
            w - Inches(0.28),
            Inches(0.24),
            claim_text,
            fill=PALE_GOLD,
            accent=GOLD,
        )
        for index, label in enumerate(rung_labels[:3]):
            row_y = ladder_y + index * (ladder_h + Inches(0.08))
            fill = PALE_CORAL if index == 0 else PAPER if index == 1 else PALE_SAGE
            accent_color = CORAL if index == 0 else TEAL if index == 1 else SAGE
            add_slot_box(
                slide,
                x + Inches(0.14),
                row_y,
                w - Inches(0.28),
                ladder_h,
                label,
                fill=fill,
                accent=accent_color,
            )
        return True

    if layout == "decision_tree":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Follow the path before you commit to a strategy.",
            fill=PAPER,
            accent=TEAL,
            title_size=15.0,
            body_size=11.5,
        )
        step_y = y + Inches(0.72)
        step_h = max(int((h - Inches(0.88)) / 3), int(Inches(0.42)))
        labels = items[:3] or ["Step 1", "Step 2", "Step 3"]
        connector_x = x + int(float(w) / 2)
        for index, label in enumerate(labels):
            box_y = step_y + index * (step_h + Inches(0.08))
            add_slot_box(slide, x + Inches(0.14), box_y, w - Inches(0.28), step_h, label, fill=PALE_BLUE, accent=TEAL)
            if index < len(labels) - 1:
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    connector_x,
                    box_y + step_h,
                    connector_x,
                    box_y + step_h + Inches(0.08),
                )
                line.line.color.rgb = TEAL
                line.line.width = Pt(1.2)
        return True

    if layout == "create_your_own":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Create a new example that still matches the lesson target.",
            fill=PALE_GOLD,
            accent=GOLD,
            title_size=15.2,
            body_size=11.6,
        )
        box_y = y + Inches(0.76)
        box_h = max(int((h - Inches(0.92)) / 2), int(Inches(0.52)))
        box_w = int((w - Inches(0.24)) / 2)
        labels = items[:4] or ["Create", "Show the Math", "Check It", "Explain It"]
        for index, label in enumerate(labels[:4]):
            row = index // 2
            col = index % 2
            add_slot_box(
                slide,
                x + col * (box_w + Inches(0.12)),
                box_y + row * (box_h + Inches(0.10)),
                box_w,
                box_h,
                label,
                fill=PAPER,
                accent=GOLD if row == 0 else TEAL,
            )
        return True

    if layout == "real_world_transfer":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Carry the lesson into a new situation and explain what still works.",
            fill=PALE_SAGE,
            accent=SAGE,
            title_size=15.0,
            body_size=11.5,
        )
        add_slot_box(
            slide,
            x + Inches(0.14),
            y + Inches(0.74),
            w - Inches(0.28),
            Inches(0.34),
            items[0] if items else "Start from the source idea.",
            fill=PALE_BLUE,
            accent=TEAL,
        )
        box_y = y + Inches(1.18)
        box_h = max(int((h - Inches(1.66)) / 2), int(Inches(0.44)))
        box_w = int((w - Inches(0.40)) / 2)
        add_slot_box(
            slide,
            x + Inches(0.14),
            box_y,
            box_w,
            box_h,
            items[1] if len(items) > 1 else "What stays true?",
            fill=PAPER,
            accent=SAGE,
        )
        add_slot_box(
            slide,
            x + Inches(0.26) + box_w,
            box_y,
            box_w,
            box_h,
            items[2] if len(items) > 2 else "What changes?",
            fill=PALE_GOLD,
            accent=GOLD,
        )
        add_lined_area(
            slide,
            x + Inches(0.14),
            box_y + box_h + Inches(0.12),
            w - Inches(0.28),
            max(h - (box_y - y) - box_h - Inches(0.26), int(Inches(0.44))),
            "Transfer Justification",
            items[3] if len(items) > 3 else "How would you justify the transfer to a new situation?",
            lines=2,
            fill=PAPER,
        )
        return True

    if layout == "mastery_tracker":
        add_card(
            slide,
            x,
            y,
            w,
            h,
            title,
            body or "Rate your confidence, name what you understand, and prove one idea before you finish.",
            fill=PALE_SAGE,
            accent=SAGE,
            title_size=15.0,
            body_size=11.8,
        )
        chip_y = y + Inches(0.72)
        chip_w = int((w - Inches(0.52)) / 3)
        for index, label in enumerate(("Not yet", "Getting there", "Ready to teach")):
            add_chip(slide, x + Inches(0.14) + index * (chip_w + Inches(0.06)), chip_y, chip_w, Inches(0.18), label, fill=PAPER, accent=SAGE)
        add_slot_box(slide, x + Inches(0.14), y + Inches(1.05), w - Inches(0.28), Inches(0.28), items[1] if len(items) > 1 else "What I understand now", fill=PAPER, accent=SAGE)
        add_slot_box(slide, x + Inches(0.14), y + Inches(1.42), w - Inches(0.28), Inches(0.28), items[2] if len(items) > 2 else "What I still need help with", fill=PAPER, accent=TEAL)
        if float(h) > float(Inches(1.9)):
            add_lined_area(
                slide,
                x + Inches(0.14),
                y + Inches(1.80),
                w - Inches(0.28),
                h - Inches(1.94),
                "Prove-It Mini Check",
                items[3] if len(items) > 3 else "Write one short proof, check, or example that shows what you know.",
                lines=2,
                fill=PAPER,
            )
        return True

    if layout == "turn_and_teach":
        if not plan_slide.get("discussion_questions"):
            plan_slide["discussion_questions"] = generated_discussion_questions(plan_slide)
        add_discussion_panel(
            slide,
            plan_slide=plan_slide,
            x=x,
            y=y,
            w=w,
            h=h,
            fill=PALE_BLUE,
            accent=TEAL,
            max_questions=2,
            include_prompt=True,
        )
        return True

    return False


def premium_layout_uses_full_spread(plan_slide: dict[str, Any]) -> bool:
    return normalize_whitespace(plan_slide.get("premium_layout", "")) in PREMIUM_FULL_SPREAD_LAYOUTS


def render_premium_feature_spread(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    accent: RGBColor,
    image_lookup: dict[int, dict[str, Any]] | None = None,
) -> bool:
    layout = normalize_whitespace(plan_slide.get("premium_layout", ""))
    if layout not in PREMIUM_FULL_SPREAD_LAYOUTS:
        return False

    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=accent,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)

    problem_cards = problem_display_cards(plan_slide, variant="practice" if plan_slide.get("kind") != "exit_ticket" else "exit")
    interactive_slide = activity_slide_for_render(plan_slide)

    if layout == "evidence_ladder":
        add_problem_prompt_stack(
            slide,
            x=Inches(0.5),
            y=Inches(2.0),
            w=Inches(4.35),
            h=Inches(3.95),
            title="Source Problem",
            problems=problem_cards,
            accent=GOLD,
            fill=PAPER,
            subtitle="Pull the strongest evidence before you write.",
            kicker="Evidence Path",
        )
        add_card(
            slide,
            Inches(0.5),
            Inches(6.08),
            Inches(4.35),
            Inches(0.70),
            "Use the Ladder",
            truncate_display_copy(
                plan_slide.get("response_prompt")
                or plan_slide.get("secondary_text")
                or "Start with a clear claim, then add the evidence and reasoning that prove it.",
                118,
            ),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=13.6,
            body_size=10.6,
        )
        render_premium_panel(
            slide,
            plan_slide=plan_slide,
            x=Inches(5.08),
            y=Inches(2.0),
            w=Inches(7.67),
            h=Inches(2.72),
            accent=CORAL,
        )
        add_lined_area(
            slide,
            Inches(5.08),
            Inches(4.96),
            Inches(7.67),
            Inches(1.82),
            "Evidence-Based Explanation",
            plan_slide.get("response_prompt")
            or "Write the final explanation that connects your claim, evidence, and reasoning.",
            lines=3,
            fill=PALE_GOLD,
        )
        return True

    asset = choose_image_asset(plan_slide, image_lookup or {})
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(3.78),
        Inches(2.18),
        asset,
        fill=PALE_GOLD,
        label=plan_slide.get("image_caption") or "Source context",
    )
    render_premium_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(4.52),
        y=Inches(2.0),
        w=Inches(8.23),
        h=Inches(2.18),
        accent=SAGE,
    )
    add_problem_prompt_stack(
        slide,
        x=Inches(0.5),
        y=Inches(4.46),
        w=Inches(5.05),
        h=Inches(2.32),
        title="Source to New Situation",
        problems=problem_cards,
        accent=GOLD,
        fill=PAPER,
        subtitle="Hold onto what still works, then stretch the idea.",
        kicker="Transfer Check",
    )
    if has_activity(interactive_slide):
        render_activity_board(
            slide,
            plan_slide=interactive_slide,
            x=Inches(5.80),
            y=Inches(4.46),
            w=Inches(6.95),
            h=Inches(2.32),
        )
    else:
        add_lined_area(
            slide,
            Inches(5.80),
            Inches(4.46),
            Inches(6.95),
            Inches(2.32),
            "Transfer Workspace",
            plan_slide.get("response_prompt")
            or plan_slide.get("secondary_text")
            or "Explain how the lesson idea changes when you move it into a new real-world situation.",
            lines=4,
            fill=PAPER,
        )
    return True


def render_activity_board(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
) -> None:
    accent = activity_accent(plan_slide)
    fill = activity_fill(plan_slide)
    support_h = flagship_activity_support_panel_height(plan_slide, h)
    if support_h:
        render_flagship_activity_supports(
            slide,
            plan_slide=plan_slide,
            x=x,
            y=y,
            w=w,
            h=support_h,
            accent=accent,
        )
        y += support_h + Inches(0.08)
        h = max(int(h - support_h - Inches(0.08)), int(Inches(0.72)))
    pieces = activity_pieces(plan_slide, limit=4) or ["Card 1", "Card 2", "Card 3", "Card 4"]
    instruction_text = plan_slide.get("activity_instructions") or "Use the interactive pieces to engage with the source lesson."
    answer_text = normalize_whitespace(plan_slide.get("answer_check", ""))
    add_rect(slide, x, y, w, h, fill, line_color=accent)
    add_activity_header(
        slide,
        x + Inches(0.14),
        y + Inches(0.10),
        w - Inches(0.28),
        activity_display_title(plan_slide),
        plan_slide.get("activity_family", ""),
        accent,
    )
    micro_layout = float(w) <= float(Inches(3.0)) or float(h) <= float(Inches(1.15))
    if micro_layout:
        compact_pieces = pieces[:2] or ["Move", "Explain"]
        add_text(
            slide,
            x + Inches(0.14),
            y + Inches(0.34),
            w - Inches(0.28),
            Inches(0.12),
            trim_dangling_display_text(truncate_text(instruction_text or answer_text or "Move a piece and explain.", 48)),
            size=10.4,
            color=MUTED,
            font=FONT_BODY,
            margin=0.01,
        )
        piece_gap = Inches(0.08)
        piece_w = int((float(w - Inches(0.28)) - float(piece_gap) * (len(compact_pieces) - 1)) / max(len(compact_pieces), 1))
        piece_y = y + h - Inches(0.24)
        for index, piece in enumerate(compact_pieces):
            add_draggable_piece_box(
                slide,
                x + Inches(0.14) + index * (piece_w + piece_gap),
                piece_y,
                piece_w,
                Inches(0.18),
                piece,
                fill=PAPER,
                accent=accent,
            )
        return
    if float(h) <= float(Inches(0.9)):
        compact_pieces = pieces[:2] or ["Move", "Explain"]
        if float(h) >= float(Inches(0.62)):
            add_text(
                slide,
                x + Inches(0.14),
                y + Inches(0.30),
                w - Inches(0.28),
                Inches(0.12),
                trim_dangling_display_text(
                    truncate_text(instruction_text or answer_text or "Move a piece and explain your choice.", 78)
                ),
                size=10.4,
                color=MUTED,
                font=FONT_BODY,
                margin=0.01,
            )
        piece_gap = Inches(0.08)
        piece_w = int((float(w - Inches(0.28)) - float(piece_gap) * (len(compact_pieces) - 1)) / max(len(compact_pieces), 1))
        piece_y = y + h - Inches(0.22)
        for index, piece in enumerate(compact_pieces):
            add_draggable_piece_box(
                slide,
                x + Inches(0.14) + index * (piece_w + piece_gap),
                piece_y,
                piece_w,
                Inches(0.18),
                piece,
                fill=PAPER,
                accent=accent,
            )
        return
    if float(h) <= float(Inches(1.9)):
        add_text(
            slide,
            x + Inches(0.14),
            y + Inches(0.38),
            w - Inches(0.28),
            Inches(0.18),
            trim_dangling_display_text(truncate_text(instruction_text, 96)),
            size=11.8,
            color=MUTED,
            font=FONT_BODY,
            margin=0.02,
        )
        compact_pieces = pieces[:3] or ["Move", "Match", "Explain"]
        chip_w = int(float(w - Inches(0.52)) / len(compact_pieces))
        for index, piece in enumerate(compact_pieces):
            add_draggable_piece_box(
                slide,
                x + Inches(0.14) + index * chip_w,
                y + h - Inches(0.30),
                chip_w - Inches(0.04),
                Inches(0.20),
                piece,
                fill=PAPER,
                accent=accent,
            )
        return
    add_text(
        slide,
        x + Inches(0.14),
            y + Inches(0.40),
            w - Inches(0.28),
            Inches(0.22),
            trim_dangling_display_text(truncate_text(instruction_text, 104)),
            size=12.6,
            color=MUTED,
            font=FONT_BODY,
            margin=0.02,
        )
    inner_y = y + Inches(0.92)
    if answer_text:
        add_text(
            slide,
            x + Inches(0.14),
            y + Inches(0.66),
            w - Inches(0.28),
            Inches(0.20),
            f"Check: {trim_dangling_display_text(truncate_text(answer_text, 96))}",
            size=10.9,
            color=MUTED,
            font=FONT_BODY,
            margin=0.01,
        )
        inner_y = y + Inches(1.00)

    family = plan_slide.get("activity_family", "")
    zone_labels = activity_zone_labels(plan_slide)
    inner_x = x + Inches(0.14)
    inner_w = w - Inches(0.28)
    chip_rows = max(1, (min(len(pieces), 4) + 1) // 2)
    chip_y = y + h - Inches(0.10) - Inches(0.22) * chip_rows
    workspace_h = max(int(Inches(0.56)), int(chip_y - inner_y - Inches(0.10)))
    tall_slot_h = max(int(Inches(0.64)), min(int(Inches(0.92)), workspace_h))
    short_slot_h = max(int(Inches(0.50)), min(int(Inches(0.64)), workspace_h))

    if family == "sort_classify":
        if len(zone_labels) >= 4:
            box_w = int(float(inner_w - Inches(0.12)) / 2)
            box_h = max(int(Inches(0.44)), int((workspace_h - Inches(0.10)) / 2))
            for index, label in enumerate(zone_labels[:4]):
                row = index // 2
                col = index % 2
                add_slot_box(
                    slide,
                    inner_x + col * (box_w + Inches(0.12)),
                    inner_y + row * (box_h + Inches(0.10)),
                    box_w,
                    box_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
        else:
            labels = zone_labels[:2] or ["Category A", "Category B"]
            box_w = int(float(inner_w - Inches(0.16)) / 2)
            add_slot_box(slide, inner_x, inner_y, box_w, tall_slot_h, labels[0], fill=PAPER, accent=accent)
            add_slot_box(slide, inner_x + box_w + Inches(0.16), inner_y, box_w, tall_slot_h, labels[1], fill=PAPER, accent=accent)
    elif family == "match_pair":
        slot_h = max(int(Inches(0.52)), tall_slot_h - Inches(0.10))
        if len(zone_labels) >= 4:
            box_w = int(float(inner_w - Inches(0.12)) / 2)
            box_h = max(int(Inches(0.44)), int((workspace_h - Inches(0.10)) / 2))
            for index, label in enumerate(zone_labels[:4]):
                row = index // 2
                col = index % 2
                add_slot_box(
                    slide,
                    inner_x + col * (box_w + Inches(0.12)),
                    inner_y + row * (box_h + Inches(0.10)),
                    box_w,
                    box_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
        else:
            labels = zone_labels[:2] or ["Source cards", "Match here"]
            box_w = int(float(inner_w - Inches(0.30)) / 2)
            add_slot_box(slide, inner_x, inner_y, box_w, slot_h, labels[0], fill=PAPER, accent=accent)
            add_slot_box(slide, inner_x + box_w + Inches(0.30), inner_y, box_w, slot_h, labels[1], fill=PAPER, accent=accent)
        if workspace_h > int(Inches(0.70)):
            add_text(
                slide,
                inner_x,
                inner_y + slot_h + Inches(0.04),
                inner_w,
                Inches(0.16),
                "Draw or drag to connect each pair.",
                size=10.6,
                color=MUTED,
                font=FONT_BODY,
                align=PP_ALIGN.CENTER,
                margin=0.02,
            )
    elif family == "sequence_order":
        slot_w = int(float(inner_w - Inches(0.20)) / 3)
        for index in range(3):
            add_slot_box(
                slide,
                inner_x + index * (slot_w + Inches(0.10)),
                inner_y,
                slot_w,
                short_slot_h,
                f"Step {index + 1}",
                fill=PAPER,
                accent=accent,
            )
    elif family == "build_construct":
        labels = zone_labels[:2] or ["Build zone", "Use these pieces"]
        add_slot_box(slide, inner_x, inner_y, int(float(inner_w) * 0.58), tall_slot_h, labels[0], fill=PAPER, accent=accent)
        add_slot_box(
            slide,
            inner_x + int(float(inner_w) * 0.62),
            inner_y,
            int(float(inner_w) * 0.38) - Inches(0.02),
            tall_slot_h,
            labels[1] if len(labels) > 1 else "Use these pieces",
            fill=PAPER,
            accent=accent,
        )
    elif family == "plot_place":
        grid_w = int(float(inner_w) * 0.58)
        add_slot_box(slide, inner_x, inner_y, grid_w, tall_slot_h, "Place labels on the source model", fill=PAPER, accent=accent)
        add_slot_box(
            slide,
            inner_x + grid_w + Inches(0.14),
            inner_y,
            inner_w - grid_w - Inches(0.14),
            max(int(Inches(0.30)), tall_slot_h // 2 - Inches(0.04)),
            zone_labels[0] if zone_labels else "What to label",
            fill=PAPER,
            accent=accent,
        )
        add_slot_box(
            slide,
            inner_x + grid_w + Inches(0.14),
            inner_y + max(int(Inches(0.34)), tall_slot_h // 2 + Inches(0.02)),
            inner_w - grid_w - Inches(0.14),
            max(int(Inches(0.30)), tall_slot_h // 2 - Inches(0.04)),
            zone_labels[1] if len(zone_labels) > 1 else "Why it matters",
            fill=PAPER,
            accent=accent,
        )
    elif family == "detect_justify":
        if len(zone_labels) >= 4:
            box_w = int(float(inner_w - Inches(0.12)) / 2)
            box_h = max(int(Inches(0.44)), int((workspace_h - Inches(0.10)) / 2))
            for index, label in enumerate(zone_labels[:4]):
                row = index // 2
                col = index % 2
                add_slot_box(
                    slide,
                    inner_x + col * (box_w + Inches(0.12)),
                    inner_y + row * (box_h + Inches(0.10)),
                    box_w,
                    box_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
        else:
            labels = zone_labels[:2] or ["Error or evidence", "Fix it / justify it"]
            add_slot_box(slide, inner_x, inner_y, int(float(inner_w) * 0.42), tall_slot_h, labels[0], fill=PAPER, accent=accent)
            add_slot_box(
                slide,
                inner_x + int(float(inner_w) * 0.46),
                inner_y,
                inner_w - int(float(inner_w) * 0.46),
                tall_slot_h,
                labels[1],
                fill=PAPER,
                accent=accent,
            )
    elif family == "compare_rank":
        if len(zone_labels) >= 4:
            box_w = int(float(inner_w - Inches(0.12)) / 2)
            box_h = max(int(Inches(0.44)), int((workspace_h - Inches(0.10)) / 2))
            for index, label in enumerate(zone_labels[:4]):
                row = index // 2
                col = index % 2
                add_slot_box(
                    slide,
                    inner_x + col * (box_w + Inches(0.12)),
                    inner_y + row * (box_h + Inches(0.10)),
                    box_w,
                    box_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
        else:
            slot_w = int(float(inner_w - Inches(0.20)) / 3)
            labels = zone_labels[:3] or ["High", "Middle", "Low"]
            for index, label in enumerate(labels):
                add_slot_box(
                    slide,
                    inner_x + index * (slot_w + Inches(0.10)),
                    inner_y,
                    slot_w,
                    short_slot_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
    elif family == "reveal_discuss":
        if len(zone_labels) >= 4:
            box_w = int(float(inner_w - Inches(0.12)) / 2)
            box_h = max(int(Inches(0.44)), int((workspace_h - Inches(0.10)) / 2))
            for index, label in enumerate(zone_labels[:4]):
                row = index // 2
                col = index % 2
                add_slot_box(
                    slide,
                    inner_x + col * (box_w + Inches(0.12)),
                    inner_y + row * (box_h + Inches(0.10)),
                    box_w,
                    box_h,
                    label,
                    fill=PAPER,
                    accent=accent,
                )
        else:
            slot_w = int(float(inner_w - Inches(0.20)) / 3)
            for index in range(3):
                add_slot_box(
                    slide,
                    inner_x + index * (slot_w + Inches(0.10)),
                    inner_y,
                    slot_w,
                    short_slot_h,
                    zone_labels[index] if index < len(zone_labels) else f"Reveal {index + 1}",
                    fill=PAPER,
                    accent=accent,
                )
    else:
        add_slot_box(slide, inner_x, inner_y, inner_w, tall_slot_h, "Interactive workspace", fill=PAPER, accent=accent)

    chip_w = int(float(inner_w - Inches(0.30)) / 2)
    for index, piece in enumerate(pieces[:4]):
        chip_x = inner_x + (index % 2) * (chip_w + Inches(0.10))
        row_y = chip_y + Inches(0.22) * (index // 2)
        add_draggable_piece_box(slide, chip_x, row_y, chip_w, Inches(0.20), piece, fill=PAPER, accent=accent)


def choose_image_asset(plan_slide: dict[str, Any], image_lookup: dict[int, dict[str, Any]]) -> dict[str, Any] | None:
    image_source_slide = plan_slide.get("image_source_slide", 0)
    if image_source_slide and image_source_slide in image_lookup:
        return image_lookup[image_source_slide]
    for slide_number in plan_slide.get("source_slide_numbers", []):
        if slide_number in image_lookup:
            return image_lookup[slide_number]
    return None


def choose_image_assets(plan_slide: dict[str, Any], image_lookup: dict[int, dict[str, Any]], limit: int = 4) -> list[dict[str, Any]]:
    assets: list[dict[str, Any]] = []
    seen: set[int] = set()
    image_source_slide = int(plan_slide.get("image_source_slide", 0) or 0)
    if image_source_slide and image_source_slide in image_lookup:
        assets.append(image_lookup[image_source_slide])
        seen.add(image_source_slide)
    for slide_number in plan_slide.get("source_slide_numbers", []):
        slide_number = int(slide_number)
        if slide_number in seen or slide_number not in image_lookup:
            continue
        assets.append(image_lookup[slide_number])
        seen.add(slide_number)
        if len(assets) >= limit:
            break
    return assets[:limit]


def should_merge_fragment(previous: str, current: str) -> bool:
    prev = normalize_whitespace(previous)
    curr = normalize_whitespace(current)
    if not prev or not curr:
        return False
    if prev.endswith((".", "?", "!")):
        return False
    if prev.endswith("…"):
        return True
    if prev.split()[-1].lower() in DANGLING_ENDING_WORDS:
        return True
    if curr[:1].islower() or len(curr.split()) <= 3:
        return True
    return False


def merge_fragmented_items(items: Iterable[str], limit: int | None = None) -> list[str]:
    merged: list[str] = []
    for item in items:
        cleaned = normalize_whitespace(item)
        if not cleaned:
            continue
        if merged and should_merge_fragment(merged[-1], cleaned):
            merged[-1] = trim_dangling_display_text(normalize_whitespace(f"{merged[-1]} {cleaned}"))
        else:
            merged.append(cleaned)
        if limit is not None and len(merged) >= limit:
            break
    return merged


def source_problem_cards(plan_slide: dict[str, Any], *, limit: int = 4, max_len: int = 180) -> list[str]:
    cards: list[str] = []
    seen: list[str] = []
    raw_items = merge_fragmented_items(
        list(plan_slide.get("source_problem_cards", []))
        + list(plan_slide.get("tasks", []))
        + list(plan_slide.get("bullets", []))
        + [plan_slide.get("primary_text", "")],
        limit=limit * 3,
    )
    candidates: list[tuple[int, str]] = []
    for raw_item in raw_items:
        raw_cleaned = clean_source_prompt(raw_item)
        fragments = [raw_cleaned] if raw_cleaned else []
        fragments.extend(split_problem_card_fragments(raw_item))
        for fragment in fragments:
            cleaned = rewrite_problem_card_text(fragment)
            if not cleaned:
                continue
            if is_dense_numeric_table_text(cleaned):
                continue
            if is_generic_slide_text(cleaned):
                continue
            if cleaned.lower() in {plan_slide.get("title", "").lower(), plan_slide.get("section", "").lower()}:
                continue
            if cleaned == cleaned.title() and not is_problem_like_text(cleaned) and len(cleaned.split()) <= 5:
                continue
            lowered = cleaned.lower()
            if any(lowered == prior or lowered in prior or prior in lowered for prior in seen):
                continue
            seen.append(lowered)
            candidates.append((problem_card_priority(cleaned), cleaned))
    for _score, candidate in sorted(candidates, key=lambda item: (-item[0], len(item[1]))):
        cards.append(trim_dangling_display_text(truncate_display_copy(candidate, max_len)))
        if len(cards) >= limit:
            break
    return cards


def render_exact_cover_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    session_label: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    set_background(slide, PAPER_WARM)
    add_cover_background_decor(slide)
    add_tag(slide, Inches(0.82), Inches(0.90), session_label.upper(), GOLD)
    add_text(
        slide,
        Inches(0.82),
        Inches(1.54),
        Inches(5.8),
        Inches(2.12),
        truncate_text(plan_slide.get("title", ""), 120),
        size=31.0,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        line_spacing=0.92,
    )
    subtitle = first_distinct_text(
        [
            plan_slide.get("primary_text", ""),
            f"{plan_slide.get('unit_label', '')}  {plan_slide.get('standard_label', '')}",
        ],
        excluded=[plan_slide.get("title", ""), session_label],
    ) or "Track the key idea, explain your reasoning, and show your strategy clearly."
    add_card(
        slide,
        Inches(0.82),
        Inches(4.36),
        Inches(2.60),
        Inches(1.10),
        "Unit",
        truncate_text(plan_slide.get("unit_label", "") or "Source lesson", 54),
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=13.0,
        body_size=12.0,
    )
    add_card(
        slide,
        Inches(3.62),
        Inches(4.36),
        Inches(2.98),
        Inches(1.10),
        "Standard",
        truncate_text(plan_slide.get("standard_label", "") or "Lesson standard", 64),
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.0,
        body_size=11.9,
    )
    add_text(
        slide,
        Inches(0.82),
        Inches(3.98),
        Inches(5.5),
        Inches(0.34),
        truncate_text(subtitle, 180),
        size=14.3,
        color=TEAL,
        font=FONT_BODY,
        margin=0.01,
    )
    title_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(5.74), Inches(2.05), Inches(0.06))
    title_rule.fill.solid()
    title_rule.fill.fore_color.rgb = GOLD
    title_rule.line.fill.background()
    add_image_panel(
        slide,
        Inches(7.46),
        Inches(1.08),
        Inches(4.62),
        Inches(2.72),
        choose_image_asset(plan_slide, image_lookup),
        fill=PAPER,
        label=plan_slide.get("image_caption", "") or "Lesson visual",
    )
    toc_text = "\n".join(f"- {truncate_text(item, 34)}" for item in plan_slide.get("toc_items", [])[:6]) or "- Preview the lesson visuals\n- Notice the key quantities\n- Use the notebook to explain your thinking"
    add_card(
        slide,
        Inches(7.15),
        Inches(4.26),
        Inches(2.42),
        Inches(1.58),
        "What's Inside",
        toc_text,
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.1,
        body_size=11.0,
    )
    add_card(
        slide,
        Inches(9.84),
        Inches(4.26),
        Inches(2.92),
        Inches(1.58),
        "Notebook Promise",
        truncate_text(plan_slide.get("primary_text", "") or "Use the source lesson and solve in the notebook.", 112),
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=13.1,
        body_size=11.1,
    )
    add_name_bar(slide)
def render_exact_learning_objectives_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide.get("subtitle", "") or "My Goal Today",
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(6.0),
        Inches(1.44),
        "Content Objective",
        f"{plan_slide.get('primary_text', '')}\n{plan_slide.get('content_objective_b', '')}",
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=14.6,
        body_size=13.0,
    )
    add_card(
        slide,
        Inches(6.78),
        Inches(2.0),
        Inches(5.97),
        Inches(1.44),
        "Language Objective",
        plan_slide.get("secondary_text", ""),
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=14.6,
        body_size=12.7,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(3.70),
        Inches(2.52),
        Inches(0.84),
        "Standard",
        plan_slide.get("standard_label", ""),
        fill=PAPER,
        accent=NAVY,
        title_size=12.8,
        body_size=9.6,
    )
    add_card(
        slide,
        Inches(3.28),
        Inches(3.70),
        Inches(9.47),
        Inches(1.18),
        "Today's 4-Step Path",
        "Move through the notebook in order so each page builds the next part of the lesson.",
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=13.2,
        body_size=11.0,
    )
    for index, step in enumerate(plan_slide.get("session_map_steps", [])[:4] or ["Be Curious", "Vocabulary", "Try It", "Discuss"]):
        add_chip(
            slide,
            Inches(3.56) + index * Inches(2.22),
            Inches(5.10),
            Inches(1.88),
            Inches(0.36),
            step,
            fill=PAPER if index % 2 == 0 else PAPER_WARM,
            accent=TEAL if index < 2 else GOLD,
        )
    add_card(
        slide,
        Inches(0.5),
        Inches(5.54),
        Inches(12.25),
        Inches(1.14),
        "Language Move",
        plan_slide.get("language_frame_support", "") or plan_slide.get("secondary_text", ""),
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.0,
        body_size=11.2,
    )


def render_exact_be_curious_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    notice_kernels, wonder_kernels = be_curious_sentence_kernels(plan_slide)
    vocab_items = be_curious_vocabulary_items(plan_slide, session_plan)
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(4.15),
        Inches(4.35),
        choose_image_asset(plan_slide, image_lookup),
        fill=PALE_BLUE,
        label=plan_slide.get("image_caption", "") or "Notice + Wonder image",
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.08),
        Inches(3.70),
        Inches(0.64),
        plan_slide.get("reference_panel_title", "Bridge"),
        plan_slide.get("bridge_sentence", "") or plan_slide.get("response_prompt", "") or f"📐 {plan_slide.get('prior_formula', '')}",
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=12.6,
        body_size=10.8,
    )
    add_lined_area(
        slide,
        Inches(5.05),
        Inches(2.0),
        Inches(3.62),
        Inches(3.28),
        "Notice",
        be_curious_panel_prompt(
            " ".join(plan_slide.get("notice_lines", [])),
            notice_kernels,
            fallback=plan_slide.get("primary_text", "") or "Record details you notice in the image before solving.",
        ),
        lines=4,
        fill=PAPER,
    )
    add_lined_area(
        slide,
        Inches(9.13),
        Inches(2.0),
        Inches(3.62),
        Inches(3.28),
        "Wonder",
        be_curious_panel_prompt(
            plan_slide.get("wonder_prompt", "") or plan_slide.get("secondary_text", ""),
            wonder_kernels,
            fallback="Write a question or prediction the image makes you think about.",
        ),
        lines=4,
        fill=PALE_BLUE,
    )
    if has_activity(plan_slide):
        add_vocabulary_snapshot(
            slide,
            Inches(5.05),
            Inches(5.58),
            Inches(4.95),
            Inches(1.10),
            vocab_items,
        )
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(10.24),
            y=Inches(5.58),
            w=Inches(2.51),
            h=Inches(1.10),
        )
    else:
        add_vocabulary_snapshot(
            slide,
            Inches(5.05),
            Inches(5.58),
            Inches(7.70),
            Inches(1.10),
            vocab_items,
        )


def render_exact_vocabulary_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    vocab_items = plan_slide.get("vocabulary", [])[:4]
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.54),
        "Word | Definition | Example | Visual",
        plan_slide.get("primary_text", "") or "Keep the words, examples, and visual clues together while you solve.",
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.4,
        body_size=11.5,
    )
    table_data = [["Word", "Definition", "Example", "Visual"]]
    for item in vocab_items:
        table_data.append(
            [
                truncate_text(display_term_label(item.get("word", "")), 22),
                truncate_text(item.get("definition", ""), 52),
                truncate_text(item.get("example", ""), 52),
                truncate_text(item.get("visual_cue", ""), 36),
            ]
        )
    add_table(
        slide,
        Inches(0.5),
        Inches(2.78),
        Inches(12.25),
        Inches(3.38),
        table_data,
        column_widths=[Inches(2.10), Inches(3.34), Inches(3.42), Inches(3.39)],
        row_heights=[Inches(0.44)] + [Inches(0.72)] * max(len(table_data) - 1, 0),
        header_font_size=11.8,
        body_font_size=11.2,
        body_align=PP_ALIGN.LEFT,
        column_alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(6.32),
        Inches(12.25),
        Inches(0.50),
        plan_slide.get("reference_flow_title", "Reference Tool"),
        "\n".join(plan_slide.get("reference_flow_lines", [])),
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=12.6,
        body_size=10.7,
    )


def render_exact_vocabulary_activity_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.56),
        "Word Match",
        plan_slide.get("primary_text", "") or "Move each word card to the best clue.",
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.8,
        body_size=11.9,
    )
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.82),
        Inches(5.02),
        Inches(3.16),
        choose_image_asset(plan_slide, image_lookup),
        fill=PALE_BLUE,
        label=plan_slide.get("image_caption") or "Source vocabulary visual",
        picture_padding=Inches(0.18),
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(6.18),
        Inches(5.02),
        Inches(0.60),
        "Say It",
        plan_slide.get("secondary_text", "") or "Say: '___ means ___.'",
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=12.6,
        body_size=11.4,
    )
    render_activity_board(
        slide,
        plan_slide=plan_slide,
        x=Inches(5.76),
        y=Inches(2.82),
        w=Inches(6.99),
        h=Inches(3.96),
    )


def render_exact_guided_practice_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    source_problem = primary_source_problem_targets(plan_slide, limit=1)
    source_problem_text = source_problem[0] if source_problem else plan_slide.get("primary_text", "")
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.62),
        "Source Problem",
        source_problem_text or plan_slide.get("context_hook", ""),
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.0,
        body_size=11.3,
    )
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.90),
        Inches(4.02),
        Inches(2.56),
        choose_image_asset(plan_slide, image_lookup),
        fill=PALE_BLUE,
        label=plan_slide.get("image_caption", "") or "Labeled source model",
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(5.72),
        Inches(4.02),
        Inches(0.66),
        plan_slide.get("formula_label", "Formula"),
        plan_slide.get("formula_text", ""),
        fill=PALE_NAVY,
        accent=TEAL,
        title_size=12.5,
        body_size=12.0,
    )
    add_table(
        slide,
        Inches(4.82),
        Inches(2.90),
        Inches(3.54),
        Inches(2.88),
        plan_slide.get("guided_table", [["A", "B", "C"], ["___", "___", "___"]]),
    )
    add_lined_area(
        slide,
        Inches(8.62),
        Inches(2.90),
        Inches(4.13),
        Inches(2.88),
        "TWR Frame",
        "\n".join(plan_slide.get("twr_frames", [])) or "Think: ___.\nWrite: ___.\nReason: ___.",
        lines=4,
        fill=PALE_BLUE,
    )
    add_lined_area(
        slide,
        Inches(4.82),
        Inches(5.96),
        Inches(7.93),
        Inches(0.72),
        "Student Work Area",
        plan_slide.get("response_prompt", "") or "Use the source problem, the reference table, and the TWR frame to show your thinking.",
        lines=2,
        fill=PAPER,
    )
REFERENCE_EXACT_NAV_ITEMS = ("EduWonderLab", "Be Curious", "Vocabulary", "Try It", "Practice", "Discuss")
REFERENCE_EXACT_SECTION_COLORS = {
    "cover": NAVY,
    "be_curious": TEAL,
    "vocabulary": GOLD,
    "worked_example": TEAL,
    "practice": CORAL,
    "quick_review": SAGE,
}


def reference_exact_session_key(session_label: str) -> str:
    return "session_2" if "2" in str(session_label) else "session_1"


def reference_exact_source_numbers(plan_slide: dict[str, Any]) -> list[int]:
    numbers: list[int] = []
    for key in ("source_numbers", "problem_numbers", "examples", "source_slide_numbers"):
        values = plan_slide.get(key) or []
        if isinstance(values, int):
            values = [values]
        for value in values:
            try:
                number = int(value)
            except (TypeError, ValueError):
                continue
            if number not in numbers:
                numbers.append(number)
    return numbers


REFERENCE_EXACT_SECTION_TITLES = {
    "objectives + session map",
    "be curious",
    "vocabulary",
    "try it: guided problem",
    "interactive activity",
    "find the error",
    "partner practice: volume of rectangular prism",
    "partner practice: volume of a rectangular prism",
}

REFERENCE_EXACT_GENERIC_PROBLEMS = {
    "determine the prism's length, width, and height from the source model.",
    "use the measurements from the source problem to determine the volume.",
}

REFERENCE_EXACT_ROLE_FALLBACKS = {
    "be_curious": "Is it always true? A cube is a rectangular prism.",
    "guided": "What strategies have you used in the past to find volume?",
    "practice": "How can you use this information to find the volume of the rectangular prism?",
    "error": "How can you use this information to find the volume of the rectangular prism?",
}

REFERENCE_EXACT_ROLE_PATTERNS = {
    "be_curious": (r"Is it always true\??\s*A cube is a rectangular prism\.?",),
    "guided": (r"What strategies have you used in the past to find volume\?",),
    "practice": (r"How can you use this information to find the volume of the rectangular prism\?",),
    "error": (r"How can you use this information to find the volume of the rectangular prism\?",),
}

REFERENCE_EXACT_FRACTIONAL_PRISM_OBJECTIVES = (
    "I can determine the volume of a right rectangular prism with fractional edge lengths by packing it with unit cubes.",
    "I can show that volume is the same as multiplying the edge lengths of a rectangular prism.",
    "I can use formulas V = l × w × h and V = B × h to find volumes of right rectangular prisms.",
)


def reference_exact_source_texts(deck: dict[str, Any]) -> list[str]:
    texts: list[str] = []
    for source_slide in deck.get("slides", []):
        if not isinstance(source_slide, dict):
            continue
        for key in ("title", "text", "body", "raw_text", "content", "speaker_notes", "notes", "text_joined"):
            value = source_slide.get(key)
            if isinstance(value, str):
                text = normalize_whitespace(value.replace("\u200b", " "))
                if text:
                    texts.append(text)
        for key in ("text_blocks", "texts", "shapes"):
            fragments = source_slide.get(key) or []
            if not isinstance(fragments, list):
                continue
            for item in fragments:
                if isinstance(item, str):
                    text = normalize_whitespace(item.replace("\u200b", " "))
                elif isinstance(item, dict):
                    text = normalize_whitespace(str(item.get("text", "")).replace("\u200b", " "))
                else:
                    text = ""
                if text:
                    texts.append(text)
    return texts


def reference_exact_source_blob(deck: dict[str, Any]) -> str:
    return " ".join(reference_exact_source_texts(deck))


def reference_exact_is_fractional_prism_deck(deck: dict[str, Any]) -> bool:
    blob = reference_exact_source_blob(deck).lower()
    return (
        "rectangular prism" in blob
        and "fractional edge lengths" in blob
        and ("volume formulas" in blob or "volume formula" in blob)
    )


def reference_exact_is_generic_text(text: str) -> bool:
    lower = normalize_whitespace(str(text)).lower()
    return (
        lower in REFERENCE_EXACT_SECTION_TITLES
        or lower in REFERENCE_EXACT_GENERIC_PROBLEMS
        or lower.startswith("determine the prism's length")
    )


def reference_exact_extract_prompt(deck: dict[str, Any], role: str) -> str:
    for pattern in REFERENCE_EXACT_ROLE_PATTERNS.get(role, ()):
        for text in reference_exact_source_texts(deck):
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                return normalize_whitespace(match.group(0))
    return REFERENCE_EXACT_ROLE_FALLBACKS.get(role, "")


def reference_exact_learning_target(deck: dict[str, Any]) -> str:
    targets: list[str] = []
    for text in reference_exact_source_texts(deck):
        if "I can" not in text or "Learning Target" not in text:
            continue
        cleaned = re.sub(r"^Learning Targets?\s*", "", text, flags=re.IGNORECASE)
        cleaned = normalize_whitespace(cleaned)
        for match in re.finditer(r"\bI can\b.*?(?=(?:\s+\bI can\b)|$)", cleaned, flags=re.IGNORECASE):
            statement = normalize_whitespace(match.group(0))
            if not statement:
                continue
            if not statement.endswith((".", "?", "!")):
                statement = f"{statement}."
            if statement not in targets:
                targets.append(statement)
    if reference_exact_is_fractional_prism_deck(deck):
        return " ".join(REFERENCE_EXACT_FRACTIONAL_PRISM_OBJECTIVES)
    if targets:
        return " ".join(targets[:4])
    return ""


def reference_exact_source_standard(deck: dict[str, Any]) -> str:
    for text in reference_exact_source_texts(deck):
        match = re.search(r"\b\d+(?:\.[A-Z0-9]+){1,3}\b", text)
        if match:
            return match.group(0)
    blob = reference_exact_source_blob(deck).lower()
    if "rectangular prism" in blob and "volume" in blob:
        return "6.G.2"
    return ""


def reference_exact_title(plan_slide: dict[str, Any], deck: dict[str, Any]) -> str:
    for candidate in (deck.get("lesson_title", ""), deck.get("title", ""), plan_slide.get("title", "")):
        text = normalize_whitespace(str(candidate).replace("\u200b", " "))
        text = re.sub(r"^Session\s+\d+\s+", "", text, flags=re.IGNORECASE)
        if text and not reference_exact_is_generic_text(text):
            return text
    return "Student Notebook"


def reference_exact_unit_title(deck: dict[str, Any]) -> str:
    return normalize_whitespace(deck.get("unit_title", "") or deck.get("unit", "") or "Volume of 3-D Figures")


def reference_exact_standard_text(deck: dict[str, Any]) -> str:
    source_standard = reference_exact_source_standard(deck)
    if source_standard:
        return source_standard
    standard = first_standard_text(deck)
    match = re.search(r"\b\d+(?:\.[A-Z]+)?\.[A-Z]?\d+\b", standard)
    return match.group(0) if match else (standard or "6.G.2")


def reference_exact_formula_text(deck: dict[str, Any], session_label: str) -> str:
    formula = normalize_formula_text(formula_for_session(deck, reference_exact_session_key(session_label)))
    if not formula or formula.lower() == "show every step.":
        return "V = l × w × h"
    if formula.lower().startswith("volume ="):
        return "V = l × w × h"
    return formula.replace(" x ", " × ")


def reference_exact_objective(plan_slide: dict[str, Any], deck: dict[str, Any]) -> str:
    source_target = reference_exact_learning_target(deck)
    if source_target:
        return source_target
    candidates = [
        plan_slide.get("learning_target", ""),
        plan_slide.get("objective", ""),
        plan_slide.get("primary_text", ""),
        deck.get("objective", ""),
        deck.get("summary", ""),
    ]
    for candidate in candidates:
        text = normalize_whitespace(str(candidate))
        if text and not reference_exact_is_generic_text(text):
            return truncate_text(text, 245)
    return "I can determine the volume of a rectangular prism using today's formula and explain my reasoning."


def reference_exact_problem_text(plan_slide: dict[str, Any], deck: dict[str, Any] | None = None, role: str = "") -> str:
    source_problem = primary_source_problem_targets(plan_slide, limit=1)
    if source_problem:
        text = normalize_whitespace(source_problem[0])
        if text and not reference_exact_is_generic_text(text):
            return text
    if deck and role:
        prompt = reference_exact_extract_prompt(deck, role)
        if prompt:
            return truncate_text(prompt, 190)
    for key in ("primary_text", "context_hook", "directions", "secondary_text", "title"):
        text = normalize_whitespace(str(plan_slide.get(key, "")))
        if text and not reference_exact_is_generic_text(text):
            return text
    if role:
        return REFERENCE_EXACT_ROLE_FALLBACKS.get(role, "Use the source prompt to solve.")
    return "Use the measurements from the source problem to determine the volume."


def reference_exact_context_text(plan_slide: dict[str, Any], deck: dict[str, Any] | None = None, role: str = "") -> str:
    if deck and role:
        prompt = reference_exact_extract_prompt(deck, role)
        if prompt:
            return truncate_text(prompt, 190)
    for key in ("context_hook", "secondary_text", "notes", "directions"):
        text = normalize_whitespace(str(plan_slide.get(key, "")))
        if text and not reference_exact_is_generic_text(text):
            return truncate_text(text, 190)
    if role:
        return REFERENCE_EXACT_ROLE_FALLBACKS.get(role, "Study the model and record what you notice.")
    return "Study the model, labels, and measurements. Record what you notice before solving."


def reference_exact_vocab_rows(deck: dict[str, Any], plan_slide: dict[str, Any]) -> list[list[str]]:
    source_numbers = reference_exact_source_numbers(plan_slide)
    vocab_deck = deck
    if not vocab_deck.get("keyword_candidates"):
        source_records = source_slides_from_numbers(vocab_deck, source_numbers)
        keyword_candidates = source_term_candidates(source_records or vocab_deck.get("slides", []), limit=8)
        if not keyword_candidates:
            keyword_candidates = [
                str(item.get("word", ""))
                for item in plan_slide.get("vocabulary", [])
                if isinstance(item, dict) and item.get("word")
            ]
        vocab_deck = {**vocab_deck, "keyword_candidates": keyword_candidates}
    vocab = session_esol_vocabulary(vocab_deck, source_numbers, limit=5)
    rows: list[list[str]] = [["Word", "Definition", "Example", "Visual"]]
    for item in vocab[:5]:
        rows.append(
            [
                truncate_text(item.get("word", ""), 20),
                truncate_text(item.get("definition", ""), 74),
                truncate_text(item.get("example", ""), 48),
                truncate_text(item.get("visual_cue", ""), 38),
            ]
        )
    fallback = [
        ["Volume", "Amount of space inside a 3-D figure.", "Count cubic units.", "inside space"],
        ["Rectangular Prism", "A box-shaped solid with rectangular faces.", "A cereal box.", "box shape"],
        ["Cube", "A prism with all edges the same length.", "A number cube.", "equal edges"],
        ["Unit Cube", "A cube with side lengths of 1 unit.", "One cubic unit.", "1 × 1 × 1"],
        ["Dimensions", "Length, width, and height measurements.", "8 by 3 by 4.", "l, w, h"],
    ]
    while len(rows) < 6:
        rows.append(fallback[len(rows) - 1])
    return rows[:6]


def add_reference_nav(slide: Any, active: str, accent: RGBColor) -> None:
    set_background(slide, BG)
    left = Inches(0.42)
    top = Inches(0.20)
    item_widths = [Inches(1.84), Inches(1.74), Inches(1.92), Inches(1.28), Inches(1.50), Inches(1.42)]
    x = left
    for item, w in zip(REFERENCE_EXACT_NAV_ITEMS, item_widths):
        is_active = item.lower() in active.lower() or active.lower() in item.lower()
        fill = accent if is_active else PAPER
        text_color = RGBColor(255, 255, 255) if is_active else NAVY
        add_rect(slide, x, top, w, Inches(0.42), fill, line_color=accent if is_active else SOFT_LINE, rounded=True)
        add_text(
            slide,
            x + Inches(0.06),
            top + Inches(0.08),
            w - Inches(0.12),
            Inches(0.23),
            item,
            size=10.6,
            bold=item == "EduWonderLab",
            color=text_color,
            align=PP_ALIGN.CENTER,
            margin=0.0,
            min_size=10.4,
        )
        x += w + Inches(0.10)


def add_reference_footer(slide: Any, session_label: str, formula_text: str, deck: dict[str, Any]) -> None:
    footer = f"{session_label} | {formula_text} | Unit: {reference_exact_unit_title(deck)}"
    add_footer_bar(slide, footer)


def add_reference_page_title(
    slide: Any,
    title: str,
    *,
    subtitle: str = "",
    accent: RGBColor = NAVY,
    y: int | None = None,
    size: float = 22.5,
) -> None:
    top = y or Inches(0.72)
    title_h = Inches(0.50 if size >= 24 else 0.44)
    add_text(slide, Inches(0.55), top, Inches(8.92), title_h, title, size=size, bold=True, color=NAVY, min_size=16.0)
    if subtitle:
        add_text(slide, Inches(0.58), top + Inches(0.50), Inches(8.9), Inches(0.30), subtitle, size=10.8, color=MUTED, min_size=10.4)
    add_rect(slide, Inches(0.55), top + Inches(0.84), Inches(2.35), Inches(0.055), accent, line_color=accent, rounded=False)


def add_reference_box(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    body: str = "",
    *,
    fill: RGBColor = PAPER,
    accent: RGBColor = NAVY,
    title_size: float = 11.6,
    body_size: float = 10.6,
    body_color: RGBColor = INK,
) -> None:
    add_rect(slide, x, y, w, h, fill, line_color=SOFT_LINE, rounded=True)
    add_text(slide, x + Inches(0.14), y + Inches(0.11), w - Inches(0.28), Inches(0.23), title, size=title_size, bold=True, color=accent, min_size=10.4)
    if body:
        add_text(
            slide,
            x + Inches(0.14),
            y + Inches(0.41),
            w - Inches(0.28),
            h - Inches(0.52),
            body,
            size=body_size,
            color=body_color,
            min_size=10.4,
            margin=0.02,
        )


def add_reference_rule(slide: Any, x: int, y: int, w: int, *, color: RGBColor = SOFT_LINE, height_in: float = 0.012) -> None:
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, max(int(w), 1), Inches(height_in))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def add_reference_field(slide: Any, x: int, y: int, w: int, label: str, *, color: RGBColor = MUTED) -> None:
    add_text(slide, x, y, Inches(0.75), Inches(0.22), label, size=10.4, bold=True, color=color, margin=0.0, min_size=10.4)
    add_reference_rule(slide, x + Inches(0.82), y + Inches(0.20), w - Inches(0.82), height_in=0.012)


def add_reference_lines(slide: Any, x: int, y: int, w: int, count: int, *, gap: float = 0.34) -> None:
    for index in range(count):
        yy = y + Inches(0.16 + gap * index)
        add_reference_rule(slide, x, yy, w, height_in=0.010)


def add_reference_prism_visual(slide: Any, x: int, y: int, w: int, h: int, *, accent: RGBColor = TEAL) -> None:
    back = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.58), y + Inches(0.24), w - Inches(1.10), h - Inches(0.76))
    back.fill.solid()
    back.fill.fore_color.rgb = RGBColor(232, 245, 247)
    back.line.color.rgb = accent
    front = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.25), y + Inches(0.60), w - Inches(1.10), h - Inches(0.76))
    front.fill.solid()
    front.fill.fore_color.rgb = RGBColor(247, 253, 252)
    front.line.color.rgb = accent
    for start_x, start_y in (
        (x + Inches(0.25), y + Inches(0.60)),
        (x + w - Inches(0.85), y + Inches(0.60)),
        (x + Inches(0.25), y + h - Inches(0.16)),
        (x + w - Inches(0.85), y + h - Inches(0.16)),
    ):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x,
            start_y,
            start_x + Inches(0.33),
            start_y - Inches(0.36),
        )
        connector.line.color.rgb = accent
    add_text(slide, x + Inches(0.42), y + h - Inches(0.04), Inches(1.42), Inches(0.24), "length", size=10.4, color=MUTED, align=PP_ALIGN.CENTER, margin=0.0, min_size=10.4)
    add_text(slide, x + w - Inches(1.42), y + Inches(0.30), Inches(1.24), Inches(0.24), "height", size=10.4, color=MUTED, align=PP_ALIGN.RIGHT, margin=0.0, min_size=10.4)
    add_text(slide, x + w - Inches(1.76), y + h - Inches(0.43), Inches(1.28), Inches(0.24), "width", size=10.4, color=MUTED, align=PP_ALIGN.CENTER, margin=0.0, min_size=10.4)


def render_reference_exact_cover(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    is_fractional_prism = reference_exact_is_fractional_prism_deck(deck)
    formula_display = f"{formula}\nV = B × h" if is_fractional_prism and "V = B" not in formula else formula
    path_body = "Be Curious -> Vocabulary -> Try It -> Practice -> Discuss"
    title = reference_exact_title(plan_slide, deck)
    accent = REFERENCE_EXACT_SECTION_COLORS["cover"]
    add_reference_nav(slide, "EduWonderLab", accent)
    add_reference_page_title(slide, title, subtitle=f"Standard {reference_exact_standard_text(deck)}", accent=accent, size=25.5)
    add_reference_field(slide, Inches(8.95), Inches(0.90), Inches(3.65), "Name")
    add_reference_field(slide, Inches(8.95), Inches(1.26), Inches(3.65), "Date")
    add_reference_field(slide, Inches(8.95), Inches(1.62), Inches(3.65), "Class")
    if is_fractional_prism:
        goal_body = "\n".join(REFERENCE_EXACT_FRACTIONAL_PRISM_OBJECTIVES)
        unit_body = f"{reference_exact_unit_title(deck)}\nBe Curious -> Vocabulary -> Try It\nPractice -> Discuss"
        add_reference_box(
            slide,
            Inches(0.58),
            Inches(1.92),
            Inches(7.84),
            Inches(1.88),
            "Today's Goals",
            goal_body,
            fill=PAPER,
            accent=TEAL,
            title_size=12.6,
            body_size=10.4,
        )
        add_reference_box(
            slide,
            Inches(8.68),
            Inches(1.92),
            Inches(3.82),
            Inches(1.88),
            "Unit + Path",
            unit_body,
            fill=PALE_GOLD,
            accent=GOLD,
            title_size=12.6,
            body_size=10.8,
        )
    else:
        add_reference_box(
            slide,
            Inches(0.58),
            Inches(1.92),
            Inches(4.26),
            Inches(1.72),
            "Today's Goal",
            reference_exact_objective(plan_slide, deck),
            fill=PAPER,
            accent=TEAL,
            title_size=12.6,
            body_size=10.4,
        )
        add_reference_box(
            slide,
            Inches(5.04),
            Inches(1.92),
            Inches(3.38),
            Inches(1.72),
            "Unit",
            f"{reference_exact_unit_title(deck)}\nFormula focus: {formula}",
            fill=PALE_BLUE,
            accent=NAVY,
            title_size=12.6,
            body_size=11.0,
        )
        add_reference_box(
            slide,
            Inches(8.68),
            Inches(1.92),
            Inches(3.82),
            Inches(1.72),
            "Learning Path",
            path_body,
            fill=PALE_GOLD,
            accent=GOLD,
            title_size=12.6,
            body_size=10.7,
        )
    add_reference_box(
        slide,
        Inches(0.58),
        Inches(3.90),
        Inches(4.55),
        Inches(1.82),
        "Before We Start",
        "What do you already know about measuring space inside a box?",
        fill=PAPER,
        accent=CORAL,
        title_size=12.3,
        body_size=10.5,
    )
    add_reference_lines(slide, Inches(0.86), Inches(4.68), Inches(3.95), 3)
    add_reference_box(
        slide,
        Inches(5.42),
        Inches(3.90),
        Inches(3.22),
        Inches(1.82),
        "Formula",
        formula_display,
        fill=PALE_NAVY,
        accent=NAVY,
        title_size=12.3,
        body_size=18.0 if "\n" in formula_display else 20.0,
    )
    add_reference_prism_visual(slide, Inches(9.00), Inches(3.88), Inches(3.18), Inches(1.55), accent=TEAL)
    tracker = [["Skill", "Before", "After"], ["I can name dimensions.", "1 2 3 4", "1 2 3 4"], ["I can find volume.", "1 2 3 4", "1 2 3 4"]]
    add_table(
        slide,
        Inches(0.58),
        Inches(5.92),
        Inches(11.92),
        Inches(0.96),
        tracker,
        column_widths=[Inches(3.70), Inches(4.10), Inches(4.12)],
        row_heights=[Inches(0.30), Inches(0.33), Inches(0.33)],
        header_font_size=10.6,
        body_font_size=10.4,
    )
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_be_curious(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    accent = REFERENCE_EXACT_SECTION_COLORS["be_curious"]
    add_reference_nav(slide, "Be Curious", accent)
    prompt = reference_exact_context_text(plan_slide, deck, "be_curious")
    add_reference_page_title(slide, "Be Curious", subtitle=prompt, accent=accent)
    add_reference_box(slide, Inches(0.58), Inches(1.72), Inches(4.72), Inches(3.15), "Look Closely", "", fill=PALE_BLUE, accent=TEAL)
    add_reference_prism_visual(slide, Inches(1.00), Inches(2.22), Inches(3.70), Inches(2.02), accent=TEAL)
    add_reference_box(slide, Inches(5.58), Inches(1.72), Inches(2.80), Inches(0.76), "Formula", formula, fill=PALE_NAVY, accent=NAVY, title_size=11.4, body_size=14.0)
    add_reference_box(
        slide,
        Inches(8.66),
        Inches(1.72),
        Inches(3.83),
        Inches(0.76),
        "Keywords",
        "volume | length | width | height | cubic units",
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=11.4,
        body_size=10.4,
    )
    add_reference_box(slide, Inches(5.58), Inches(2.82), Inches(3.18), Inches(2.05), "I Notice", "Write 2 facts you can see in the model.", fill=PAPER, accent=TEAL)
    add_reference_lines(slide, Inches(5.88), Inches(3.60), Inches(2.56), 3)
    add_reference_box(slide, Inches(9.10), Inches(2.82), Inches(3.40), Inches(2.05), "I Wonder", "Write 1 question you have before solving.", fill=PAPER, accent=CORAL)
    add_reference_lines(slide, Inches(9.40), Inches(3.60), Inches(2.76), 3)
    add_reference_box(
        slide,
        Inches(0.58),
        Inches(5.20),
        Inches(5.72),
        Inches(1.10),
        "Visual Prompt",
        prompt,
        fill=PAPER,
        accent=NAVY,
        title_size=11.6,
        body_size=10.4,
    )
    add_reference_box(slide, Inches(6.62), Inches(5.20), Inches(5.87), Inches(1.10), "My First Idea", "", fill=PAPER, accent=TEAL)
    add_reference_lines(slide, Inches(6.92), Inches(5.88), Inches(5.25), 2)
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_vocabulary(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    accent = REFERENCE_EXACT_SECTION_COLORS["vocabulary"]
    add_reference_nav(slide, "Vocabulary", accent)
    add_reference_page_title(slide, "Vocabulary", subtitle="Use precise math words when you explain your reasoning.", accent=accent)
    add_table(
        slide,
        Inches(0.58),
        Inches(1.72),
        Inches(11.92),
        Inches(3.90),
        reference_exact_vocab_rows(deck, plan_slide),
        column_widths=[Inches(2.10), Inches(4.45), Inches(2.90), Inches(2.47)],
        row_heights=[Inches(0.42), Inches(0.70), Inches(0.70), Inches(0.70), Inches(0.70), Inches(0.68)],
        header_font_size=10.5,
        body_font_size=10.4,
        column_alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
    )
    add_reference_box(slide, Inches(0.58), Inches(5.92), Inches(5.75), Inches(0.84), "Sentence Frame", "The volume is ___ cubic units because ___.", fill=PALE_BLUE, accent=TEAL, body_size=11.0)
    add_reference_box(slide, Inches(6.60), Inches(5.92), Inches(5.90), Inches(0.84), "My Sketch", "Draw or label one vocabulary word.", fill=PAPER, accent=GOLD, body_size=10.5)
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_guided_problem(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    accent = REFERENCE_EXACT_SECTION_COLORS["worked_example"]
    problem = reference_exact_problem_text(plan_slide, deck, "guided")
    add_reference_nav(slide, "Try It", accent)
    add_reference_page_title(slide, "Try It: Guided Problem", subtitle=truncate_text(problem, 150), accent=accent)
    add_reference_box(slide, Inches(0.58), Inches(1.72), Inches(4.72), Inches(1.32), "Source Problem", problem, fill=PALE_BLUE, accent=TEAL, body_size=10.5)
    add_reference_prism_visual(slide, Inches(0.98), Inches(3.30), Inches(3.60), Inches(1.42), accent=TEAL)
    organizer = [
        ["Read", "What measurements do I have?"],
        ["Plan", "Which formula will I use?"],
        ["Solve", "Substitute and multiply."],
        ["Explain", "What does the answer mean?"],
    ]
    add_table(
        slide,
        Inches(5.58),
        Inches(1.72),
        Inches(3.08),
        Inches(2.25),
        organizer,
        column_widths=[Inches(1.00), Inches(2.08)],
        row_heights=[Inches(0.56), Inches(0.56), Inches(0.56), Inches(0.57)],
        header_font_size=10.6,
        body_font_size=10.4,
        column_alignments=[PP_ALIGN.CENTER, PP_ALIGN.LEFT],
    )
    add_reference_box(slide, Inches(9.00), Inches(1.72), Inches(3.50), Inches(2.25), "Set Up", f"{formula}\nV = ___ × ___ × ___\nV = ___ cubic units", fill=PAPER, accent=NAVY, body_size=13.2)
    add_reference_box(slide, Inches(5.58), Inches(4.34), Inches(6.92), Inches(1.88), "Show Your Work", "", fill=PAPER, accent=TEAL)
    add_reference_lines(slide, Inches(5.88), Inches(5.02), Inches(6.30), 4, gap=0.30)
    add_reference_box(slide, Inches(0.58), Inches(5.20), Inches(4.72), Inches(1.02), "Answer Sentence", "The volume is ___ cubic units because ___.", fill=PALE_GOLD, accent=GOLD, body_size=10.6)
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_practice(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    accent = REFERENCE_EXACT_SECTION_COLORS["practice"]
    title = normalize_whitespace(str(plan_slide.get("title", "")))
    if reference_exact_is_generic_text(title) or not title or "rectangular prism" not in title.lower():
        title = "Partner Practice: Volume of a Rectangular Prism"
    problem = reference_exact_problem_text(plan_slide, deck, "practice")
    add_reference_nav(slide, "Practice", accent)
    add_reference_page_title(slide, truncate_text(title, 62), subtitle=truncate_text(problem, 150), accent=accent)
    add_reference_box(slide, Inches(0.58), Inches(1.72), Inches(4.35), Inches(1.18), "Problem", problem, fill=PALE_CORAL, accent=CORAL, body_size=10.5)
    add_reference_box(slide, Inches(0.58), Inches(3.18), Inches(4.35), Inches(2.40), "Model / Label", "", fill=PAPER, accent=CORAL)
    add_reference_prism_visual(slide, Inches(1.05), Inches(3.78), Inches(3.36), Inches(1.38), accent=CORAL)
    add_reference_box(slide, Inches(5.28), Inches(1.72), Inches(3.18), Inches(1.18), "Formula", formula, fill=PALE_NAVY, accent=NAVY, body_size=15.0)
    add_reference_box(slide, Inches(8.78), Inches(1.72), Inches(3.72), Inches(1.18), "Substitute", "V = ___ × ___ × ___", fill=PAPER, accent=NAVY, body_size=14.0)
    add_reference_box(slide, Inches(5.28), Inches(3.18), Inches(7.22), Inches(2.40), "Work Space", "", fill=PAPER, accent=TEAL)
    add_reference_lines(slide, Inches(5.60), Inches(3.95), Inches(6.60), 5, gap=0.30)
    add_reference_box(slide, Inches(0.58), Inches(5.90), Inches(11.92), Inches(0.74), "Explain to Your Partner", "My answer is ___ cubic units. I know because ___.", fill=PALE_BLUE, accent=TEAL, body_size=10.8)
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_error_analysis(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> None:
    formula = reference_exact_formula_text(deck, session_label)
    prompt = reference_exact_problem_text(plan_slide, deck, "error")
    problem = prompt
    source_targets = [
        target
        for target in primary_source_problem_targets(plan_slide, limit=2)
        if not reference_exact_is_generic_text(target)
    ]
    if source_targets and not source_problem_text_overlap(prompt, source_targets):
        problem = source_targets[0]
    student_work = f"A student's work -- something is wrong!\n{problem}"
    accent = SAGE
    add_reference_nav(slide, "Discuss", accent)
    add_reference_page_title(slide, "Find the Error", subtitle="A student's work -- something is wrong!", accent=accent)
    add_reference_box(slide, Inches(0.58), Inches(1.62), Inches(5.82), Inches(1.70), "Student Work", student_work, fill=PALE_SAGE, accent=SAGE, body_size=10.4)
    add_reference_box(slide, Inches(6.65), Inches(1.62), Inches(2.40), Inches(1.70), "What Went Wrong?", "Circle or describe the error.", fill=PAPER, accent=CORAL, body_size=10.5)
    add_reference_lines(slide, Inches(6.94), Inches(2.55), Inches(1.78), 2)
    add_reference_box(slide, Inches(9.32), Inches(1.62), Inches(3.18), Inches(1.70), "Correct Formula", formula, fill=PALE_NAVY, accent=NAVY, body_size=14.0)
    add_reference_box(slide, Inches(0.58), Inches(3.58), Inches(5.78), Inches(1.88), "Correct the Work", "", fill=PAPER, accent=TEAL)
    add_reference_lines(slide, Inches(0.90), Inches(4.24), Inches(5.15), 4, gap=0.27)
    add_reference_box(slide, Inches(6.70), Inches(3.58), Inches(5.80), Inches(1.88), "Explain Your Reasoning", "The mistake was ___. The correct volume is ___ cubic units.", fill=PAPER, accent=SAGE, body_size=10.5)
    add_reference_lines(slide, Inches(7.02), Inches(4.36), Inches(5.12), 3, gap=0.27)
    add_reference_box(slide, Inches(0.58), Inches(5.80), Inches(11.92), Inches(0.84), "Discussion Check", "What does cubic units tell us that square units does not?", fill=PALE_GOLD, accent=GOLD, body_size=10.4)
    add_reference_footer(slide, session_label, formula, deck)


def render_reference_exact_slide(
    slide: Any,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    page: int,
) -> bool:
    kind = plan_slide.get("kind", "")
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    title = normalize_whitespace(plan_slide.get("title", "")).lower()
    if kind in {"cover", "learning_target"} or template_role in {"cover_page", "learning_objectives"}:
        render_reference_exact_cover(slide, plan_slide, session_plan, session_label, deck, page)
    elif kind == "be_curious":
        render_reference_exact_be_curious(slide, plan_slide, session_plan, session_label, deck, page)
    elif kind == "vocabulary":
        render_reference_exact_vocabulary(slide, plan_slide, session_plan, session_label, deck, page)
    elif kind == "worked_example" or "guided problem" in title or title.startswith("try it"):
        render_reference_exact_guided_problem(slide, plan_slide, session_plan, session_label, deck, page)
    elif template_role == "best_fit_review" or "find the error" in title or "error" in title:
        render_reference_exact_error_analysis(slide, plan_slide, session_plan, session_label, deck, page)
    elif kind == "practice":
        render_reference_exact_practice(slide, plan_slide, session_plan, session_label, deck, page)
    else:
        return False
    return True


def reference_exact_slide_engagement_modes(plan_slide: dict[str, Any]) -> set[str]:
    kind = plan_slide.get("kind", "")
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    title = normalize_whitespace(plan_slide.get("title", "")).lower()
    modes: set[str] = set()
    if kind in {"cover", "learning_target"} or template_role in {"cover_page", "learning_objectives"}:
        modes.add("reflect")
    if kind == "be_curious":
        modes.update({"notice", "wonder"})
    elif kind == "vocabulary":
        modes.update({"match", "visual"})
    elif kind == "worked_example" or "guided problem" in title or title.startswith("try it"):
        modes.update({"guided_solve", "annotate"})
    elif template_role == "best_fit_review" or "find the error" in title or "error" in title:
        modes.update({"error_analysis", "discuss"})
    elif kind == "practice":
        modes.update({"partner", "apply"})
    return modes


def reference_exact_engagement_modes(session: dict[str, Any]) -> set[str]:
    modes: set[str] = set()
    for plan_slide in session.get("slides", []):
        modes.update(reference_exact_slide_engagement_modes(plan_slide))
    return modes


def reference_exact_engagement_slide_count(session: dict[str, Any]) -> int:
    return sum(1 for plan_slide in session.get("slides", []) if reference_exact_slide_engagement_modes(plan_slide))


def render_cover_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    session_label: str,
    deck: dict[str, Any],
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    if normalize_whitespace(plan_slide.get("template_role", "")) == "cover_page":
        render_exact_cover_slide(slide, plan_slide=plan_slide, session_label=session_label, image_lookup=image_lookup)
        return
    session_key = "session_1" if "1" in session_label else "session_2"
    hero_subtitle = first_distinct_text(
        [plan_slide.get("subtitle", ""), plan_slide.get("secondary_text", ""), default_cover_subtitle(deck, session_key)],
        excluded=[plan_slide.get("title", ""), session_label],
    ) or default_cover_subtitle(deck, session_key)
    focus_text = first_distinct_text(
        [
            plan_slide.get("primary_text", ""),
            cover_focus_statement(deck, session_key, plan_slide.get("title", "")),
            bullet_text(plan_slide.get("bullets", []), limit=2, max_len=64),
        ],
        excluded=[plan_slide.get("title", ""), session_label, hero_subtitle],
    ) or cover_focus_statement(deck, session_key, plan_slide.get("title", ""))
    right_body = first_distinct_text(
        [
            plan_slide.get("secondary_text", ""),
            bullet_text(plan_slide.get("bullets", []), limit=3, max_len=62),
            first_standard_text(deck),
        ],
        excluded=[plan_slide.get("title", ""), session_label, hero_subtitle, focus_text],
    ) or "Use the notebook to track the lesson's ideas, models, and source problems."

    set_background(slide, PAPER_WARM)
    add_cover_background_decor(slide)
    add_tag(slide, Inches(0.82), Inches(0.90), session_label.upper(), GOLD)
    add_text(
        slide,
        Inches(0.82),
        Inches(1.54),
        Inches(5.85),
        Inches(2.26),
        truncate_text(plan_slide["title"], 120),
        size=31.4,
        color=NAVY,
        bold=True,
        font=FONT_HEAD,
        line_spacing=0.92,
    )
    add_text(
        slide,
        Inches(0.82),
        Inches(4.02),
        Inches(5.55),
        Inches(0.40),
        truncate_text(hero_subtitle, 220),
        size=14.5,
        color=TEAL,
        font=FONT_BODY,
        margin=0.01,
    )
    title_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(5.72), Inches(2.05), Inches(0.06))
    title_rule.fill.solid()
    title_rule.fill.fore_color.rgb = GOLD
    title_rule.line.fill.background()
    add_image_panel(
        slide,
        Inches(7.46),
        Inches(1.08),
        Inches(4.62),
        Inches(2.72),
        choose_image_asset(plan_slide, image_lookup),
        fill=PAPER,
    )
    add_card(
        slide,
        Inches(7.15),
        Inches(4.26),
        Inches(2.42),
        Inches(1.52),
        "Look For",
        focus_text,
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=13.2,
        body_size=11.7,
    )
    add_card(
        slide,
        Inches(9.84),
        Inches(4.26),
        Inches(2.92),
        Inches(1.52),
        "Lesson Thread",
        right_body or "Use the notebook to track important lesson ideas.",
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=13.2,
        body_size=11.4,
    )
    terms = flagship_cover_terms(deck, plan_slide.get("source_slide_numbers", []), limit=3)
    if terms:
        add_text(
            slide,
            Inches(0.82),
            Inches(5.22),
            Inches(2.6),
            Inches(0.18),
            "Lesson Words",
            size=10.8,
            color=TEAL,
            bold=True,
            font=FONT_HEAD,
            margin=0.01,
        )
        for index, term in enumerate(terms[:3]):
            add_chip(
                slide,
                Inches(0.82) + index * Inches(1.8),
                Inches(5.48),
                Inches(1.64),
                Inches(0.28),
                truncate_text(term, 18),
                fill=PALE_GOLD if index == 0 else PAPER,
                accent=GOLD if index == 0 else TEAL,
            )
    add_name_bar(slide)


def render_be_curious_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    session_plan: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    if normalize_whitespace(plan_slide.get("template_role", "")) == "prior_session_review":
        render_exact_be_curious_slide(
            slide,
            plan_slide=plan_slide,
            session_plan=session_plan,
            page=page,
            footer_text=footer_text,
            image_lookup=image_lookup,
        )
        return
    notice_kernels, wonder_kernels = be_curious_sentence_kernels(plan_slide)
    vocab_items = be_curious_vocabulary_items(plan_slide, session_plan)
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=CORAL,
        footer_text=footer_text,
    )
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(4.15),
        Inches(4.35),
        choose_image_asset(plan_slide, image_lookup),
        fill=PALE_BLUE,
        label=plan_slide.get("image_caption", "") or "Notice + Wonder image",
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.08),
        Inches(3.70),
        Inches(0.64),
        "Math Lens",
        plan_slide.get("response_prompt", "") or "What clue in the source image or prompt should shape your first move?",
        fill=PALE_CORAL,
        accent=CORAL,
        title_size=12.8,
        body_size=11.7,
    )
    add_lined_area(
        slide,
        Inches(5.05),
        Inches(2.0),
        Inches(3.62),
        Inches(2.40),
        "Notice",
        be_curious_panel_prompt(
            plan_slide["primary_text"] or "",
            notice_kernels,
            fallback="Record details you notice in the image before solving anything.",
        ),
        lines=3,
    )
    add_lined_area(
        slide,
        Inches(9.13),
        Inches(2.0),
        Inches(3.62),
        Inches(2.40),
        "Wonder",
        be_curious_panel_prompt(
            plan_slide["secondary_text"] or "",
            wonder_kernels,
            fallback="Ask a question or make a prediction that could deepen your understanding.",
        ),
        lines=3,
    )
    starters_text = "  ".join(f"• {k}" for k in notice_kernels[:1] + wonder_kernels[:1])
    add_card(
        slide,
        Inches(5.05),
        Inches(4.56),
        Inches(7.70),
        Inches(0.86),
        "Sentence Starters",
        starters_text,
        fill=PAPER_WARM,
        accent=CORAL,
        title_size=12.8,
        body_size=11.7,
    )
    if has_activity(plan_slide):
        add_vocabulary_snapshot(
            slide,
            Inches(5.05),
            Inches(5.58),
            Inches(4.95),
            Inches(1.10),
            vocab_items,
        )
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(10.24),
            y=Inches(5.58),
            w=Inches(2.51),
            h=Inches(1.10),
        )
    else:
        add_vocabulary_snapshot(
            slide,
            Inches(5.05),
            Inches(5.58),
            Inches(7.70),
            Inches(1.10),
            vocab_items,
        )


def render_learning_target_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    if normalize_whitespace(plan_slide.get("template_role", "")) == "learning_objectives":
        render_exact_learning_objectives_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_checkmark_tracker(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(2.00),
        plan_slide["primary_text"],
        plan_slide["secondary_text"],
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(4.28),
        Inches(4.12),
        Inches(2.42),
        "Success Moves",
        bullet_text(plan_slide["bullets"], limit=3, max_len=72)
        or "- Meet the content objective.\n- Meet the language objective.\n- Show evidence in notes and practice.",
        fill=PALE_SAGE,
        accent=SAGE,
        title_size=16.8,
        body_size=13.2,
    )
    add_sentence_starters(
        slide,
        Inches(4.9),
        Inches(4.28),
        Inches(7.85),
        Inches(2.42),
        plan_slide["sentence_starters"] or [
            "The content of today's lesson is ___.",
            "The objective is to ___.",
            "Before the lesson, I could ___.",
            "After the lesson, I can ___ because ___.",
        ],
        accent=TEAL,
    )


def render_vocabulary_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    if template_role == "vocabulary_table":
        render_exact_vocabulary_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text, image_lookup=image_lookup)
        return
    if template_role == "vocabulary_activity":
        render_exact_vocabulary_activity_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text, image_lookup=image_lookup)
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=SAGE,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.56),
        "Vocabulary in Context",
        "Use the source words, examples, and visuals together so the vocabulary feels useful, visible, and ready for your discussion and writing.",
        fill=PALE_SAGE,
        accent=SAGE,
        title_size=14.4,
        body_size=11.7,
    )
    vocab_items = plan_slide["vocabulary"][:4] or [
        {"word": "term", "definition": "student-friendly definition", "example": "example from the lesson", "visual_cue": "visual clue"}
    ]
    vocab_assets = choose_image_assets(plan_slide, image_lookup, limit=4)
    card_specs = [
        (Inches(0.5), Inches(2.78), PALE_BLUE, TEAL),
        (Inches(6.39), Inches(2.78), PALE_GOLD, GOLD),
        (Inches(0.5), Inches(4.42), PALE_CORAL, CORAL),
        (Inches(6.39), Inches(4.42), PALE_SAGE, SAGE),
    ]
    for index, item in enumerate(vocab_items[:4]):
        asset = vocab_assets[index % len(vocab_assets)] if vocab_assets else None
        card_x, card_y, fill, accent = card_specs[index]
        add_vocabulary_feature_card(
            slide,
            x=card_x,
            y=card_y,
            w=Inches(5.86),
            h=Inches(1.48),
            word=item.get("word", ""),
            definition=item.get("definition", ""),
            example=item.get("example", ""),
            visual_cue=item.get("visual_cue", ""),
            asset=asset,
            accent=accent,
            fill=fill,
        )
    add_slot_box(
        slide,
        Inches(0.5),
        Inches(6.12),
        Inches(5.86),
        Inches(0.46),
        "My strongest word today: ____________________",
        fill=PALE_SAGE,
        accent=SAGE,
    )
    if has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(6.39),
            y=Inches(5.94),
            w=Inches(6.36),
            h=Inches(0.84),
        )
    else:
        add_card(
            slide,
            Inches(6.39),
            Inches(5.94),
            Inches(6.36),
            Inches(0.84),
            "Use the Words",
            plan_slide["primary_text"] or "Use these words when you talk, write, and explain your thinking.",
            fill=PALE_SAGE,
            accent=SAGE,
            title_size=12.7,
            body_size=11.9,
        )


def render_guided_notes_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(4.95),
        Inches(4.85),
        "Source Summary",
        plan_slide["primary_text"] or "Summarize the core idea from the source lesson.",
        fill=PAPER,
        accent=CORAL,
        title_size=17,
        body_size=14.1,
    )
    asset = choose_image_asset(plan_slide, image_lookup)
    if asset:
        add_image_panel(
            slide,
            Inches(0.82),
            Inches(4.42),
            Inches(4.3),
            Inches(1.85),
            asset,
            fill=PALE_BLUE,
            label=plan_slide["image_caption"] or "Source image",
        )
    add_card(
        slide,
        Inches(5.72),
        Inches(2.0),
        Inches(7.03),
        Inches(2.1),
        "Key Takeaways",
        bullet_text(plan_slide["bullets"], limit=4, max_len=95)
        or "- Record the important details.\n- Track how the lesson ideas connect.",
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=16.2,
        body_size=14.2,
    )
    if render_premium_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(5.72),
        y=Inches(4.32),
        w=Inches(7.03),
        h=Inches(2.53),
        accent=TEAL,
    ):
        return
    if has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(5.72),
            y=Inches(4.32),
            w=Inches(7.03),
            h=Inches(2.53),
        )
    else:
        add_lined_area(
            slide,
            Inches(5.72),
            Inches(4.42),
            Inches(7.03),
            Inches(2.43),
            "Notes in Your Own Words",
            plan_slide["response_prompt"] or "Write the main idea in your own words.",
            lines=4,
        )


def render_interactive_apply_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    accent = KIND_ACCENT.get(plan_slide.get("kind", "practice"), GOLD)
    add_header(
        slide,
        section=plan_slide.get("section", "Practice Extension") or "Practice Extension",
        title=plan_slide.get("title", "Practice Extension") or "Practice Extension",
        subtitle=plan_slide.get("subtitle", "Use draggable pieces to plan, solve, and justify the full source problem."),
        page=page,
        accent=accent,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    problem_cards = problem_display_cards(plan_slide, variant="practice")
    add_problem_prompt_stack(
        slide,
        x=Inches(0.5),
        y=Inches(2.0),
        w=Inches(12.25),
        h=Inches(1.34),
        title="Practice Extension",
        problems=problem_cards,
        accent=accent,
        fill=PAPER,
        subtitle="Keep the full source problem visible while you build your solve path.",
        kicker="Apply + Explain",
    )
    asset = choose_image_asset(plan_slide, image_lookup)
    if asset:
        add_image_panel(
            slide,
            Inches(0.5),
            Inches(3.62),
            Inches(4.85),
            Inches(2.22),
            asset,
            fill=PALE_BLUE,
            label=plan_slide.get("image_caption", "") or "Source model",
        )
    else:
        add_card(
            slide,
            Inches(0.5),
            Inches(3.62),
            Inches(4.85),
            Inches(2.22),
            "Source Focus",
            truncate_display_copy(rendered_source_problem_statement(plan_slide), 148),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=15.0,
            body_size=12.4,
        )
    interactive_slide = activity_slide_for_render(plan_slide)
    render_activity_board(
        slide,
        plan_slide=interactive_slide,
        x=Inches(5.72),
        y=Inches(3.62),
        w=Inches(7.03),
        h=Inches(2.22),
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(6.06),
        Inches(12.25),
        Inches(0.76),
        "Apply + Explain",
        plan_slide.get("response_prompt", "") or "Explain how the interactive plan helped you solve the full source problem.",
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=14.2,
        body_size=11.7,
    )


def render_worked_example_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    render_mode = problem_render_mode(plan_slide)
    if render_mode == PROBLEM_RENDER_MODE_INTERACTIVE:
        render_interactive_apply_slide(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            image_lookup=image_lookup,
        )
        return
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    if template_role == "guided_practice":
        render_exact_guided_practice_slide(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            image_lookup=image_lookup,
        )
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=CORAL,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    asset = choose_image_asset(plan_slide, image_lookup)
    problem_cards = problem_display_cards(plan_slide, variant="guided")
    add_problem_prompt_stack(
        slide,
        x=Inches(0.5),
        y=Inches(2.0),
        w=Inches(7.05),
        h=Inches(1.34),
        title="Worked Example",
        problems=problem_cards,
        accent=CORAL,
        fill=PAPER,
        subtitle=workbook_prompt_stack_subtitle("guided"),
        kicker=workbook_prompt_stack_kicker("guided"),
    )
    if asset:
        add_image_panel(
            slide,
            Inches(7.82),
            Inches(2.0),
            Inches(4.93),
            Inches(1.34),
            asset,
            fill=PALE_BLUE,
            label=plan_slide["image_caption"] or "Source model",
        )
    else:
        add_card(
            slide,
            Inches(7.82),
            Inches(2.0),
            Inches(4.93),
            Inches(1.34),
            "Look For",
            bullet_text(plan_slide["bullets"], limit=2, max_len=78)
            or "Look for the measurements, labels, and relationships that help you model the source problem.",
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=15.4,
            body_size=12.4,
        )
    render_problem_workbook_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(0.5),
        y=Inches(3.52),
        w=Inches(7.12),
        h=Inches(3.10),
        accent=CORAL,
        variant="guided",
    )
    if render_mode == PROBLEM_RENDER_MODE_FOCUS:
        add_card(
            slide,
            Inches(7.82),
            Inches(3.52),
            Inches(4.93),
            Inches(3.10),
            "Next Step",
            "Use the next slide to apply this same full problem with draggable pieces and a written explanation.",
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=17,
            body_size=13.0,
        )
    else:
        interactive_slide = activity_slide_for_render(plan_slide)
        if render_premium_panel(
            slide,
            plan_slide=plan_slide,
            x=Inches(7.82),
            y=Inches(3.52),
            w=Inches(4.93),
            h=Inches(3.10),
            accent=CORAL,
        ):
            pass
        elif has_activity(interactive_slide):
            render_activity_board(
                slide,
                plan_slide=interactive_slide,
                x=Inches(7.82),
                y=Inches(3.52),
                w=Inches(4.93),
                h=Inches(3.10),
            )
        else:
            add_discussion_panel(
                slide,
                plan_slide=plan_slide,
                x=Inches(7.82),
                y=Inches(3.52),
                w=Inches(4.93),
                h=Inches(3.10),
                fill=PALE_BLUE,
                accent=TEAL,
                max_questions=2,
                include_prompt=True,
            )
    add_sentence_starters(
        slide,
        Inches(0.5),
        Inches(6.76),
        Inches(12.25),
        Inches(0.40),
        plan_slide["sentence_starters"] or ["First, I ...", "Next, I ...", "So, the answer is ..."],
        accent=GOLD,
    )


def render_drag_sort_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    source_problem = primary_source_problem_targets(plan_slide, limit=1)
    source_problem_text = source_problem[0] if source_problem else plan_slide.get("primary_text", "")
    intro_body = normalize_whitespace(
        "\n".join(
            part
            for part in [
                truncate_display_copy(source_problem_text, 188),
                truncate_display_copy(
                    plan_slide.get("response_prompt", "") or "Move each card into the category that fits best, then explain one choice.",
                    74,
                ),
            ]
            if normalize_whitespace(part)
        )
    )
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.94),
        "Source Problem + Sort",
        intro_body,
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=14.6,
        body_size=11.6,
    )
    labels = unique_nonempty(plan_slide.get("tasks", []), limit=2)
    if len(labels) < 2 or any(len(label) > 44 or is_problem_like_text(label) for label in labels):
        labels = ["Matches the source problem", "Needs a second look"]
    add_slot_box(slide, Inches(0.5), Inches(3.18), Inches(5.88), Inches(2.08), labels[0], fill=PAPER, accent=TEAL)
    add_slot_box(slide, Inches(6.87), Inches(3.18), Inches(5.88), Inches(2.08), labels[1], fill=PAPER, accent=GOLD)
    pieces = activity_pieces(plan_slide, limit=6) or ["Card 1", "Card 2", "Card 3", "Card 4"]
    chip_w = Inches(3.78)
    chip_h = Inches(0.26)
    for index, piece in enumerate(pieces[:6]):
        row = index // 3
        col = index % 3
        add_draggable_piece_box(
            slide,
            Inches(0.68) + col * Inches(4.06),
            Inches(5.54) + row * Inches(0.36),
            chip_w,
            chip_h,
            piece,
            fill=PAPER,
            accent=TEAL if row == 0 else GOLD,
        )
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(6.46),
        Inches(12.25),
        Inches(0.42),
        "Explain One Sort",
        "One card belongs in ___ because ___.",
        lines=1,
        fill=PALE_GOLD,
    )


def render_error_analysis_workbook_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=CORAL,
        footer_text=footer_text,
    )
    steps = unique_nonempty(plan_slide.get("error_steps", []), limit=4)
    if steps:
        for index, step in enumerate(steps):
            add_card(
                slide,
                Inches(0.5),
                Inches(2.0) + index * Inches(0.78),
                Inches(5.48),
                Inches(0.64),
                f"Step {index + 1}",
                step,
                fill=PALE_CORAL if "❌" in step else PAPER,
                accent=CORAL if "❌" in step else TEAL,
                title_size=12.8,
                body_size=11.2,
            )
    else:
        add_card(
            slide,
            Inches(0.5),
            Inches(2.0),
            Inches(5.48),
            Inches(3.10),
            "Student Work",
            plan_slide.get("primary_text", ""),
            fill=PALE_CORAL,
            accent=CORAL,
            title_size=16.0,
            body_size=12.4,
        )
    add_lined_area(
        slide,
        Inches(6.28),
        Inches(2.0),
        Inches(6.47),
        Inches(1.18),
        "✏️ The mistake is:",
        plan_slide.get("response_prompt", "") or "The mistake is: ___.",
        lines=2,
        fill=PAPER,
    )
    add_card(
        slide,
        Inches(6.28),
        Inches(3.38),
        Inches(6.47),
        Inches(0.84),
        "✅ Fix It",
        plan_slide.get("fix_it_text", "") or "Fix it: ___.",
        fill=PALE_SAGE,
        accent=SAGE,
        title_size=13.2,
        body_size=11.4,
    )
    add_lined_area(
        slide,
        Inches(6.28),
        Inches(4.44),
        Inches(6.47),
        Inches(1.42),
        "💬 Explain",
        plan_slide.get("why_prompt", "") or "Why does this fix work?",
        lines=3,
        fill=PAPER,
    )
    if has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(0.5),
            y=Inches(6.02),
            w=Inches(12.25),
            h=Inches(0.66),
        )
    else:
        add_sentence_starters(
            slide,
            Inches(0.5),
            Inches(6.02),
            Inches(12.25),
            Inches(0.66),
            plan_slide["sentence_starters"] or ["The mistake is ___.", "The fix is ___.", "This works because ___."],
            accent=CORAL,
        )


def render_two_column_compare_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=GOLD,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.52),
        "Directions",
        "Each partner solves one tile. Then answer Q3 together.",
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=12.6,
        body_size=10.8,
    )
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(2.76),
        Inches(5.88),
        Inches(3.08),
        "Partner A",
        plan_slide.get("partner_a_problem", "") or source_problem_statement(plan_slide),
        lines=4,
        fill=PAPER,
    )
    add_lined_area(
        slide,
        Inches(6.87),
        Inches(2.76),
        Inches(5.88),
        Inches(3.08),
        "Partner B",
        plan_slide.get("partner_b_problem", "") or similar_problem_statement(plan_slide),
        lines=4,
        fill=PALE_GOLD,
    )
    if has_activity(plan_slide):
        add_lined_area(
            slide,
            Inches(0.5),
            Inches(5.48),
            Inches(7.06),
            Inches(1.20),
            "Q3 Compare the Thinking",
            plan_slide.get("compare_frame", "") or plan_slide.get("response_prompt", "") or "It is the same because ___. It is different because ___.",
            lines=2,
            fill=PALE_BLUE,
        )
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(7.74),
            y=Inches(5.48),
            w=Inches(5.01),
            h=Inches(1.20),
        )
    else:
        add_lined_area(
            slide,
            Inches(0.5),
            Inches(5.48),
            Inches(12.25),
            Inches(1.20),
            "Q3 Compare the Thinking",
            plan_slide.get("compare_frame", "") or plan_slide.get("response_prompt", "") or "It is the same because ___. It is different because ___.",
            lines=2,
            fill=PALE_BLUE,
        )


def render_choice_board_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    paths = plan_slide.get("choice_paths", [])
    if not paths:
        paths = [
            {"label": "Path A", "problem": source_problem_statement(plan_slide), "hint": "Hint: start with the formula.", "answer": "My answer: ___"},
            {"label": "Path B", "problem": similar_problem_statement(plan_slide), "hint": "Hint: show one extra step.", "answer": "My answer: ___"},
            {"label": "Path C", "problem": third_problem_statement(plan_slide) or build_problem_creation_prompt(plan_slide), "hint": "Hint: keep the same rule.", "answer": "My answer: ___"},
        ]
    for index, item in enumerate(paths[:3]):
        x = Inches(0.5) + index * Inches(4.12)
        fill = PAPER if index != 1 else PALE_GOLD
        add_lined_area(
            slide,
            x,
            Inches(2.0),
            Inches(3.94),
            Inches(2.98),
            item.get("label", f"Path {index + 1}"),
            item.get("problem", ""),
            lines=2,
            fill=fill,
        )
        add_card(
            slide,
            x,
            Inches(5.00),
            Inches(3.94),
            Inches(0.54),
            "Hint",
            item.get("hint", ""),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=12.2,
            body_size=10.5,
        )
        add_lined_area(
            slide,
            x,
            Inches(5.66),
            Inches(3.94),
            Inches(0.62),
            "My Answer",
            item.get("answer", "My answer: ___"),
            lines=1,
            fill=PAPER,
        )
    if has_activity(plan_slide):
        add_lined_area(
            slide,
            Inches(0.5),
            Inches(6.40),
            Inches(7.06),
            Inches(0.48),
            "My Choice",
            plan_slide.get("response_prompt", "") or "I chose Path ___ because ___.",
            lines=1,
            fill=PALE_SAGE,
        )
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(7.74),
            y=Inches(6.40),
            w=Inches(5.01),
            h=Inches(0.48),
        )
    else:
        add_lined_area(
            slide,
            Inches(0.5),
            Inches(6.40),
            Inches(12.25),
            Inches(0.48),
            "My Choice",
            plan_slide.get("response_prompt", "") or "I chose Path ___ because ___.",
            lines=1,
            fill=PALE_SAGE,
        )


def render_collaborative_practice_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=GOLD,
        footer_text=footer_text,
    )
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(5.9),
        Inches(3.52),
        "Partner A",
        plan_slide.get("partner_a_problem", "") or source_problem_statement(plan_slide),
        lines=4,
        fill=PAPER,
    )
    add_lined_area(
        slide,
        Inches(6.85),
        Inches(2.0),
        Inches(5.9),
        Inches(3.52),
        "Partner B",
        plan_slide.get("partner_b_problem", "") or similar_problem_statement(plan_slide),
        lines=4,
        fill=PALE_GOLD,
    )
    discussion = plan_slide.get("discussion_questions", [])
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(5.72),
        Inches(6.05),
        Inches(0.96),
        "Q1",
        discussion[0] if discussion else "How is the new formula like the old one?",
        lines=2,
        fill=PALE_BLUE,
    )
    if has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(6.70),
            y=Inches(5.72),
            w=Inches(6.05),
            h=Inches(0.96),
        )
    else:
        add_lined_area(
            slide,
            Inches(6.70),
            Inches(5.72),
            Inches(6.05),
            Inches(0.96),
            "Q2",
            discussion[1] if len(discussion) > 1 else "The ___ changes because ___.",
            lines=2,
            fill=PALE_SAGE,
        )


def render_independent_practice_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=GOLD,
        footer_text=footer_text,
    )
    problems = plan_slide.get("independent_problems", [])
    if not problems:
        prompts = problem_display_cards(plan_slide, variant="practice")
        problems = [{"label": f"P{index + 1}", "prompt": prompt, "answer": "Answer: ___"} for index, prompt in enumerate(prompts[:3])]
    profile = math_profile_for_text(normalize_whitespace(" ".join([plan_slide.get("title", ""), plan_slide.get("primary_text", ""), " ".join(plan_slide.get("source_problem_cards", []))])))
    for index, item in enumerate(problems[:3]):
        x = Inches(0.5) + index * Inches(4.12)
        fill = PAPER if index != 1 else PALE_GOLD
        add_lined_area(
            slide,
            x,
            Inches(2.0),
            Inches(3.94),
            Inches(2.52),
            item.get("label", f"P{index + 1}"),
            item.get("prompt", ""),
            lines=2,
            fill=fill,
        )
        add_small_math_visual(
            slide,
            x + Inches(0.12),
            Inches(4.70),
            Inches(1.26),
            Inches(1.00),
            profile=profile,
            accent=TEAL,
        )
        add_lined_area(
            slide,
            x + Inches(1.54),
            Inches(4.58),
            Inches(2.28),
            Inches(1.12),
            "Answer",
            item.get("answer", "Answer: ___"),
            lines=2,
            fill=PAPER,
        )
    if has_activity(plan_slide):
        add_card(
            slide,
            Inches(0.5),
            Inches(5.92),
            Inches(7.06),
            Inches(0.76),
            "Helpful Hints",
            "\n".join(f"- {line}" for line in plan_slide.get("helpful_tips", [])[:3]) or problem_tip_body(plan_slide),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=12.8,
            body_size=10.5,
        )
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(7.74),
            y=Inches(5.92),
            w=Inches(5.01),
            h=Inches(0.76),
        )
    else:
        add_card(
            slide,
            Inches(0.5),
            Inches(5.92),
            Inches(12.25),
            Inches(0.76),
            "Helpful Hints",
            "\n".join(f"- {line}" for line in plan_slide.get("helpful_tips", [])[:3]) or problem_tip_body(plan_slide),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=12.8,
            body_size=10.5,
        )


def render_twr_frame_workbook_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=SAGE,
        footer_text=footer_text,
    )
    prompts = plan_slide.get("twr_frames", []) or [
        "The area is ___ cm².",
        "Base = ___. Height = ___.",
        "I use the formula because ___.",
        "Same: ___. Different: ___.",
    ]
    labels = ["Frame 1 Identify", "Frame 2 Explain", "Frame 3 Justify", "Frame 4 Compare"]
    fills = [PAPER, PALE_BLUE, PALE_GOLD, PALE_SAGE]
    for index, label in enumerate(labels):
        row = index // 2
        col = index % 2
        add_lined_area(
            slide,
            Inches(0.5) + col * Inches(6.18),
            Inches(2.0) + row * Inches(2.28),
            Inches(5.93),
            Inches(2.04),
            label,
            prompts[index] if index < len(prompts) else prompts[-1],
            lines=3,
            fill=fills[index],
        )


def render_tiered_exit_workbook_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=CORAL,
        footer_text=footer_text,
    )
    tiers = plan_slide.get("exit_tiers", [])
    if not tiers:
        tiers = [
            {"label": "⭐ Tier 1", "prompt": source_problem_statement(plan_slide), "frame": "Say: 'The answer is ___.'"},
            {"label": "⭐⭐ Tier 2", "prompt": similar_problem_statement(plan_slide), "frame": "Say: 'I use ___ to find ___.'"},
            {"label": "⭐⭐⭐ Tier 3", "prompt": third_problem_statement(plan_slide) or plan_slide.get("response_prompt", ""), "frame": "Say: 'They are alike because ___.'"},
        ]
    for index, item in enumerate(tiers[:3]):
        x = Inches(0.5) + index * Inches(4.12)
        add_lined_area(
            slide,
            x,
            Inches(2.0),
            Inches(3.94),
            Inches(3.66),
            item.get("label", f"Tier {index + 1}"),
            item.get("prompt", ""),
            lines=3,
            fill=PAPER if index != 1 else PALE_CORAL,
        )
        add_lined_area(
            slide,
            x,
            Inches(5.84),
            Inches(3.94),
            Inches(0.82),
            "💬 Say It",
            item.get("frame", "Say: '___.'"),
            lines=1,
            fill=PALE_GOLD,
        )


def render_goal_tracker_workbook_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=SAGE,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(12.25),
        Inches(0.82),
        "Today's Goal",
        plan_slide.get("primary_text", "") or "Rate your confidence and be honest about what you know right now.",
        fill=PALE_SAGE,
        accent=SAGE,
        title_size=14.4,
        body_size=11.8,
    )
    levels = plan_slide.get("goal_levels", []) or [
        "1. I need more help. I am not sure yet.",
        "2. I am getting there. I can do it with help.",
        "3. I got it. I can do it alone.",
        "4. I can teach it. I can connect both sessions.",
    ]
    for index, text in enumerate(levels[:4]):
        add_card(
            slide,
            Inches(0.5) + index * Inches(3.10),
            Inches(3.10),
            Inches(2.82),
            Inches(2.18),
            f"Level {index + 1}",
            text,
            fill=PAPER if index % 2 == 0 else PALE_GOLD,
            accent=SAGE if index < 2 else TEAL,
            title_size=14.2,
            body_size=10.8,
        )
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(5.70),
        Inches(12.25),
        Inches(0.98),
        "One Thing I Learned Today",
        plan_slide.get("secondary_text", "") or "One thing I learned today:",
        lines=2,
        fill=PALE_BLUE,
    )


def render_practice_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    render_mode = problem_render_mode(plan_slide)
    if render_mode == PROBLEM_RENDER_MODE_INTERACTIVE:
        render_interactive_apply_slide(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            image_lookup=image_lookup,
        )
        return
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    if template_role == "interactive_activity":
        plan_slide = {
            **plan_slide,
            "template_role": normalize_whitespace(plan_slide.get("interactive_render_role", "")) or "drag_sort",
        }
        template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    elif template_role == "best_fit_review":
        plan_slide = {
            **plan_slide,
            "template_role": normalize_whitespace(plan_slide.get("review_render_role", "")) or "turn_and_teach",
        }
        template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    if template_role == "drag_sort":
        render_drag_sort_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "error_analysis":
        render_error_analysis_workbook_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "two_column_compare":
        render_two_column_compare_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "choice_board":
        render_choice_board_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "collaborative_practice":
        render_collaborative_practice_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "independent_practice":
        render_independent_practice_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if premium_layout_uses_full_spread(plan_slide):
        render_premium_feature_spread(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            accent=GOLD,
            image_lookup=image_lookup,
        )
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=GOLD,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    asset = choose_image_asset(plan_slide, image_lookup)
    problem_cards = problem_display_cards(plan_slide, variant="practice")
    add_problem_prompt_stack(
        slide,
        x=Inches(0.5),
        y=Inches(2.0),
        w=Inches(7.55),
        h=Inches(1.56),
        title="Solve the Problems",
        problems=problem_cards,
        accent=GOLD,
        fill=PAPER,
        subtitle=workbook_prompt_stack_subtitle("practice"),
        kicker=workbook_prompt_stack_kicker("practice"),
    )
    if plan_slide.get("premium_layout") == "turn_and_teach" or plan_slide.get("partner_prompt"):
        render_premium_panel(
            slide,
            plan_slide=plan_slide,
            x=Inches(8.30),
            y=Inches(2.0),
            w=Inches(4.45),
            h=Inches(1.56),
            accent=TEAL,
        )
    elif plan_slide.get("discussion_questions"):
        add_discussion_panel(
            slide,
            plan_slide=plan_slide,
            x=Inches(8.30),
            y=Inches(2.0),
            w=Inches(4.45),
            h=Inches(1.56),
            fill=PALE_BLUE,
            accent=TEAL,
            max_questions=2,
            include_prompt=True,
        )
    else:
        add_card(
            slide,
            Inches(8.30),
            Inches(2.0),
            Inches(4.45),
            Inches(1.56),
            "Helpful Hints",
            problem_tip_body(plan_slide, include_visual_tip=bool(asset)),
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=15.5,
            body_size=11.8,
        )
    render_problem_workbook_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(0.5),
        y=Inches(3.80),
        w=Inches(7.55),
        h=Inches(2.98),
        accent=GOLD,
        variant="practice",
    )
    if render_mode == PROBLEM_RENDER_MODE_FOCUS:
        add_card(
            slide,
            Inches(8.30),
            Inches(3.80),
            Inches(4.45),
            Inches(2.98),
            "Next Step",
            "Use the next slide to apply this full problem with draggable pieces and written justification.",
            fill=PALE_GOLD,
            accent=SAGE,
            title_size=15.5,
            body_size=12.0,
        )
    else:
        interactive_slide = activity_slide_for_render(plan_slide)
        if has_activity(interactive_slide):
            render_activity_board(
                slide,
                plan_slide=interactive_slide,
                x=Inches(8.30),
                y=Inches(3.80),
                w=Inches(4.45),
                h=Inches(2.98),
            )
        else:
            add_card(
                slide,
                Inches(8.30),
                Inches(3.80),
                Inches(4.45),
                Inches(2.98),
                "Try a Similar One",
                truncate_display_copy(similar_problem_statement(plan_slide), 168),
                fill=PALE_GOLD,
                accent=SAGE,
                title_size=15.5,
                body_size=12.4,
            )


def render_quick_review_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=GOLD,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(5.4),
        Inches(4.85),
        "Recall the Big Idea",
        plan_slide["primary_text"] or "Use this space for key review notes from Session 1.",
        fill=PALE_GOLD,
        accent=GOLD,
        title_size=17,
        body_size=14.8,
    )
    review_tasks = merge_fragmented_items(plan_slide["tasks"][:2] or plan_slide["bullets"][:2], limit=2)
    first_prompt = review_tasks[0] if review_tasks else "What is one thing you remember clearly?"
    second_prompt = review_tasks[1] if len(review_tasks) > 1 else "What is one question you still have?"
    add_lined_area(
        slide,
        Inches(6.15),
        Inches(2.0),
        Inches(3.0),
        Inches(2.15),
        "Quick Check 1",
        first_prompt,
        lines=2,
    )
    add_lined_area(
        slide,
        Inches(9.4),
        Inches(2.0),
        Inches(3.35),
        Inches(2.15),
        "Quick Check 2",
        second_prompt,
        lines=2,
    )
    if render_premium_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(6.15),
        y=Inches(4.46),
        w=Inches(6.6),
        h=Inches(2.29),
        accent=GOLD,
    ):
        pass
    elif has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(6.15),
            y=Inches(4.46),
            w=Inches(6.6),
            h=Inches(2.29),
        )
    else:
        add_lined_area(
            slide,
            Inches(6.15),
            Inches(4.52),
            Inches(6.6),
            Inches(2.23),
            "Connect to Today",
            plan_slide["response_prompt"] or "How will this review help you with today's work?",
            lines=4,
        )


def render_challenge_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
    image_lookup: dict[int, dict[str, Any]],
) -> None:
    render_mode = problem_render_mode(plan_slide)
    if render_mode == PROBLEM_RENDER_MODE_INTERACTIVE:
        render_interactive_apply_slide(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            image_lookup=image_lookup,
        )
        return
    if premium_layout_uses_full_spread(plan_slide):
        render_premium_feature_spread(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            accent=TEAL,
            image_lookup=image_lookup,
        )
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=TEAL,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    add_image_panel(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(4.55),
        Inches(1.56),
        choose_image_asset(plan_slide, image_lookup),
        fill=PALE_GOLD,
        label=plan_slide["image_caption"] or "Source challenge image",
    )
    challenge_cards = problem_display_cards(plan_slide, variant="practice")
    add_problem_prompt_stack(
        slide,
        x=Inches(5.30),
        y=Inches(2.0),
        w=Inches(7.45),
        h=Inches(1.56),
        title="Challenge Problems",
        problems=challenge_cards,
        accent=GOLD,
        fill=PAPER,
        subtitle=workbook_prompt_stack_subtitle("practice"),
        kicker="Challenge Page",
    )
    render_problem_workbook_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(0.5),
        y=Inches(3.82),
        w=Inches(7.18),
        h=Inches(2.96),
        accent=TEAL,
        variant="practice",
    )
    if render_mode == PROBLEM_RENDER_MODE_FOCUS:
        add_card(
            slide,
            Inches(7.95),
            Inches(3.82),
            Inches(4.80),
            Inches(2.96),
            "Next Step",
            "Use the next slide to apply this challenge with draggable pieces and a written explanation.",
            fill=PALE_BLUE,
            accent=TEAL,
            title_size=15.8,
            body_size=12.4,
        )
    else:
        interactive_slide = activity_slide_for_render(plan_slide)
        if has_activity(interactive_slide):
            render_activity_board(
                slide,
                plan_slide=interactive_slide,
                x=Inches(7.95),
                y=Inches(3.82),
                w=Inches(4.80),
                h=Inches(2.96),
            )
        else:
            add_discussion_panel(
                slide,
                plan_slide=plan_slide,
                x=Inches(7.95),
                y=Inches(3.82),
                w=Inches(4.80),
                h=Inches(2.96),
                fill=PALE_BLUE,
                accent=TEAL,
                max_questions=2,
                include_prompt=True,
            )


def render_reflection_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    template_role = normalize_whitespace(plan_slide.get("template_role", ""))
    if template_role == "twr_frame":
        render_twr_frame_workbook_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if template_role == "goal_tracker":
        render_goal_tracker_workbook_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=SAGE,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    add_lined_area(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(5.85),
        Inches(3.0),
        "Reflection 1",
        plan_slide["primary_text"] or "What idea feels strongest to you right now?",
        lines=5,
    )
    add_lined_area(
        slide,
        Inches(6.62),
        Inches(2.0),
        Inches(6.13),
        Inches(3.0),
        "Reflection 2",
        plan_slide["secondary_text"] or "What do you still want to review or apply next?",
        lines=5,
    )
    if render_premium_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(0.5),
        y=Inches(5.12),
        w=Inches(12.25),
        h=Inches(1.73),
        accent=SAGE,
    ):
        pass
    elif has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(0.5),
            y=Inches(5.12),
            w=Inches(12.25),
            h=Inches(1.73),
        )
    else:
        add_discussion_panel(
            slide,
            plan_slide=plan_slide,
            x=Inches(0.5),
            y=Inches(5.12),
            w=Inches(12.25),
            h=Inches(1.73),
            fill=PALE_SAGE,
            accent=SAGE,
            max_questions=2,
            include_prompt=True,
        )


def render_exit_ticket_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    if normalize_whitespace(plan_slide.get("template_role", "")) == "tiered_exit":
        render_tiered_exit_workbook_slide(slide, plan_slide=plan_slide, page=page, footer_text=footer_text)
        return
    if premium_layout_uses_full_spread(plan_slide):
        render_premium_feature_spread(
            slide,
            plan_slide=plan_slide,
            page=page,
            footer_text=footer_text,
            accent=CORAL,
        )
        return
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=CORAL,
        footer_text=footer_text,
    )
    add_premium_meta_chips(slide, plan_slide)
    problem_cards = problem_display_cards(plan_slide, variant="exit")
    add_problem_prompt_stack(
        slide,
        x=Inches(0.5),
        y=Inches(2.0),
        w=Inches(4.35),
        h=Inches(4.85),
        title="Final Check",
        problems=problem_cards,
        accent=CORAL,
        fill=PAPER,
        subtitle=workbook_prompt_stack_subtitle("exit"),
        kicker=workbook_prompt_stack_kicker("exit"),
    )
    render_problem_workbook_panel(
        slide,
        plan_slide=plan_slide,
        x=Inches(5.10),
        y=Inches(2.0),
        w=Inches(4.18),
        h=Inches(4.85),
        accent=CORAL,
        variant="exit",
    )
    interactive_slide = activity_slide_for_render(plan_slide)
    if has_activity(interactive_slide):
        render_activity_board(
            slide,
            plan_slide=interactive_slide,
            x=Inches(9.52),
            y=Inches(2.0),
            w=Inches(3.23),
            h=Inches(4.85),
        )
    else:
        add_lined_area(
            slide,
            Inches(9.52),
            Inches(2.0),
            Inches(3.23),
            Inches(4.85),
            "Teacher Check",
            plan_slide["response_prompt"] or plan_slide["secondary_text"] or "Add one last sentence that shows why your answer makes sense.",
            lines=6,
        )


def render_generic_slide(
    slide: Any,
    *,
    plan_slide: dict[str, Any],
    page: int,
    footer_text: str,
) -> None:
    accent = KIND_ACCENT.get(plan_slide["kind"], TEAL)
    add_header(
        slide,
        section=plan_slide["section"],
        title=plan_slide["title"],
        subtitle=plan_slide["subtitle"],
        page=page,
        accent=accent,
        footer_text=footer_text,
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(2.0),
        Inches(5.7),
        Inches(4.85),
        "Main Idea",
        plan_slide["primary_text"] or bullet_text(plan_slide["bullets"], limit=4, max_len=95),
        fill=PAPER,
        accent=accent,
        title_size=17,
        body_size=14.4,
    )
    add_card(
        slide,
        Inches(6.45),
        Inches(2.0),
        Inches(6.3),
        Inches(2.8),
        "Details",
        bullet_text(plan_slide["bullets"], limit=4, max_len=95)
        or plan_slide["secondary_text"]
        or "Use this space for supporting details.",
        fill=PALE_BLUE,
        accent=TEAL,
        title_size=17,
        body_size=14.2,
    )
    if has_activity(plan_slide):
        render_activity_board(
            slide,
            plan_slide=plan_slide,
            x=Inches(6.45),
            y=Inches(5.0),
            w=Inches(6.3),
            h=Inches(1.85),
        )
    else:
        add_lined_area(
            slide,
            Inches(6.45),
            Inches(5.05),
            Inches(6.3),
            Inches(1.8),
            "Respond",
            plan_slide["response_prompt"] or "Write one thoughtful response here.",
            lines=3,
        )


def build_image_lookup(deck: dict[str, Any]) -> dict[int, dict[str, Any]]:
    lookup: dict[int, dict[str, Any]] = {}
    for slide in deck["slides"]:
        images = slide.get("images") or []
        slide_number = slide.get("slide_number")
        if not images or not slide_number:
            continue
        lookup[slide_number] = max(
            images,
            key=lambda asset: asset["pixel_width"] * asset["pixel_height"],
        )
    return lookup


def rgb_to_hex(color: RGBColor) -> str:
    return "#{:02x}{:02x}{:02x}".format(int(color[0]), int(color[1]), int(color[2]))


def htmlize_text(text: str) -> str:
    return "<br>".join(html_escape(line) for line in (text or "").split("\n"))


def render_html_list(items: list[str], *, title: str, class_name: str = "bullet-list") -> str:
    cleaned = unique_nonempty(items, limit=6)
    if not cleaned:
        return ""
    rows = "".join(f"<li>{html_escape(item)}</li>" for item in cleaned)
    return (
        f'<section class="info-card">'
        f'<h4>{html_escape(title)}</h4>'
        f'<ul class="{class_name}">{rows}</ul>'
        f"</section>"
    )


def render_html_chips(items: list[str], *, class_name: str = "chip-row") -> str:
    cleaned = unique_nonempty(items, limit=8)
    if not cleaned:
        return ""
    return (
        f'<div class="{class_name}">'
        + "".join(f'<span class="mini-chip">{html_escape(item)}</span>' for item in cleaned)
        + "</div>"
    )


def render_html_image(plan_slide: dict[str, Any], image_lookup: dict[int, dict[str, Any]]) -> str:
    asset = choose_image_asset(plan_slide, image_lookup)
    if not asset:
        return ""
    caption = plan_slide.get("image_caption") or f"Source image from slide {asset['slide_number']}"
    return (
        '<figure class="slide-figure">'
        f'<img src="{html_escape(asset["relative_path"])}" alt="{html_escape(caption)}" />'
        f'<figcaption>{html_escape(caption)}</figcaption>'
        "</figure>"
    )


def html_activity_zone_labels(family: str) -> list[str]:
    return {
        "sort_classify": ["Category A", "Category B"],
        "match_pair": ["Match A", "Match B"],
        "sequence_order": ["Step 1", "Step 2", "Step 3"],
        "build_construct": ["Build Zone", "Support Pieces"],
        "plot_place": ["Placement Zone", "Label Zone"],
        "detect_justify": ["Spot It", "Justify It"],
        "compare_rank": ["Most Convincing", "Middle", "Least Convincing"],
        "reveal_discuss": ["Reveal 1", "Reveal 2", "Reveal 3"],
    }.get(family, ["Workspace"])


def render_html_activity_board(plan_slide: dict[str, Any], slide_id: str) -> str:
    if not has_activity(plan_slide):
        return ""
    family = plan_slide.get("activity_family", "")
    accent = rgb_to_hex(activity_accent(plan_slide))
    fill = rgb_to_hex(activity_fill(plan_slide))
    help_id = f"{slide_id}-help"
    status_id = f"{slide_id}-status"
    zones = "".join(
        f'<div class="dropzone" data-dropzone data-zone-label="{html_escape(label)}" tabindex="0" role="button" aria-label="Place the selected piece in {html_escape(label)}"><span>{html_escape(label)}</span></div>'
        for label in html_activity_zone_labels(family)
    )
    pieces = activity_pieces(plan_slide, limit=6) or ["Move", "Match", "Explain"]
    bank_id = f"{slide_id}-bank"
    chips = "".join(
        f'<button class="activity-chip" type="button" draggable="true" data-home="{bank_id}" aria-pressed="false" aria-label="Movable piece: {html_escape(piece)}">{html_escape(piece)}</button>'
        for piece in pieces
    )
    answer_check = plan_slide.get("answer_check", "")
    return f"""
    <section class="activity-shell family-{html_escape(family or 'default')}" style="--activity-accent:{accent}; --activity-fill:{fill};" id="{html_escape(slide_id)}" aria-describedby="{html_escape(help_id)} {html_escape(status_id)}">
      <div class="activity-top">
        <div>
          <p class="activity-kicker">Interactive Notebook Move</p>
          <h4>{html_escape(plan_slide["activity_name"])}</h4>
        </div>
        <button class="ghost-button" type="button" data-reset-target="{html_escape(slide_id)}">Reset Pieces</button>
      </div>
      <p class="activity-copy">{html_escape(plan_slide.get("activity_instructions") or "Use the pieces to interact with the lesson content.")}</p>
      <p class="activity-help" id="{html_escape(help_id)}">Drag pieces, or select a piece and press Enter on a drop zone to place it.</p>
      <div class="activity-zones">{zones}</div>
      <div class="chip-bank" id="{html_escape(bank_id)}" data-chip-bank tabindex="0" aria-label="Piece bank. Move the selected piece back here to reset it manually.">{chips}</div>
      <p class="sr-only" id="{html_escape(status_id)}" data-activity-status aria-live="polite"></p>
      {f'<details class="answer-panel"><summary>Check your thinking</summary><p>{html_escape(answer_check)}</p></details>' if answer_check else ''}
    </section>
    """


def render_html_vocabulary_table(vocabulary: list[dict[str, str]]) -> str:
    if not vocabulary:
        return ""
    rows = "".join(
        "<tr>"
        f"<td>{html_escape(item.get('word', ''))}</td>"
        f"<td>{html_escape(item.get('definition', ''))}</td>"
        f"<td>{html_escape(item.get('example', ''))}</td>"
        f"<td>{html_escape(item.get('visual_cue', ''))}</td>"
        "</tr>"
        for item in vocabulary[:5]
    )
    return f"""
    <section class="info-card vocab-card">
      <h4>Vocabulary in Context</h4>
      <div class="table-wrap">
        <table>
          <thead>
            <tr>
              <th>Word</th>
              <th>Definition</th>
              <th>Example</th>
              <th>Visual Cue</th>
            </tr>
          </thead>
          <tbody>{rows}</tbody>
        </table>
      </div>
    </section>
    """


def render_html_content_stack(plan_slide: dict[str, Any]) -> str:
    blocks: list[str] = []
    if plan_slide.get("primary_text"):
        blocks.append(
            f'<section class="info-card lead-card"><h4>Main Idea</h4><p>{htmlize_text(plan_slide["primary_text"])}</p></section>'
        )
    if plan_slide.get("secondary_text"):
        blocks.append(
            f'<section class="info-card"><h4>Look For</h4><p>{htmlize_text(plan_slide["secondary_text"])}</p></section>'
        )
    if plan_slide.get("vocabulary"):
        blocks.append(render_html_vocabulary_table(plan_slide["vocabulary"]))
    if plan_slide.get("bullets"):
        blocks.append(render_html_list(plan_slide["bullets"], title="Key Moves"))
    if plan_slide.get("tasks"):
        blocks.append(render_html_list(plan_slide["tasks"], title="Try It Next"))
    if plan_slide.get("response_prompt"):
        blocks.append(
            f'<section class="info-card prompt-card"><h4>Respond</h4><p>{html_escape(plan_slide["response_prompt"])}</p></section>'
        )
    if plan_slide.get("sentence_starters"):
        blocks.append(
            '<section class="info-card starter-card"><h4>Kernel Sentences</h4>'
            + render_html_chips(plan_slide["sentence_starters"])
            + "</section>"
        )
    return "".join(blocks)


def render_html_cover_slide(
    plan_slide: dict[str, Any],
    *,
    session_label: str,
    image_lookup: dict[int, dict[str, Any]],
) -> str:
    image_markup = render_html_image(plan_slide, image_lookup)
    bullet_markup = render_html_list(plan_slide.get("bullets", []), title="Session Focus")
    return f"""
    <article class="slide-card cover-slide">
      <div class="cover-copy">
        <p class="slide-kicker">{html_escape(session_label)} Student Notebook</p>
        <h3>{html_escape(plan_slide["title"])}</h3>
        <p class="slide-subtitle">{html_escape(plan_slide["subtitle"] or plan_slide["primary_text"])}</p>
        <div class="cover-panels">
          <section class="info-card lead-card">
            <h4>Focus</h4>
            <p>{htmlize_text(plan_slide.get("primary_text", ""))}</p>
          </section>
          <section class="info-card">
            <h4>Look For</h4>
            <p>{htmlize_text(plan_slide.get("secondary_text", ""))}</p>
          </section>
        </div>
        {bullet_markup}
      </div>
      <div class="cover-visual">
        {image_markup or '<div class="empty-visual">Source image will appear here when available.</div>'}
      </div>
    </article>
    """


def render_html_slide(
    plan_slide: dict[str, Any],
    *,
    session_label: str,
    page_index: int,
    image_lookup: dict[int, dict[str, Any]],
) -> str:
    if plan_slide["kind"] == "cover":
        return render_html_cover_slide(plan_slide, session_label=session_label, image_lookup=image_lookup)

    slide_id = slugify(f"{session_label}-{page_index}-{plan_slide['title']}")
    image_markup = render_html_image(plan_slide, image_lookup)
    content_stack = render_html_content_stack(plan_slide)
    activity_markup = render_html_activity_board(plan_slide, slide_id)
    source_refs = ", ".join(str(number) for number in plan_slide.get("source_slide_numbers", []))
    activity_badge = (
        f'<span class="meta-pill accent">{html_escape(plan_slide["activity_name"])}</span>'
        if has_activity(plan_slide)
        else ""
    )
    return f"""
    <article class="slide-card kind-{html_escape(plan_slide['kind'])}">
      <header class="slide-top">
        <div class="slide-top-left">
          <span class="meta-pill">{html_escape(plan_slide["section"])}</span>
          <span class="meta-pill muted">{html_escape(plan_slide["kind"].replace("_", " ").title())}</span>
          {activity_badge}
        </div>
        <span class="page-pill">Page {page_index}</span>
      </header>
      <h3>{html_escape(plan_slide["title"])}</h3>
      <p class="slide-subtitle">{html_escape(plan_slide["subtitle"])}</p>
      <div class="slide-layout {'with-image' if image_markup else 'no-image'}">
        {f'<div class="visual-column">{image_markup}</div>' if image_markup else ''}
        <div class="content-column">{content_stack}</div>
      </div>
      {activity_markup}
    </article>
    """


def render_html_session(
    session_plan: dict[str, Any],
    *,
    session_label: str,
    image_lookup: dict[int, dict[str, Any]],
) -> str:
    slide_html = "".join(
        render_html_slide(
            plan_slide,
            session_label=session_label,
            page_index=index,
            image_lookup=image_lookup,
        )
        for index, plan_slide in enumerate(session_plan["slides"], start=1)
    )
    session_id = slugify(session_label)
    return f"""
    <section class="session-shell" id="{html_escape(session_id)}">
      <header class="session-header">
        <div>
          <p class="session-kicker">{html_escape(session_label)}</p>
          <h2>{html_escape(session_plan["session_title"])}</h2>
          <p>{html_escape(session_plan["session_subtitle"])}</p>
        </div>
      </header>
      <div class="slides-column">{slide_html}</div>
    </section>
    """


def collect_activity_names(plan: dict[str, Any], limit: int = 10) -> list[str]:
    names: list[str] = []
    for session_key in planned_session_keys(plan):
        for slide in plan.get(session_key, {}).get("slides", []):
            if slide.get("activity_name"):
                names.append(slide["activity_name"])
    return unique_nonempty(names, limit=limit)


def render_html_notebook(plan: dict[str, Any], deck: dict[str, Any], output_dir: Path) -> Path:
    output_path = output_dir / "Flagship Student Notebook.html"
    image_lookup = build_image_lookup(deck)
    standards_markup = render_html_chips(plan.get("standards", []), class_name="standards-row")
    activity_markup = render_html_chips(collect_activity_names(plan), class_name="hero-activity-row")
    session_links = "".join(
        f'<a href="#{label.lower().replace(" ", "-")}">{label}</a>'
        for _session_key, label, _output_key in planned_session_specs(plan)
    )
    session_html = "".join(
        render_html_session(plan[session_key], session_label=label, image_lookup=image_lookup)
        for session_key, label, _output_key in planned_session_specs(plan)
    )
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <meta name="description" content="{html_escape(plan.get('topic_summary', plan['lesson_title']))}" />
  <title>{html_escape(plan["lesson_title"])} - Flagship Student Notebook</title>
  <style>
    :root {{
      --bg:#fffaf1;
      --bg-2:#f9efd8;
      --panel:#ffffff;
      --panel-2:#fff6e8;
      --text:#27404a;
      --muted:#5c7a87;
      --navy:{rgb_to_hex(NAVY)};
      --teal:{rgb_to_hex(TEAL)};
      --gold:{rgb_to_hex(GOLD)};
      --coral:{rgb_to_hex(CORAL)};
      --sage:{rgb_to_hex(SAGE)};
      --line:rgba(39,63,74,.10);
      --shadow:0 24px 70px rgba(30,95,116,.10);
      --radius-xl:30px;
      --radius-lg:22px;
      --radius-md:16px;
      --font-ui:"Arial","Helvetica Neue",sans-serif;
      --font-display:"Arial","Helvetica Neue",sans-serif;
    }}
    * {{ box-sizing:border-box; }}
    html {{ scroll-behavior:smooth; }}
    body {{
      margin:0;
      font-family:var(--font-ui);
      color:var(--text);
      background:
        radial-gradient(circle at 84% 8%, rgba(249,168,38,.20), transparent 22%),
        radial-gradient(circle at 91% 78%, rgba(249,168,38,.18), transparent 14%),
        linear-gradient(180deg, #fffef9 0%, var(--bg) 42%, var(--bg-2) 100%);
    }}
    a, button {{
      transition:transform .16s ease, box-shadow .16s ease, background-color .16s ease;
    }}
    a:focus-visible, button:focus-visible {{
      outline:3px solid rgba(56,127,132,.45);
      outline-offset:3px;
    }}
    .skip-link {{
      position:absolute;
      left:12px;
      top:-48px;
      padding:10px 14px;
      border-radius:14px;
      background:var(--navy);
      color:#fff;
      text-decoration:none;
      z-index:10;
    }}
    .skip-link:focus {{
      top:12px;
    }}
    .page-shell {{
      width:min(1240px, calc(100% - 28px));
      margin:0 auto;
      padding:26px 0 42px;
    }}
    .hero {{
      position:relative;
      overflow:hidden;
      border-radius:var(--radius-xl);
      border:1px solid var(--line);
      background:linear-gradient(140deg, rgba(255,253,248,.96), rgba(248,237,220,.98));
      box-shadow:var(--shadow);
      padding:28px;
    }}
    .hero::after {{
      content:"";
      position:absolute;
      width:320px;
      height:320px;
      right:-72px;
      top:-36px;
      border-radius:40px;
      background:
        linear-gradient(145deg, rgba(198,91,61,.18), rgba(56,127,132,.12)),
        radial-gradient(circle at 32% 26%, rgba(255,255,255,.78), transparent 36%);
      transform:rotate(18deg);
      border:1px solid rgba(255,255,255,.52);
    }}
    .hero-top {{
      display:flex;
      justify-content:space-between;
      gap:18px;
      align-items:flex-start;
      position:relative;
      z-index:1;
    }}
    .eyebrow {{
      display:inline-flex;
      align-items:center;
      gap:8px;
      padding:8px 12px;
      border-radius:999px;
      background:rgba(56,127,132,.09);
      border:1px solid rgba(56,127,132,.16);
      color:var(--teal);
      font-size:12px;
      font-weight:800;
      letter-spacing:.08em;
      text-transform:uppercase;
    }}
    h1 {{
      margin:16px 0 12px;
      font-family:var(--font-display);
      font-size:clamp(40px, 7vw, 78px);
      line-height:.92;
      letter-spacing:-.05em;
      max-width:10ch;
    }}
    .hero-copy {{
      max-width:64ch;
      color:var(--muted);
      line-height:1.65;
      margin:0;
    }}
    .hero-meta, .standards-row, .chip-row {{
      display:flex;
      flex-wrap:wrap;
      gap:10px;
    }}
    .hero-meta {{ margin-top:18px; }}
    .hero-activity-row {{
      display:flex;
      flex-wrap:wrap;
      gap:10px;
      margin-top:14px;
      position:relative;
      z-index:1;
    }}
    .meta-pill, .mini-chip {{
      display:inline-flex;
      align-items:center;
      gap:6px;
      padding:8px 12px;
      border-radius:999px;
      border:1px solid var(--line);
      background:rgba(255,255,255,.72);
      font-size:12px;
      font-weight:800;
      letter-spacing:.05em;
      text-transform:uppercase;
    }}
    .meta-pill.accent {{ color:var(--coral); }}
    .meta-pill.muted {{ color:var(--muted); }}
    .session-nav {{
      display:flex;
      gap:12px;
      flex-wrap:wrap;
      margin-top:18px;
      position:relative;
      z-index:1;
    }}
    .session-nav a {{
      text-decoration:none;
      color:var(--text);
      font-weight:800;
      padding:12px 16px;
      border-radius:16px;
      background:rgba(255,255,255,.76);
      border:1px solid var(--line);
    }}
    .session-nav a:hover {{
      transform:translateY(-1px);
    }}
    .session-shell {{
      margin-top:26px;
    }}
    .session-header {{
      margin-bottom:16px;
      padding:18px 4px 2px;
    }}
    .session-kicker {{
      margin:0 0 6px;
      color:var(--coral);
      font-size:12px;
      font-weight:800;
      letter-spacing:.08em;
      text-transform:uppercase;
    }}
    .session-header h2 {{
      margin:0 0 8px;
      font-family:var(--font-display);
      font-size:clamp(28px, 4vw, 44px);
      letter-spacing:-.04em;
    }}
    .session-header p {{
      margin:0;
      color:var(--muted);
      max-width:62ch;
      line-height:1.6;
    }}
    .slides-column {{
      display:grid;
      gap:20px;
    }}
    .slide-card {{
      position:relative;
      border-radius:var(--radius-lg);
      border:1px solid var(--line);
      background:linear-gradient(180deg, rgba(255,253,248,.97), rgba(255,248,239,.98));
      box-shadow:var(--shadow);
      padding:20px;
    }}
    .slide-card::before {{
      content:"";
      position:absolute;
      left:0;
      right:0;
      top:0;
      height:5px;
      border-radius:22px 22px 0 0;
      background:linear-gradient(90deg, var(--navy), rgba(255,255,255,0));
    }}
    .slide-card.kind-be_curious::before,
    .slide-card.kind-worked_example::before,
    .cover-slide::before {{
      background:linear-gradient(90deg, var(--coral), rgba(255,255,255,0));
    }}
    .slide-card.kind-learning_target::before,
    .slide-card.kind-guided_notes::before,
    .slide-card.kind-challenge::before {{
      background:linear-gradient(90deg, var(--teal), rgba(255,255,255,0));
    }}
    .slide-card.kind-practice::before,
    .slide-card.kind-quick_review::before {{
      background:linear-gradient(90deg, var(--gold), rgba(255,255,255,0));
    }}
    .slide-card.kind-reflection::before,
    .slide-card.kind-vocabulary::before {{
      background:linear-gradient(90deg, var(--sage), rgba(255,255,255,0));
    }}
    .cover-slide {{
      display:grid;
      grid-template-columns:1.05fr .95fr;
      gap:18px;
      align-items:start;
    }}
    .cover-copy {{
      display:grid;
      gap:14px;
      align-content:start;
    }}
    .cover-panels {{
      display:grid;
      gap:12px;
    }}
    .cover-visual,
    .visual-column {{
      display:grid;
      gap:14px;
      align-content:start;
    }}
    .slide-kicker {{
      margin:0;
      color:var(--coral);
      font-size:12px;
      font-weight:800;
      letter-spacing:.08em;
      text-transform:uppercase;
    }}
    .slide-top {{
      display:flex;
      justify-content:space-between;
      gap:10px;
      align-items:flex-start;
    }}
    .slide-top-left {{
      display:flex;
      gap:8px;
      flex-wrap:wrap;
    }}
    .page-pill {{
      display:inline-flex;
      align-items:center;
      justify-content:center;
      min-width:86px;
      padding:10px 14px;
      border-radius:999px;
      background:rgba(232,237,243,.95);
      color:var(--navy);
      font-size:12px;
      font-weight:800;
      text-transform:uppercase;
      letter-spacing:.06em;
    }}
    .slide-card h3 {{
      margin:14px 0 8px;
      font-family:var(--font-display);
      font-size:clamp(28px, 3vw, 40px);
      letter-spacing:-.04em;
    }}
    .slide-subtitle {{
      margin:0 0 14px;
      color:var(--muted);
      line-height:1.6;
    }}
    .slide-layout {{
      display:grid;
      gap:16px;
      align-items:start;
    }}
    .slide-layout.with-image {{
      grid-template-columns:minmax(280px, .9fr) minmax(0, 1.1fr);
    }}
    .content-column {{
      display:grid;
      gap:14px;
    }}
    .slide-figure {{
      margin:0;
      border-radius:22px;
      overflow:hidden;
      border:1px solid var(--line);
      background:rgba(255,255,255,.84);
    }}
    .slide-figure img {{
      display:block;
      width:100%;
      height:auto;
      max-height:420px;
      object-fit:cover;
      background:#fff;
    }}
    .slide-figure figcaption {{
      padding:10px 12px 12px;
      font-size:12px;
      color:var(--muted);
    }}
    .info-card {{
      border-radius:18px;
      border:1px solid var(--line);
      background:rgba(255,255,255,.82);
      padding:16px;
    }}
    .lead-card {{
      background:linear-gradient(180deg, rgba(234,242,247,.88), rgba(255,255,255,.82));
    }}
    .prompt-card {{
      background:linear-gradient(180deg, rgba(252,245,224,.95), rgba(255,255,255,.82));
    }}
    .starter-card {{
      background:linear-gradient(180deg, rgba(252,236,230,.95), rgba(255,255,255,.82));
    }}
    .info-card h4 {{
      margin:0 0 8px;
      font-size:13px;
      text-transform:uppercase;
      letter-spacing:.08em;
      color:var(--muted);
    }}
    .info-card p {{
      margin:0;
      line-height:1.65;
    }}
    .bullet-list {{
      margin:0;
      padding-left:18px;
      display:grid;
      gap:8px;
      line-height:1.55;
    }}
    .table-wrap {{
      overflow:auto;
    }}
    table {{
      width:100%;
      border-collapse:collapse;
      font-size:14px;
    }}
    th, td {{
      border:1px solid rgba(31,39,49,.10);
      padding:10px 12px;
      vertical-align:top;
    }}
    th {{
      background:var(--navy);
      color:#fff;
      text-align:left;
      font-size:12px;
      letter-spacing:.06em;
      text-transform:uppercase;
    }}
    .activity-shell {{
      margin-top:16px;
      border-radius:22px;
      border:1px solid rgba(31,39,49,.10);
      background:var(--activity-fill);
      box-shadow:inset 0 1px 0 rgba(255,255,255,.65);
      padding:16px;
    }}
    .activity-top {{
      display:flex;
      justify-content:space-between;
      gap:12px;
      align-items:flex-start;
    }}
    .activity-kicker {{
      margin:0 0 6px;
      color:var(--activity-accent);
      font-size:11px;
      font-weight:800;
      letter-spacing:.08em;
      text-transform:uppercase;
    }}
    .activity-top h4 {{
      margin:0;
      font-size:22px;
      letter-spacing:-.03em;
    }}
    .activity-copy {{
      margin:10px 0 14px;
      color:var(--muted);
      line-height:1.55;
    }}
    .activity-help {{
      margin:0 0 12px;
      font-size:13px;
      color:var(--muted);
    }}
    .activity-zones {{
      display:grid;
      grid-template-columns:repeat(auto-fit, minmax(160px, 1fr));
      gap:12px;
    }}
    .dropzone {{
      min-height:110px;
      border-radius:18px;
      border:1.5px dashed rgba(31,39,49,.16);
      background:rgba(255,255,255,.72);
      padding:12px;
      display:flex;
      flex-direction:column;
      gap:8px;
    }}
    .dropzone:hover,
    .dropzone:focus-visible,
    .dropzone.is-active,
    .chip-bank:focus-visible,
    .chip-bank.is-active {{
      border-color:rgba(56,127,132,.34);
      box-shadow:0 0 0 3px rgba(56,127,132,.12);
    }}
    .dropzone span {{
      font-size:11px;
      font-weight:800;
      text-transform:uppercase;
      letter-spacing:.08em;
      color:var(--muted);
    }}
    .chip-bank {{
      margin-top:12px;
      display:flex;
      flex-wrap:wrap;
      gap:10px;
      min-height:52px;
      padding:12px;
      border-radius:18px;
      border:1px solid rgba(31,39,49,.08);
      background:rgba(255,255,255,.72);
    }}
    .activity-chip {{
      appearance:none;
      border:1px solid rgba(31,39,49,.12);
      background:#fff;
      color:var(--text);
      border-radius:999px;
      padding:10px 14px;
      font:inherit;
      cursor:grab;
      box-shadow:0 10px 24px rgba(46,40,26,.08);
    }}
    .activity-chip:active {{
      cursor:grabbing;
    }}
    .activity-chip:hover {{
      transform:translateY(-1px);
    }}
    .activity-chip.is-selected,
    .activity-chip[aria-pressed="true"] {{
      border-color:rgba(56,127,132,.44);
      background:rgba(234,242,247,.94);
      box-shadow:0 14px 28px rgba(56,127,132,.18);
    }}
    .ghost-button {{
      appearance:none;
      border:1px solid rgba(31,39,49,.12);
      background:rgba(255,255,255,.7);
      padding:10px 14px;
      border-radius:14px;
      font:inherit;
      font-weight:800;
      cursor:pointer;
    }}
    .answer-panel {{
      margin-top:12px;
      border-radius:16px;
      border:1px solid rgba(31,39,49,.10);
      background:rgba(255,255,255,.72);
      padding:12px 14px;
    }}
    .answer-panel summary {{
      cursor:pointer;
      font-weight:800;
      color:var(--navy);
    }}
    .empty-visual {{
      min-height:280px;
      border-radius:24px;
      border:1px dashed rgba(31,39,49,.16);
      display:grid;
      place-items:center;
      background:rgba(255,255,255,.74);
      color:var(--muted);
      text-align:center;
      padding:20px;
    }}
    .sr-only {{
      position:absolute;
      width:1px;
      height:1px;
      padding:0;
      margin:-1px;
      overflow:hidden;
      clip:rect(0, 0, 0, 0);
      white-space:nowrap;
      border:0;
    }}
    @media (prefers-reduced-motion: reduce) {{
      html {{ scroll-behavior:auto; }}
      a, button {{ transition:none; }}
    }}
    @media (max-width: 960px) {{
      .cover-slide,
      .slide-layout.with-image {{
        grid-template-columns:1fr;
      }}
    }}
  </style>
</head>
<body>
  <a class="skip-link" href="#main">Skip to notebook</a>
  <div class="page-shell">
    <section class="hero">
      <div class="hero-top">
        <div>
          <div class="eyebrow">Flagship Student Notebook</div>
          <h1>{html_escape(plan["lesson_title"])}</h1>
          <p class="hero-copy">{html_escape(plan.get("topic_summary", ""))}</p>
        </div>
        <div class="hero-meta">
          <span class="meta-pill">{html_escape(plan.get("subject", "Lesson"))}</span>
          <span class="meta-pill">{html_escape(plan.get("grade_level", "Grade Level Unspecified"))}</span>
        </div>
      </div>
      {standards_markup}
      {activity_markup}
      <nav class="session-nav" aria-label="Session navigation">{session_links}</nav>
    </section>
    <main id="main">
      {session_html}
    </main>
  </div>
  <script>
    const chips = document.querySelectorAll('.activity-chip');
    let draggingChip = null;
    let selectedChip = null;
    function setSelectedChip(chip) {{
      chips.forEach((item) => {{
        const isSelected = item === chip;
        item.classList.toggle('is-selected', isSelected);
        item.setAttribute('aria-pressed', String(isSelected));
      }});
      selectedChip = chip;
    }}
    function clearSelectedChip() {{
      chips.forEach((item) => {{
        item.classList.remove('is-selected');
        item.setAttribute('aria-pressed', 'false');
      }});
      selectedChip = null;
    }}
    function announce(zone, message) {{
      const activity = zone.closest('.activity-shell');
      const status = activity?.querySelector('[data-activity-status]');
      if (status) {{
        status.textContent = message;
      }}
    }}
    function moveChip(zone, chip) {{
      if (!zone || !chip) return;
      zone.appendChild(chip);
      zone.classList.add('is-active');
      window.setTimeout(() => zone.classList.remove('is-active'), 220);
      announce(zone, `${{chip.textContent.trim()}} moved to ${{zone.dataset.zoneLabel || 'the selected area'}}.`);
      clearSelectedChip();
    }}
    chips.forEach((chip) => {{
      chip.addEventListener('dragstart', () => {{
        draggingChip = chip;
        chip.classList.add('is-dragging');
        setSelectedChip(chip);
      }});
      chip.addEventListener('dragend', () => {{
        chip.classList.remove('is-dragging');
        draggingChip = null;
      }});
      chip.addEventListener('click', () => {{
        setSelectedChip(selectedChip === chip ? null : chip);
      }});
    }});
    document.querySelectorAll('[data-dropzone], [data-chip-bank]').forEach((zone) => {{
      zone.addEventListener('dragover', (event) => {{
        event.preventDefault();
      }});
      zone.addEventListener('drop', (event) => {{
        event.preventDefault();
        if (draggingChip) {{
          moveChip(zone, draggingChip);
        }}
      }});
      zone.addEventListener('click', () => {{
        if (selectedChip) {{
          moveChip(zone, selectedChip);
        }}
      }});
      zone.addEventListener('keydown', (event) => {{
        if ((event.key === 'Enter' || event.key === ' ') && selectedChip) {{
          event.preventDefault();
          moveChip(zone, selectedChip);
        }}
      }});
    }});
    document.querySelectorAll('[data-reset-target]').forEach((button) => {{
      button.addEventListener('click', () => {{
        const target = document.getElementById(button.dataset.resetTarget);
        if (!target) return;
        const bank = target.querySelector('[data-chip-bank]');
        if (!bank) return;
        target.querySelectorAll('.activity-chip').forEach((chip) => {{
          if (chip.dataset.home === bank.id) {{
            bank.appendChild(chip);
          }}
        }});
        announce(bank, 'All pieces returned to the piece bank.');
        clearSelectedChip();
      }});
    }});
  </script>
</body>
</html>
"""
    output_path.write_text(html, encoding="utf-8")
    return output_path


def paragraph_line_spacing_value(paragraph: Any) -> float:
    value = getattr(paragraph, "line_spacing", None)
    if isinstance(value, (int, float)):
        return float(value)
    return DEFAULT_TEXT_LINE_SPACING


def paragraph_font_point_sizes(paragraph: Any) -> list[float]:
    sizes: list[float] = []
    paragraph_font = getattr(paragraph, "font", None)
    paragraph_size = getattr(paragraph_font, "size", None)
    if paragraph_size is not None and getattr(paragraph_size, "pt", 0):
        sizes.append(float(paragraph_size.pt))
    for run in getattr(paragraph, "runs", []):
        font_size = getattr(getattr(run, "font", None), "size", None)
        if font_size is not None and getattr(font_size, "pt", 0):
            sizes.append(float(font_size.pt))
    return sizes


def shape_text_value(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return normalize_whitespace(shape.text_frame.text)


def is_notebook_write_line(shape: Any) -> bool:
    return (
        getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE
        and emu_to_inches(int(getattr(shape, "height", 0))) < 0.04
        and emu_to_inches(int(getattr(shape, "width", 0))) >= 1.8
    )


def remove_shape(shape: Any) -> None:
    element = getattr(shape, "_element", None)
    if element is None:
        return
    parent = element.getparent()
    if parent is None:
        return
    parent.remove(element)


def strip_notebook_write_lines(prs: Presentation) -> int:
    removed = 0
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not is_notebook_write_line(shape):
                continue
            remove_shape(shape)
            removed += 1
    return removed


def run_session_design_review(prs: Presentation, session_plan: dict[str, Any], session_label: str) -> None:
    issues: list[str] = []
    session_slide_count = len(session_plan.get("slides", []))
    exact_template = uses_exact_esol_template(session_plan)
    activity_floor = min(4 if exact_template else 5, session_slide_count)
    activity_cap = min(
        MAX_ACTIVITY_SLIDES_PER_SESSION,
        activity_floor + 3 + (1 if session_slide_count >= PREMIUM_TARGET_SESSION_SLIDES else 0),
    )
    activity_names = [
        normalize_whitespace(slide.get("activity_name", ""))
        for slide in session_plan.get("slides", [])
        if normalize_whitespace(slide.get("activity_name", ""))
    ]
    activity_count = (
        reference_exact_engagement_slide_count(session_plan) if exact_template else len(activity_names)
    )
    if activity_count < activity_floor:
        issues.append(f"{session_label}: not enough named activity slides to support the premium default")
    if activity_count > activity_cap:
        issues.append(f"{session_label}: too many activity slides for a clean premium layout")
    engagement_count = session_engagement_slide_count(session_plan)
    engagement_modes = set(session_engagement_modes(session_plan))
    if exact_template:
        engagement_count = max(engagement_count, reference_exact_engagement_slide_count(session_plan))
        engagement_modes.update(reference_exact_engagement_modes(session_plan))
    if engagement_count < engagement_slide_target(session_plan):
        issues.append(f"{session_label}: not enough high-agency engagement pages across the session")
    if len(engagement_modes) < engagement_mode_target(session_plan):
        issues.append(f"{session_label}: engagement variety is too narrow across the session")
    premium_features = session_plan.get("premium_features", [])
    if premium_features and not 2 <= len(premium_features) <= 4:
        issues.append(f"{session_label}: premium decision layer selected too many or too few features")
    duplicate_activity_names = [name for name, count in Counter(activity_names).items() if count > 1]
    if duplicate_activity_names:
        issues.append(f"{session_label}: repeated activity names reduce premium variety")
    for issue in validate_flagship_activities(session_plan):
        issues.append(f"{session_label}: {issue}")
    required_kinds = required_kinds_for_session(session_plan)
    present_kinds = {slide.get("kind") for slide in session_plan.get("slides", [])}
    missing_kinds = sorted(required_kinds - present_kinds)
    if missing_kinds:
        issues.append(f"{session_label}: missing locked-architecture sections ({', '.join(missing_kinds)})")
    required_roles = required_template_roles_for_session(session_plan)
    if required_roles and template_role_signature(session_plan) != required_roles:
        issues.append(f"{session_label}: exact ESOL workbook sequence drifted")
    objective_slides = [slide for slide in session_plan.get("slides", []) if slide.get("kind") == "learning_target"]
    if not objective_slides:
        issues.append(f"{session_label}: missing learning_target slide")
    for slide in objective_slides:
        if not normalize_whitespace(slide.get("primary_text", "")):
            issues.append(f"{session_label}: missing content objective text")
        language_objective = normalize_whitespace(slide.get("secondary_text", ""))
        if not exact_template and not language_objective:
            issues.append(f"{session_label}: missing language objective text")
        if not is_i_can_objective(slide.get("primary_text", "")):
            issues.append(f"{session_label}: content objective must start with 'I can'")
        if language_objective and not is_i_can_objective(language_objective):
            issues.append(f"{session_label}: language objective must start with 'I can'")
    for page_index, (ppt_slide, plan_slide) in enumerate(zip(prs.slides, session_plan.get("slides", [])), start=1):
        template_role = normalize_whitespace(plan_slide.get("template_role", ""))
        picture_count = sum(1 for shape in ppt_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        wide_write_lines = sum(1 for shape in ppt_slide.shapes if is_notebook_write_line(shape))
        rendered_texts = [
            shape_text_value(shape)
            for shape in ppt_slide.shapes
            if getattr(shape, "has_text_frame", False) and shape_text_value(shape)
        ]
        rendered_text_blob = normalize_whitespace(" ".join(rendered_texts))
        if (
            plan_slide.get("kind") == "vocabulary"
            and not exact_template
            and plan_slide.get("image_source_slide")
            and picture_count == 0
        ):
            issues.append(f"{session_label} page {page_index}: vocabulary slide is missing the source visual")
        if wide_write_lines:
            issues.append(f"{session_label} page {page_index}: writing area still uses worksheet-style lines instead of boxed workspace")
        if plan_slide.get("kind") in PROBLEM_SOLVING_KINDS:
            workbook_prompts = workbook_problem_prompts(plan_slide, limit=3)
            full_targets = primary_source_problem_targets(plan_slide, limit=2)
            if action_prompt_count(workbook_prompts) < 2:
                issues.append(f"{session_label} page {page_index}: problem-solving slide is not solve-first enough")
            if len(normalize_whitespace(plan_slide.get("subtitle", ""))) > 76:
                issues.append(f"{session_label} page {page_index}: problem-solving subtitle is too long for a clean layout")
            if full_targets and not any(source_problem_text_overlap(rendered_text_blob, [target]) for target in full_targets):
                issues.append(f"{session_label} page {page_index}: rendered slide does not preserve the full source problem")
        if plan_slide.get("kind") == "cover":
            duplicate_targets = {
                display_text_key(session_label),
                display_text_key(plan_slide.get("title", "")),
            }
            for field_name in ("subtitle", "primary_text", "secondary_text"):
                key = display_text_key(plan_slide.get(field_name, ""))
                if key and key in duplicate_targets:
                    issues.append(f"{session_label} page {page_index}: cover {field_name} duplicates the title or session label")
                    break
        draggable_texts: set[str] = set()
        if has_activity(plan_slide) and plan_slide.get("kind") not in {"be_curious", "cover"}:
            draggable_texts = set(unique_nonempty(plan_slide.get("movable_pieces", []), limit=12))
        if plan_slide.get("kind") == "learning_target":
            draggable_texts.add(CHECKMARK_CHIP)
        non_textbox_shape_texts = {
            shape_text_value(shape)
            for shape in ppt_slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX
            and shape_text_value(shape)
        }
        non_textbox_shape_text_raws = {normalize_whitespace(text) for text in non_textbox_shape_texts if normalize_whitespace(text)}
        non_textbox_shape_text_keys = {display_text_key(text) for text in non_textbox_shape_texts if display_text_key(text)}
        text_shape_count = 0
        for shape in ppt_slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_value = shape_text_value(shape)
            if not text_value:
                continue
            text_shape_count += 1
            if int(shape.width) <= 0 or int(shape.height) <= 0:
                issues.append(f"{session_label} page {page_index}: invalid text shape geometry")
                continue
            if (
                plan_slide.get("kind") != "vocabulary"
                and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                and text_value in draggable_texts
                and text_value not in non_textbox_shape_texts
            ):
                issues.append(
                    f"{session_label} page {page_index}: draggable text '{truncate_text(text_value, 40)}' rendered as a separate text box"
                )
            text_issues = publisher_copyedit_issues(text_value)
            if text_issues:
                issues.append(f"{session_label} page {page_index}: copyedit issue in rendered text ({text_issues[0]})")
            if len(text_value) > 260 and plan_slide.get("kind") not in {"learning_target", "vocabulary"}:
                issues.append(f"{session_label} page {page_index}: oversized text block is likely to crowd the layout")
            paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if normalize_whitespace(paragraph.text)]
            if len(paragraphs) >= 2:
                for paragraph in paragraphs:
                    if paragraph_line_spacing_value(paragraph) < 1.08:
                        issues.append(f"{session_label} page {page_index}: multi-line text has cramped line spacing")
                        break
            small_font_found = False
            for paragraph in paragraphs:
                font_sizes = paragraph_font_point_sizes(paragraph)
                if font_sizes and min(font_sizes) < FORMAL_REVIEW_MIN_FONT_PT:
                    issues.append(f"{session_label} page {page_index}: text rendered below the minimum readable size")
                    small_font_found = True
                if small_font_found:
                    break
        if plan_slide.get("kind") == "learning_target" and not exact_template:
            if CHECKMARK_CHIP not in non_textbox_shape_text_raws:
                issues.append(f"{session_label} page {page_index}: learning-target checkmarks are not embedded in draggable shapes")
        elif not exact_template and has_activity(plan_slide) and draggable_texts and not normalize_whitespace(plan_slide.get("premium_layout", "")) and (
            not template_role or template_role in TEMPLATE_ROLES_WITH_DRAGGABLES
        ):
            if not any(display_text_key(piece) in non_textbox_shape_text_keys for piece in draggable_texts if display_text_key(piece)):
                issues.append(f"{session_label} page {page_index}: draggable text is not embedded in movable shapes")
        text_shape_cap = 52 if plan_slide.get("kind") in PROBLEM_SOLVING_KINDS else 44
        if has_activity(plan_slide) or normalize_whitespace(plan_slide.get("premium_layout", "")):
            text_shape_cap += 4
        if flagship_activity_spec(plan_slide):
            text_shape_cap += 4
        if text_shape_count > text_shape_cap and plan_slide.get("kind") not in {"learning_target", "vocabulary"}:
            issues.append(f"{session_label} page {page_index}: too many text objects for a clean premium page")

    if issues:
        raise RuntimeError("Notebook design review failed: " + " | ".join(issues[:8]))


def run_formal_release_render_review(prs: Presentation, session_plan: dict[str, Any], session_label: str) -> None:
    run_session_design_review(prs, session_plan, session_label)


def slide_number_from_pptx_name(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def extract_pptx_slide_texts(path: Path) -> list[str]:
    slide_texts: list[str] = []
    with zipfile.ZipFile(str(path)) as archive:
        slide_names = sorted(
            [
                name
                for name in archive.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ],
            key=slide_number_from_pptx_name,
        )
        for name in slide_names:
            xml_text = archive.read(name).decode("utf-8", "ignore")
            slide_texts.append(normalize_whitespace(" ".join(re.findall(r"<a:t>(.*?)</a:t>", xml_text))))
    return slide_texts


def rendered_problem_target_status(rendered_text: str, plan_slide: dict[str, Any]) -> tuple[bool, str]:
    targets = primary_source_problem_targets(plan_slide, limit=2)
    if not targets:
        return True, ""
    for target in targets:
        if source_problem_text_overlap(rendered_text, [target]):
            return True, target
    return False, targets[0]


def build_rendered_quality_report(plan: dict[str, Any], outputs: dict[str, Path]) -> dict[str, Any]:
    report: dict[str, Any] = {
        "passed": True,
        "style_version": PUBLISHER_STYLE_VERSION,
        "sessions": {},
    }
    session_specs = planned_session_specs(plan)
    for session_key, label, output_key in session_specs:
        session_plan = plan.get(session_key, {})
        output_path = outputs.get(output_key)
        slide_texts = extract_pptx_slide_texts(output_path) if output_path else []
        slides = session_plan.get("slides", [])
        session_report: dict[str, Any] = {
            "slide_count": len(slides),
            "problem_slide_count": sum(1 for slide in slides if slide.get("kind") in PROBLEM_SOLVING_KINDS),
            "activity_slide_count": sum(1 for slide in slides if has_activity(slide)),
            "flagship_activity_count": sum(1 for slide in slides if flagship_activity_spec(slide)),
            "flagship_activity_fallbacks": list(session_plan.get("flagship_activity_fallbacks", []))[:4],
            "engagement_slide_count": session_engagement_slide_count(session_plan),
            "engagement_modes": session_engagement_modes(session_plan),
            "long_problem_stack_pages": [],
            "problem_checks": [],
            "copyedit_flags": [],
            "issues": [],
        }
        if output_path is None:
            report["passed"] = False
            session_report["issues"].append("missing rendered output for planned session")
        elif len(slide_texts) != len(slides):
            report["passed"] = False
            session_report["issues"].append("slide count mismatch between plan and rendered deck")
        for page_index, rendered_text in enumerate(slide_texts, start=1):
            style_issues = [
                issue
                for issue in publisher_copyedit_issues(rendered_text)
                if issue != "repeated opening phrase"
            ]
            if style_issues:
                report["passed"] = False
                session_report["copyedit_flags"].append(
                    {
                        "page": page_index,
                        "issue": style_issues[0],
                    }
                )
                session_report["issues"].append(f"page {page_index}: rendered copy issue ({style_issues[0]})")
        for page_index, plan_slide in enumerate(slides, start=1):
            if plan_slide.get("kind") not in PROBLEM_SOLVING_KINDS:
                continue
            rendered_text = slide_texts[page_index - 1] if page_index - 1 < len(slide_texts) else ""
            preserved, main_target = rendered_problem_target_status(rendered_text, plan_slide)
            if prompt_stack_layout_mode(problem_display_cards(plan_slide, variant="practice")) == "focus":
                session_report["long_problem_stack_pages"].append(page_index)
            session_report["problem_checks"].append(
                {
                    "page": page_index,
                    "kind": plan_slide.get("kind", ""),
                    "full_source_problem_preserved": preserved,
                    "main_source_problem": truncate_display_copy(main_target, 220),
                }
            )
            if not preserved:
                report["passed"] = False
                session_report["issues"].append(f"page {page_index}: full source problem not found in rendered text")
        report["sessions"][label] = session_report
    return report


def render_session_notebook(
    session_plan: dict[str, Any],
    *,
    session_label: str,
    lesson_title: str,
    deck: dict[str, Any],
    output_path: Path,
) -> Path:
    prs = make_presentation()
    blank = prs.slide_layouts[6]
    image_lookup = build_image_lookup(deck)
    footer_text = ""
    exact_reference_template = uses_exact_esol_template(session_plan)

    for page_index, plan_slide in enumerate(session_plan["slides"], start=1):
        slide = prs.slides.add_slide(blank)
        kind = plan_slide["kind"]
        if exact_reference_template and render_reference_exact_slide(
            slide,
            plan_slide,
            session_plan,
            session_label,
            deck,
            page_index,
        ):
            continue
        if kind == "cover":
            render_cover_slide(
                slide,
                plan_slide=plan_slide,
                session_label=session_label,
                deck=deck,
                image_lookup=image_lookup,
            )
        elif kind == "be_curious":
            render_be_curious_slide(
                slide,
                plan_slide=plan_slide,
                session_plan=session_plan,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "learning_target":
            render_learning_target_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
            )
        elif kind == "vocabulary":
            render_vocabulary_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "guided_notes":
            render_guided_notes_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "worked_example":
            render_worked_example_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "practice":
            render_practice_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "quick_review":
            render_quick_review_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
            )
        elif kind == "challenge":
            render_challenge_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
                image_lookup=image_lookup,
            )
        elif kind == "reflection":
            render_reflection_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
            )
        elif kind == "exit_ticket":
            render_exit_ticket_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
            )
        else:
            render_generic_slide(
                slide,
                plan_slide=plan_slide,
                page=page_index,
                footer_text=footer_text,
            )

    strip_notebook_write_lines(prs)
    run_formal_release_render_review(prs, session_plan, session_label)
    prs.save(str(output_path))
    return output_path


def render_plan(
    plan: dict[str, Any],
    deck: dict[str, Any],
    output_dir: Path,
) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    outputs: dict[str, Path] = {}
    for session_key, session_label, output_key in planned_session_specs(plan):
        output_path = output_dir / f"{session_label} - Student Notebook.pptx"
        render_session_notebook(
            plan[session_key],
            session_label=session_label,
            lesson_title=plan["lesson_title"],
            deck=deck,
            output_path=output_path,
        )
        outputs[output_key] = output_path

    html_output = render_html_notebook(plan, deck, output_dir)
    outputs["html_notebook"] = html_output

    return outputs


def generate_notebook_artifacts(
    source_pptx: Path | str,
    *,
    output_dir: Path | str | None = None,
    model: str = DEFAULT_MODEL,
    offline: bool = False,
    custom_guidance: str = "",
    api_key: str = "",
) -> dict[str, Any]:
    source_path, resolved_output_dir = notebook_runtime_preflight(source_pptx, output_dir=output_dir)
    effective_guidance = enforce_runtime_quality_guidance(custom_guidance)
    deck, deck_path = extract_source_deck(source_path, resolved_output_dir)
    if offline:
        plan, plan_path = generate_heuristic_plan(
            deck,
            resolved_output_dir,
            custom_guidance=effective_guidance,
        )
    else:
        plan, plan_path = generate_plan_with_openai(
            deck,
            resolved_output_dir,
            model,
            custom_guidance=effective_guidance,
            api_key=api_key,
        )
    outputs = render_plan(plan, deck, resolved_output_dir)
    quality_report = build_rendered_quality_report(plan, outputs)
    quality_report_path = resolved_output_dir / "quality_report.json"
    write_json(quality_report_path, quality_report)
    if not quality_report.get("passed", False):
        issues: list[str] = []
        for session_label, session_report in quality_report.get("sessions", {}).items():
            issues.extend(f"{session_label} {issue}" for issue in session_report.get("issues", []))
        raise RuntimeError("Rendered notebook quality report failed: " + " | ".join(issues[:8]))
    return {
        "source_pptx": source_path,
        "output_dir": resolved_output_dir,
        "effective_guidance": effective_guidance,
        "deck": deck,
        "deck_path": deck_path,
        "plan": plan,
        "plan_path": plan_path,
        "outputs": outputs,
        "quality_report": quality_report,
        "quality_report_path": quality_report_path,
    }


def resolve_output_dir(source_pptx: Path, output_dir: Path | None) -> Path:
    if output_dir:
        return output_dir.resolve()
    return (ROOT / f"{slugify(source_pptx.stem)}-notebook-build").resolve()


def cmd_extract(args: argparse.Namespace) -> int:
    source_pptx, output_dir = notebook_runtime_preflight(args.source, output_dir=args.output_dir)
    _deck, deck_path = extract_source_deck(source_pptx, output_dir)
    print(f"Extracted source deck -> {deck_path}")
    return 0


def cmd_plan(args: argparse.Namespace) -> int:
    source_pptx, output_dir = notebook_runtime_preflight(args.source, output_dir=args.output_dir)
    deck, _deck_path = extract_source_deck(source_pptx, output_dir)
    effective_guidance = enforce_runtime_quality_guidance(args.guidance)
    if args.offline:
        _plan, plan_path = generate_heuristic_plan(deck, output_dir, custom_guidance=effective_guidance)
    else:
        _plan, plan_path = generate_plan_with_openai(deck, output_dir, args.model, custom_guidance=effective_guidance)
    print(f"Notebook plan -> {plan_path}")
    return 0


def cmd_render(args: argparse.Namespace) -> int:
    plan_path = validate_json_artifact_path(args.plan, label="Notebook plan JSON")
    output_dir = validate_output_dir_path(plan_path.parent if not args.output_dir else args.output_dir)
    deck_path = validate_json_artifact_path(args.deck, label="Source deck JSON")
    plan = read_json(plan_path)
    deck = read_json(deck_path)
    effective_guidance = enforce_runtime_quality_guidance(getattr(args, "guidance", ""))
    plan = enforce_plan_requirements(plan, deck, custom_guidance=effective_guidance)
    outputs = render_plan(plan, deck, output_dir)
    quality_report = build_rendered_quality_report(plan, outputs)
    quality_report_path = output_dir / "quality_report.json"
    write_json(quality_report_path, quality_report)
    if not quality_report.get("passed", False):
        issues: list[str] = []
        for session_label, session_report in quality_report.get("sessions", {}).items():
            issues.extend(f"{session_label} {issue}" for issue in session_report.get("issues", []))
        raise RuntimeError("Rendered notebook quality report failed: " + " | ".join(issues[:8]))
    for label, path in outputs.items():
        print(f"{label}: {path}")
    print(f"quality_report: {quality_report_path}")
    return 0


def cmd_run(args: argparse.Namespace) -> int:
    result = generate_notebook_artifacts(
        args.source,
        output_dir=args.output_dir,
        model=args.model,
        offline=args.offline,
        custom_guidance=args.guidance,
    )
    print(f"Extracted deck: {result['deck_path']}")
    print(f"Notebook plan: {result['plan_path']}")
    for label, path in result["outputs"].items():
        print(f"{label}: {path}")
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Generate a compact Session 1 student notebook from a source PPTX deck."
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    extract_parser = subparsers.add_parser("extract", help="Extract a PPTX source deck into structured JSON.")
    extract_parser.add_argument("source", help="Path to the source .pptx file")
    extract_parser.add_argument("--output-dir", help="Directory for extracted artifacts")
    extract_parser.set_defaults(func=cmd_extract)

    plan_parser = subparsers.add_parser("plan", help="Create a compact Session 1 notebook plan JSON.")
    plan_parser.add_argument("source", help="Path to the source .pptx file")
    plan_parser.add_argument("--output-dir", help="Directory for generated artifacts")
    plan_parser.add_argument("--model", default=DEFAULT_MODEL, help="OpenAI model to use")
    plan_parser.add_argument("--offline", action="store_true", help="Use the local heuristic planner instead of the OpenAI API")
    plan_parser.add_argument("--guidance", default="", help="Optional extra instructions for notebook style or constraints")
    plan_parser.set_defaults(func=cmd_plan)

    render_parser = subparsers.add_parser("render", help="Render notebook PPTX files from a saved plan.")
    render_parser.add_argument("plan", help="Path to notebook_plan.json")
    render_parser.add_argument("--deck", required=True, help="Path to source_deck.json")
    render_parser.add_argument("--output-dir", help="Directory for rendered notebooks")
    render_parser.add_argument("--guidance", default="", help="Optional extra instructions for notebook style or constraints")
    render_parser.set_defaults(func=cmd_render)

    run_parser = subparsers.add_parser("run", help="Extract, plan, and render in one command.")
    run_parser.add_argument("source", help="Path to the source .pptx file")
    run_parser.add_argument("--output-dir", help="Directory for generated artifacts")
    run_parser.add_argument("--model", default=DEFAULT_MODEL, help="OpenAI model to use")
    run_parser.add_argument("--offline", action="store_true", help="Use the local heuristic planner instead of the OpenAI API")
    run_parser.add_argument("--guidance", default="", help="Optional extra instructions for notebook style or constraints")
    run_parser.set_defaults(func=cmd_run)

    return parser


def main() -> int:
    load_dotenv()
    parser = build_parser()
    args = parser.parse_args()
    try:
        return args.func(args)
    except KeyboardInterrupt:
        print("Notebook engine run cancelled.", file=sys.stderr)
        return 130
    except Exception as exc:
        if os.getenv(DEBUG_TRACEBACK_ENV, "").strip().lower() in {"1", "true", "yes", "on"}:
            raise
        print(f"Error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
