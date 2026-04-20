from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from utils import (
    LessonPlanError,
    clean_line,
    clean_source_line,
    normalize_whitespace,
    parse_standards,
    sanitize_artifact_path,
    unique_preserve,
    write_json,
)


LEARNING_TARGET_MARKERS = (
    "learning target",
    "objective",
    "i can",
    "we will",
    "success criteria",
)


def iter_leaf_shapes(shapes: Any):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def extract_table_text(table: Any) -> list[str]:
    rows: list[str] = []
    for row in table.rows:
        cells = [clean_source_line(cell.text) for cell in row.cells]
        cells = [cell for cell in cells if cell]
        if cells:
            rows.append(" | ".join(cells))
    return rows


def extract_text_items(slide: Any) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for shape in iter_leaf_shapes(slide.shapes):
        text_lines: list[str] = []
        if getattr(shape, "has_text_frame", False):
            text_lines = [clean_source_line(line) for line in shape.text_frame.text.splitlines()]
        elif getattr(shape, "has_table", False):
            text_lines = extract_table_text(shape.table)
        text_lines = [line for line in text_lines if line]
        if not text_lines:
            continue

        is_title_placeholder = False
        if getattr(shape, "is_placeholder", False):
            try:
                is_title_placeholder = shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                )
            except Exception:
                is_title_placeholder = False

        items.append(
            {
                "text": " ".join(text_lines),
                "lines": text_lines,
                "top": int(getattr(shape, "top", 0)),
                "left": int(getattr(shape, "left", 0)),
                "is_title_placeholder": is_title_placeholder,
            }
        )
    items.sort(key=lambda item: (item["top"], item["left"], len(item["text"])))
    return items


def pick_slide_title(text_items: list[dict[str, Any]]) -> str:
    for item in text_items:
        if item["is_title_placeholder"]:
            return item["text"][:140]
    for item in text_items:
        if len(item["text"]) <= 140:
            return item["text"]
    return text_items[0]["text"][:140] if text_items else ""


def detect_learning_target_candidate(title: str, full_text: str) -> bool:
    lowered = f"{title} {full_text}".lower()
    return any(marker in lowered for marker in LEARNING_TARGET_MARKERS) or bool(parse_standards(full_text))


def extract_speaker_notes(slide: Any) -> list[str]:
    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        return []

    lines = [
        clean_source_line(line)
        for line in notes_text.splitlines()
        if clean_source_line(line)
    ]
    filtered = [
        line
        for line in lines
        if "click to add notes" not in line.lower()
    ]
    return unique_preserve(filtered, cleaner=clean_source_line)


def extract_slide_deck(
    pptx_path: Path,
    output_path: Path,
    *,
    base_dir: Path | None = None,
) -> dict[str, Any]:
    if not pptx_path.exists():
        raise LessonPlanError(f"Slide deck not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise LessonPlanError(f"Expected a .pptx file, found: {pptx_path.name}")

    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:  # pragma: no cover - depends on malformed binary input
        raise LessonPlanError(f"Could not read PowerPoint deck: {pptx_path}") from exc

    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_items = extract_text_items(slide)
        if not text_items:
            full_text = ""
            title = ""
            visible_lines: list[str] = []
        else:
            title = pick_slide_title(text_items)
            visible_lines = unique_preserve(
                [line for item in text_items for line in item["lines"]],
                cleaner=clean_source_line,
            )
            full_text = normalize_whitespace(" ".join(visible_lines))
        speaker_notes = extract_speaker_notes(slide)

        slide_payload = {
            "slide_number": slide_number,
            "title": title,
            "text_items": visible_lines,
            "full_text": full_text,
            "speaker_notes": speaker_notes,
            "is_learning_target_candidate": detect_learning_target_candidate(title, full_text),
        }
        slides.append(slide_payload)

    candidate_numbers = [slide["slide_number"] for slide in slides if slide["is_learning_target_candidate"]]
    raw_payload = {
        "source_file": sanitize_artifact_path(pptx_path, base_dir),
        "source_filename": pptx_path.name,
        "slide_count": len(slides),
        "slides": slides,
        "learning_target_candidate_numbers": candidate_numbers,
    }
    write_json(output_path, raw_payload)
    return raw_payload
